"""
app.py — ImpactProof: Pre-submission evidence quality scorer for MEL teams.

Run with:  streamlit run app.py

Screen flow driven by st.session_state["screen"] (0-3):
  0  Landing & Onboarding
  1  Reported Result Submission
  2  Confidence Snapshot & Next Steps
  3  Portfolio / Framework Dashboard (multi-indicator logframe upload)

Scoring engine: deterministic, no API calls — see evaluator.py. Same inputs always produce the same output.
Audit My Report (Screen 3): transmits document content to Anthropic's Claude API for result extraction.
"""

import base64
import json
import os
import pathlib
import re
import threading
import time
import urllib.parse
from datetime import datetime, date

import streamlit as st
import evaluator as _evaluator
import metrics
from diagnostics import (
    _DIAGNOSTIC_BADGE, _READINESS_BAND, _READINESS_STYLE, _LIMITS_DISCLAIMER,
    _readiness_banner_html, get_diagnostic_state,
    _BRAND_BADGE, _VERDICT_CSS, _DIRECTNESS_TIPS, _VERIFICATION_TIPS,
    _RECENCY_TIPS, _CLARITY_TIPS, _SCORING_GUIDE, _axis_badge_html,
    _overview_score_values, _build_overview_chart_b64,
    DONOR_PROFILES, build_donor_crosswalk_html,
)
from prompts import (
    TOOLTIP_DEFINITION, TOOLTIP_MEASUREMENT, TOOLTIP_INTEGRITY,
    TOOLTIP_SCOPE, TOOLTIP_GOVERNANCE,
    BENEFICIARY_VOICE_TOOLTIP, BENEFICIARY_VOICE_WHATTOFIX,
    METHODOLOGY_STACK,
)
from donor_templates import DONOR_DIAGNOSTICS

# --- Payment / auth / DB utilities ---
try:
    from utils.db import (
        get_user, upsert_user, mark_paid,
        is_still_paid, save_example, get_examples,
        save_user_draft, load_user_draft, clear_user_draft,
        get_payment_history, delete_wa_conversations,
    )
    from utils.paystack import (
        initialize_payment, verify_payment, last_payment_error,
        initialize_subscription_payment, disable_subscription,
    )
    from utils.anonymize import anonymize as _anonymize_value
    from utils.auth import (
        send_login_email, verify_magic_link_token, redeem_magic_link_token,
        issue_session_token, verify_session_token,
        list_sessions, revoke_session, revoke_all_sessions,
    )
    from utils.metering import check_access, record_check, FREE_CHECKS_LIMIT
    _UTILS_AVAILABLE = True
except ImportError:
    _UTILS_AVAILABLE = False
    import logging as _logging
    _logging.exception(
        "Payment/auth/DB utils failed to import -- app is running in degraded "
        "stub mode (no login, no payments, no usage tracking) until this is fixed."
    )
    def get_user(e): return None
    def upsert_user(e): return None
    def mark_paid(e, days=30): pass
    def is_still_paid(u): return False
    def save_example(f, s, v): pass
    def save_user_draft(email, json_str): pass
    def load_user_draft(email): return None
    def clear_user_draft(email): pass
    def get_examples(f, s, k=5): return []
    def get_payment_history(e, limit=50): return []
    def delete_wa_conversations(e): pass
    def initialize_payment(e, a, p="per_use"): return ""
    def verify_payment(r): return {"status": "error", "amount": 0, "plan": ""}
    def last_payment_error(): return ""
    def initialize_subscription_payment(e, a, plan_code, plan_label): return ""
    def disable_subscription(subscription_code, email_token): return False, "Billing is not configured."
    def _anonymize_value(v): return None
    def send_login_email(e, base_url): return False, "Login is not configured.", ""
    def verify_magic_link_token(t): return None
    def redeem_magic_link_token(t): return None
    def issue_session_token(e, user_agent=""): return ""
    def verify_session_token(t): return None
    def list_sessions(e): return []
    def revoke_session(token_hash, e): pass
    def revoke_all_sessions(e): pass
    FREE_CHECKS_LIMIT = 3
    def check_access(e):
        return {"is_paid": False, "plan": "free", "checks_used": 0,
                "checks_remaining": 0, "allowed": False}
    def record_check(e): pass
# --- End utils imports ---

# --- OTP email verification ---
try:
    from utils.email_otp import otp_enabled
except ImportError:
    def otp_enabled(): return False
# --- End OTP email verification ---

# --- Opt-in audit persistence / Logframe Library / benchmark ---
# Isolated in its own try/except (not bundled with the payment/auth/DB block
# above) so a failure here degrades only these opt-in features, not login or
# payments -- exactly the kind of hard-to-diagnose bundled failure that
# turned out to be the root cause of an earlier production incident this
# session (see CLAUDE.md/git history: a single shared except ImportError
# silently masked which of five unrelated modules had actually failed).
try:
    from utils.audits import (
        save_audit, list_audits, get_audit, delete_audit,
        create_logframe_library, list_logframe_libraries,
        add_library_items, get_library_items, delete_logframe_library,
        get_benchmark, MIN_BENCHMARK_SAMPLE, check_rate_limit, log_access,
        purge_account_audit_content, last_audit_error,
    )
    _AUDITS_AVAILABLE = True
except ImportError:
    _AUDITS_AVAILABLE = False
    import logging as _logging
    _logging.exception(
        "utils.audits failed to import -- saved audit history, Logframe "
        "Library, and the comparison benchmark are unavailable until fixed. "
        "Scoring, login, and payments are unaffected."
    )
    def save_audit(email, submissions, evaluations, ref_id): return None
    def list_audits(email, limit=50): return []
    def get_audit(email, audit_id): return None
    def delete_audit(email, audit_id): pass
    def create_logframe_library(email, name): return None
    def list_logframe_libraries(email): return []
    def add_library_items(library_id, email, items): pass
    def get_library_items(library_id, email): return []
    def delete_logframe_library(library_id, email): pass
    def get_benchmark(donor, sector, org_type, my_confidence, my_clarity): return None
    MIN_BENCHMARK_SAMPLE = 10
    def check_rate_limit(email, action, max_count, window_seconds): return True  # fail open
    def log_access(email, action, resource_type=None, resource_id=None, ip_address=None): pass
    def purge_account_audit_content(email): return {"audits_deleted": 0, "libraries_deleted": 0}
    def last_audit_error(): return "utils.audits failed to import."
# --- End opt-in audit persistence ---


def _safe_rate_limit_ok(email: str, action: str, max_count: int, window_seconds: int) -> bool:
    """check_rate_limit() already fails open (returns True) on any DB error --
    this is a second, belt-and-suspenders layer at the call site so this
    optional subsystem can never crash a core feature (scoring, IRC,
    portfolio upload, saving) no matter what goes wrong underneath it."""
    try:
        return check_rate_limit(email, action, max_count, window_seconds)
    except Exception:
        return True


def _safe_log_access(email: str, action: str, resource_type: str | None = None,
                      resource_id=None, ip_address: str | None = None) -> None:
    """Same reasoning as _safe_rate_limit_ok -- logging must never be able to
    break the feature it's observing."""
    try:
        log_access(email, action, resource_type=resource_type,
                   resource_id=resource_id, ip_address=ip_address)
    except Exception:
        pass

# --- UX: INSTANT REPORT CHECK IMPORTS (v3.2) ---
import anthropic as _anthropic
try:
    import fitz as _fitz
    _HAS_FITZ = True
except ImportError:
    _fitz = None
    _HAS_FITZ = False
try:
    import docx as _docx
    _HAS_DOCX = True
except ImportError:
    _docx = None
    _HAS_DOCX = False
try:
    import pdfplumber as _pdfplumber
    _HAS_PDFPLUMBER = True
except ImportError:
    _pdfplumber = None
    _HAS_PDFPLUMBER = False
try:
    import pandas as _pd
    _HAS_PANDAS = True
except ImportError:
    _pd = None
    _HAS_PANDAS = False
try:
    import pptx as _pptx
    _HAS_PPTX = True
except ImportError:
    _pptx = None
    _HAS_PPTX = False
# --- END UX: INSTANT REPORT CHECK IMPORTS (v3.2) ---

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# --- Payment / usage constants (edit amounts here) ---
# FREE_CHECKS_LIMIT now lives in utils/metering.py (imported above) — the
# module that enforces it is the canonical place to define it.
PRICE_PER_CHECK_GHS   = 500        # pesewas  (GHS 5.00)
PRICE_MONTHLY_GHS     = 5000       # pesewas  (GHS 50.00/month)
PRICE_AGENCY_GHS      = 20000      # pesewas  (GHS 200.00/month — Agency tier)
PRICE_ANNUAL_GHS      = 50000      # pesewas  (GHS 500.00/year — Professional annual)

# Canonical app URL — used in reports, emails, payment callbacks, share links.
# Override by setting APP_BASE_URL in Streamlit secrets or environment variable.
def _get_app_url() -> str:
    _env = os.environ.get("APP_BASE_URL", "")
    try:
        _cfg = st.secrets.get("APP_BASE_URL") or _env
        if _cfg:
            return _cfg.rstrip("/")
        # Try to detect live URL from request context (Streamlit 1.30+)
        try:
            _host = st.context.headers.get("Host") or st.context.headers.get("host", "")
            if _host and "localhost" not in _host:
                return f"https://{_host}"
        except Exception:
            pass
    except Exception:
        if _env:
            return _env.rstrip("/")
    return "https://impact-integrity-diagnostic.streamlit.app"

APP_URL = _get_app_url()
# --- End payment constants ---

# --- Disposable / temporary email domains (blocked at the email gate) ---
# Prevents users from cycling through throwaway addresses to reset the
# free-checks counter. Not exhaustive, but covers the most common services.
DISPOSABLE_EMAIL_DOMAINS = {
    "mailinator.com", "10minutemail.com", "10minutemail.net", "guerrillamail.com",
    "guerrillamail.net", "guerrillamail.org", "guerrillamail.biz", "sharklasers.com",
    "tempmail.com", "temp-mail.org", "tempmail.net", "tempmail.dev", "throwawaymail.com",
    "yopmail.com", "yopmail.net", "fakeinbox.com", "trashmail.com", "trashmail.net",
    "getnada.com", "mailnesia.com", "maildrop.cc", "mintemail.com", "mohmal.com",
    "dispostable.com", "moakt.com", "emailondeck.com", "mailcatch.com", "spambog.com",
    "33mail.com", "discard.email", "mytemp.email", "tempinbox.com", "burnermail.io",
    "mail-temp.com", "fakemailgenerator.com", "inboxkitten.com", "tempm.com",
}


def _is_disposable_email(email: str) -> bool:
    """Return True if the email's domain is a known disposable/temp-mail provider."""
    domain = email.strip().lower().rsplit("@", 1)[-1]
    return domain in DISPOSABLE_EMAIL_DOMAINS

EVIDENCE_TYPES = [
    "Select evidence type...",
    "Attendance sheets / participant registers",
    "Raw datasets or survey exports",
    "Partner verification letters",
    "Photos with metadata",
    "Tracer survey results",
    "Baseline and endline study",
    "Financial records",
    "Third-party audits",
    "Case study",
    "Outcome harvesting",
    "Beneficiary narrative or testimony",
    "Community register / village book",
    "Community scorecard / participatory assessment",
    "Participatory Rural Appraisal (PRA) output",
    "Other",
]

# Evidence types scored on sourcing rigor, triangulation, and bias
# mitigation (Qualitative Evidence Track) instead of measurement precision.
QUALITATIVE_EVIDENCE_TYPES = (
    "Case study",
    "Outcome harvesting",
    "Beneficiary narrative or testimony",
    "Community register / village book",
    "Community scorecard / participatory assessment",
    "Participatory Rural Appraisal (PRA) output",
)

# Per-type wording for the five Qualitative Rigor checkboxes (qual_sourcing,
# qual_triangulated, qual_bias feed Measurement/"Sourcing & Triangulation";
# qual_voice, qual_consent feed Definition/"Narrative Definition"). Same
# keys/scoring across types — only the label is tailored so it matches what
# the officer is actually looking at.
QUAL_RIGOR_CHECKLIST = {
    "Case study": (
        "Case/respondent selection method documented "
        "(not just convenience — explain how and why these cases were chosen)",
        "Cross-checked against another source or method (triangulation)",
        "Recall, social-desirability, or selection bias considered and addressed",
        "The case(s) reflect a representative range of beneficiaries, not just one outlier story",
        "Consent to share this account was obtained, and identifying details handled appropriately",
    ),
    "Outcome harvesting": (
        "Outcome description and contribution claim documented "
        "(what changed, and how the program plausibly contributed)",
        "Outcome substantiated with the people/organizations who experienced or observed it (triangulation)",
        "Confirmation bias or post-hoc rationalization considered and addressed",
        "The people/organizations consulted reflect a representative range of those affected",
        "Consent to share this account was obtained, and identifying details handled appropriately",
    ),
    "Beneficiary narrative or testimony": (
        "Selection of who was interviewed/recorded is documented "
        "(not just the most positive or available story)",
        "Account cross-checked against another source or method (triangulation)",
        "Social-desirability or success-story bias considered and addressed",
        "The voices represented reflect a range of beneficiaries, not just one outlier story",
        "Consent to share this account was obtained, and identifying details handled appropriately",
    ),
}

# Per-evidence-type quality checks rendered inline (no expander wrapper).
# Ordered most-specific-first to avoid false positives on substring matches.
_EV_TYPE_KEYWORDS: dict[str, list[str]] = {
    # Specific first — prevent "tracer survey" matching generic "survey"
    "Tracer survey results": [
        "tracer survey", "tracer study", "tracer interview", "cohort follow-up",
        "follow-up survey", "follow up survey",
    ],
    # Third-party requires multi-word match — avoid "audit" alone matching anything
    "Third-party audits": [
        "third-party audit", "third party audit", "independent audit",
        "external auditor", "external evaluation firm", "independent evaluator",
        "external verifier", "external verification report",
    ],
    "Outcome harvesting": [
        "outcome harvest", "outcome harvesting", "outcome mapping",
    ],
    "Case study": [
        "case study", "case-study", "success story", "learning story",
        "interview transcript", "focus group discussion", "focus group", "fgd",
        "most significant change", " msc ", "qualitative interview",
        "key informant interview", "kii", "narrative report",
    ],
    "Beneficiary narrative or testimony": [
        "testimony", "beneficiary narrative", "beneficiary story",
        "quote from", "beneficiary quote", "voice of beneficiar",
        "direct testimony", "beneficiary feedback",
    ],
    "Partner verification letters": [
        "partner letter", "verification letter", "confirmation letter",
        "letter from partner", "signed letter", "partner confirmation",
        "letter of verification", "letter of confirmation",
    ],
    "Photos with metadata": [
        "photo", "photograph", "image with", "gps metadata", "gps coordinate",
        "geotagged", "timestamped photo", "photos of", "site photo",
    ],
    "Attendance sheets / participant registers": [
        "attendance sheet", "attendance register", "participant register",
        "participant list", "sign-in sheet", "sign in sheet", "sign-up sheet",
        "training register", "beneficiary roster", "roster of",
        "attendance record",
    ],
    "Financial records": [
        "financial record", "receipt", "bank statement", "payment record",
        "invoice", "bank transfer", "ledger", "expenditure record",
        "petty cash", "financial report", "transaction record",
    ],
    "Raw datasets or survey exports": [
        "dataset", "questionnaire", "kobo", "kobotoolbox", "csv export",
        "excel export", "survey data", "survey export", "household survey",
        "baseline survey", "endline survey", "midline survey",
        "pre/post test", "pre-test", "post-test", "pretest", "posttest",
        "assessment score", "test score", "exam result",
        "data collection form", "interview schedule", "collected from",
        # broad fallback keywords — placed last so specific types match first
        "survey", "dataset", "data collected",
    ],
}


def _smart_extract_ev_type(ev_desc: str, ev_type_key: str) -> None:
    """Auto-set evidence type from description keywords if still on default.
    Iterates _EV_TYPE_KEYWORDS in order (most-specific first) and returns on
    first match, so general keywords like 'survey' never beat specific ones."""
    if not ev_desc or len(ev_desc) < 15:
        return
    current = st.session_state.get(ev_type_key, "")
    if current and current not in ("Select evidence type...", EVIDENCE_TYPES[0], EVIDENCE_TYPES[1], ""):
        return  # user already made an explicit choice — don't overwrite
    desc_lower = ev_desc.lower()
    for ev_type, keywords in _EV_TYPE_KEYWORDS.items():
        if ev_type not in EVIDENCE_TYPES:
            continue  # skip any dict keys that don't map to a real option
        if any(kw in desc_lower for kw in keywords):
            st.session_state[ev_type_key] = ev_type
            return


# Maps fix dimension → IRC field keys whose page numbers are relevant to that fix.
# Used to annotate Priority fixes with "found on page N" when IRC was used.
_FIX_FIELD_SOURCE_MAP: dict[str, list[str]] = {
    "confidence":   ["result_statement", "evidence_description"],
    "clarity":      ["result_statement", "logframe_indicator"],
    "directness":   ["evidence_description"],
    "verification": ["evidence_description"],
    "recency":      ["evidence_description"],
    "definition":   ["result_statement"],
    "measurement":  ["logframe_indicator"],
    "integrity":    ["evidence_description"],
    "scope":        ["result_statement"],
    "governance":   ["result_statement"],
}


_EV_QUALITY_CHECKS: dict[str, list[tuple[str, str]]] = {
    "Attendance sheets / participant registers": [
        ("signatures_verified", "Signatures verified against ID list"),
        ("date_stamped",        "Sheets dated and stamped"),
        ("cross_ref",           "Cross-referenced with another source (e.g., facilitator notes)"),
    ],
    "Raw datasets or survey exports": [
        ("sample_doc",   "Sampling method documented"),
        ("clean_data",   "Dataset cleaned and de-duplicated"),
        ("version_ctrl", "Original raw export retained for audit"),
    ],
    "Partner verification letters": [
        ("letterhead",       "Letter on official partner letterhead"),
        ("authority_signed", "Signed by authorized partner representative"),
        ("recent_letter",    "Letter dated within 6 months of reporting period"),
    ],
    "Photos with metadata": [
        ("gps_meta",       "Photos contain GPS metadata"),
        ("timestamp_photo","Timestamps visible/verifiable"),
        ("consent_photo",  "Beneficiary consent obtained for photos"),
    ],
    "Tracer survey results": [
        ("followup_tracer", "Follow-up conducted at appropriate interval (3+ months)"),
        ("response_rate",   "Response rate documented (target: 60%+)"),
        ("bias_ack",        "Sampling bias / non-response acknowledged"),
    ],
    "Financial records": [
        ("receipts_dated",  "Receipts/transactions dated"),
        ("reconciled_ev",   "Reconciled with bank/MoMo statements"),
        ("audit_trail_ev",  "Audit trail intact (request → approval → payment)"),
    ],
    "Third-party audits": [
        ("independent_ev",    "Auditor independent from implementer"),
        ("signed_audit",      "Audit report signed and dated"),
        ("recommendations_ev","Audit recommendations addressed/disclosed"),
    ],
}

# Provenance questions most relevant per evidence type — show only these (not all 6).
# Keys must match the provenance_* session state keys used in _render_tab3_slot().
_PROVENANCE_FOR_EV_TYPE: dict[str, list[str]] = {
    "Attendance sheets / participant registers": ["double_counting_checked", "collector_independent"],
    "Raw datasets or survey exports":           ["sampling_documented", "double_counting_checked", "collection_tool_named"],
    "Partner verification letters":             ["auditor_traceable"],
    "Photos with metadata":                     ["collector_independent", "auditor_traceable"],
    "Tracer survey results":                    ["sampling_documented", "recall_period_ok", "collector_independent"],
    "Financial records":                        ["double_counting_checked", "auditor_traceable"],
    "Third-party audits":                       ["collector_independent", "auditor_traceable"],
}
# Fall-back (no evidence type selected or unlisted): show all 5
_PROVENANCE_ALL = ["sampling_documented", "double_counting_checked", "collection_tool_named",
                   "collector_independent", "recall_period_ok"]

# Smart semantic defaults per (evidence_type, provenance_key).
# Phase A: deterministic, semantics-based. Only applied if the user hasn't already answered.
# Phase B: once ≥50 real answers accumulate per type in Supabase, modal answers will override.
_PROVENANCE_DEFAULTS: dict[tuple, str] = {
    # Financial records: exhaustive enumeration — sampling and recall not relevant
    ("Financial records", "sampling_documented"):   "Not applicable",
    ("Financial records", "collector_independent"): "Not applicable",
    ("Financial records", "recall_period_ok"):      "Not applicable",
    # Third-party audits: no sampling methodology needed
    ("Third-party audits", "sampling_documented"):  "Not applicable",
    # Partner verification letters: simple attestation, no sampling or recall
    ("Partner verification letters", "sampling_documented"):     "Not applicable",
    ("Partner verification letters", "recall_period_ok"):        "Not applicable",
    ("Partner verification letters", "double_counting_checked"): "Not applicable",
    # Photos: no recall period risk
    ("Photos with metadata", "recall_period_ok"): "Not applicable",
}

# Human-readable labels for provenance keys
_PROVENANCE_LABELS = {
    "sampling_documented":     "Sampling or selection method documented",
    "double_counting_checked": "Checked for double-counting across activities or periods",
    "collection_tool_named":   "Data-collection tool/method identified (e.g. KoboToolbox, paper form, admin records)",
    "collector_independent":   "Data collected by someone independent of those reporting the result (enumerator-bias risk)",
    "recall_period_ok":        "Recall-period risk assessed (data collected close to the event, or recall bias mitigated)",
}

# --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
PII_EVIDENCE_TYPES = [
    "Attendance sheets / participant registers",
    "Photos with metadata",
    "Raw datasets or survey exports",
    "Tracer survey results",
]

# Evidence types that may carry beneficiary stories, photos, or testimony —
# trigger the do-no-harm / safeguarding check, not just the PII/data-law checks.
SAFEGUARDING_EVIDENCE_TYPES = [
    "Photos with metadata",
    "Case study",
    "Outcome harvesting",
    "Beneficiary narrative or testimony",
]

# Governance checklist display maps: value → (icon, description, pts_earned)
CONSENT_CHECKLIST_MAP = {
    "Yes — written consent forms on file":    ("✓", "Written consent on file", 5),
    "Yes — verbal consent documented":         ("✓", "Verbal consent documented", 3),
    "Partial — some beneficiaries consented":  ("⚠", "Partial consent", 1),
    "Not applicable (no personal data)":       ("✓", "Not applicable", 3),
    "No — consent not obtained":               ("✗", "Consent not obtained", 0),
}
ANON_CHECKLIST_MAP = {
    "Yes — fully anonymized":   ("✓", "Fully anonymized", 4),
    "Partially anonymized":     ("⚠", "Partially anonymized", 2),
    "No — not anonymized":      ("✗", "Not anonymized", 0),
    "Not applicable":           ("✓", "Not applicable", 3),
}
LAW_CHECKLIST_MAP = {
    "Yes — compliant (e.g. Ghana Act 843, Nigeria NDPA, Kenya DPA)": ("✓", "Compliant", 3),
    "Unsure — we haven't checked": ("⚠", "Unsure — needs verification", 1),
    "No — we are not compliant":   ("✗", "Not compliant", 0),
    "Not applicable":              ("✓", "Not applicable", 0),
}
SAFEGUARDING_CHECKLIST_MAP = {
    "Yes — reviewed, no concerns identified":         ("✓", "Do-no-harm review completed", 3),
    "Yes — reviewed, identifying details removed":    ("✓", "Reviewed & de-identified", 3),
    "Partial — some content not yet reviewed":        ("⚠", "Partial review", 1),
    "Not applicable (no beneficiary stories/photos)": ("✓", "Not applicable", 3),
    "No — not yet reviewed":                          ("✗", "Do-no-harm review not completed", 0),
}

# Keywords suggesting the result statement/target group may involve minors —
# used only to decide whether to prompt the child-safeguarding question and
# show a warning. Never used to auto-answer the checklist.
CHILD_SAFEGUARDING_KEYWORDS = (
    "child", "children", "minor", "minors", "student", "students", "pupil", "pupils",
    "adolescent", "adolescents", "youth", "girl", "girls", "boy", "boys", "orphan", "orphans",
)

CHILD_SAFEGUARDING_CHECKLIST_MAP = {
    "Yes — child safeguarding policy applied, guardian consent obtained where applicable":
        ("✓", "Child safeguarding policy applied", 3),
    "Partial — some steps taken but not complete":
        ("⚠", "Partial child safeguarding review", 1),
    "Not applicable (no minors involved)":
        ("✓", "Not applicable", 3),
    "No — not yet reviewed":
        ("✗", "Child safeguarding review not completed", 0),
}

SECURE_HANDLING_CHECKLIST_MAP = {
    "Yes — stored securely, access restricted to authorised staff":
        ("✓", "Securely stored, access-restricted", 3),
    "Partial — some identifiable material not yet secured":
        ("⚠", "Partial secure handling", 1),
    "Not applicable (no identifiable testimony/photos)":
        ("✓", "Not applicable", 3),
    "No — not yet secured":
        ("✗", "Secure handling not yet confirmed", 0),
}
# --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---

# --- UX: INSTANT REPORT CHECK (v3.2) ---
INSTANT_CHECK_SYSTEM_PROMPT = r'''You are an expert MEL (Monitoring, Evaluation, and Learning) data extraction engine for ImpactProof. Your job is to read donor-funded project progress reports and extract structured data to pre-fill a submission verification form.

## YOUR TASK

Extract data from the progress report provided by the user and return a single, valid JSON object — no preamble, no markdown fences, no explanation. Return only the JSON.

## EXTRACTION RULES

### Rule 1 — Always Extract, Never Invent
Extract only what is explicitly stated or can be directly inferred from the document text. Do not fabricate data. If a field is genuinely absent from the document, return "Not found" as the value (string, not null).

### Rule 2 — Infer Intelligently
For fields not explicitly labelled, infer from context. Examples:
- If the report mentions "USAID / Feed the Future", the Primary Donor is "USAID / Feed the Future Ghana".
- If activities span January–March 2026, the Reporting Period Start is "2026/01/01" and End is "2026/03/31".
- If the report says "reviewed by MEL Officer", Internal Review is "Reviewed by MEL Officer".

### Rule 3 — Date Format
All dates must be formatted as YYYY/MM/DD.

### Rule 4 — Geographic Scope
Return as a JSON array of strings, one entry per district/region/country/location mentioned
in relation to the programme — including countries, regions, provinces, districts, cities,
or communities. Do NOT require an explicit "Geographic scope:" label — extract location
names from narrative sentences too, e.g. "implemented in Northern Region", "across Tamale
and Yendi districts", "the Ghana programme", "operating in Kenya and Uganda". If multiple
administrative levels are mentioned for the same place, include the most specific
(e.g. "Tamale, Yendi (Northern Region, Ghana)" rather than just "Ghana").

### Rule 5 — Evidence Type Selection
Map the described evidence to the closest standard type from this list:
- "Attendance sheets / participant registers"
- "Raw datasets or survey exports"
- "Partner verification letters"
- "Photos with metadata"
- "Tracer survey results"
- "Financial records"
- "Third-party audits"
- "Other"
If multiple types apply, list the most dominant one.

### Rule 6 — Submission Type Selection
Choose the closest match from this list, based on the document's title, framing, and content:
"Quarterly progress report", "Annual progress report", "Baseline report", "Mid-term review",
"End-line evaluation", "Final/closeout report", "Project proposal", "Financial report",
"MEL plan", "Others (Special/Ad-hoc reports)"

### Rule 7 — Compliance Flags
For the six compliance fields (consent, anonymisation, data protection, safeguarding, child safeguarding, secure data handling), if the document does not explicitly address them, return "Not found" — do NOT assume compliance. For safeguarding_measures, return text evidence if the document mentions "do no harm" protocols, referral pathways, or participant safety measures. For child_safeguarding, return text evidence if child safeguarding policies are mentioned. For secure_data_handling, return text evidence if secure file handling (e.g., password-protected files, restricted access) is described.

### Rule 8 — Logframe Linkage
Extract the PRIMARY indicator that the main result statement reports against. If multiple indicators are listed, select the one with the highest strategic prominence (usually the reach/beneficiary count indicator at Output level).

### Rule 9 — Evidence Description
For "evidence_description", extract and quote the most detailed evidence description directly from the document — the passage that best explains what data was collected, how it was collected, and by whom. Prefer verbatim or near-verbatim quotes over paraphrase. If multiple passages describe the evidence, combine them into a single coherent paragraph without adding or inferring anything new.

### Rule 10 — Result Statement
The result statement should be the single clearest achievement sentence from the Executive Summary or KPI table. It must contain: (a) a number, (b) a target group, (c) a timeframe, and (d) a % achievement or comparison to target if available.

### Rule 11 — Sector Selection
Choose the closest match from this list based on the document's subject matter:
"WASH", "Health & Nutrition", "Education & Skills", "Agriculture & Livelihoods",
"Youth Employment & TVET", "Climate Resilience", "Governance & Accountability",
"Digital Economy & Technology", "Energy & Clean Energy", "Gender & Social Inclusion",
"Nutrition & Food Security", "Private Sector Development", "Other"

### Rule 12 — Primary Donor
If the document names a specific donor/funder (e.g., USAID, FCDO, GIZ, RVO, World Bank, AfDB,
EU/EuropeAid), return that donor's name. If a donor is mentioned but not one of these, return
the donor's name as written. If no donor is mentioned anywhere, return "Not found".

### Rule 13 — Donor Readiness Inputs
Read the ENTIRE document — including annexes, appendices, lessons-learned sections, M&E/MEL
sections, and any image captions or figure/table descriptions — for the following:
- "learning_and_adaptation": A 1–3 sentence synthesis of what the implementing team learned
  and how the program adapted as a result. Only include this if the document explicitly
  describes a lesson learned, adaptation, course-correction, or change in approach.
- "limitations": A 1–2 sentence synthesis of what the reported data does NOT show, cannot
  confirm, or cannot be generalized to (e.g., sample limitations, geographic scope limits,
  self-reported data caveats). Only include this if the document explicitly states a
  limitation or caveat.
- "result_owner_and_decision": If the document names a person, role, or unit responsible for
  this result (e.g., "MEL Lead", "Project Manager") AND/OR describes a decision the result
  will inform (e.g., "will inform the Q3 budget review"), synthesise both into one sentence.
- "attribution_vs_contribution": Return "Yes" if the document explicitly distinguishes its
  own contribution from other actors/factors (e.g., "alongside government and other NGOs"),
  "No" if it claims sole credit without acknowledging other factors, or "Not found" if
  attribution/contribution isn't discussed.
- "disaggregation_status": Return "Yes — fully disaggregated" if beneficiary data is broken
  down by sex, age, disability, AND location; "Partially disaggregated" if only some of these
  dimensions are present; "No" if beneficiary numbers are reported only as totals; or
  "Not found" if no beneficiary data is reported at all.
For any of the above not found in the document, return "Not found".

### Rule 14 — Documents Referenced
Return a JSON array of short strings naming any standard report components that this
document itself contains, references as attached, or refers to as available annexes —
e.g., "Logframe", "Budget", "Financial report", "M&E plan", "Audit report", "Beneficiary
list", "Disaggregated data", "Case studies", "Sustainability plan", "Action plan". Base
this only on what is explicitly present or referenced in the text — do not guess.

### Rule 15 — Beneficiary Voice
Choose the closest match from this list, based on whether and how beneficiaries
contributed to or validated the evidence in this document:
- "No beneficiary voice captured"
- "Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS)"
- "Beneficiary representatives consulted (community leaders, beneficiary committees)"
- "Anecdotal beneficiary quotes only (uncollected, not systematic)"
- "Not applicable to this result type"
If the document doesn't address this at all, return "Not found".

### Rule 16 — Evidence Strengthening Checks
For the `evidence_strengthening_checks` array, check whether each of the following items is EXPLICITLY confirmed in the document. Return ONLY items confirmed; use EXACTLY the label text shown below (copy it verbatim). Do not paraphrase.

For "Attendance sheets / participant registers": "Signatures verified against ID list" | "Sheets dated and stamped" | "Cross-referenced with another source (e.g., facilitator notes)"
For "Raw datasets or survey exports": "Sampling method documented (who you included, and how)" | "Dataset cleaned and de-duplicated" | "Original raw export retained for audit"
For "Partner verification letters": "Letter on official partner letterhead" | "Signed by authorized partner representative" | "Letter dated within 6 months of reporting period"
For "Photos with metadata": "Photos contain GPS metadata" | "Timestamps visible/verifiable" | "Beneficiary consent obtained for photos"
For "Tracer survey results": "Follow-up conducted at appropriate interval (3+ months)" | "Response rate documented (target: 60%+)" | "Sampling bias / non-response acknowledged"
For "Financial records": "Receipts/transactions dated" | "Reconciled with bank/MoMo statements" | "Audit trail intact (request → approval → payment)"
For "Third-party audits": "Auditor independent from implementer" | "Audit report signed and dated" | "Audit recommendations addressed/disclosed"

Return an empty array [] if none are confirmed or the evidence type does not match any category above.

### Rule 17 — Independent Verifier
For "independent_verifier", extract the name or role of the person or organisation that independently reviewed, verified, or validated the reported data — NOT the implementing organisation itself. Examples: "District Health Information Officer", "External MEL consultant", "Third-party auditor (KPMG)". If no independent verification is mentioned, return "Not found".

## REQUIRED JSON OUTPUT STRUCTURE

Return exactly this structure. Do not add or remove keys.

{
  "result_basics": {
    "result_statement": "<string>",
    "target_group": "<string>",
    "timeframe": "<string>",
    "geographic_scope": ["<string>"],
    "sector": "<string>",
    "primary_donor": "<string>",
    "submission_type": "<string>",
    "beneficiary_voice": "<string>"
  },
  "logframe_linkage": {
    "indicator_name": "<string>",
    "original_target": "<string>",
    "actual_achievement": "<string>"
  },
  "evidence_verification": {
    "evidence_description": "<string>",
    "evidence_type": "<string>",
    "internal_review": "<string>",
    "external_review": "<string>",
    "reporting_period_start": "<YYYY/MM/DD>",
    "reporting_period_end": "<YYYY/MM/DD>",
    "evidence_collection_date": "<YYYY/MM/DD>",
    "consent_documented": "<string>",
    "data_anonymised": "<string>",
    "data_protection_compliant": "<string>",
    "safeguarding_measures": "<string>",
    "child_safeguarding": "<string>",
    "secure_data_handling": "<string>",
    "independent_verifier": "<string>"
  },
  "funder_readiness_inputs": {
    "learning_and_adaptation": "<string>",
    "limitations": "<string>",
    "result_owner_and_decision": "<string>",
    "attribution_vs_contribution": "<string>",
    "disaggregation_status": "<string>"
  },
  "documents_referenced": ["<string>"],
  "evidence_strengthening_checks": ["<string>"],
  "extraction_metadata": {
    "implementing_org": "<string>",
    "report_prepared_by": "<string>",
    "project_name": "<string — the name of the programme/project as stated in the document, or 'Not found'>",
    "confidence_note": "<one sentence describing extraction confidence and any gaps>"
  },
  "field_sources": {
    "result_statement":   {"page": 0, "confidence": "high|medium|low"},
    "target_group":       {"page": 0, "confidence": "high|medium|low"},
    "timeframe":          {"page": 0, "confidence": "high|medium|low"},
    "evidence_description": {"page": 0, "confidence": "high|medium|low"},
    "logframe_indicator": {"page": 0, "confidence": "high|medium|low"}
  }
}

IMPORTANT: For each field in field_sources, set "page" to the 1-based page number in the document where that content was found. If the page cannot be determined, use 0. Use "confidence": "high" if the value was explicitly stated, "medium" if inferred, "low" if uncertain.'''

_UX_TAB_NAMES = ["Your Result", "Logframe", "Evidence", "Review & Score"]

# ---------------------------------------------------------------------------
# BATCH EXTRACTION PROMPT — "Score My Report" pipeline (council XVII)
# Extracts ALL results from a donor report and maps to the 29-column
# portfolio schema, including review_status metadata per field.
# ---------------------------------------------------------------------------
BATCH_EXTRACTION_SYSTEM_PROMPT = r'''You are an expert MEL evidence quality extraction engine.
Your job: read a donor-funded project progress report and extract EVERY reportable result
(KPIs, indicators, outcomes) as structured JSON ready for batch quality scoring.

## OUTPUT FORMAT
Return a JSON object with one key "results" — an array, one object per result found.
Return ONLY the JSON. No preamble, no markdown fences, no explanation.

Each result object must follow this schema exactly:

{
  "indicator_name": "<indicator code + name from logframe — e.g. KPI 2.1: Number of households reached>",
  "result_statement": "<achievement sentence: number + target group + timeframe + % of target if available>",
  "target_group": "<defined population — age, gender, role, occupation>",
  "timeframe": "<bounded date range — e.g. January–June 2025>",
  "geographic_scope": "<location(s) — districts, regions, country>",
  "sector": "<one of: WASH | Health & Nutrition | Education & Skills | Agriculture & Livelihoods | Youth Employment & TVET | Climate Resilience | Governance & Accountability | Digital Economy & Technology | Energy & Clean Energy | Gender & Social Inclusion | Nutrition & Food Security | Private Sector Development | Other>",
  "primary_donor": "<donor name or Not found>",
  "evidence_type": "<one of: Attendance sheets / participant registers | Raw datasets or survey exports | Partner verification letters | Photos with metadata | Tracer survey results | Baseline and endline study | Financial records | Third-party audits | Case study | Outcome harvesting | Beneficiary narrative or testimony | Other>",
  "evidence_description": "<verbatim or near-verbatim quote from document describing what data was collected, how, and by whom>",
  "evidence_date": "<YYYY/MM/DD — date evidence was collected, or Not found>",
  "internal_review": "<one of: Reviewed by MEL Officer | Reviewed by Program Manager | Collected only (no review) | Not reviewed | Not found>",
  "external_review": "<one of: Verified by independent third party | External partner review | No external review | Not found>",
  "verifier": "<name or role of independent reviewer — NOT the implementing org, or Not found>",
  "logframe_indicator": "<exact indicator text from logframe or Technical Proposal, or Not found>",
  "logframe_baseline": "<pre-intervention value with unit and year — e.g. 32% (2022 baseline) — or Not found>",
  "logframe_target": "<approved target from logframe — e.g. 85% by Dec 2025 — or Not found>",
  "logframe_achievement": "<actual reported achievement — e.g. 91% by June 2025 — 107% of target — or Not found>",
  "learning_notes": "<what the team learned and how the programme adapted — or Not found>",
  "limitations_notes": "<what the data does NOT show or cannot confirm — or Not found>",
  "additional_context": "<result owner name/role AND what decision this result informs — or Not found>",
  "beneficiary_voice": "<one of: Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS) | Beneficiary representatives consulted (community leaders, beneficiary committees) | Anecdotal beneficiary quotes only (uncollected, not systematic) | No beneficiary voice captured | Not applicable to this result type | Not found>",
  "bv_method_detail": "<describe how beneficiary feedback was collected, when, and n= if stated — or Not found>",
  "provenance_sampling": "<Yes | No | Not applicable | Not found — was sampling/selection method documented?>",
  "provenance_dedup": "<Yes | No | Not applicable | Not found — was double-counting explicitly checked?>",
  "provenance_tool": "<Yes | No | Not applicable | Not found — is the data collection tool named (e.g., KoboToolbox)?>",
  "provenance_independent": "<Yes | No | Not applicable | Not found — was data collected by someone independent of programme staff?>",
  "provenance_recall": "<Yes | No | Not applicable | Not found — was recall-period risk acknowledged or mitigated?>",
  "provenance_traceable": "<Yes or 'Yes — an auditor could retrieve the original records' | No | Not applicable | Not found — could an auditor trace back to original records?>",
  "field_confidence": {
    "<field_name>": "high|medium|low|not_found"
  }
}

## EXTRACTION RULES

1. EXTRACT ALL RESULTS — do not limit to one. A quarterly progress report may have 5–20 reportable results.
   Each KPI row in a logframe table = one result. Each "Output X.Y" or "Indicator X.Y" row = one result.

2. NEVER INVENT — extract only what is explicitly stated or directly inferable. If absent, return "Not found".
   Do NOT return null or omit fields — always return the string "Not found".

3. PROVENANCE FIELDS — these are often buried in methodology sections.
   - provenance_sampling: look for "random sampling", "purposive sampling", "census", "convenience sample"
   - provenance_dedup: look for "double-counting", "de-duplication", "unique beneficiaries"
   - provenance_tool: look for "KoboToolbox", "ODK", "CommCare", "survey form", "data collection tool"
   - provenance_independent: look for "independent enumerators", "external data collectors", "third-party data collection"
   - provenance_recall: look for "recall period", "recall bias", "data collected same day"
   - provenance_traceable: look for "original records", "source documents", "audit trail", "auditor", "verifiable"

4. LOGFRAME BASELINE — look in annexes, baseline reports, or logframe tables for pre-intervention values.
   Format: "<value> (<year> baseline)" e.g. "32% stunting prevalence (2022 Ghana Health Service baseline)".

5. FIELD CONFIDENCE — for each extracted field, return "high" if explicitly stated, "medium" if inferred,
   "low" if uncertain, "not_found" if the field was absent. Include ALL field names in field_confidence.

6. SECTOR — infer from programme focus, not document title. A health programme reporting
   on water access should use "WASH", not "Health & Nutrition".

7. RESULT STATEMENT — must contain: (a) a number, (b) a defined population, (c) a timeframe.
   Include % achievement vs target if stated. Prefer Executive Summary or logframe achievement rows.

8. ORDER results by logframe/KPI numbering if available (e.g., KPI 2.1 before KPI 2.2).
   If no numbering, use document order.

Return ONLY the JSON. Start your response with { and end with }.'''

# IRC field map: extracted key → session_state key
# Excludes selectbox widgets (sector, donor_selected) that render before tab1
_IRC_FIELD_MAP = {
    # Result Basics tab
    "result_statement":   "result_statement",
    "target_group":       "target_group",
    "timeframe":          "timeframe",
    "geographic_scope":   "geographic_scope",
    # Logframe Linkage tab
    "logframe_indicator": "logframe_indicator",
    "logframe_target":    "logframe_target",
    "logframe_achievement": "logframe_achievement",
    # Evidence & Verification tab
    "evidence_description": "evidence_description",
    "verifier":           "verifier",
}

_IRC_PATTERNS = {
    # --- Result Basics ---
    "result_statement": [
        r"(?:key\s+)?result\s+statement\s*[:\-]\s*(.+)",
        r"reported\s+result\s*[:\-]\s*(.+)",
        r"(?:key\s+)?result\s*[:\-]\s*(.+)",
        r"output\s+(?:statement|achieved)\s*[:\-]\s*(.+)",
        r"outcome\s+(?:statement|achieved)\s*[:\-]\s*(.+)",
        r"achievement\s*[:\-]\s*(.+)",
        r"project\s+(?:name|title)\s*[:\-]\s*(.+)",
        r"programme\s+(?:name|title)\s*[:\-]\s*(.+)",
    ],
    "target_group": [
        r"target\s+(?:group|population|beneficiaries|community)\s*[:\-]\s*(.+)",
        r"beneficiar(?:y|ies)\s*[:\-]\s*(.+)",
        r"primary\s+(?:beneficiar(?:y|ies)|target)\s*[:\-]\s*(.+)",
        r"direct\s+beneficiar(?:y|ies)\s*[:\-]\s*(.+)",
    ],
    "timeframe": [
        r"reporting\s+period\s*[:\-]\s*(.+)",
        r"period\s+covered\s*[:\-]\s*(.+)",
        r"implementation\s+period\s*[:\-]\s*(.+)",
        r"timeframe\s*[:\-]\s*(.+)",
        r"report\s+period\s*[:\-]\s*(.+)",
        r"(?:project|programme)\s+(?:duration|period)\s*[:\-]\s*(.+)",
    ],
    "geographic_scope": [
        r"geographic(?:al)?\s+(?:scope|coverage|area)\s*[:\-]\s*(.+)",
        r"country\s*[:\-]\s*(.+)",
        r"(?:project\s+)?location\s*[:\-]\s*(.+)",
        r"region\s*[:\-]\s*(.+)",
        r"(?:target\s+)?(?:district|county|province|state)s?\s*[:\-]\s*(.+)",
        # Natural-language mentions without explicit labels
        r"(?:implemented|delivered|conducted|operating|operated)\s+(?:in|across|within)\s+([A-Z][^,.\n]+(?:,\s*[A-Z][^,.\n]+)*)",
        r"([A-Z][a-z]+ (?:Region|District|Province|State|Municipality)(?:,\s*[A-Z][a-z]+)*)",
    ],
    # --- Logframe Linkage ---
    "logframe_indicator": [
        r"(?:logframe\s+)?indicator\s*(?:name|description)?\s*[:\-]\s*(.+)",
        r"(?:key\s+)?performance\s+indicator\s*[:\-]\s*(.+)",
        r"output\s+indicator\s*[:\-]\s*(.+)",
        r"outcome\s+indicator\s*[:\-]\s*(.+)",
        r"KPI\s*[:\-]\s*(.+)",
        r"M&E\s+indicator\s*[:\-]\s*(.+)",
    ],
    "logframe_target": [
        r"(?:annual|cumulative|indicator)?\s*target\s*[:\-]\s*(.+)",
        r"planned\s+(?:result|output|target)\s*[:\-]\s*(.+)",
        r"(?:project|programme)\s+target\s*[:\-]\s*(.+)",
    ],
    "logframe_achievement": [
        r"(?:actual|cumulative)?\s*achievement\s*[:\-]\s*(.+)",
        r"actual\s+(?:result|output|figure)\s*[:\-]\s*(.+)",
        r"(?:result|output)\s+achieved\s*[:\-]\s*(.+)",
        r"delivered\s*[:\-]\s*(.+)",
        r"number\s+(?:reached|trained|served|treated|supported)\s*[:\-]\s*(.+)",
    ],
    # --- Evidence & Verification ---
    "evidence_description": [
        r"(?:supporting\s+)?evidence\s*(?:description|type|source)?\s*[:\-]\s*(.+)",
        r"data\s+source\s*[:\-]\s*(.+)",
        r"verification\s+(?:source|method|means)\s*[:\-]\s*(.+)",
        r"means\s+of\s+verification\s*[:\-]\s*(.+)",
        r"data\s+collection\s+(?:method|tool|instrument)\s*[:\-]\s*(.+)",
        r"(?:key\s+)?evidence\s+collected\s*[:\-]\s*(.+)",
    ],
    "verifier": [
        r"verified\s+by\s*[:\-]\s*(.+)",
        r"verification\s+by\s*[:\-]\s*(.+)",
        r"implementing\s+(?:organization|organisation|partner|agency)\s*[:\-]\s*(.+)",
        r"submitted\s+by\s*[:\-]\s*(.+)",
        r"prepared\s+by\s*[:\-]\s*(.+)",
        r"(?:MEL|M&E)\s+officer\s*[:\-]\s*(.+)",
    ],
}


def _extract_text_from_file(fname_lower, raw):
    """Extract plain text from a PDF, DOCX, TXT, CSV, PPTX, or XLSX file's raw bytes.

    Returns (text, error_message). On success error_message is "".
    On failure text is "" and error_message describes why.
    """
    import io as _io
    text = ""
    if fname_lower.endswith(".pdf"):
        if not _HAS_PDFPLUMBER:
            return "", "pdfplumber not installed. Run: pip install pdfplumber"
        with _pdfplumber.open(_io.BytesIO(raw)) as _pdf:
            for _pg in _pdf.pages:
                _pt = _pg.extract_text()
                if _pt:
                    text += _pt + "\n"
    elif fname_lower.endswith(".docx"):
        if not _HAS_DOCX:
            return "", "python-docx not installed. Run: pip install python-docx"
        _dobj = _docx.Document(_io.BytesIO(raw))
        text = "\n".join(p.text for p in _dobj.paragraphs)
    elif fname_lower.endswith(".txt") or fname_lower.endswith(".csv"):
        text = raw.decode("utf-8", errors="replace")
    elif fname_lower.endswith(".pptx"):
        if not _HAS_PPTX:
            return "", "python-pptx not installed. Run: pip install python-pptx"
        _prs = _pptx.Presentation(_io.BytesIO(raw))
        text = "\n".join(
            shape.text
            for slide in _prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    elif fname_lower.endswith(".xlsx") or fname_lower.endswith(".xls"):
        if not _HAS_PANDAS:
            return "", "pandas not installed. Run: pip install pandas openpyxl"
        try:
            _sheets = _pd.read_excel(_io.BytesIO(raw), sheet_name=None)
        except Exception as _xl_exc:
            return "", f"Could not read Excel file: {_xl_exc}"
        for _sheet_name, _df in _sheets.items():
            text += f"\n--- Sheet: {_sheet_name} ---\n"
            text += _df.to_string(index=False) + "\n"
    else:
        return "", "Unsupported file type. Upload a PDF, DOCX, TXT, CSV, PPTX, or XLSX."
    return text, ""


def _extract_report_fields(uploaded_file):
    """Rule-based extraction. No AI. Returns (fields, found_list, not_found_list) or (None, error_str, [])."""
    import re as _re
    fname = uploaded_file.name.lower()
    raw = uploaded_file.read()
    text, _err = _extract_text_from_file(fname, raw)
    if _err:
        return None, _err, []
    if not text.strip():
        return None, "Could not extract text. The file may be scanned/image-based.", []
    fields, found, not_found = {}, [], []
    for field, pats in _IRC_PATTERNS.items():
        matched = False
        for pat in pats:
            m = _re.search(pat, text, _re.IGNORECASE)
            if m:
                fields[field] = m.group(1).strip()[:120]
                found.append(field)
                matched = True
                break
        if not matched:
            fields[field] = ""
            not_found.append(field)
    return fields, found, not_found


def _irc_extract_combined(doc_files):
    """Extract and combine text + rule-based fields from one or more uploaded files.

    Returns (full_text, raw_fields, error_message). error_message is "" on success.
    """
    if not doc_files:
        return "", {}, ("No readable text found in this document. Please upload a "
                         "text-based PDF, DOCX, TXT, PPTX, or XLSX — scanned image "
                         "files cannot be extracted.")
    parts = []
    raw_fields = {}
    for f in doc_files:
        f.seek(0)
        text, err = _extract_text_from_file(f.name.lower(), f.read())
        if err:
            return "", {}, f"{f.name}: {err}"
        if not text.strip():
            return "", {}, (f"No readable text found in {f.name}. Please upload a "
                             f"text-based PDF, DOCX, TXT, PPTX, or XLSX — scanned image "
                             f"files cannot be extracted.")
        parts.append(f"--- Document: {f.name} ---\n\n{text}")
        f.seek(0)
        fields, _, _ = _extract_report_fields(f)
        for fk, fv in (fields or {}).items():
            if fv and not raw_fields.get(fk):
                raw_fields[fk] = fv
    return "\n\n".join(parts), raw_fields, ""


def _irc_parse_date(s):
    """Parse YYYY/MM/DD string to date, return None on failure."""
    if not s or s == "Not found":
        return None
    try:
        return date.fromisoformat(s.replace("/", "-"))
    except (ValueError, AttributeError):
        return None


def _irc_match_option(value, options):
    """Fuzzy-match extracted string to a selectbox options list. Returns matched option or None."""
    if not value or value == "Not found":
        return None
    vl = value.lower().strip()
    for opt in options:
        if vl == opt.lower().strip():
            return opt
    for opt in options:
        ol = opt.lower()
        if vl in ol or ol in vl:
            return opt
    vwords = vl.split()
    for opt in options:
        owords = opt.lower().split()
        if vwords and (vwords[0] in owords or (owords and owords[0] in vwords)):
            return opt
    return None
# --- END UX: INSTANT REPORT CHECK (v3.2) ---

EVIDENCE_TYPE_HELP = (
    "Choose the type that best describes your primary evidence document.\n\n"
    "• Attendance sheets / participant registers — Signed records of participants by name, date, and session. "
    "Example: 'Signed attendance sheets from 12 training sessions, names + signatures + dates'\n\n"
    "• Raw datasets or survey exports — Unprocessed data files exported from a survey or data collection tool. "
    "Example: 'KoboToolbox CSV export of 487 farmer surveys; SPSS dataset from baseline survey'\n\n"
    "• Partner verification letters — Formal letters from partner organizations confirming they witnessed or validated the activity. "
    "Example: 'Letter from District Agriculture Officer confirming attendance at all 12 training sessions'\n\n"
    "• Photos with metadata — Photos with embedded GPS, timestamps, and EXIF data proving where and when they were taken. "
    "Example: 'Geotagged photos of borehole installation with date stamps'\n\n"
    "• Tracer survey results — Follow-up surveys conducted weeks/months after the activity to measure actual outcomes. "
    "Example: '3-month tracer survey results showing 65% of trained farmers adopted climate-smart techniques'\n\n"
    "• Financial records — Receipts, payment confirmations, payroll records that prove transactions occurred. "
    "Example: 'Mobile money transfer receipts to 250 farmer cash transfer recipients'\n\n"
    "• Third-party audits — Independent audits or verification reports from external organizations. "
    "Example: 'External audit by SGS Ghana of distribution logistics and beneficiary lists'\n\n"
    "• Case study — An in-depth account of one participant, site, or community, used to illustrate how change happened. "
    "Example: 'Case study of a women's savings group in Tamale showing how loan access changed household decision-making'\n\n"
    "• Outcome harvesting — Outcomes identified after the fact, then worked backwards to assess the program's contribution. "
    "Example: 'Outcome harvesting exercise with district officials identifying 8 policy changes influenced by the program'\n\n"
    "• Beneficiary narrative or testimony — First-person accounts from participants describing the change they experienced. "
    "Example: 'Recorded interviews with 12 farmers describing changes in income and food security after training'\n\n"
    "• Other (specify) — Evidence that doesn't fit any category above. Use sparingly; most evidence fits one of the above."
)

SECTOR_OPTIONS = [
    "(No sector selected)",
    "WASH",
    "Health & Nutrition",
    "Education & Skills",
    "Agriculture & Livelihoods",
    "Youth Employment & TVET",
    "Climate Resilience",
    "Governance & Accountability",
    "Digital Economy & Technology",
    "Energy & Clean Energy",
    "Gender & Social Inclusion",
    "Nutrition & Food Security",
    "Private Sector Development",
    "Other",
]

# Sector-specific beneficiary voice HOW-TO guidance (Council XXVII)
# Rendered in Screen 2 when BV score < 0.5 and a sector is selected.
_BV_SECTOR_GUIDANCE: dict[str, str] = {
    "Health & Nutrition": (
        "Patient exit interviews (5 min at clinic exit), community health worker "
        "focus groups (10–15 participants), or Ghana Health Service facility satisfaction "
        "survey templates — accepted by FCDO and MCF as structured feedback evidence."
    ),
    "WASH": (
        "Water point committee satisfaction surveys (quarterly, 10 questions), or "
        "CWSA community assessment tools. Document collection date, sample size, "
        "and who administered the survey."
    ),
    "Education & Skills": (
        "Student satisfaction surveys at term end, parent perception interviews at "
        "community level, or GES-administered student assessment forms independently "
        "scored. Beneficiary voice must come from students or parents — not teachers "
        "reporting on their own delivery."
    ),
    "Youth Employment & TVET": (
        "End-of-training satisfaction survey (at programme completion) plus 6-month "
        "employment follow-up phone interview. MCF requires BOTH 6-month AND 12-month "
        "tracer calls. Use a 5-question script with gender and employment status fields. "
        "N=30+ is sufficient for most TVET programmes."
    ),
    "Agriculture & Livelihoods": (
        "Post-season farmer adoption survey (% applying technique + estimated yield change), "
        "women farmer focus group discussion, or market buyer satisfaction data from "
        "offtake partners. USAID FtF and GIZ accept post-season yield comparison "
        "(this season vs. pre-programme) as direct beneficiary outcome evidence."
    ),
    "Climate Resilience": (
        "Community risk perception survey before and after intervention, early warning "
        "system utilisation logs, or community committee meeting minutes documenting "
        "adaptation decisions. EPA Ghana and NADMO accept community-led risk monitoring data."
    ),
    "Governance & Accountability": (
        "Citizen satisfaction survey with district services (pre/post), community "
        "accountability meeting attendance + action point tracking, or NCCE-administered "
        "civic knowledge assessment. Must include marginalised groups (women, PWDs, "
        "youth) — not only elected officials or community leaders."
    ),
    "Nutrition & Food Security": (
        "Household dietary diversity survey (24-hour recall method, N=30+ per site), "
        "caregiver knowledge assessment on complementary feeding, or CMAM exit interview "
        "conducted by an enumerator independent of the nutrition programme team. "
        "GHS CHPS compound records and UNICEF nutrition monitoring tools are accepted "
        "by UNICEF, WFP, and USAID Food for Peace as structured beneficiary voice evidence."
    ),
    "Digital Economy & Technology": (
        "User satisfaction survey after digital tool onboarding (N=20+ per site), "
        "utilisation log showing return sessions beyond registration, or structured "
        "interview on whether the tool changed a workflow or income-generating activity. "
        "Evidence must show active use — not just account creation. GSMA Mobile for "
        "Humanitarian Innovation and FCDO Digital accepts utilisation-based feedback."
    ),
    "Energy & Clean Energy": (
        "Household energy consumption diary (weekly, before/after clean energy access), "
        "post-installation satisfaction survey covering reliability, safety, and household "
        "saving, or community meeting minutes on cookstove or solar adoption decisions. "
        "GOGLA and SE4ALL frameworks accept structured post-installation feedback as BV. "
        "Include a 3-month follow-up to capture sustained use vs. adoption-day reporting."
    ),
    "Gender & Social Inclusion": (
        "Women's agency and decision-making survey (household and community level, pre/post), "
        "GBV incident reporting rate change via community safety committee records, or "
        "structured interview with women from marginalised groups on access change. "
        "Evidence must come from the rights-holder group directly — not programme staff "
        "reporting on observed change. FCDO EQuALS 2 and Sida HRBA require this distinction."
    ),
    "Private Sector Development": (
        "Business owner satisfaction survey on BDS quality (post-support, N=20+ per cohort), "
        "market linkage outcome report from offtake partners (actual purchase records), or "
        "structured interview with SME beneficiaries on revenue or employment change "
        "6–12 months post-support. GIZ Value Chain, IFC DOTS, and USAID CDCS frameworks "
        "accept structured SME-level outcome interviews as verifiable beneficiary voice."
    ),
}

DONOR_GUIDANCE = {
    "USAID": {
        "key_emphasis": "USAID is governed by ADS 201. Five data-quality standards apply to every reported result: Validity, Integrity, Accuracy, Completeness, and Timeliness. Activity MEL Plans must be submitted within 90 days of project start and updated annually.",
        "common_rejection": "Results not tied to PIRS (Performance Indicator Reference Sheets), or missing sex/age disaggregation. USAID auditors specifically flag results without baseline data or without a documented data-collection methodology.",
        "tip": "USAID requires evidence collected within 12 months for a full Recency score. Older evidence must be explicitly flagged and justified in your MEL Plan. Always document your data source, collection date, and sample size in the result narrative.",
    },
    "FCDO": {
        "key_emphasis": "FCDO's January 2025 Evaluation Policy requires all evaluations to be registered in the UK Government Evaluation Registry before fieldwork begins. Results must connect activities to outcomes via an explicit Theory of Change — outputs alone are insufficient.",
        "common_rejection": "Outputs reported without contribution analysis — showing what happened without explaining how the programme contributed. Evaluations that lack documented evaluator independence or a safeguarding review are also routinely rejected.",
        "tip": "FCDO accepts qualitative evidence if triangulated (at least two independent sources). Register your evaluation in the Government Evaluation Registry immediately — mandatory for all evaluations approved after 1 April 2024.",
    },
    "GIZ": {
        "key_emphasis": "GIZ uses the Capacity WORKS framework and requires SMART indicators with documented baselines. Evaluation reports are scored on a mandatory 100-point quality grid — a minimum score of 60 is required before acceptance.",
        "common_rejection": "Insufficient evaluability: weak results model, missing SMART indicators, or no baseline data. Methodological quality scores below 60/100 and poor documentation of partner and target-group perspectives are common rejection triggers.",
        "tip": "GIZ values qualitative learning narratives alongside quantitative KPIs. Use the KOMPASS qualitative method to document stakeholder feedback systematically, and include a lessons-learned section — GIZ evaluators score this explicitly.",
    },
    "RVO": {
        "key_emphasis": "RVO (Netherlands Enterprise Agency) requires logframe-anchored reporting with a clear link to the original Technical Proposal. Every result MUST tie to an approved indicator — deviations must be pre-approved.",
        "common_rejection": "Missing M&E data tied to the original logframe — the #1 RVO rejection cause. Always include a logframe progress table in every progress report.",
        "tip": "RVO final reports require four components: narrative + financial + audit + logframe update. Confirm all four are in your submission package before sending. Missing any one component triggers an automatic return.",
    },
    "World Bank": {
        "key_emphasis": "World Bank projects use a Results Framework with PDO (Project Development Objective) indicators and intermediate outcomes. Baseline data must be collected at project effectiveness — not retrospectively. A PLR (Performance & Learning Review) at mid-term is mandatory for all IDA operations.",
        "common_rejection": "Outcome indicators that are actually outputs (e.g., 'number of trainings held' instead of 'skills improvement rate'). Missing baseline data is the #1 cause of unsatisfactory M&E ratings in World Bank Implementation Status Reports.",
        "tip": "World Bank IDA projects above $5M require third-party verification for Tier-1 indicators. Identify proxy indicators early where direct outcome data is unavailable — document the data gap and mitigation plan in your MEL annex.",
    },
    "AfDB": {
        "key_emphasis": "AfDB's 2024–2029 Country Strategy Paper for Ghana prioritises industrialisation and sustainable transport. Results must align to AfDB's High 5s priorities and explicitly reference the country strategic pillar they support.",
        "common_rejection": "Results not linked to AfDB strategic pillars or missing regional development context. Insufficient evidence of private sector engagement or domestic resource mobilisation is a growing trigger under the 2024–2029 strategy.",
        "tip": "AfDB values African-led monitoring and evaluation. Reference AfrEA or the African Evidence Network methodology where possible, and track private sector participation and business-growth metrics alongside social outcomes.",
    },
    "EU / EuropeAid": {
        "key_emphasis": "EU follows DG INTPA reporting standards and the PRAG (Practical Guide) framework. The Logical Framework Approach (LFA) is the foundation — all results must trace to the logframe with updated assumptions and risks columns.",
        "common_rejection": "Assumptions and risks columns not updated when reporting deviations, and procurement documentation with incomplete audit trails. Missing beneficiary feedback or gender mainstreaming evidence are also common flags.",
        "tip": "EU expects gender mainstreaming and rights-based analysis to be explicitly visible in the results narrative — not just mentioned in passing. Use EU core indicators in your M&E plan and conduct monthly data quality reviews against them.",
    },
    "Global Fund": {
        "key_emphasis": "The Global Fund requires country-owned disease surveillance data that meets its DQA (Data Quality Assessment) minimum standards before each tranche is disbursed. All HIV, TB, and malaria data must be reconciled from facility to district to national level.",
        "common_rejection": "Poor data quality in disease surveillance registers — missing patient IDs, duplicate records, or facility-level data not reconciled with district totals. Weak domestic government co-financing commitment is a growing rejection trigger as Ghana transitions to self-funding by 2030.",
        "tip": "Conduct an internal DQA at facility level every quarter before the Global Fund's external DQA visit. Track the percentage of domestic health budget allocated to programme costs — Global Fund reviewers now weight this heavily in continuation decisions.",
    },
    "Mastercard Foundation": {
        "key_emphasis": "Mastercard Foundation's Young Africa Works strategy targets 3 million dignified jobs in Ghana by 2030, with a minimum 70% women-participation target across all programmes. Results must measure employment outcomes — not just training completion — at 6 and 12 months post-programme.",
        "common_rejection": "Training headcount reported as the primary result without 6-month follow-up employment tracking. Gender equity gaps (women below 70% of participants) and weak digital or innovation integration are the most common programme-design flags.",
        "tip": "Track tracer outcomes at 6 and 12 months post-programme — Mastercard Foundation considers employment retention at 12 months a core outcome, not just initial job placement. Disaggregate all data by gender, age, region, and disability status.",
    },
    "KOICA": {
        "key_emphasis": "KOICA (Korea International Cooperation Agency) focuses on technology transfer and capacity building. Results must demonstrate farmer or community adoption rates of transferred technologies — not just delivery of inputs or training sessions.",
        "common_rejection": "Completion reports showing inputs delivered or training conducted without evidence of technology adoption or productivity change. Weak linkage to market demand and poor post-project maintenance planning are common rejection triggers.",
        "tip": "KOICA evaluators look for adoption sustainability — document who maintains the technology after project end and how community institutions have been strengthened. Track yield or productivity improvement at least two seasons after training, not just immediately after.",
    },
    "SIDA": {
        "key_emphasis": "Sida (Swedish International Development Cooperation Agency) applies a Human Rights-Based Approach (HRBA) to all programming. Results must demonstrate change experienced by rights-holders — disaggregated by gender, age, disability, and marginalisation status. Independently verified outcome data is required under Sida's Results Strategy.",
        "common_rejection": "Outputs reported without evidence of rights-holder outcome change. Gender disaggregation missing or tokenistic — Sida expects women's voices in evidence, not just participant head counts. Evaluations without documented evaluator independence also fail review.",
        "tip": "Sida evaluators explicitly look for participation of marginalised groups in evidence collection. Commission an independent outcome assessment, and document how rights-holders validated the reported outcomes. File your evaluation plan with Sida's country office before fieldwork begins.",
    },
    "SDC": {
        "key_emphasis": "SDC (Swiss Agency for Development and Cooperation) requires outcome-level evidence anchored in its Results Measurement Framework, with a specific emphasis on sustainability and systemic change. Outputs delivered at programme end are insufficient without evidence of continued use or institutional change post-closure.",
        "common_rejection": "Project-level output reporting without outcome evidence. Missing sustainability indicators — SDC evaluators look for evidence the change persists after programme closure. Results that cannot be traced to the SDC Programme Agreement indicator are routinely returned.",
        "tip": "Reference the specific outcome indicator from your SDC Programme Agreement. Conduct a follow-up assessment 6–12 months after key activities to generate sustainability evidence — this is the most frequently cited gap in SDC programme evaluations reviewed by the Swiss Embassy.",
    },
}

SECTOR_PLACEHOLDERS = {
    "WASH": {
        "result":               "e.g., Constructed 25 boreholes serving 12,000 people across 5 districts in Northern Region between January and June 2025",
        "target_group":         "e.g., Rural households without access to safe drinking water; women and children primarily responsible for water collection",
        "geographic_scope":     "e.g., Tamale, Yendi, Savelugu, Karaga, and Kumbungu districts",
        "evidence_description": "e.g., Borehole functionality reports from 25 sites + water quality test results from district lab + GPS-tagged photos of completed structures",
        "logframe_indicator":   "e.g., Indicator 2.1: Number of households with access to safely managed drinking water",
        "logframe_baseline":    "e.g., 3,200 households with access (2021 CWSA district baseline survey)",
        "logframe_target":      "e.g., 12,000 households with access by Q4 2025",
        "logframe_achievement": "e.g., 12,000 people reached by June 2025 — 100% of target",
        "verifier":             "e.g., District Water and Sanitation Officer, Water Resource Commission inspector",
    },
    "Health & Nutrition": {
        "result":               "e.g., Vaccinated 8,500 children under 5 against measles across 3 health districts in Eastern Region between July and September 2025",
        "target_group":         "e.g., Children aged 6 months to 5 years and pregnant or lactating women residing in target communities",
        "geographic_scope":     "e.g., New Juaben, Suhum, and Akropong health districts",
        "evidence_description": "e.g., Patient records from 3 health facilities + immunization registers signed by district health officer + cold chain monitoring logs + MUAC screening results",
        "logframe_indicator":   "e.g., Indicator 1.3: % of children under 5 fully immunized in target districts",
        "logframe_baseline":    "e.g., 62% immunization coverage (2022 DHIMS2 district baseline)",
        "logframe_target":      "e.g., 85% immunization coverage in 3 districts by Dec 2025",
        "logframe_achievement": "e.g., 8,500 children vaccinated by Sept 2025 — 100% of district target",
        "verifier":             "e.g., District Health Officer, Regional Health Directorate field supervisor, UNICEF field monitor",
    },
    "Education & Skills": {
        "result":               "e.g., Improved literacy scores by 35% among 1,200 primary school students across 15 schools in Central Region between September 2024 and June 2025",
        "target_group":         "e.g., Primary school students grades 3–6, ages 8–12, in selected public schools; majority from low-income households",
        "geographic_scope":     "e.g., Cape Coast, Mfantsiman, and Ekumfi districts (15 schools)",
        "evidence_description": "e.g., Pre/post standardized EGRA test results + enrollment registers + teacher observation logs certified by GES inspector + sample of student work",
        "logframe_indicator":   "e.g., Indicator 3.2: % of students achieving minimum reading proficiency (EGRA standard)",
        "logframe_baseline":    "e.g., 28% of students at grade-level literacy (2023 GES baseline EGRA assessment)",
        "logframe_target":      "e.g., 60% of students at grade-level literacy by June 2025",
        "logframe_achievement": "e.g., 1,200 students with improved literacy scores by June 2025 — 100% of target",
        "verifier":             "e.g., Ghana Education Service district inspector, headteacher certification, external EGRA assessor",
    },
    "Agriculture & Livelihoods": {
        "result":               "e.g., Trained 487 smallholder farmers in climate-smart agriculture across 3 districts in Northern Ghana between January and June 2025",
        "target_group":         "e.g., Smallholder farmers (18–60 years), majority women, with land holdings under 2 hectares in target districts",
        "geographic_scope":     "e.g., Tamale, Yendi, and Savelugu districts (Northern Region)",
        "evidence_description": "e.g., Signed attendance sheets from 12 training sessions across 3 districts + District Agriculture Officer verification + farmer cooperative records + pre/post yield measurement data",
        "logframe_indicator":   "e.g., Indicator 2.4: Number of smallholder farmers trained in climate-smart agricultural practices",
        "logframe_baseline":    "e.g., 0 farmers trained (programme start, Jan 2024); avg yield 0.42 t/ha (MOFA 2022 baseline)",
        "logframe_target":      "e.g., 400 farmers trained by Q4 2025",
        "logframe_achievement": "e.g., 487 farmers trained by June 2025 — 122% of target",
        "verifier":             "e.g., District Agriculture Officer, USAID Feed the Future field monitor, partner org M&E lead",
    },
    "Youth Employment & TVET": {
        "result":               "e.g., Provided vocational training to 250 unemployed youth in IT and entrepreneurship across Accra and Kumasi between January and June 2025, with 68% employed within 3 months",
        "target_group":         "e.g., Unemployed youth aged 18–35, with secondary school qualifications, predominantly women (70%+), residing in urban and peri-urban areas",
        "geographic_scope":     "e.g., Accra (Greater Accra Region) and Kumasi (Ashanti Region) — 4 training centres",
        "evidence_description": "e.g., Signed attendance sheets for all training modules + COTVET certificates issued to 245 graduates + 3-month tracer survey (n=200) + employer sign-off letters confirming placement",
        "logframe_indicator":   "e.g., Indicator 1.2: Number of unemployed youth completing accredited vocational training and securing employment within 3 months",
        "logframe_baseline":    "e.g., 18% youth employment rate in target districts (2023 GSS Labour Force Report baseline)",
        "logframe_target":      "e.g., 250 youth trained and 60% employed within 3 months by Q2 2025",
        "logframe_achievement": "e.g., 250 trained by June 2025; 170 (68%) employed within 3 months — exceeds 60% target",
        "verifier":             "e.g., COTVET assessor, Mastercard Foundation programme officer, employer partner HR records",
    },
    "Climate Resilience": {
        "result":               "e.g., Established 50 community-managed early-warning systems across 10 coastal communities in Volta Region between March and December 2025",
        "target_group":         "e.g., Coastal fishing and farming communities (approx. 25,000 people) vulnerable to flooding, drought, and sea-level rise in Volta Region",
        "geographic_scope":     "e.g., Keta, Anloga, Ada East, and Ada West districts (Volta Region)",
        "evidence_description": "e.g., Installation completion logs + GPS coordinates of 50 stations + community management committee minutes + EPA certification of each station + monthly rainfall data reports",
        "logframe_indicator":   "e.g., Indicator 4.1: Number of community-managed climate early-warning systems established and operational",
        "logframe_baseline":    "e.g., 3 functional early-warning stations in target area (2022 EPA district assessment baseline)",
        "logframe_target":      "e.g., 50 weather stations operational by Dec 2025",
        "logframe_achievement": "e.g., 50 weather stations operational by Dec 2025 — 100% of target",
        "verifier":             "e.g., Environmental Protection Agency inspector, Meteorological Service Department officer, community committee chair",
    },
    "Governance & Accountability": {
        "result":               "e.g., Trained 180 district-level officials on participatory budgeting processes across 6 districts in Ghana between April and August 2025",
        "target_group":         "e.g., Elected district assembly members, district planning officers, and civil society representatives in 6 selected districts",
        "geographic_scope":     "e.g., 6 selected districts in Ashanti, Eastern, and Western regions",
        "evidence_description": "e.g., Training attendance records + pre/post knowledge assessments + signed certificates of completion + post-training participatory budget reports from 4 districts (verified by District Coordinating Directors)",
        "logframe_indicator":   "e.g., Indicator 3.3: Number of district officials demonstrating improved participatory budgeting skills (post-test score ≥70%)",
        "logframe_baseline":    "e.g., 12 officials with participatory budgeting training (2023 pre-programme assessment baseline)",
        "logframe_target":      "e.g., 150 district officials trained by Aug 2025",
        "logframe_achievement": "e.g., 180 officials trained by Aug 2025 — 120% of target",
        "verifier":             "e.g., District Coordinating Director, NCCE observer, civil society partner M&E lead, auditor-general representative",
    },
    "Digital Economy & Technology": {
        "result":               "e.g., Connected 12,000 micro-enterprises to digital payment platforms across 4 regions in Ghana between January and June 2025",
        "target_group":         "e.g., Micro and small enterprise owners (18–45 years), majority women, operating in informal markets in urban and peri-urban areas",
        "geographic_scope":     "e.g., Greater Accra, Ashanti, Northern, and Volta regions (48 districts)",
        "evidence_description": "e.g., Merchant onboarding records from MoMo/GhIPSS platform data + quarterly transaction logs + enumerator-administered post-adoption survey (n=1,200) verified by GIFEC officer",
        "logframe_indicator":   "e.g., Indicator 2.3: Number of MSEs actively using digital financial or business services at least once per month",
        "logframe_baseline":    "e.g., 1,400 MSEs with active digital payment accounts (2022 BoG GIFEC survey baseline)",
        "logframe_target":      "e.g., 12,000 MSEs with active digital payment accounts by Q4 2025",
        "logframe_achievement": "e.g., 12,000 MSEs onboarded by June 2025 — 100% of target",
        "verifier":             "e.g., Ghana Investment Fund for Electronic Communications (GIFEC) officer, GhIPSS data verification, World Bank Ghana Digital Acceleration Project field monitor",
    },
    "Energy & Clean Energy": {
        "result":               "e.g., Installed solar home systems for 3,500 off-grid households across 12 communities in Upper West Region between March and December 2025",
        "target_group":         "e.g., Rural off-grid households earning under GHS 500/month, with no access to the national electricity grid",
        "geographic_scope":     "e.g., Sissala East, Sissala West, Wa East, and Wa West districts (Upper West Region)",
        "evidence_description": "e.g., Installation completion certificates signed by certified NABCB technicians + GPS-tagged photos of each system + 3-month post-installation functionality audit by district energy officer + household income surveys (n=350)",
        "logframe_indicator":   "e.g., Indicator 1.4: Number of off-grid households with access to reliable clean energy (Tier 3+ on SE4ALL Multi-Tier Framework)",
        "logframe_baseline":    "e.g., 0 households with Tier 3+ energy access in target communities (2023 GEC off-grid baseline)",
        "logframe_target":      "e.g., 3,500 households with solar home systems by Dec 2025",
        "logframe_achievement": "e.g., 3,500 systems installed and operational by Dec 2025 — 100% of target",
        "verifier":             "e.g., Ghana Energy Commission inspector, Electricity Company of Ghana off-grid officer, GIZ Renewable Energy officer",
    },
    "Gender & Social Inclusion": {
        "result":               "e.g., Increased women's participation in community decision-making bodies by 42 percentage points across 30 communities in Western North Region between January and June 2025",
        "target_group":         "e.g., Women aged 18–55 from marginalised communities, including women with disabilities and single mothers, in 30 selected communities",
        "geographic_scope":     "e.g., Sefwi Wiawso, Bibiani-Anhwiaso-Bekwai, and Suaman districts (Western North Region)",
        "evidence_description": "e.g., Pre/post community governance meeting attendance registers + baseline and endline participatory assessments (n=300 women) + key informant interviews with 30 community leaders (triangulated with civil society observer data)",
        "logframe_indicator":   "e.g., Indicator 3.1: Percentage of community governance decision-making roles held by women in target communities",
        "logframe_baseline":    "e.g., 11% of governance roles held by women (2022 NCCE community assessment baseline)",
        "logframe_target":      "e.g., 40% of governance roles held by women in 30 communities by Q2 2025",
        "logframe_achievement": "e.g., Women's participation increased by 42 pp across 30 communities by June 2025 — exceeds 40% target",
        "verifier":             "e.g., District Gender Desk Officer, National Commission for Civic Education (NCCE) observer, FCDO programme officer",
    },
    "Nutrition & Food Security": {
        "result":               "e.g., Reduced stunting prevalence by 8 percentage points among children under 5 across 5 districts in Northern Ghana between January 2024 and December 2025",
        "target_group":         "e.g., Children aged 6–59 months and pregnant or lactating women in smallholder farming households in 5 target districts",
        "geographic_scope":     "e.g., Tolon, Kumbungu, Nanton, Mion, and Sagnarigu districts (Northern Region)",
        "evidence_description": "e.g., MUAC measurements from 1,200 children at baseline and endline + health facility nutrition register data + mother/caregiver 24-hour dietary recall surveys + Ghana Health Service district nutrition reports",
        "logframe_indicator":   "e.g., Indicator 4.2: Prevalence of stunting (height-for-age < −2 SD) among children 6–59 months in target districts",
        "logframe_baseline":    "e.g., 32% stunting prevalence in target districts (2022 GHS-NHS district nutrition baseline)",
        "logframe_target":      "e.g., Reduce stunting prevalence from 32% to 24% by Dec 2025 (8 pp reduction)",
        "logframe_achievement": "e.g., Stunting reduced from 32% to 24% across 5 districts by Dec 2025 — meets 8 pp target",
        "verifier":             "e.g., District Nutrition Officer, Ghana Health Service Regional Nutrition Directorate, UNICEF field monitoring officer",
    },
    "Private Sector Development": {
        "result":               "e.g., 184 entrepreneurs across 8 cohorts increased aggregate net profit from a pre-programme baseline of net losses to GHS 313,609 post-evaluation, with 608 jobs created across Cohorts 1–8 between 2019 and 2024",
        "target_group":         "e.g., Early-stage entrepreneurs and SME owners (Cohorts 1–8); majority youth-led businesses (under 35) with at least one product or service at market-entry stage",
        "geographic_scope":     "e.g., Accra, Kumasi, and Takoradi (Greater Accra, Ashanti, and Western regions)",
        "evidence_description": "e.g., Pre-evaluation (Annex 5a) and post-evaluation (Annex 6) datasets collected by independent evaluation firm across 184 participants + financial records from partner financial institutions + employment rosters verified by business mentors",
        "logframe_indicator":   "e.g., KPI 2.3: Increase in profit levels among programme participants (target: 60% of participants reporting positive net profit post-programme)",
        "logframe_baseline":    "e.g., -€1,314,576.78 net loss (pre-evaluation, 2019)",
        "logframe_target":      "e.g., 60% of participants achieving positive net profit by programme end",
        "logframe_achievement": "e.g., 50% of participants achieved positive net profit — GHS 313,609 aggregate (vs. pre-evaluation net loss)",
        "verifier":             "e.g., Independent evaluation firm (lead evaluator), programme MEL officer, Dutch Embassy country representative",
    },
    "Other": {
        "result":               "e.g., [Action verb] [number] [target population] in [location] between [start date] and [end date]",
        "target_group":         "e.g., Women 18–35, rural community health workers in 3 districts",
        "geographic_scope":     "e.g., Ashanti Region — Kumasi, Obuasi, and Bekwai districts",
        "evidence_description": "e.g., Type of records + who collected them + how they were verified + any third-party validation",
        "logframe_indicator":   "e.g., Indicator [X.X]: [Indicator name from approved Technical Proposal or logframe]",
        "logframe_baseline":    "e.g., [Baseline value] ([year] [data source] baseline)",
        "logframe_target":      "e.g., [Number + unit + deadline from logframe]",
        "logframe_achievement": "e.g., [Actual delivered number] by [date] — [%] of original target",
        "verifier":             "e.g., [Implementing partner M&E lead], [government line ministry], [external evaluator]",
    },
}

def _render_readiness_banner(diag_state: str):
    band = _READINESS_BAND.get(diag_state, "Needs Work")
    style = _READINESS_STYLE[band]
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
    st.markdown(
        f"<div class='readiness-banner' style='background:{style['bg']};{_pca}'>"
        f"{style['icon']} {band} &mdash; {style['caption']}"
        f"</div>",
        unsafe_allow_html=True,
    )
    st.caption(_LIMITS_DISCLAIMER)


INTERNAL_REVIEW_OPTIONS = [
    "Choose an option...",
    "Reviewed by MEL Officer",
    "Collected only (no review)",
    "Not reviewed",
    "Other",
]

EXTERNAL_REVIEW_OPTIONS = [
    "Choose an option...",
    "Verified by independent third party",
    "External partner review",
    "No external review",
    "Other",
]

TRACEABILITY_OPTIONS = [
    "Choose an option...",
    "Yes — an auditor could retrieve the original records",
    "Partially — some records would take effort to locate",
    "No / not sure",
]

PROVENANCE_YES_NO_NA_OPTIONS = [
    "Choose an option...",
    "Yes",
    "No",
    "Not applicable",
]

# ---------------------------------------------------------------------------
# Portfolio / Framework Dashboard — CSV/Excel column schema
# ---------------------------------------------------------------------------

# (column_name, required, example value for the downloadable template)
_PORTFOLIO_COLUMNS = [
    # --- Required fields ---
    ("indicator_name",       True,  "Indicator 2.1: Number of rural households with access to safely managed drinking water (WASH)"),
    ("result_statement",     True,  "Constructed 12 boreholes serving 3,400 rural households across 4 districts in Northern Region between January and June 2025"),
    ("target_group",         True,  "Rural households without access to safe drinking water; women and children primarily responsible for water collection"),
    ("timeframe",            True,  "January–June 2025"),
    ("geographic_scope",     True,  "Tamale, Yendi, Savelugu, and Karaga districts (Northern Region, Ghana)"),
    ("evidence_type",        True,  "Attendance sheets / participant registers"),
    ("evidence_description", True,  "Borehole functionality reports from 12 sites + water quality test results from district lab + GPS-tagged photos of completed structures, verified by DWSO"),
    # --- Verification fields ---
    ("evidence_date",        False, "June 2025"),
    ("internal_review",      False, "Reviewed by MEL Officer"),
    ("external_review",      False, "No external review"),
    ("verifier",             False, "District Water and Sanitation Officer, Water Resource Commission inspector"),
    # --- Logframe linkage ---
    ("logframe_indicator",   False, "Indicator 2.1: Number of households with access to safely managed drinking water (SDG 6.1 aligned)"),
    ("logframe_baseline",    False, "1,200 households with access (2019 baseline)"),
    ("logframe_target",      False, "3,000 households with access by Q4 2025"),
    ("logframe_achievement", False, "3,400 households by June 2025 — 113% of target"),
    # --- Advisory flags ---
    ("learning_notes",       False, ""),
    ("limitations_notes",    False, ""),
    # --- Qualitative rigor (TRUE/FALSE) ---
    ("qual_sourcing_documented",           False, "FALSE"),
    ("qual_triangulated",                  False, "FALSE"),
    ("qual_bias_considered",               False, "FALSE"),
    ("qual_beneficiary_voice_represented", False, "FALSE"),
    ("qual_consent_ethics_addressed",      False, "FALSE"),
    # --- Beneficiary voice — affects Confidence bonus (+0–0.5) ---
    ("beneficiary_voice",      False, "No beneficiary voice captured"),
    # --- Governance context — affects Clarity > Governance sub-score ---
    ("additional_context",     False, "MEL Lead owns this result. Informs Q3 budget reallocation for WASH component."),
    # --- Provenance checklist — each affects Confidence > Verification score ---
    # Values: Yes / No / Not applicable  (blank defaults to 'Not applicable' — neutral)
    ("provenance_sampling",    False, "Not applicable"),
    ("provenance_dedup",       False, "Not applicable"),
    ("provenance_tool",        False, "Not applicable"),
    ("provenance_independent", False, "Not applicable"),
    ("provenance_recall",      False, "Not applicable"),
    ("provenance_traceable",   False, "Not applicable"),
]

# Accepted values for internal_review / external_review (unrecognized values
# default to 0, same as evaluate_submission()'s own fallback).
_PORTFOLIO_REVIEW_HINT = (
    "internal_review accepts: 'Reviewed by MEL Officer', 'Collected only (no review)', "
    "'Not reviewed'. external_review accepts: 'Verified by independent third party', "
    "'External partner review', 'No external review'. Leave blank to default to "
    "'Not reviewed' / 'No external review'. "
    "provenance_sampling / provenance_dedup / provenance_tool / provenance_independent / "
    "provenance_recall / provenance_traceable accept: 'Yes', 'No', 'Not applicable' — "
    "leave blank to default to 'Not applicable' (neutral, no penalty). "
    "beneficiary_voice accepts: 'No beneficiary voice captured', "
    "'Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS)', "
    "'Beneficiary representatives consulted (community leaders, beneficiary committees)', "
    "'Anecdotal beneficiary quotes only (uncollected, not systematic)', "
    "'Not applicable to this result type'."
)

SUBMISSION_CHECKLIST = {
    "Project proposal": [
        ("cl_proposal",       "Narrative / technical proposal"),
        ("cl_budget",         "Budget / financial plan"),
        ("cl_logframe",       "Logframe / results framework"),
        ("cl_annexes",        "Annexes (supporting docs)"),
        ("cl_implementation", "Implementation plan"),
        ("cl_org_docs",       "Organisational documents"),
    ],
    "Baseline report": [
        ("cl_methodology",    "Baseline methodology & tools"),
        ("cl_findings",       "Baseline findings report"),
        ("cl_disaggregated",  "Disaggregated / beneficiary data"),
        ("cl_logframe",       "Updated indicators / logframe"),
        ("cl_annexes",        "Annexes with tools or datasets"),
    ],
    "Quarterly progress report": [
        ("cl_narrative",      "Narrative / technical report"),
        ("cl_logframe",       "Updated logframe with achievements"),
        ("cl_disaggregated",  "Beneficiary / disaggregated data"),
        ("cl_annexes",        "Annexes / evidence"),
        ("cl_variance",       "Budget variance report (if requested)"),
    ],
    "Annual progress report": [
        ("cl_narrative",      "Narrative / technical report"),
        ("cl_financial",      "Financial report"),
        ("cl_logframe",       "Updated logframe with achievements"),
        ("cl_disaggregated",  "Beneficiary / disaggregated data"),
        ("cl_annexes",        "Annexes / evidence"),
        ("cl_sustainability", "Sustainability / next-step plan"),
    ],
    "Mid-term review": [
        ("cl_findings",       "Review / field verification report"),
        ("cl_logframe",       "Progress against logframe"),
        ("cl_annexes",        "Annexes / evidence"),
        ("cl_action_plan",    "Action plan (if required)"),
    ],
    "End-line evaluation": [
        ("cl_evaluation",     "Evaluation report"),
        ("cl_methodology",    "Methodology documentation"),
        ("cl_findings",       "Findings with disaggregated outcomes"),
        ("cl_annexes",        "Annexes / evidence"),
        ("cl_learning",       "Learning / dissemination summary"),
    ],
    "Final/closeout report": [
        ("cl_narrative",      "Final narrative report"),
        ("cl_financial",      "Final financial report"),
        ("cl_audit",          "Audit report (if required)"),
        ("cl_logframe",       "Updated logframe / achievements"),
        ("cl_disaggregated",  "Beneficiary / disaggregated data"),
        ("cl_sustainability", "Sustainability / exit plan"),
        ("cl_annexes",        "Annexes"),
    ],
    "Financial report": [
        ("cl_expense_summary","Expense summary"),
        ("cl_variance",       "Budget vs actual / variance explanation"),
        ("cl_schedules",      "Supporting financial schedules"),
        ("cl_bank_recon",     "Bank / ledger reconciliation"),
    ],
    "MEL plan": [
        ("cl_indicators",     "Indicators & targets"),
        ("cl_data_sources",   "Data sources & collection methods"),
        ("cl_reporting_cal",  "Reporting calendar"),
        ("cl_roles",          "Roles & responsibilities"),
        ("cl_data_quality",   "Data quality procedures"),
        ("cl_methodology",    "Data collection tools / methodology"),
    ],
    "Others (Special/Ad-hoc reports)": [
        ("cl_narrative",      "Narrative / technical report"),
        ("cl_annexes",        "Donor-specific annexes"),
        ("cl_case_studies",   "Case studies / special studies"),
        ("cl_safeguarding",   "Safeguarding / compliance reports"),
    ],
}

# "Strengthen this evidence" checklist items, keyed by evidence type (v3.6 IRC auto-fill)
EVIDENCE_STRENGTHEN_CHECKLIST = {
    "Attendance sheets / participant registers": [
        ("signatures_verified", "Signatures verified against ID list"),
        ("date_stamped",        "Sheets dated and stamped"),
        ("cross_ref",           "Cross-referenced with another source (e.g., facilitator notes)"),
    ],
    "Raw datasets or survey exports": [
        ("sample_doc",   "Sampling method documented"),
        ("clean_data",   "Dataset cleaned and de-duplicated"),
        ("version_ctrl", "Original raw export retained for audit"),
    ],
    "Partner verification letters": [
        ("letterhead",       "Letter on official partner letterhead"),
        ("authority_signed", "Signed by authorized partner representative"),
        ("recent_letter",    "Letter dated within 6 months of reporting period"),
    ],
    "Photos with metadata": [
        ("gps_meta",       "Photos contain GPS metadata"),
        ("timestamp_photo","Timestamps visible/verifiable"),
        ("consent_photo",  "Beneficiary consent obtained for photos"),
    ],
    "Tracer survey results": [
        ("followup_tracer", "Follow-up conducted at appropriate interval (3+ months)"),
        ("response_rate",   "Response rate documented (target: 60%+)"),
        ("bias_ack",        "Sampling bias / non-response acknowledged"),
    ],
    "Financial records": [
        ("receipts_dated", "Receipts/transactions dated"),
        ("reconciled_ev",  "Reconciled with bank/MoMo statements"),
        ("audit_trail_ev", "Audit trail intact (request → approval → payment)"),
    ],
    "Third-party audits": [
        ("independent_ev",     "Auditor independent from implementer"),
        ("signed_audit",       "Audit report signed and dated"),
        ("recommendations_ev", "Audit recommendations addressed/disclosed"),
    ],
}

# ---------------------------------------------------------------------------
# CSS — injected once at app load
# ---------------------------------------------------------------------------

CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;700&family=JetBrains+Mono:wght@400&display=swap');

:root {
  --brand-green: #1B5E20;
  --gold:        #8A6500;
  --body-text:   #212121;
  --muted:       #424242;  /* raised from #616161 — WCAG 1.4.3 AA: 5.3:1 on white */
  --bg-card:     #F5F5F5;
  --border:      rgba(27,90,32,0.15);
}


html, body, [class*="css"] {
  font-family: 'Inter', sans-serif;
  color: var(--body-text);
}
/* Prevent pull-to-refresh on mobile (Chrome/Android).
   Stops accidental page refresh when user scrolls past the top. */
html { overscroll-behavior-y: contain; }
body { overscroll-behavior-y: contain; }

h1, h2, h3, h4 {
  font-family: 'Inter', sans-serif;
  font-weight: 700;
  color: var(--brand-green);
}

/* Primary button -> Impact Green */
.stButton > button[kind="primary"],
.stFormSubmitButton > button[kind="primary"] {
  background-color: #1B5E20 !important;
  border-color: #1B5E20 !important;
  color: white !important;
  font-family: 'Inter', sans-serif;
  font-weight: 700;
  border-radius: 8px;
}

/* Secondary button -> Trust Gold outline */
.stButton > button[kind="secondary"],
.stFormSubmitButton > button[kind="secondary"] {
  border-color: #8A6500 !important;
  color: #8A6500 !important;
  font-family: 'Inter', sans-serif;
  background: transparent !important;
}

/* Card container */
.result-card {
  background: var(--bg-card);
  border: 1px solid var(--border);
  border-radius: 8px;
  padding: 24px 28px;
  margin-bottom: 20px;
}

/* Axis score badge */
.axis-badge {
  padding: 10px 14px;
  border-radius: 8px;
  text-align: center;
  font-family: 'Inter', sans-serif;
  font-weight: 700;
  font-size: 0.85rem;
  margin-bottom: 6px;
}

/* Verdict banner */
.verdict-banner {
  background: #1B5E20;
  color: white;
  border-radius: 8px;
  padding: 14px 20px;
  font-family: 'Inter', sans-serif;
  font-weight: 700;
  text-align: center;
  margin: 16px 0;
  font-size: 1rem;
}
.verdict-banner.misleading { background: #E65100; }
.verdict-banner.weak-conf  { background: #F57F17; }
.verdict-banner.high-risk  { background: #B71C1C; }


/* Hero section */
.hero-block {
  padding: 12px 0 20px 0;
  border-bottom: 1px solid #8A6500;
  margin-bottom: 20px;
}
.hero-block h1 {
  font-size: 1.85rem;
  line-height: 1.25;
  margin-bottom: 8px;
}
.hero-tagline {
  font-style: normal;
  font-weight: 600;
  color: #8A6500;
  font-size: 1.05rem;
  margin: 4px 0 14px 0;
}

/* CTA call button */
.cta-call-btn a {
  display: inline-block;
  background: #8A6500;
  color: white !important;
  font-family: 'Inter', sans-serif;
  font-weight: 700;
  padding: 10px 20px;
  border-radius: 8px;
  text-decoration: none;
  font-size: 0.95rem;
}

/* Trust Gold tagline footer */
.trust-tagline {
  font-style: italic;
  color: #8A6500;
  font-size: 0.82rem;
  text-align: center;
  padding: 12px 0 4px 0;
  border-top: 1px solid rgba(138,101,0,0.2);
  margin-top: 24px;
}

/* GTM conversion hook card */
.gtm-card {
  border: 1px solid #8A6500;
  border-radius: 8px;
  padding: 20px 24px;
  margin: 24px 0;
  background: #FFFEF7;
}
.gtm-card p { color: #212121; font-size: 0.95rem; margin: 0 0 4px 0; }
.gtm-card .gtm-sub { color: #616161; font-size: 0.85rem; margin-bottom: 14px; }

/* GTM buttons */
.gtm-btn-gold a {
  display: inline-block;
  border: 2px solid #8A6500;
  color: #8A6500 !important;
  padding: 8px 18px;
  border-radius: 8px;
  text-decoration: none;
  font-weight: 700;
  font-size: 0.9rem;
  font-family: 'Inter', sans-serif;
  margin-right: 10px;
}

/* Gold info box */
.gold-info-box {
  background: #FFFEF7;
  border-left: 4px solid #8A6500;
  padding: 10px 16px;
  border-radius: 8px;
  font-size: 0.9rem;
  color: #212121;
  margin: 12px 0;
}
.gold-info-box a { color: #1B5E20; }

/* Score metric font */
[data-testid="stMetricValue"] { font-family: 'JetBrains Mono', monospace; }

/* Diagnostic state badge */
.diagnostic-badge {
  padding: 10px 16px;
  border-radius: 8px;
  font-weight: 700;
  font-size: 1rem;
  margin-bottom: 12px;
  display: inline-block;
  letter-spacing: 0.02em;
}

/* Top-line "is this good enough to submit?" banner */
.readiness-banner {
  padding: 14px 20px;
  border-radius: 8px;
  font-weight: 700;
  font-size: 1.1rem;
  margin-bottom: 12px;
  color: #FFFFFF;
  text-align: center;
}

/* Mobile-first improvements */
@media (max-width: 768px) {
  .stButton button {
    min-height: 48px !important;
    font-size: 16px !important;
    padding: 12px 16px !important;
  }
  .stTextInput input, .stTextArea textarea,
  .stSelectbox div[role="combobox"] {
    min-height: 44px !important;
    font-size: 16px !important;
  }
  .main .block-container {
    padding-left: 1rem !important;
    padding-right: 1rem !important;
  }
  .stTabs [data-baseweb="tab-list"] {
    overflow-x: auto !important;
    flex-wrap: nowrap !important;
  }
  .stCheckbox label {
    min-height: 36px !important;
    display: flex !important;
    align-items: center !important;
  }
  /* Pitch strip: reduce padding on mobile so it stays compact */
  .md-pitch { padding: 6px 10px !important; }
  .md-pitch-stages { max-width: 100% !important; }
  /* Sidebar: reduce push-down to match smaller strip height */
  [data-testid="stSidebarUserContent"] { padding-top: 72px !important; }
}
/* Very narrow phones: hide pitch strip text labels, show dots only */
@media (max-width: 420px) {
  .md-pitch .lbl { display: none !important; }
  .md-pitch { padding: 8px 8px !important; }
  .md-pstage .dot { width: 24px !important; height: 24px !important; font-size: 10px !important; }
  /* Accessibility D6: prevent dots overlapping at high display scale */
  .md-pitch-stages { gap: 2px; }
  .md-pstage { min-width: 22px; }
}
/* Active tab: bold + underline */
.stTabs [data-baseweb="tab"][aria-selected="true"] button {
    font-weight: 700;
    text-decoration: underline;
}
/* Mobile: Vega/Altair chart tooltips — keep within viewport, above bars */
.vega-tooltip {
    z-index: 10000 !important;
    max-width: 90vw !important;
    word-wrap: break-word !important;
    font-size: 0.875rem !important;
}
@media (max-width: 768px) {
  .vega-tooltip { position: fixed !important; bottom: 8px !important; top: auto !important; left: 5% !important; right: 5% !important; }
}
/* Form labels: consistent weight */
.stTextInput label, .stTextArea label, .stSelectbox label,
.stDateInput label, .stFileUploader label, .stNumberInput label {
    font-size: 0.875rem !important;
    font-weight: 600 !important;
}
/* Expander headers: lighter weight than section headings */
details summary, .stExpander summary {
    font-weight: 500 !important;
    font-size: 0.9rem !important;
}
/* Breathing room between expander groups */
.stExpander {
    margin-bottom: 1.25rem !important;
}
/* Captions: clearly subordinate to body text */
.stCaptionContainer p, [data-testid="stCaptionContainer"] p {
    font-size: 0.78rem !important;
    color: var(--muted) !important;
}
/* Accessibility — WCAG 1.4.4 Resize Text: minimum readable floor for small text */
.stCaptionContainer p, [data-testid="stCaptionContainer"] p,
.stCaption, small, .trust-tagline { font-size: max(0.875rem, 14px) !important; }
/* Accessibility — WCAG 2.4.7 Focus Visible: keyboard focus ring using brand green */
*:focus-visible {
    outline: 3px solid #1B5E20 !important;
    outline-offset: 3px !important;
    border-radius: 4px;
}
/* Screen-reader-only utility class for visually-hidden text */
.sr-only {
    position: absolute; width: 1px; height: 1px; overflow: hidden;
    clip: rect(0,0,0,0); white-space: nowrap; border: 0;
}
/* Accessibility — WCAG 2.5.5 Target Size: HTML anchors styled as buttons */
.cta-call-btn a, .gtm-btn-gold a {
    min-height: 44px !important;
    display: flex !important;
    align-items: center !important;
    justify-content: center !important;
}
/* Smooth navigation transitions — reduces ghost/shadow flash between tabs and screens */
[data-testid="stMainBlockContainer"] > div {
    animation: fadeInContent 0.12s ease-in;
}
@keyframes fadeInContent {
    from { opacity: 0.5; }
    to   { opacity: 1; }
}
/* Push sidebar user content below the fixed pitch strip */
[data-testid="stSidebarUserContent"] {
    padding-top: 88px !important;
}
/* Keep sidebar collapse/expand toggle above pitch strip */
[data-testid="stSidebarCollapsedControl"],
[data-testid="stSidebarCollapseButton"] {
    z-index: 9999999 !important;
}
</style>
"""

# ---------------------------------------------------------------------------
# Session state helpers
# ---------------------------------------------------------------------------

def _init_session_state():
    defaults = {
        "screen":              0,
        "error_message":       None,
        "evaluations":         None,
        "submissions_snapshot": None,
        "active_slots":        1,
        "has_seen_tutorial":   False,
        "tutorial_step":       0,
        "sector":              SECTOR_OPTIONS[0],
        "confirm_reset":       False,
        "gov_dpp_uploaded":    False,
        "lite_mode":           False,
        # --- UX (v3.2) ---
        "current_tab":         0,
        "remembered_donor":    "",
        "remembered_sector":   "",
        # --- END UX (v3.2) ---
        # --- v3.3 additions ---
        "donor_other":         "",
        "donor_framework":     "Generic",
        "org_type":            "International NGO (INGO)",
        # --- auth / payment ---
        "user_email":          "",
        "is_paid":             False,
        "consent_examples":    False,
        "_pay_once_url":       "",
        "_pay_monthly_url":    "",
        # --- end auth ---
        "report_level":        "(Not specified)",
        "report_prepared_by":  "",
        "report_status":       _REPORT_STATUS_OPTIONS[0],
        "report_notes":        "",
        "_tab1_auto_advanced": False,
        "_tab2_auto_advanced": False,
        # --- chat help (council XV) ---
        "chat_messages":       [],
        # --- END v3.3 ---
    }
    for key, default in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = default


_BASE_FORM_KEYS = [
    "result_statement", "target_group", "timeframe", "geographic_scope",
    "evidence_description", "evidence_type", "evidence_type_other",
    "internal_review", "internal_review_other",
    "external_review", "external_review_other",
    "verifier", "sector", "sector_other", "beneficiary_voice",
    "logframe_indicator", "logframe_baseline", "logframe_target", "logframe_achievement",
    "logframe_data_forthcoming",
    # evidence sub-prompt checkboxes (informational, v3.3)
    "signatures_verified", "date_stamped", "cross_ref",
    "sample_doc", "clean_data", "version_ctrl",
    "letterhead", "authority_signed", "recent_letter",
    "gps_meta", "timestamp_photo", "consent_photo",
    "followup_tracer", "response_rate", "bias_ack",
    "receipts_dated", "reconciled_ev", "audit_trail_ev",
    "independent_ev", "signed_audit", "recommendations_ev",
    # governance & compliance (v3.2)
    "gov_consent_status", "gov_anonymization_status", "gov_compliance_law_status",
    "gov_safeguarding_status", "gov_child_safeguarding_status", "gov_secure_handling_status",
    # v3.3
    "donor_other", "report_level",
    # v3.3 checklist keys
    "cl_proposal", "cl_implementation", "cl_org_docs", "cl_methodology",
    "cl_findings", "cl_disaggregated", "cl_variance", "cl_action_plan",
    "cl_evaluation", "cl_learning", "cl_expense_summary", "cl_schedules",
    "cl_bank_recon", "cl_indicators", "cl_data_sources", "cl_reporting_cal",
    "cl_roles", "cl_data_quality", "cl_case_studies", "cl_safeguarding",
    "cl_beneficiary",
    # v3.4 advisory checklist (optional, score-neutral)
    "attribution_contribution", "disaggregation_status",
    # v3.5 learning / limitations / decision-ownership notes
    "learning_notes", "limitations_notes", "additional_context",
    # v3.6 review-handoff layer (per-submission, no auth)
    "review_status", "reviewer_name", "reviewer_role", "reviewer_date",
    "reviewer_decision", "reviewer_notes",
]

# v3.6 — Review-handoff layer: per-submission status and reviewer decision.
# No user accounts / authentication — reviewer identity is free text within
# the current session, consistent with the rest of the app's single-session model.
_SUBMISSION_STATUS_OPTIONS = ["Draft", "In review", "Approved", "Returned"]
_REVIEW_DECISION_OPTIONS = ["", "Approve", "Return for revision"]
_SUBMISSION_STATUS_COLORS = {
    "Draft":      ("#F5F5F5", "#616161"),
    "In review":  ("#FFF9C4", "#F57F17"),
    "Approved":   ("#C8E6C9", "#1B5E20"),
    "Returned":   ("#FFE0B2", "#E65100"),
}

_BV_OPTIONS = [
    "Choose an option...",
    "No beneficiary voice captured",
    "Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS)",
    "Beneficiary representatives consulted (community leaders, beneficiary committees)",
    "Anecdotal beneficiary quotes only (uncollected, not systematic)",
    "Not applicable to this result type",
]


def _slot_suffix(slot: int) -> str:
    return "" if slot == 1 else f"_{slot}"


def _reset_all_slots():
    active = st.session_state.get("active_slots", 1)
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        for k in _BASE_FORM_KEYS:
            st.session_state.pop(f"{k}{s}", None)
        for k in ["evidence_date", "uploaded_files", "draft_uploaded_filenames"]:
            st.session_state.pop(f"{k}{s}", None)
    for k in ["active_slots", "evaluations", "submissions_snapshot",
              "evaluation", "submission_snapshot", "error_message", "active_slots_run",
              "_tab1_auto_advanced", "_tab2_auto_advanced",
              "_results_email_sent", "current_tab", "chat_messages",
              "smr_chat_messages"]:
        st.session_state.pop(k, None)


def _init_from_query_params() -> None:
    """Restore screen/tab from URL on first load. Skips Paystack callback URLs."""
    if st.session_state.get("_nav_initialized"):
        return
    st.session_state["_nav_initialized"] = True
    _p = st.query_params
    if any(k in _p for k in ("paystack_ref", "reference", "trxref", "login_token")):
        return
    if "screen" in _p:
        try:
            _s = int(_p["screen"])
            if 0 <= _s <= 3:
                st.session_state["screen"] = _s
        except (ValueError, TypeError):
            pass
    if "tab" in _p:
        try:
            _t = int(_p["tab"])
            if 0 <= _t <= 3:
                st.session_state["current_tab"] = _t
        except (ValueError, TypeError):
            pass
    if _p.get("demo") == "1":
        for _k, _v in _DEMO_SUBMISSION.items():
            st.session_state[_k] = _v
        for _k, _v in _DEMO_SELECT_FIELDS.items():
            st.session_state[_k] = _v
        if not st.session_state.get("_demo_viewed_logged"):
            metrics.log_event("demo_viewed", _metrics_session_id())
            st.session_state["_demo_viewed_logged"] = True
    # Track referral source for conversion analytics
    if "ref" in _p:
        st.session_state["_referral_source"] = _p["ref"]
    if _p.get("billing") == "1":
        st.session_state["_show_billing"] = True
    if _p.get("my_audits") == "1":
        st.session_state["_show_my_audits"] = True


def _restore_session_from_query_param() -> None:
    """Silently re-authenticate a returning visitor from a bookmarked
    ?session=... URL. Deliberately lighter-weight than _complete_email_login,
    which does one-time things (welcome email, pending-payment resolution,
    draft restore, an unconditional rerun) that shouldn't fire on every normal
    page load just because a valid session token was found."""
    if st.session_state.get("user_email"):
        return
    _tok = st.query_params.get("session", "")
    if not _tok:
        return
    _email = verify_session_token(_tok)
    if _email:
        st.session_state["user_email"] = _email
        _u = get_user(_email)
        if _u and is_still_paid(_u):
            st.session_state["is_paid"] = True


def _go_to_screen(screen: int, reset: bool = False):
    if reset:
        _reset_all_slots()
    st.session_state["screen"] = screen
    st.session_state["_scroll_to_content"] = True
    st.query_params["screen"] = str(screen)
    if screen != 1:
        st.query_params.pop("tab", None)
        st.query_params.pop("demo", None)
    st.rerun()


# ---------------------------------------------------------------------------
# Shared UI helpers
# ---------------------------------------------------------------------------

def _metrics_session_id() -> str:
    """A stable identifier for this browser session, for metrics hashing only.
    Prefers the logged-in email (already hashed one-way by metrics.log_event);
    falls back to a random id kept in session_state for anonymous visitors."""
    email = st.session_state.get("user_email", "")
    if email:
        return email
    sid = st.session_state.get("_anon_session_id", "")
    if not sid:
        import uuid
        sid = str(uuid.uuid4())
        st.session_state["_anon_session_id"] = sid
    return sid


def _render_tagline_footer():
    st.markdown(
        '<div class="trust-tagline">ImpactProof · Upload your report. Get a determination for every result. Submit with confidence. · Built in Accra 🇬🇭 · Deterministic scoring, AI-assisted, fabrication-proof</div>',
        unsafe_allow_html=True,
    )
    st.caption("We log anonymous usage counts only — never your results or documents.")


def _format_date(d) -> str:
    """Convert date/datetime to 'Month YYYY' string for evaluator."""
    if d is None:
        return ""
    if hasattr(d, "strftime"):
        return d.strftime("%B %Y")
    return str(d)


def _ss_str(key: str, default: str = "") -> str:
    """Read a session_state value as a string, tolerating non-string values
    (e.g. a number or list accidentally written by IRC extraction)."""
    val = st.session_state.get(key, default)
    if isinstance(val, str):
        return val
    return default if val is None else str(val)


_DRAFT_PATH = os.path.join("inputs", "draft.json")


def _save_draft():
    active = st.session_state.get("active_slots", 1)
    draft = {"active_slots": active}
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        for k in _BASE_FORM_KEYS:
            draft[f"{k}{s}"] = st.session_state.get(f"{k}{s}", "")
        ed = st.session_state.get(f"evidence_date{s}")
        draft[f"evidence_date{s}"] = ed.isoformat() if hasattr(ed, "isoformat") else ""
        raw_files = st.session_state.get(f"uploaded_files_widget{s}") or []
        draft[f"uploaded_filenames{s}"] = [f.name for f in raw_files if hasattr(f, "name")]
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        for dk in ("reporting_start", "reporting_end"):
            d = st.session_state.get(f"{dk}{s}")
            draft[f"{dk}{s}"] = d.isoformat() if hasattr(d, "isoformat") else ""
    # --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
    for _gs in range(1, active + 1):
        _s = _slot_suffix(_gs)
        for _gk in ["gov_consent_status", "gov_anonymization_status", "gov_compliance_law_status"]:
            draft[f"{_gk}{_s}"] = st.session_state.get(f"{_gk}{_s}", "")
    draft["gov_dpp_uploaded"] = st.session_state.get("gov_dpp_uploaded", False)
    # --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
    for gk in ("submission_type", "cl_narrative", "cl_financial", "cl_audit",
               "cl_logframe", "cl_annexes", "cl_beneficiary", "cl_sustainability",
               "cl_budget", "donor_selected", "donor_other"):
        draft[gk] = st.session_state.get(gk, "")
    _draft_json = json.dumps(draft, indent=2, ensure_ascii=False)
    st.session_state["_draft_bytes"] = _draft_json.encode("utf-8")
    st.session_state["_last_saved_time"] = datetime.now().strftime("%H:%M")


def _load_draft():
    if not os.path.exists(_DRAFT_PATH):
        return
    try:
        with open(_DRAFT_PATH, encoding="utf-8") as f:
            draft = json.load(f)
    except Exception:
        return
    active = int(draft.get("active_slots", 1))
    st.session_state["active_slots"] = active
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        for k in _BASE_FORM_KEYS:
            key = f"{k}{s}"
            if key in draft:
                st.session_state[key] = draft[key]
        raw_date = draft.get(f"evidence_date{s}", "")
        if raw_date:
            try:
                st.session_state[f"evidence_date{s}"] = date.fromisoformat(raw_date)
            except (ValueError, TypeError):
                pass
        st.session_state[f"draft_uploaded_filenames{s}"] = draft.get(f"uploaded_filenames{s}", [])
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        for dk in ("reporting_start", "reporting_end"):
            raw = draft.get(f"{dk}{s}", "")
            if raw:
                try:
                    st.session_state[f"{dk}{s}"] = date.fromisoformat(raw)
                except (ValueError, TypeError):
                    pass
    for gk in ("submission_type", "cl_narrative", "cl_financial", "cl_audit",
               "cl_logframe", "cl_annexes", "cl_beneficiary", "cl_sustainability",
               "cl_budget", "donor_selected", "donor_other"):
        if gk in draft:
            st.session_state[gk] = draft[gk]


def _clear_draft():
    st.session_state.pop("_draft_bytes", None)
    st.session_state.pop("_last_saved_time", None)
    if os.path.exists(_DRAFT_PATH):  # clean up any legacy file
        os.remove(_DRAFT_PATH)


# ---------------------------------------------------------------------------
# Tutorial renderer
# ---------------------------------------------------------------------------

_TUTORIAL_COPY = {
    1: {
        "title": "📝 Each field below contributes to your score.",
        "body": (
            "Watch the **Submission Summary** panel (sidebar) update as you fill in each field.\n"
            "The **Review & Submit** tab shows your live Confidence and Clarity scores with fix buttons — check it before submitting."
        ),
    },
    2: {
        "title": "🎯 Your result is now scored on two axes:",
        "body": (
            "• **Confidence:** How much we should trust the evidence\n"
            "• **Clarity:** How clearly the result is defined\n\n"
            "Both must meet your organisation's evidence standard to reach this tool's top band.\n\n"
            "The **What to Fix** section tells you exactly how to improve."
        ),
    },
    3: {
        "title": "📄 Download your report and submit it to your donor.",
        "body": (
            "Use the **Download PDF Report** button above "
            "to get a shareable copy of your results.\n\n"
            "Used your free checks? Upgrade options (pay-per-check or monthly) appear "
            "wherever you've reached the free-check limit."
        ),
    },
}

_TUTORIAL_LAST_STEP = max(_TUTORIAL_COPY)


def _render_tutorial(step: int):
    if st.session_state.get("has_seen_tutorial") or st.session_state.get("tutorial_step", 0) > step:
        return
    copy = _TUTORIAL_COPY.get(step)
    if not copy:
        return
    st.caption(f"💡 **{copy['title']}** — {copy['body']}")
    if st.button("Got it →", key=f"tutorial_got_{step}"):
        if step == _TUTORIAL_LAST_STEP:
            st.session_state["has_seen_tutorial"] = True
        st.session_state["tutorial_step"] = step + 1
        st.rerun()


# ---------------------------------------------------------------------------
# Live score preview
# ---------------------------------------------------------------------------

# --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
def _minors_possibly_involved(slot: int) -> bool:
    """True if the result statement or target group mentions minors — used only
    to decide whether to prompt the child-safeguarding question, never to
    auto-answer it."""
    s = _slot_suffix(slot)
    text = " ".join([
        st.session_state.get(f"result_statement{s}", "") or "",
        st.session_state.get(f"target_group{s}", "") or "",
    ]).lower()
    return any(kw in text for kw in CHILD_SAFEGUARDING_KEYWORDS)


def _compute_governance_score(slot: int):
    """Returns (governance_score 0-24, pii_selected bool, gaps list,
    safeguarding_triggered bool, minors_possibly_involved bool)."""
    s = _slot_suffix(slot)
    ev_type = st.session_state.get(f"evidence_type{s}", "")
    pii_selected = ev_type in PII_EVIDENCE_TYPES
    safeguarding_triggered = ev_type in SAFEGUARDING_EVIDENCE_TYPES
    minors_involved = _minors_possibly_involved(slot)

    consent    = st.session_state.get(f"gov_consent_status{s}", "")
    anon       = st.session_state.get(f"gov_anonymization_status{s}", "")
    law        = st.session_state.get(f"gov_compliance_law_status{s}", "")
    safeguard  = st.session_state.get(f"gov_safeguarding_status{s}", "")
    child_safe = st.session_state.get(f"gov_child_safeguarding_status{s}", "")
    secure     = st.session_state.get(f"gov_secure_handling_status{s}", "")
    dpp        = st.session_state.get("gov_dpp_uploaded", False)

    score = 0
    gaps  = []

    if consent in ("Choose an option...", "Select consent status..."):
        pass  # 0 pts, no gap — user has not answered yet
    elif consent == "Yes — written consent forms on file":
        score += 5
    elif consent == "Yes — verbal consent documented":
        score += 3
    elif consent.startswith("Partial"):
        score += 1
    elif consent.startswith("Not applicable"):
        score += 3
    else:
        gaps.append("Consent not obtained")

    if anon in ("Choose an option...", "Select anonymization status..."):
        pass  # 0 pts, no gap
    elif anon == "Yes — fully anonymized":
        score += 4
    elif anon == "Partially anonymized":
        score += 2
    elif anon == "Not applicable":
        score += 3
    else:
        gaps.append("Evidence not anonymized")

    if law in ("Choose an option...", "Select compliance status..."):
        pass  # 0 pts, no gap
    elif law.startswith("Yes"):
        score += 3
    elif law.startswith("Unsure"):
        score += 1
    elif law.startswith("No"):
        gaps.append("Data law compliance not confirmed")

    if safeguard in ("Choose an option...", "Select safeguarding status..."):
        pass  # 0 pts, no gap
    elif safeguard.startswith("Yes"):
        score += 3
    elif safeguard.startswith("Partial"):
        score += 1
    elif safeguard.startswith("Not applicable"):
        score += 3
    elif safeguard.startswith("No"):
        gaps.append("Do-no-harm review not completed (Core Humanitarian Standard, Commitment 4)")

    if child_safe in ("Choose an option...", "Select child safeguarding status..."):
        pass  # 0 pts, no gap
    elif child_safe.startswith("Yes"):
        score += 3
    elif child_safe.startswith("Partial"):
        score += 1
    elif child_safe.startswith("Not applicable"):
        score += 3
    elif child_safe.startswith("No"):
        gaps.append("Child safeguarding review not completed (Core Humanitarian Standard, Commitment 4 — Keeping Children Safe)")

    if secure in ("Choose an option...", "Select secure handling status..."):
        pass  # 0 pts, no gap
    elif secure.startswith("Yes"):
        score += 3
    elif secure.startswith("Partial"):
        score += 1
    elif secure.startswith("Not applicable"):
        score += 3
    elif secure.startswith("No"):
        gaps.append("Secure handling of identifiable testimony not confirmed (Bond Evidence Principles 2024 — Transparency/Data Protection)")

    if dpp:
        score += 5

    return min(24, score), pii_selected, gaps, safeguarding_triggered, minors_involved
# --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---


def _plan_code(secret_name: str) -> str:
    """Read a PAYSTACK_PLAN_* code from secrets/env (see
    scripts/setup_paystack_plans.py). Empty string if not configured yet --
    callers fall back to a plain one-off transaction in that case, so a
    Subscribe button never breaks just because a plan code hasn't been
    pasted into secrets."""
    return (
        st.secrets.get(secret_name, "")
        if hasattr(st, "secrets") else
        os.environ.get(secret_name, "")
    )


def _render_paywall(irc_context: bool = False, custom_message: str | None = None,
                     prompt_context: str = "limit_hit"):
    """Show upgrade/payment options. irc_context=True suppresses the free-checks header.
    custom_message overrides the default re-scoring copy for call sites gating a
    different feature (e.g. Audit My Report, CSV Portfolio) so the paywall accurately
    describes what the user actually got blocked on. prompt_context identifies which
    high-intent moment triggered this paywall, for upgrade_prompt_shown/_clicked metrics."""
    email = st.session_state.get("user_email", "")
    metrics.log_event("upgrade_prompt_shown", _metrics_session_id(), context=prompt_context)
    if custom_message is not None:
        st.markdown(custom_message)
    elif not irc_context:
        st.markdown("### You've used your 3 free checks.")
        st.markdown(
            "Improved your evidence? Upgrade to re-score and see the impact:\n\n"
            "- **Re-score after every fix** — see exactly how much each change moves your score\n"
            "- **⚡ Instant Report Check** — upload your draft report and auto-fill all fields in seconds\n"
            "- **Downloadable PDF report** — shareable with your supervisor or donor\n\n"
            f"*GHS {PRICE_PER_CHECK_GHS/100:.0f} per check · or GHS {PRICE_MONTHLY_GHS/100:.0f}/month for unlimited*\n\n"
            "💡 *The ROI is immediate: GHS 50/month vs. GHS 12,000–17,000 in rework costs "
            "from a donor-queried report.*"
        )
    _c1, _c2, _c3 = st.columns(3)
    with _c1:
        st.markdown(f"**Pay-per-use:** GHS {PRICE_PER_CHECK_GHS/100:.0f}")
        if st.session_state.get("_pay_once_url"):
            st.link_button("Complete Payment →", st.session_state["_pay_once_url"],
                           use_container_width=True, type="primary")
        elif st.button("Pay for 1 Check", key="pay_once", use_container_width=True):
            with st.spinner("Preparing payment link…"):
                _url = initialize_payment(email, PRICE_PER_CHECK_GHS, "per_use")
            if _url:
                st.session_state["_pay_once_url"] = _url
                metrics.log_event("payment_initiated", _metrics_session_id())
                metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context=prompt_context)
                st.rerun()
            else:
                _detail = last_payment_error()
                st.error(f"Payment service unavailable. Try again shortly.{' (' + _detail + ')' if _detail else ''}")
    with _c2:
        st.markdown(f"**Professional:** GHS {PRICE_MONTHLY_GHS/100:.0f}/month")
        st.caption("Unlimited · Readiness Card PDF")
        if st.session_state.get("_pay_monthly_url"):
            st.link_button("Complete Payment →", st.session_state["_pay_monthly_url"],
                           use_container_width=True, type="primary")
        elif st.button("Subscribe Professional", key="pay_monthly", use_container_width=True, type="primary"):
            with st.spinner("Preparing payment link…"):
                _plan_monthly = _plan_code("PAYSTACK_PLAN_PROFESSIONAL_MONTHLY")
                _url = (
                    initialize_subscription_payment(email, PRICE_MONTHLY_GHS, _plan_monthly, "monthly")
                    if _plan_monthly else
                    initialize_payment(email, PRICE_MONTHLY_GHS, "monthly")
                )
            if _url:
                st.session_state["_pay_monthly_url"] = _url
                metrics.log_event("payment_initiated", _metrics_session_id())
                metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context=prompt_context)
                st.rerun()
            else:
                _detail = last_payment_error()
                st.error(f"Payment service unavailable. Try again shortly.{' (' + _detail + ')' if _detail else ''}")
        if st.session_state.get("_pay_annual_url"):
            st.link_button(f"Complete annual payment →", st.session_state["_pay_annual_url"],
                           use_container_width=True)
        elif st.button(f"Or pay GHS {PRICE_ANNUAL_GHS/100:.0f}/year (2 months free)",
                       key="pay_annual", use_container_width=True):
            with st.spinner("Preparing payment link…"):
                _plan_annual = _plan_code("PAYSTACK_PLAN_PROFESSIONAL_ANNUAL")
                _url = (
                    initialize_subscription_payment(email, PRICE_ANNUAL_GHS, _plan_annual, "annual")
                    if _plan_annual else
                    initialize_payment(email, PRICE_ANNUAL_GHS, "annual")
                )
            if _url:
                st.session_state["_pay_annual_url"] = _url
                metrics.log_event("payment_initiated", _metrics_session_id())
                metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context=prompt_context)
                st.rerun()
            else:
                _detail = last_payment_error()
                st.error(f"Payment service unavailable. Try again shortly.{' (' + _detail + ')' if _detail else ''}")
    with _c3:
        st.markdown(f"**Agency:** GHS {PRICE_AGENCY_GHS/100:.0f}/month")
        st.caption("Portfolio analysis · 5 seats")
        if st.session_state.get("_pay_agency_url"):
            st.link_button("Complete Payment →", st.session_state["_pay_agency_url"],
                           use_container_width=True)
        elif st.button("Subscribe Agency", key="pay_agency", use_container_width=True):
            with st.spinner("Preparing payment link…"):
                _plan_agency = _plan_code("PAYSTACK_PLAN_AGENCY_MONTHLY")
                _url = (
                    initialize_subscription_payment(email, PRICE_AGENCY_GHS, _plan_agency, "agency")
                    if _plan_agency else
                    initialize_payment(email, PRICE_AGENCY_GHS, "agency")
                )
            if _url:
                st.session_state["_pay_agency_url"] = _url
                metrics.log_event("payment_initiated", _metrics_session_id())
                metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context=prompt_context)
                st.rerun()
            else:
                _detail = last_payment_error()
                st.error(f"Payment service unavailable. Try again shortly.{' (' + _detail + ')' if _detail else ''}")

    st.caption(
        "Paid securely via Paystack — MTN MoMo, Telecel, Visa/Mastercard. "
        "If you are charged but not unlocked, contact us within 24 hours."
    )
    # Payment support WhatsApp CTA — server-side notification (council XXIV)
    _ps_email  = st.session_state.get("user_email", "")
    _ps_wa_key = "wa_payment_support_clicked"
    _ps_col1, _ps_col2 = st.columns([2, 1])
    with _ps_col1:
        st.caption("info@impact-receipts.com · WhatsApp +233 50 364 8195")
    with _ps_col2:
        if st.button("WhatsApp support →", key="wa_payment_support_btn", use_container_width=True):
            from utils.whatsapp import notify_founder
            notify_founder("payment_support", user_email=_ps_email)
            st.session_state[_ps_wa_key] = True
    if st.session_state.get(_ps_wa_key):
        from utils.whatsapp import build_wa_url
        st.link_button("Open WhatsApp →",
                       build_wa_url("payment_support", _ps_email),
                       use_container_width=True)
        st.success("✓ Payment support notified — we'll resolve within 4 hours.")


def _subscore_chart(items):
    """Build an interactive horizontal bar chart of sub-scores with hover tooltips.

    items: list of (label, score, max_val, detail) tuples.
    """
    import pandas as pd
    import altair as alt

    rows = []
    for label, score, max_val, detail in items:
        pct = round(min(score / max_val, 1.0) * 100, 1) if max_val else 0.0
        status = "Strong" if pct >= 75 else ("Acceptable" if pct >= 50 else "Below target")
        rows.append({
            "Component": label, "Score": score, "Max": max_val,
            "% of target": pct, "Status": status,
            "Detail": (detail or "").split("\n\n")[0],
        })
    df = pd.DataFrame(rows)

    return (
        alt.Chart(df)
        .mark_bar()
        .encode(
            x=alt.X("% of target:Q", scale=alt.Scale(domain=[0, 100]), title="% of target"),
            y=alt.Y("Component:N", sort=None, title=None),
            color=alt.Color(
                "Status:N",
                scale=alt.Scale(
                    domain=["Strong", "Acceptable", "Below target"],
                    range=["#1B5E20", "#8A6500", "#C62828"],
                ),
                legend=alt.Legend(title=None, orient="bottom"),
            ),
            tooltip=[
                alt.Tooltip("Component:N", title="Component"),
                alt.Tooltip("Score:Q", title="Score", format=".2f"),
                alt.Tooltip("Max:Q", title="Max", format=".2f"),
                alt.Tooltip("% of target:Q", title="% of target"),
                alt.Tooltip("Status:N", title="Status"),
                alt.Tooltip("Detail:N", title="Why"),
            ],
        )
        .properties(width="container", height=alt.Step(28))
    )


def _render_subscore_chart(items, key: str):
    """Render an interactive horizontal bar chart of sub-scores with hover tooltips.

    In low-bandwidth mode, render plain progress bars instead of an Altair chart.
    """
    if st.session_state.get("lite_mode", False):
        for label, score, max_val, _detail in items:
            pct = min(score / max_val, 1.0) if max_val else 0.0
            st.progress(pct, text=f"{label}: {score}/{max_val}")
        return
    st.altair_chart(_subscore_chart(items), use_container_width=True, key=key)


_EVIDENCE_LADDER_TIER_DESCRIPTIONS = {
    "Basic": "Attendance, registration, logs, photos",
    "Moderate": "Follow-up surveys, testimonials",
    "Stronger": "Business/regulatory records, mentor verification, "
                "baseline/endline, external evaluation, comparison groups",
}


def _evidence_ladder_chart(ladder):
    """Build an interactive 'climb the ladder' chart: one rung per evidence tier,
    with the rung you're currently standing on highlighted and hover detail on
    what was matched."""
    import pandas as pd
    import altair as alt

    counts   = ladder.get("tier_counts", {})
    matches  = ladder.get("tier_matches", {})
    dominant = ladder.get("dominant_tier")

    rows = []
    for rung, tier in enumerate(_evaluator.EVIDENCE_LADDER_TIERS, start=1):
        rows.append({
            "Tier": tier,
            "Rung": f"Rung {rung} — {tier}",
            "Sources detected": counts.get(tier, 0),
            "You are here": "👈 Your evidence currently sits here" if tier == dominant else "",
            "Matched keywords": ", ".join(matches.get(tier, [])) or "None detected yet",
            "What counts": _EVIDENCE_LADDER_TIER_DESCRIPTIONS[tier],
        })
    df = pd.DataFrame(rows)

    if dominant:
        opacity_enc = alt.condition(
            alt.FieldEqualPredicate(field="Tier", equal=dominant),
            alt.value(1.0), alt.value(0.35),
        )
    else:
        opacity_enc = alt.value(0.85)

    bars = (
        alt.Chart(df)
        .mark_bar(cornerRadiusEnd=4)
        .encode(
            x=alt.X("Sources detected:Q", title="Sources detected"),
            y=alt.Y("Rung:N", sort=["Rung 3 — Stronger", "Rung 2 — Moderate", "Rung 1 — Basic"], title=None),
            color=alt.Color(
                "Tier:N",
                scale=alt.Scale(
                    domain=["Basic", "Moderate", "Stronger"],
                    range=["#C62828", "#8A6500", "#1B5E20"],
                ),
                legend=alt.Legend(title=None, orient="bottom"),
            ),
            opacity=opacity_enc,
            tooltip=[
                alt.Tooltip("Tier:N", title="Tier"),
                alt.Tooltip("What counts:N", title="What counts"),
                alt.Tooltip("Sources detected:Q", title="Sources detected"),
                alt.Tooltip("Matched keywords:N", title="Matched keywords"),
                alt.Tooltip("You are here:N", title=" "),
            ],
        )
    )

    labels = (
        alt.Chart(df)
        .mark_text(align="left", dx=6, fontWeight="bold", color="#1B5E20")
        .encode(
            x=alt.X("Sources detected:Q"),
            y=alt.Y("Rung:N", sort=["Rung 3 — Stronger", "Rung 2 — Moderate", "Rung 1 — Basic"]),
            text="You are here:N",
        )
    )

    return (bars + labels).properties(width="container", height=alt.Step(40))


def _render_evidence_ladder_chart(ladder, key: str):
    """Render the Evidence Ladder chart.

    In low-bandwidth mode, render a plain markdown list instead of an Altair chart.
    """
    if st.session_state.get("lite_mode", False):
        dominant = ladder.get("dominant_tier")
        counts   = ladder.get("tier_counts", {})
        for tier in _evaluator.EVIDENCE_LADDER_TIERS:
            marker = "👉 " if tier == dominant else "◦ "
            st.markdown(f"{marker}**{tier}** — {counts.get(tier, 0)} source(s) detected")
        return
    st.altair_chart(_evidence_ladder_chart(ladder), use_container_width=True, key=key)


def _render_live_score_preview(slot: int = 1):
    # Reuse cached evaluation from Tab 3 if available (avoids triple evaluate_submission() call)
    _cache = st.session_state.get("_tab3_ev_cache") if slot == 1 else None
    if _cache:
        sub = _cache["sub"]
        ev  = _cache["ev"]
    else:
        sub = _build_submission_from_session(slot)
        try:
            ev = _evaluator.evaluate_submission(sub)
        except Exception:
            st.caption("Fill in the form fields above to see your live score.")
            return

    conf_score     = ev.get("confidence_score", 0)       # post-multiplier (gate assessment)
    raw_conf       = ev.get("raw_confidence_score", conf_score)  # pre-multiplier (matches sub-scores)
    multiplier     = ev.get("content_quality_multiplier", 1.0)
    content_issues = ev.get("content_issues", [])
    clar_score     = ev.get("clarity_score", 0)
    clar_label  = ev.get("clarity_label", "—")
    conf_comp   = ev.get("confidence_components", {})
    clar_comp   = ev.get("clarity_components", {})

    # --- v3.4: cache "what to fix" so destination tabs can show highlighted notes ---
    s = _slot_suffix(slot)
    st.session_state[f"_fixes_computed{s}"] = ev.get("fixes", [])

    # Labels derived from the raw confidence so they match the displayed number
    raw_conf_label, _ = _evaluator.interpret_score(raw_conf) if hasattr(_evaluator, "interpret_score") else (ev.get("confidence_label", "—"), "")

    # Headline: is this good enough to submit?
    _live_diag_state, _ = get_diagnostic_state(raw_conf, clar_score, content_issues, sub.get("beneficiary_voice", ""))
    if _live_diag_state != "INVALID INPUT":
        _render_readiness_banner(_live_diag_state)

    c1, c2 = st.columns(2)
    with c1:
        # Show raw score — always equals sum of sub-scores below
        st.metric("Confidence", f"{raw_conf}/5.0", delta=raw_conf_label, delta_color="off")
    with c2:
        st.metric("Clarity", f"{clar_score}/5.0", delta=clar_label, delta_color="off")

    # Penalty warning: effective score used for gate assessment
    if multiplier < 1.0:
        _penalty_lines = [f"- {ci}" for ci in (content_issues or [])]
        _penalty_body = "\n".join(_penalty_lines) if _penalty_lines else "Review your result statement and evidence description."
        st.warning(
            f"**Content quality penalty (×{multiplier}) applied** — effective score: **{conf_score}/5.0**\n\n"
            f"Issues detected:\n{_penalty_body}\n\n"
            "Fix these to remove the penalty."
        )
        # --- UX: ACTIONABLE SCORE PREVIEW (v3.2) ---
        if st.button("→ Fix: Go to Result Basics", key=f"fix_content_quality_{slot}", type="primary"):
            st.session_state["current_tab"] = 0
            st.session_state["_scroll_to_content"] = True
            st.query_params["tab"] = "0"
            st.rerun()
        # --- END UX: ACTIONABLE SCORE PREVIEW (v3.2) ---

    bd1, bd2 = st.columns(2)

    with bd1:
        st.markdown("**Confidence**")
        dl = conf_comp.get("direct_level", 0)
        vl = conf_comp.get("verify_level", 0)
        rl = conf_comp.get("recency_level", 0)
        ds = conf_comp.get("direct_score", 0)
        vs = conf_comp.get("verify_score", 0)
        rs = conf_comp.get("recency_score", 0)
        _render_subscore_chart([
            ("Directness", ds, 2.0, _DIRECTNESS_TIPS.get(dl, "How directly traceable the evidence is to the result. Target: 1.5+/2.0.")),
            ("Verification", vs, 2.0, _VERIFICATION_TIPS.get(vl, "How rigorously the evidence has been reviewed. Target: 1.5+/2.0.")),
            ("Recency", rs, 1.0, _RECENCY_TIPS.get(rl, "How recent the evidence is relative to the reporting period. Target: 0.7+/1.0.")),
        ], key=f"live_conf_chart{s}")

    with bd2:
        st.markdown("**Clarity**")
        def_s  = clar_comp.get("definition_score", 0)
        meas_s = clar_comp.get("measurement_score", 0)
        integ  = clar_comp.get("integrity_score", 0)
        scope  = clar_comp.get("scope_score", 0)
        gov    = clar_comp.get("governance_score", 0)
        is_qual    = clar_comp.get("is_qualitative", False)
        def_label  = "Narrative Definition" if is_qual else "Definition"
        def_tip    = _CLARITY_TIPS["definition_qualitative"] if is_qual else _CLARITY_TIPS["definition"]
        meas_label = "Sourcing & Triangulation" if is_qual else "Measurement"
        meas_tip   = _CLARITY_TIPS["measurement_qualitative"] if is_qual else _CLARITY_TIPS["measurement"]
        _render_subscore_chart([
            (def_label, def_s, 1.25, def_tip),
            (meas_label, meas_s, 1.25, meas_tip),
            ("Integrity", integ, 1.0, _CLARITY_TIPS["integrity"]),
            ("Scope", scope, 0.75, _CLARITY_TIPS["scope"]),
            ("Governance", gov, 0.75, _CLARITY_TIPS["governance"]),
        ], key=f"live_clar_chart{s}")

    # Gate assessment uses penalized conf_score
    state, state_sub = get_diagnostic_state(conf_score, clar_score)
    st.caption(f"Status: **{state}** — {state_sub}")

    # --- UX: ACTIONABLE SCORE PREVIEW (v3.2) ---
    if state in ("MISLEADING", "FUNDAMENTALLY WEAK"):
        if st.button("→ Fix: Sharpen Result Statement", key=f"fix_misleading_{slot}", type="primary"):
            st.session_state["current_tab"] = 0
            st.session_state["_scroll_to_content"] = True
            st.query_params["tab"] = "0"
            st.rerun()
    if state in ("UNDEREVIDENCED", "FUNDAMENTALLY WEAK"):
        if st.button("→ Fix: Strengthen Evidence", key=f"fix_underevidenced_{slot}", type="primary"):
            st.session_state["current_tab"] = 2
            st.session_state["_scroll_to_content"] = True
            st.query_params["tab"] = "2"
            st.rerun()
    if state == "NEEDS REFINEMENT":
        if st.button("→ Fix: Review Specific Gaps", key=f"fix_refinement_{slot}", type="primary"):
            st.session_state["current_tab"] = 1
            st.session_state["_scroll_to_content"] = True
            st.query_params["tab"] = "1"
            st.rerun()
    # --- END UX: ACTIONABLE SCORE PREVIEW (v3.2) ---

    # --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
    gov_score, pii_selected, gov_gaps, safeguarding_triggered, minors_involved = _compute_governance_score(slot)
    conf_100 = round(raw_conf * 20, 1)
    gov_adjustment = round(gov_score * 0.3, 1)
    if gov_score == 0 and (pii_selected or safeguarding_triggered or minors_involved):
        gov_adjustment -= 8
    adjusted_conf = min(100.0, conf_100 + gov_adjustment)

    st.session_state["_gov_score_computed"] = gov_score
    st.session_state["_gov_gaps_computed"]  = gov_gaps
    st.session_state["_gov_pii_computed"]   = pii_selected
    st.session_state["_conf_adj_computed"]  = adjusted_conf

    # Read governance field values for display (scoring unchanged — uses _compute_governance_score above)
    _disp_consent    = st.session_state.get(f"gov_consent_status{s}", "")
    _disp_anon       = st.session_state.get(f"gov_anonymization_status{s}", "")
    _disp_law        = st.session_state.get(f"gov_compliance_law_status{s}", "")
    _disp_safeguard  = st.session_state.get(f"gov_safeguarding_status{s}", "")
    _disp_child_safe = st.session_state.get(f"gov_child_safeguarding_status{s}", "")
    _disp_secure     = st.session_state.get(f"gov_secure_handling_status{s}", "")
    _disp_dpp        = st.session_state.get("gov_dpp_uploaded", False)
    _answered = sum([
        _disp_consent    not in ("", "Choose an option...", "Select consent status..."),
        _disp_anon       not in ("", "Choose an option...", "Select anonymization status..."),
        _disp_law        not in ("", "Choose an option...", "Select compliance status..."),
        _disp_safeguard  not in ("", "Choose an option...", "Select safeguarding status..."),
        _disp_child_safe not in ("", "Choose an option...", "Select child safeguarding status..."),
        _disp_secure     not in ("", "Choose an option...", "Select secure handling status..."),
    ])

    st.divider()
    st.markdown("#### 🛡️ Governance & Compliance")
    _gc1, _gc2 = st.columns([1, 2])
    with _gc1:
        st.metric("Governance readiness", f"{int(gov_score / 24 * 100)}%")
    with _gc2:
        st.metric("Governance-Adjusted Confidence", f"{adjusted_conf:.0f} / 100")

    # Single bold status line — no alarming styling
    if gov_score >= 20:
        st.markdown("**Strong governance — major requirements addressed.**")
    elif gov_score >= 11:
        st.markdown("**Partial compliance — some requirements still recommended before submission.**")
    elif _answered == 0:
        st.markdown("**Data governance checklist not yet completed — 0 of 6 questions answered.**")
    else:
        st.markdown(f"**Governance readiness: {int(gov_score / 24 * 100)}% — review items in the checklist below.**")

    # Remediation action — placed right next to the status line so the fix is one click away
    if gov_score < 20:
        if st.button("→ Fix: Governance Issues", key=f"fix_gov_btn_{slot}", type="primary"):
            st.session_state["current_tab"] = 2
            st.session_state["_scroll_to_content"] = True
            st.query_params["tab"] = "2"
            st.rerun()

    # Per-item checklist
    with st.expander("Governance checklist detail", expanded=(gov_score < 11), key=f"gov_checklist_detail_{slot}"):
        for _lbl, _val, _max, _cmap in [
            ("Beneficiary consent",                  _disp_consent,    5, CONSENT_CHECKLIST_MAP),
            ("Data anonymization",                   _disp_anon,       4, ANON_CHECKLIST_MAP),
            ("Data law compliance",                  _disp_law,        3, LAW_CHECKLIST_MAP),
            ("Do-no-harm & safeguarding",            _disp_safeguard,  3, SAFEGUARDING_CHECKLIST_MAP),
            ("Child safeguarding",                   _disp_child_safe, 3, CHILD_SAFEGUARDING_CHECKLIST_MAP),
            ("Secure handling of identifiable testimony", _disp_secure, 3, SECURE_HANDLING_CHECKLIST_MAP),
        ]:
            _icon, _desc, _earned = _cmap.get(_val, ("✗", "Not answered", 0))
            st.markdown(f"{_icon} **{_lbl}** — {_desc} ({_earned}/{_max} pts)")
        if _disp_dpp:
            st.markdown("✓ **Data protection policy** — uploaded (+5 bonus)")
        else:
            st.caption("◦ Data protection policy not uploaded (optional +5 bonus)")
    # --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---


# ---------------------------------------------------------------------------
# JSON inputs export / import
# ---------------------------------------------------------------------------

def _build_inputs_json(timestamp: str) -> str:
    active = st.session_state.get("active_slots_run", st.session_state.get("active_slots", 1))
    slots_data = []
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        slot_dict = {}
        for k in _BASE_FORM_KEYS:
            slot_dict[k] = st.session_state.get(f"{k}{s}", "")
        ed = st.session_state.get(f"evidence_date{s}")
        slot_dict["evidence_date"] = ed.isoformat() if hasattr(ed, "isoformat") else ""
        rs = st.session_state.get(f"reporting_start{s}")
        slot_dict["reporting_start"] = rs.isoformat() if hasattr(rs, "isoformat") else ""
        re_ = st.session_state.get(f"reporting_end{s}")
        slot_dict["reporting_end"] = re_.isoformat() if hasattr(re_, "isoformat") else ""
        raw_files = st.session_state.get(f"uploaded_files_widget{s}") or []
        slot_dict["uploaded_filenames"] = [f.name for f in raw_files if hasattr(f, "name")]
        # --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        for _gk in ["gov_consent_status", "gov_anonymization_status", "gov_compliance_law_status",
                     "gov_safeguarding_status", "gov_child_safeguarding_status", "gov_secure_handling_status"]:
            slot_dict[_gk] = st.session_state.get(f"{_gk}{s}", "")
        slot_dict["gov_dpp_uploaded"] = st.session_state.get("gov_dpp_uploaded", False)
        # --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        slots_data.append(slot_dict)

    payload = {
        "timestamp": timestamp,
        "session_id": f"ir-{timestamp}",
        "active_slots": active,
        "slots": slots_data,
        "consent_examples": st.session_state.get("consent_examples", False),
        "has_seen_tutorial": st.session_state.get("has_seen_tutorial", False),
        "user_email": st.session_state.get("user_email", ""),
    }
    return json.dumps(payload, indent=2, ensure_ascii=False)


def _normalize_draft_json(data: dict) -> dict:
    """Convert a flat '_save_draft()' export (the '📥 Download Draft (JSON)' format,
    e.g. {"active_slots": 1, "result_statement": "...", "logframe_indicator": "...", ...})
    into the {"slots": [...]} format produced by _build_inputs_json() / 'Save Inputs (JSON)'
    and expected by _load_from_inputs_json."""
    if "slots" in data:
        return data

    active = int(data.get("active_slots", 1))
    slots_data = []
    for slot in range(1, active + 1):
        s = _slot_suffix(slot)
        slot_dict = {}
        for k in _BASE_FORM_KEYS:
            if f"{k}{s}" in data:
                slot_dict[k] = data[f"{k}{s}"]
        for dk in ("evidence_date", "reporting_start", "reporting_end"):
            if f"{dk}{s}" in data:
                slot_dict[dk] = data[f"{dk}{s}"]
        if f"uploaded_filenames{s}" in data:
            slot_dict["uploaded_filenames"] = data[f"uploaded_filenames{s}"]
        for gk in ("gov_consent_status", "gov_anonymization_status", "gov_compliance_law_status",
                    "gov_safeguarding_status", "gov_child_safeguarding_status", "gov_secure_handling_status"):
            if f"{gk}{s}" in data:
                slot_dict[gk] = data[f"{gk}{s}"]
        slot_dict["gov_dpp_uploaded"] = data.get("gov_dpp_uploaded", False)
        slots_data.append(slot_dict)

    return {
        "timestamp": data.get("timestamp", ""),
        "active_slots": active,
        "slots": slots_data,
    }


def _load_from_inputs_json(data: dict):
    data = _normalize_draft_json(data)
    if "slots" not in data:
        st.error("Invalid file format — missing 'slots' key. Please upload a file exported by ImpactProof.")
        return

    active = int(data.get("active_slots", 1))
    st.session_state["active_slots"] = active

    for slot_idx, slot_dict in enumerate(data["slots"]):
        slot = slot_idx + 1
        s = _slot_suffix(slot)
        for k in _BASE_FORM_KEYS:
            if k in slot_dict:
                try:
                    st.session_state[f"{k}{s}"] = slot_dict[k]
                except Exception:
                    # key already backs a widget instantiated earlier in this run
                    # (e.g. the global "sector" selector) — skip, non-critical
                    pass
        raw_date = slot_dict.get("evidence_date", "")
        if raw_date:
            try:
                st.session_state[f"evidence_date{s}"] = date.fromisoformat(raw_date)
            except (ValueError, TypeError):
                pass
        for _dk2, _sk2 in [("reporting_start", f"reporting_start{s}"),
                            ("reporting_end",   f"reporting_end{s}")]:
            _rd2 = slot_dict.get(_dk2, "")
            if _rd2:
                try:
                    st.session_state[_sk2] = date.fromisoformat(_rd2)
                except (ValueError, TypeError):
                    pass
        # restore governance fields not in _BASE_FORM_KEYS loop
        _gov_dpp = slot_dict.get("gov_dpp_uploaded")
        if _gov_dpp is not None:
            st.session_state["gov_dpp_uploaded"] = bool(_gov_dpp)
        st.session_state[f"draft_uploaded_filenames{s}"] = slot_dict.get("uploaded_filenames", [])

    ts = data.get("timestamp", "unknown")
    # --- UX: SMART DEFAULTS (v3.2) ---
    _prefill_count = sum(
        1 for _sd in data.get("slots", [{}])
        for _k, _v in _sd.items() if _v and _k in _BASE_FORM_KEYS
    )
    st.success(f"✅ Draft loaded — {_prefill_count} fields pre-filled. Review and update as needed.")
    st.session_state["_form_is_resumption"] = True  # show "continuing from previous session" banner
    # Restore user preferences from draft
    if "consent_examples" in data:
        st.session_state["consent_examples"] = bool(data["consent_examples"])
    if data.get("has_seen_tutorial"):
        st.session_state["has_seen_tutorial"] = True
        st.session_state["tutorial_step"] = 99  # skip all steps
    if data.get("user_email") and not st.session_state.get("user_email"):
        # Only fill in an email if no session is already authenticated -- an
        # uploaded/imported JSON is user-controlled input (e.g. via the
        # Instant Report Check file uploader), so it must never be allowed to
        # overwrite an already-verified session's identity. Without this
        # guard, a crafted {"user_email": "victim@example.com"} in an
        # uploaded file would silently hijack the uploader's own session into
        # acting as that other account for every subsequent utils/audits.py
        # call (list/view/delete their saved audits, save new ones as them).
        st.session_state["user_email"] = data["user_email"]
    # --- END UX: SMART DEFAULTS (v3.2) ---
    # bump version so _irc_widget-backed fields re-seed from the freshly loaded values
    st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
    st.session_state["_tab2_auto_advanced"] = True
    st.session_state["screen"] = 1
    st.session_state["current_tab"] = 0
    st.query_params["screen"] = "1"
    st.query_params["tab"] = "0"
    st.rerun()


def _build_submission_from_session(slot: int = 1) -> dict:
    """Assemble evaluator-compatible submission dict from session_state for a given slot."""
    s = _slot_suffix(slot)

    ev_type = st.session_state.get(f"evidence_type{s}", "")
    if ev_type == "Other":
        ev_type = st.session_state.get(f"evidence_type_other{s}", "") or "Other"

    int_rev = st.session_state.get(f"internal_review{s}", "Choose an option...")
    if int_rev == "Other":
        int_rev = st.session_state.get(f"internal_review_other{s}", "") or "Other"

    ext_rev = st.session_state.get(f"external_review{s}", "Choose an option...")
    if ext_rev == "Other":
        ext_rev = st.session_state.get(f"external_review_other{s}", "") or "Other"

    donor = st.session_state.get("donor_selected", "(No donor specified)")
    if donor == "Other":
        donor = st.session_state.get("donor_other", "") or "Other"

    sector = st.session_state.get("sector", SECTOR_OPTIONS[0])
    if sector == "Other":
        sector = st.session_state.get("sector_other", "") or "Other"

    submission_type = st.session_state.get("submission_type", "Select submission type...")

    return {
        "result_statement":   st.session_state.get(f"result_statement{s}", ""),
        "target_group":       st.session_state.get(f"target_group{s}", ""),
        "timeframe":          st.session_state.get(f"timeframe{s}", ""),
        "geographic_scope":   st.session_state.get(f"geographic_scope{s}", ""),
        "additional_context": st.session_state.get(f"additional_context{s}", ""),
        "learning_notes":     st.session_state.get(f"learning_notes{s}", ""),
        "limitations_notes":  st.session_state.get(f"limitations_notes{s}", ""),
        "internal_review":    int_rev,
        "external_review":    ext_rev,
        "attached_filenames": st.session_state.get(f"uploaded_files{s}", []),
        "beneficiary_voice":    st.session_state.get(f"beneficiary_voice{s}", ""),
        "bv_method_detail":     st.session_state.get(f"bv_method_detail{s}", ""),
        "logframe_indicator":   st.session_state.get(f"logframe_indicator{s}", ""),
        "logframe_baseline":    st.session_state.get(f"logframe_baseline{s}", ""),
        "logframe_target":      st.session_state.get(f"logframe_target{s}", ""),
        "logframe_achievement": st.session_state.get(f"logframe_achievement{s}", ""),
        "logframe_data_forthcoming": bool(st.session_state.get(f"logframe_data_forthcoming{s}", False)),
        "reporting_start":      _format_date(st.session_state.get(f"reporting_start{s}")),
        "reporting_end":        _format_date(st.session_state.get(f"reporting_end{s}")),
        "provenance_checklist": {
            "sampling_documented":     st.session_state.get(f"provenance_sampling{s}", "Choose an option..."),
            "double_counting_checked": st.session_state.get(f"provenance_dedup{s}", "Choose an option..."),
            "collection_tool_named":   st.session_state.get(f"provenance_tool{s}", "Choose an option..."),
            "collector_independent":   st.session_state.get(f"provenance_independence{s}", "Choose an option..."),
            "recall_period_ok":        st.session_state.get(f"provenance_recall{s}", "Choose an option..."),
            "auditor_traceable":       st.session_state.get(f"provenance_traceability{s}", "Choose an option..."),
        },
        "qualitative_evidence": st.session_state.get(f"qualitative_evidence{s}", False),
        "qualitative_rigor_checklist": {
            "sourcing_documented":          st.session_state.get(f"qual_sourcing{s}", False),
            "triangulated":                 st.session_state.get(f"qual_triangulated{s}", False),
            "bias_considered":              st.session_state.get(f"qual_bias{s}", False),
            "beneficiary_voice_represented": st.session_state.get(f"qual_voice{s}", False),
            "consent_ethics_addressed":      st.session_state.get(f"qual_consent{s}", False),
        },
        "attribution_contribution": st.session_state.get(f"attribution_contribution{s}", "Not specified"),
        "disaggregation_status":     st.session_state.get(f"disaggregation_status{s}", "Not specified"),
        "review_status":     st.session_state.get(f"review_status{s}", _SUBMISSION_STATUS_OPTIONS[0]),
        "reviewer_name":     st.session_state.get(f"reviewer_name{s}", ""),
        "reviewer_role":     st.session_state.get(f"reviewer_role{s}", ""),
        "reviewer_date":     st.session_state.get(f"reviewer_date{s}", ""),
        "reviewer_decision": st.session_state.get(f"reviewer_decision{s}", ""),
        "reviewer_notes":    st.session_state.get(f"reviewer_notes{s}", ""),
        "donor":                     donor,
        "sector":                    sector,
        "project_name":              st.session_state.get("project_name", ""),
        "submission_type":           submission_type,
        "org_type":                  st.session_state.get("org_type", "International NGO (INGO)"),
        "evidence": [{
            "type":        ev_type,
            "description": st.session_state.get(f"evidence_description{s}", ""),
            "recency":     _format_date(st.session_state.get(f"evidence_date{s}")),
            "verified_by": st.session_state.get(f"verifier{s}", ""),
        }],
    }


def _render_slot_fields(slot: int):
    """Render all form fields for one result slot."""
    s = _slot_suffix(slot)

    for key, default in [
        (f"evidence_type{s}", EVIDENCE_TYPES[0]),
        (f"internal_review{s}", "Choose an option..."),
        (f"external_review{s}", "Choose an option..."),
        # --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        (f"gov_consent_status{s}", "Select consent status..."),
        (f"gov_anonymization_status{s}", "Select anonymization status..."),
        (f"gov_compliance_law_status{s}", "Select compliance status..."),
        (f"gov_safeguarding_status{s}", "Select safeguarding status..."),
        (f"gov_child_safeguarding_status{s}", "Select child safeguarding status..."),
        (f"gov_secure_handling_status{s}", "Select secure handling status..."),
        # --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        # --- REVIEW-HANDOFF LAYER (v3.6) ---
        (f"review_status{s}", _SUBMISSION_STATUS_OPTIONS[0]),
        (f"reviewer_decision{s}", _REVIEW_DECISION_OPTIONS[0]),
        # --- END REVIEW-HANDOFF LAYER (v3.6) ---
    ]:
        if key not in st.session_state:
            st.session_state[key] = default

    _sector = st.session_state.get("sector", SECTOR_OPTIONS[0])
    _ph_key = "Other" if _sector in ("Other", "(No sector selected)") else _sector
    _ph = SECTOR_PLACEHOLDERS.get(_ph_key, SECTOR_PLACEHOLDERS["Other"])

    st.text_area(
        "Result statement",
        key=f"result_statement{s}",
        placeholder=_ph["result"],
        height=100,
        help="What did your project achieve? Include the action verb, number, target group, location, and timeframe.",
    )
    _rs = st.session_state.get(f"result_statement{s}", "")
    if _rs and len(_rs.strip()) < 20:
        st.warning("Result statement is very short. Include: action verb + number + population + timeframe.")
    elif _rs and not any(c.isdigit() for c in _rs):
        st.caption("Tip: Add a number (e.g., '500 farmers trained') — quantified claims score higher.")

    st.markdown("#### Logframe Linkage")
    st.caption(
        "**Why this matters:** A real African consultancy had their final donor report "
        "rejected 3 times in 2024 because results weren't tied to logframe indicators. "
        "40+ hours of rework. We don't want that to happen to you."
    )

    _lf_api_key = (
        st.secrets.get("ANTHROPIC_API_KEY", "")
        if hasattr(st, "secrets") else
        os.environ.get("ANTHROPIC_API_KEY", "")
    )
    if _lf_api_key:
        with st.expander("🎯 AI Logframe Match — paste your indicators, get a suggested match", expanded=False):
            st.caption(
                "Paste your approved logframe indicators (one per line). The AI suggests which "
                "one this result reports against — it never forces a match, and it only quotes "
                "words already in your result statement or your pasted indicators."
            )
            st.text_area(
                "Your logframe indicators (one per line)",
                key=f"_lf_paste{s}", height=100,
                placeholder=(
                    "Indicator 1.2: Number of households with access to safe water\n"
                    "Indicator 2.1: % of farmers applying climate-smart practices"
                ),
            )
            if st.button("Match my result to an indicator", key=f"lf_match_btn{s}"):
                _lf_rs = st.session_state.get(f"result_statement{s}", "")
                _lf_raw = st.session_state.get(f"_lf_paste{s}", "")
                _lf_indicators = [ln.strip() for ln in _lf_raw.splitlines() if ln.strip()]
                if not _lf_rs.strip():
                    st.warning("Enter a result statement above first.")
                elif not _lf_indicators:
                    st.warning("Paste at least one logframe indicator above.")
                else:
                    with st.spinner("Matching your result to an indicator…"):
                        from council import match_logframe_indicator
                        _lf_match = match_logframe_indicator(_lf_rs, _lf_indicators, _lf_api_key)
                    st.session_state[f"_lf_match_result{s}"] = _lf_match
                    if _lf_match.get("confidence_label") != "None" and _lf_match.get("best_match"):
                        st.session_state[f"logframe_indicator{s}"] = _lf_match["best_match"]
                        st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                    st.rerun()

            _lf_result = st.session_state.get(f"_lf_match_result{s}")
            if _lf_result:
                _lf_cl = _lf_result.get("confidence_label", "None")
                if _lf_cl == "None" or not _lf_result.get("best_match"):
                    st.info("No confident match found — enter the indicator manually below.")
                else:
                    st.success(
                        f"AI-suggested match ({_lf_cl.lower()} confidence) — confirm against "
                        f"your approved logframe before submitting."
                    )
                    if _lf_result.get("justification"):
                        st.caption(_lf_result["justification"])

    _irc_widget(
        st.text_input,
        "Logframe indicator this result reports against",
        f"logframe_indicator{s}", default="",
        placeholder=_ph.get("logframe_indicator", "e.g., Indicator 1.2: Number of [target group] achieving [outcome]"),
        help=(
            "Copy the exact indicator name and code from your approved Technical Proposal or logframe. "
            "If you cannot quote it, your donor cannot match your result to your commitment."
        ),
    )
    _irc_widget(
        st.text_input,
        "Original target for this indicator (from logframe)",
        f"logframe_target{s}", default="",
        placeholder=_ph.get("logframe_target", "e.g., 250 youth trained by Q4 2025"),
        help=(
            "The target as approved in the original Technical Proposal. Donors compare achievements "
            "against approved targets — not revised internal targets."
        ),
    )
    _irc_widget(
        st.text_input,
        "Actual achievement (must match your result statement)",
        f"logframe_achievement{s}", default="",
        placeholder=_ph.get("logframe_achievement", "e.g., [Actual number] by [date] — [%] of original target"),
        help=(
            "The actual delivered number, ideally with % achievement vs original target. "
            "Must reconcile with your result statement above."
        ),
    )

    st.text_input(
        "Target group", key=f"target_group{s}",
        placeholder=_ph["target_group"],
        help="Who specifically? Age group, gender, role, occupation. Avoid 'beneficiaries' alone — name the population.",
    )

    st.text_input(
        "Timeframe", key=f"timeframe{s}",
        placeholder="e.g., January - June 2025",
        help="Specific dates or quarters. 'January–June 2025' is stronger than 'In 2025'.",
    )

    st.text_input(
        "Geographic scope", key=f"geographic_scope{s}",
        placeholder=_ph["geographic_scope"],
        help="Districts, regions, or specific sites. 'Volta Region' beats 'rural areas'.",
    )

    st.text_area(
        "Describe your supporting evidence", key=f"evidence_description{s}",
        placeholder=_ph["evidence_description"],
        height=120,
        help="Describe the actual document or data: who collected it, how, and what's in it.",
    )
    _ed = st.session_state.get(f"evidence_description{s}", "")
    if _ed and len(_ed.strip()) < 30:
        st.warning("Evidence description is brief. Specify: who collected it, how, and what it contains.")

    st.selectbox(
        "Evidence type", key=f"evidence_type{s}",
        options=EVIDENCE_TYPES,
        help=EVIDENCE_TYPE_HELP,
    )
    ev_type = st.session_state.get(f"evidence_type{s}", EVIDENCE_TYPES[0])
    ev_desc = st.session_state.get(f"evidence_description{s}", "")
    _dl = _evaluator.get_directness_level(ev_type, ev_desc)
    _ds = round((_dl / 5) * 2.0, 1)

    if ev_type == "Other":
        st.text_input("Specify evidence type", key=f"evidence_type_other{s}")

    int_rev = st.session_state.get(f"internal_review{s}", INTERNAL_REVIEW_OPTIONS[0])
    st.selectbox(
        "Internal review", key=f"internal_review{s}",
        options=INTERNAL_REVIEW_OPTIONS,
        help="Did anyone in your organization review or cross-check this data?",
    )
    int_rev = st.session_state.get(f"internal_review{s}", INTERNAL_REVIEW_OPTIONS[0])
    _int_vl = _evaluator.get_verification_level(int_rev, "No external review", "")
    _int_vs = round((_int_vl / 5) * 2.0, 1)
    if _int_vs == 0:
        st.warning("No internal review — adding a reviewer significantly strengthens Verification.")

    if int_rev == "Other":
        st.text_input("Specify internal reviewer", key=f"internal_review_other{s}")

    ext_rev = st.session_state.get(f"external_review{s}", EXTERNAL_REVIEW_OPTIONS[0])
    st.selectbox(
        "External review", key=f"external_review{s}",
        options=EXTERNAL_REVIEW_OPTIONS,
        help="Did an outside party verify the data? Government, partner, auditor, or evaluator.",
    )
    ext_rev = st.session_state.get(f"external_review{s}", EXTERNAL_REVIEW_OPTIONS[0])
    verifier_text = st.session_state.get(f"verifier{s}", "")
    _full_vl = _evaluator.get_verification_level(int_rev, ext_rev, verifier_text)
    _full_vs = round((_full_vl / 5) * 2.0, 1)
    _added   = round(_full_vs - _int_vs, 1)
    if ext_rev == "No external review":
        st.warning("No external review — independent verification raises your score significantly.")

    if ext_rev == "Other":
        st.text_input("Specify external reviewer", key=f"external_review_other{s}")

    # --- UX: CONDITIONAL FIELDS (v3.2) ---
    if st.session_state.get(f"internal_review{s}") != "Not reviewed":
        _irc_widget(
            st.text_input, "Who verified this?", f"verifier{s}", default="",
            placeholder=_ph.get("verifier", "e.g., District Agriculture Officer, partner org M&E lead, external evaluator"),
            help="The person or organization that confirmed the data is accurate.",
        )
    # --- END UX: CONDITIONAL FIELDS (v3.2) ---

    st.markdown("#### Reporting Period")
    st.caption("The period this submission covers. Evidence dates outside this range trigger a flag.")
    _rp_col_s, _rp_col_e = st.columns(2)
    with _rp_col_s:
        st.date_input(
            "Reporting period start",
            value=st.session_state.get(f"reporting_start{s}"),
            key=f"reporting_start{s}",
            help="When does the period this report covers begin?",
        )
    with _rp_col_e:
        st.date_input(
            "Reporting period end",
            value=st.session_state.get(f"reporting_end{s}"),
            key=f"reporting_end{s}",
            help="When does the period this report covers end?",
        )

    st.date_input(
        "When was this evidence collected?",
        value=st.session_state.get(f"evidence_date{s}"),
        key=f"evidence_date{s}",
        help="When was the data collected? Use the most recent date if multiple sources.",
    )
    _ed = st.session_state.get(f"evidence_date{s}")
    if _ed and hasattr(_evaluator, "get_recency_diagnostic"):
        _rec_diag = _evaluator.get_recency_diagnostic(_ed)
        if "0.4/1.0" in _rec_diag or "0.2/1.0" in _rec_diag:
            st.warning(_rec_diag)
        elif "0.6/1.0" in _rec_diag:
            st.info(_rec_diag)
        else:
            st.success(_rec_diag)
    _rp_s = st.session_state.get(f"reporting_start{s}")
    _rp_e = st.session_state.get(f"reporting_end{s}")
    if _ed and _rp_s and _rp_e and hasattr(_evaluator, "validate_reporting_period"):
        _, _rp_msg, _rp_sev = _evaluator.validate_reporting_period(_ed, _rp_s, _rp_e)
        if _rp_sev == "ERROR":
            st.error(_rp_msg)
        elif _rp_sev == "WARNING":
            st.warning(_rp_msg)
        elif _rp_msg:
            st.success(_rp_msg)

    st.markdown("#### Beneficiary Voice")
    st.caption(
        "Did the beneficiaries contribute to or validate this evidence?"
    )
    st.selectbox(
        "How were beneficiary voices captured?",
        key=f"beneficiary_voice{s}",
        options=_BV_OPTIONS,
        help=(
            "Bond Evidence Principle 1 (2024 refresh): Voice & Inclusion. "
            "The strongest evidence includes beneficiary perspectives, not just provider reports."
        ),
    )
    _bv_sel_1 = st.session_state.get(f"beneficiary_voice{s}", "")
    _BV_HIGH = {
        "Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS)",
        "Beneficiary representatives consulted (community leaders, beneficiary committees)",
    }
    if _bv_sel_1 in _BV_HIGH:
        st.text_input(
            "Briefly describe the method — when conducted and approximately how many people participated",
            key=f"bv_method_detail{s}",
            placeholder="e.g., Phone survey with 120 farmers, March 2025",
            help="Required to receive the full beneficiary voice bonus (≥20 characters).",
        )
        _bv_detail_1 = st.session_state.get(f"bv_method_detail{s}", "")
        _bv_chars_1  = len(_bv_detail_1.strip())
        if _bv_chars_1 >= 20:
            st.caption("✓ Full beneficiary voice bonus unlocked: +0.5")
        else:
            st.caption(f"Add {20 - _bv_chars_1} more characters to unlock the full +0.5 bonus (currently +0.1).")

    prev_files = st.session_state.get(f"draft_uploaded_filenames{s}", [])
    if prev_files:
        st.caption(
            f"For security, browsers don't let us keep re-uploaded files between sessions — "
            f"please re-attach: {', '.join(prev_files)}"
        )
    st.file_uploader(
        "Attach supporting documents (optional)", key=f"uploaded_files_widget{s}",
        accept_multiple_files=True,
        type=["pdf", "docx", "xlsx", "csv", "jpg", "jpeg", "png", "txt"],
        help="Attach raw evidence files — datasets, signed sheets, photos with metadata, partner letters.",
    )


# ---------------------------------------------------------------------------
# Screen 1 — Tab helper functions (v3.3)
# ---------------------------------------------------------------------------

def _tab_slot_setup(slot: int):
    s = _slot_suffix(slot)
    for key, default in [
        (f"evidence_type{s}", EVIDENCE_TYPES[0]),
        (f"internal_review{s}", "Choose an option..."),
        (f"external_review{s}", "Choose an option..."),
        # --- GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        (f"gov_consent_status{s}", "Select consent status..."),
        (f"gov_anonymization_status{s}", "Select anonymization status..."),
        (f"gov_compliance_law_status{s}", "Select compliance status..."),
        (f"gov_safeguarding_status{s}", "Select safeguarding status..."),
        (f"gov_child_safeguarding_status{s}", "Select child safeguarding status..."),
        (f"gov_secure_handling_status{s}", "Select secure handling status..."),
        # --- END GOVERNANCE & COMPLIANCE LAYER (v3.2) ---
        # --- REVIEW-HANDOFF LAYER (v3.6) ---
        (f"review_status{s}", _SUBMISSION_STATUS_OPTIONS[0]),
        (f"reviewer_decision{s}", _REVIEW_DECISION_OPTIONS[0]),
        # --- END REVIEW-HANDOFF LAYER (v3.6) ---
    ]:
        if key not in st.session_state:
            st.session_state[key] = default
    _sector = st.session_state.get("sector", SECTOR_OPTIONS[0])
    _ph_key = "Other" if _sector in ("Other", "(No sector selected)") else _sector
    _ph = SECTOR_PLACEHOLDERS.get(_ph_key, SECTOR_PLACEHOLDERS["Other"])
    return s, _ph


def _irc_widget(widget_fn, label, base_key, default=None, **kwargs):
    """Render a widget whose canonical value lives in the plain (non-widget)
    session_state key `base_key`, so the rest of the app (sidebar summary,
    scoring, _build_submission_from_session) keeps reading and writing
    `base_key` exactly as before — including Instant Report Check pre-fills
    and the "Today" date-shortcut buttons.

    The widget itself is bound to a version-suffixed key (`base_key__w{N}`).
    Streamlit resets a session_state entry to its last frontend value at the
    start of every run once that key has ever backed a widget, which would
    silently discard out-of-band writes to `base_key` if `base_key` were used
    directly as the widget key. Keeping the widget key separate avoids that,
    and bumping `N` (via `_irc_fill_version`) mints a fresh, never-instantiated
    widget key whenever IRC re-fills the form, so the new value is picked up.
    """
    ver = st.session_state.get("_irc_fill_version", 0)
    wkey = f"{base_key}__w{ver}"
    shadow_key = f"_irc_shadow_{wkey}"
    base_val = st.session_state.get(base_key, default)
    if wkey not in st.session_state or st.session_state.get(shadow_key, object()) != base_val:
        st.session_state[wkey] = base_val
    widget_fn(label, key=wkey, **kwargs)
    st.session_state[base_key] = st.session_state[wkey]
    st.session_state[shadow_key] = st.session_state[wkey]


# --- v3.3 field-validation marker sets ---
_DEMO_MARKERS = {"farmer", "women", "woman", "youth", "child", "children", "household",
    "teacher", "worker", "patient", "student", "beneficiar", "community",
    "resident", "family", "families", "men", "girl", "boy", "aged", "adult"}
_DATE_MARKERS = {"jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep",
    "oct", "nov", "dec", "q1", "q2", "q3", "q4", "quarter", "2023", "2024",
    "2025", "2026", "2027", "month", "year", "period", "week"}
_LOC_MARKERS  = {"district", "region", "province", "state", "county", "city", "town",
    "village", "community", "ward", "national", "local", "ghana", "nigeria",
    "kenya", "uganda", "ethiopia", "senegal", "africa", "northern", "southern",
    "eastern", "western", "central", "zone", "area", "site"}

# --- v3.4: maps "what to fix" messages to the tab/field where they should be addressed ---
_FIX_FIELD_MAP = [
    ("missing unit, timeframe, or target group", 0, "Result statement, Target group, Timeframe"),
    ("sites and groups included and excluded",   0, "Geographic scope / Target group"),
    ("Name an owner for this result",            1, "Logframe indicator & linkage"),
    ("primary record",                           2, "Evidence type & description"),
    ("internal reviewer or an external partner", 2, "Internal review / External review / Verifier"),
    ("evidence date is within 6 months",         2, "Evidence date"),
    ("collection method and sampling approach",  2, "Evidence description"),
    ("Close data gaps with original source records", 2, "Evidence description"),
]


def _render_fix_notes(slot: int, tab_idx: int):
    """Show highlighted notes for gaps relevant to this tab, computed on Review & Submit."""
    s = _slot_suffix(slot)
    fixes = st.session_state.get(f"_fixes_computed{s}", [])
    notes = []
    for fix in fixes:
        msg = fix.get("message", "")
        for kw, t, field in _FIX_FIELD_MAP:
            if t == tab_idx and kw in msg:
                notes.append((field, fix["message"], fix.get("score_impact", "")))
                break
    if tab_idx == 2:
        for gap in st.session_state.get("_gov_gaps_computed", []):
            notes.append(("Compliance & Data Governance", gap, "raises Governance score"))
    if notes:
        lines = "\n".join(
            f"- **{field}** — {msg} _({impact})_" for field, msg, impact in notes
        )
        st.info(f"**📌 To improve your score, address:**\n\n{lines}")


_MONTH_NAMES = (
    "january","february","march","april","may","june",
    "july","august","september","october","november","december",
)
_MONTH_ABBR = ("jan","feb","mar","apr","may","jun","jul","aug","sep","sept","oct","nov","dec")
_ALL_MONTHS_PAT = "|".join(_MONTH_NAMES + _MONTH_ABBR)

# Compiled regexes for timeframe extraction
_TF_RANGE_RE = re.compile(
    r"(?:between\s+|from\s+)?"
    r"\b(" + _ALL_MONTHS_PAT + r")\b"
    r"[\s\w,–\-]*?"
    r"(?:\band\b|\bto\b|–|-)\s*"
    r"\b(" + _ALL_MONTHS_PAT + r")\b"
    r"\s+(\d{4})\b",
    re.IGNORECASE,
)
_TF_SINGLE_RE = re.compile(
    r"\b(" + _ALL_MONTHS_PAT + r")\b\s+(\d{4})\b", re.IGNORECASE
)
_TF_QUARTER_RE = re.compile(r"\b(q[1-4])\s*[/\-–]?\s*(\d{4})\b", re.IGNORECASE)
_TF_YEAR_RANGE_RE = re.compile(r"\b(20\d{2})\s*[-–/]\s*(20\d{2})\b")

# Regex to extract geographic phrase after a preposition — most reliable for natural language
_GEO_PREP_RE = re.compile(
    r"(?:in|across|throughout|within|covering|from)\s+([A-Z][A-Za-z\s,–\-]{3,60}?)"
    r"(?=\s+(?:between|from|during|in\s+20|\d{4}|q[1-4])"
    r"|\s+(?:received|trained|vaccinated|completed|achieved|reported|served|reached|was|were|has|have|had)"
    r"|[.,;]|$)",
    re.IGNORECASE,
)

# Known Ghana/West Africa city and region names for geo extraction.
# Order matters: most specific (district/region/city) entries are checked
# before country names, so "Northern Region, Ghana" matches the more
# specific "northern region" rather than just "ghana".
_GEO_PROPER_NAMES = [
    # Ghana regions (longest/most specific first)
    "greater accra region", "western north region", "savannah region",
    "north east region", "bono east region", "ahafo region", "oti region",
    "eastern region", "western region", "central region", "volta region",
    "northern region", "upper east region", "upper west region",
    "ashanti region", "bono region",
    # Ghana short forms
    "greater accra", "western north", "upper east", "upper west",
    "ashanti", "eastern", "western", "central",
    # Ghana cities
    "accra", "kumasi", "tamale", "takoradi", "tema", "sunyani",
    "cape coast", "koforidua", "bolgatanga", "wa", "ho", "techiman",
    # Nigeria
    "lagos", "abuja", "kano", "ibadan", "kaduna", "enugu", "port harcourt",
    # Other West Africa
    "nairobi", "kampala", "dar es salaam", "kigali", "harare",
    "freetown", "monrovia", "banjul", "conakry", "dakar",
    # Country names — checked last as the broadest fallback match
    "the gambia", "ghana", "nigeria", "kenya", "uganda", "tanzania", "rwanda",
    "zimbabwe", "sierra leone", "liberia", "gambia", "guinea", "senegal",
    "ivory coast", "côte d'ivoire", "burkina faso", "mali", "niger", "chad",
    "benin", "togo", "ethiopia", "zambia", "malawi", "mozambique", "south africa",
]


def _clean_geo_chunk(chunk: str) -> str:
    """Remove leading non-alpha chars, strip to first sentence/clause."""
    chunk = re.sub(r"^[^A-Za-z]*", "", chunk)
    chunk = re.split(r"[.;]|\bbetween\b|\bfrom\b", chunk)[0]
    return chunk.strip()


def _smart_extract_from_result(result_text: str, s: str) -> None:
    """Silently pre-fill empty Tab 0 fields from the result statement.
    Never overwrites a field the user has already filled."""
    if not result_text or len(result_text.split()) < 6:
        return
    rt = result_text.lower()

    # ── TIMEFRAME ──────────────────────────────────────────────────────────
    tf_key = f"timeframe{s}"
    if not _ss_str(tf_key).strip():
        tf_m = _TF_RANGE_RE.search(rt)
        if tf_m:
            m1 = tf_m.group(1).capitalize()
            m2 = tf_m.group(2).capitalize()
            yr = tf_m.group(3)
            st.session_state[tf_key] = f"{m1}–{m2} {yr}"
        else:
            q_m = _TF_QUARTER_RE.search(rt)
            if q_m:
                st.session_state[tf_key] = f"{q_m.group(1).upper()} {q_m.group(2)}"
            else:
                s_m = _TF_SINGLE_RE.search(rt)
                if s_m:
                    st.session_state[tf_key] = f"{s_m.group(1).capitalize()} {s_m.group(2)}"
                else:
                    yr_m = _TF_YEAR_RANGE_RE.search(rt)
                    if yr_m:
                        st.session_state[tf_key] = f"{yr_m.group(1)}–{yr_m.group(2)}"

    # ── GEOGRAPHIC SCOPE ───────────────────────────────────────────────────
    gs_key = f"geographic_scope{s}"
    if not _ss_str(gs_key).strip():
        # Strategy 1: preposition-based extraction ("in X Region", "across X and Y")
        # Find ALL matches, pick the longest (most descriptive)
        geo_candidates = _GEO_PREP_RE.findall(result_text)
        if geo_candidates:
            # Remove temporal noise, then filter candidates that start with real place names
            cleaned = []
            for cand in geo_candidates:
                cand = re.split(r"\s+(?:between|from|during)\b", cand, flags=re.IGNORECASE)[0].strip().rstrip(",;.")
                # Strip leading short non-place words (e.g. "IT", "a", "the")
                # Strip leading all-caps abbreviations (e.g. "IT "), articles, or prepositions
                # but NOT the first letter of a proper name (e.g. "Northern" starts with "N")
                cand = re.sub(r"^(?:[A-Z]{2,}\s+|[Aa]n?\s+|[Tt]he\s+|across\s+|in\s+|from\s+)", "", cand).strip()
                # Second pass: strip any remaining leading lowercase prepositions
                cand = re.sub(r"^(?:across|within|throughout|in|from)\s+", "", cand, flags=re.IGNORECASE).strip()
                if len(cand) > 4:
                    cleaned.append(cand)
            if cleaned:
                best = max(cleaned, key=len)
                st.session_state[gs_key] = best

        # Strategy 2: if not found, scan for known proper place names
        if not _ss_str(gs_key).strip():
            for place in _GEO_PROPER_NAMES:
                # Word-boundary match — plain substring matching lets short names
                # like "ho" or "wa" false-positive inside "cohorts" or "want".
                _place_m = re.search(r"\b" + re.escape(place) + r"\b", rt)
                if _place_m:
                    # Expand back to capture "X districts in Y Region" patterns
                    idx = _place_m.start()
                    raw = result_text[max(0, idx - 20):idx + len(place) + 15]
                    raw = re.sub(r"^[^A-Za-z\d]*", "", raw)
                    chunk = re.split(r"[.;,]|\bbetween\b|\bfrom\b|\bduring\b", raw, flags=re.IGNORECASE)[0].strip()
                    if len(chunk) > 4:
                        st.session_state[gs_key] = chunk.strip()
                        break

    # ── TARGET GROUP ────────────────────────────────────────────────────────
    tg_key = f"target_group{s}"
    if not _ss_str(tg_key).strip():
        for marker in _DEMO_MARKERS:
            if marker in rt:
                idx = rt.find(marker)
                # Capture 20 chars before (qualifier) + marker + 50 chars after
                raw = result_text[max(0, idx - 20):idx + len(marker) + 50]
                chunk = re.sub(r"^[^A-Za-z]*", "", raw)
                # Stop at geographic/temporal keywords
                chunk = re.split(r"\b(?:across|in|within|between|from|during)\b", chunk)[0]
                chunk = chunk.split(",")[0].split(".")[0]
                if len(chunk) > 4:
                    st.session_state[tg_key] = chunk.strip()
                    break

    # ── SECTOR (programme-level — keyword inference, never overwrites) ──────
    if st.session_state.get("sector") in (None, "", "(No sector selected)"):
        _SECTOR_KWORDS = [
            ("WASH",                        ["water", "borehole", "sanitation", "latrine", "wash", "hygiene", "handwash", "clean water", "water point"]),
            ("Health & Nutrition",          ["health", "clinic", "nutrition", "patient", "hospital", "immunis", "immuniz", "malaria", "vaccination", "antenatal", "maternal", "child health", "hiv", "tuberculosis"]),
            ("Education & Skills",          ["school", "pupil", "literacy", "enrolment", "enrollment", "education", "classroom", "teacher", "reading"]),
            ("Youth Employment & TVET",     ["tvet", "vocational", "apprentice", "placed", "hired", "youth employment", "livelihood", "income generation"]),
            ("Agriculture & Livelihoods",   ["farmer", "agriculture", "crop", "yield", "harvest", "livestock", "irrigation", "cassava", "maize", "cocoa", "poultry"]),
            ("Nutrition & Food Security",   ["stunting", "wasting", "malnutrition", "food security", "food insecurity", "micronutrient", "diet diversity", "acute malnutrition"]),
            ("Climate Resilience",          ["climate", "resilience", "adaptation", "disaster", "flood", "drought", "reforestation", "early warning"]),
            ("Governance & Accountability", ["governance", "accountability", "transparency", "civic", "citizen", "district assembly", "community accountability"]),
            ("Digital Economy & Technology",["digital", "mobile money", "fintech", "internet access", "ict ", "e-learning", "tech hub"]),
            ("Energy & Clean Energy",       ["solar", "energy", "electricity", "off-grid", "clean energy", "cookstove", "biogas"]),
            ("Gender & Social Inclusion",   ["gender equality", "gbv", "disability", "social inclusion", "women's rights"]),
            ("Private Sector Development",  ["private sector", "sme", "entrepreneur", "business development", "value chain", "market linkage"]),
        ]
        for _sk_name, _sk_kws in _SECTOR_KWORDS:
            if any(_kw in rt for _kw in _sk_kws):
                st.session_state["sector"] = _sk_name
                st.session_state["_sector_auto_inferred"] = True
                break

    # ── DONOR (programme-level — keyword inference, never overwrites) ───────
    if st.session_state.get("donor_selected") in (None, "", "(No donor specified)"):
        _DONOR_KWORDS = [
            ("USAID",                 ["usaid", "u.s. agency", "american people"]),
            ("FCDO",                  ["fcdo", "foreign commonwealth", "uk aid", "dfid"]),
            ("GIZ",                   ["giz", "deutsche gesellschaft", "german agency"]),
            ("World Bank",            ["world bank", "ibrd", "ida "]),
            ("Mastercard Foundation", ["mastercard foundation", "young africa works"]),
            ("AfDB",                  ["afdb", "african development bank"]),
            ("EU / EuropeAid",        ["european union", "eu grant", "europeaid", "dg intpa"]),
            ("KOICA",                 ["koica", "korea international"]),
            ("SIDA",                  ["sida", "swedish international development", "swedish development cooperation"]),
            ("SDC",                   ["sdc ", "swiss agency for development", "swiss development cooperation"]),
            ("RVO",                   ["rvo", "netherlands enterprise"]),
        ]
        for _dk_name, _dk_sigs in _DONOR_KWORDS:
            if any(_ds in rt for _ds in _dk_sigs):
                st.session_state["donor_selected"] = _dk_name
                st.session_state["_donor_auto_inferred"] = True
                if _dk_name in DONOR_PROFILES and st.session_state.get("donor_framework", "Generic") == "Generic":
                    st.session_state["donor_framework"] = _dk_name
                break


def _smart_extract_achievement(result_text: str, s: str) -> None:
    """Pre-fill logframe_achievement from result statement when the field is empty.
    Extracts the first quantified phrase (number + noun) stopping at location/time markers."""
    key = f"logframe_achievement{s}"
    if _ss_str(key).strip():
        return
    if not result_text or len(result_text.split()) < 5:
        return
    m = re.search(
        r"\b(\d[\d,\.%]*\s+\w[\w\s\-]{2,40}?)"
        r"(?=\s+(?:in|across|between|from|during|by|across|at|among)"
        r"|\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)"
        r"|[.,;]|$)",
        result_text, re.IGNORECASE,
    )
    if m:
        st.session_state[key] = m.group(1).strip()


def _render_tab1_slot(slot: int):
    s, _ph = _tab_slot_setup(slot)
    _render_fix_notes(slot, 0)
    _irc_widget(
        st.text_area, "Result statement *", f"result_statement{s}", default="",
        placeholder=_ph["result"],
        height=100,
        help="What did your project achieve? Include the action verb, number, target group, location, and timeframe.",
    )
    _rs = st.session_state.get(f"result_statement{s}", "")
    if _rs and len(_rs.strip()) >= 30:
        _smart_extract_from_result(_rs, s)   # auto-fill target group / timeframe / geo
        _smart_extract_achievement(_rs, s)    # auto-fill logframe actual achievement
        # Live Clarity signal — shows momentum from the very first tab
        try:
            _t0_sub = _build_submission_from_session(slot)
            _t0_ev  = _evaluator.evaluate_submission(_t0_sub)
            _t0_cl  = _t0_ev.get("clarity_score", 0)
            _t0_dots = "●" * int(_t0_cl) + "○" * (5 - int(_t0_cl))
            st.caption(f"Clarity so far: **{_t0_cl}/5.0** {_t0_dots} — fill timeframe, geography, and target group to raise this.")
        except Exception:
            pass
    if _rs and len(_rs.strip()) < 20:
        st.warning("Result statement is very short. Include: action verb + number + population + timeframe.")
    elif _rs and not any(c.isdigit() for c in _rs):
        st.caption("Tip: Add a number (e.g., '500 farmers trained') — quantified claims score higher.")

    # Result statement quality checklist — 5 inline indicators (council XXIV)
    if _rs and len(_rs.strip()) > 10:
        _rs_lo = _rs.lower()
        _change_verbs = (
            "trained", "reached", "vaccinated", "supported", "reduced", "increased",
            "improved", "constructed", "delivered", "provided", "established", "enrolled",
            "graduated", "employed", "received", "achieved", "made eligible", "completed",
        )
        _ql_checks = [
            ("Contains a number",       any(c.isdigit() for c in _rs)),
            ("Contains a target group", any(m in _rs_lo for m in _DEMO_MARKERS)),
            ("Contains a timeframe",    bool(re.search(r"\b(20\d{2}|q[1-4]|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec|quarter|annual)\b", _rs_lo))),
            ("Contains a location",     any(p in _rs_lo for p in _GEO_PROPER_NAMES) or bool(re.search(r"\b(region|district|province|state|country|national|community)\b", _rs_lo))),
            ("Contains a change verb",  any(v in _rs_lo for v in _change_verbs)),
        ]
        _all_pass = all(v for _, v in _ql_checks)
        if not _all_pass:
            with st.expander("Result statement quality check", expanded=True):
                for _label, _pass in _ql_checks:
                    st.markdown(f"{'✅' if _pass else '⬜'} {_label}")

    _irc_widget(
        st.text_input, "Target group *", f"target_group{s}", default="",
        placeholder=_ph["target_group"],
        help="Who specifically? Age group, gender, role, occupation. Avoid 'beneficiaries' alone — name the population.",
    )
    _irc_widget(
        st.text_input, "Timeframe *", f"timeframe{s}", default="",
        placeholder="e.g., January - June 2025",
        help="Specific dates or quarters. 'January–June 2025' is stronger than 'In 2025'.",
    )
    _irc_widget(
        st.text_input, "Geographic scope *", f"geographic_scope{s}", default="",
        placeholder=_ph["geographic_scope"],
        help="Districts, regions, or specific sites. 'Volta Region' beats 'rural areas'.",
    )
    _tg = _ss_str(f"target_group{s}").strip()
    if len(_tg) > 5 and not any(m in _tg.lower() for m in _DEMO_MARKERS):
        st.warning("Target group should describe who you reached — include population type, age, or role.")

    _tf = _ss_str(f"timeframe{s}").strip()
    if len(_tf) > 3 and not any(m in _tf.lower() for m in _DATE_MARKERS):
        st.warning("Timeframe should include a date range or period, e.g. January–June 2025.")

    _gs = _ss_str(f"geographic_scope{s}").strip()
    _gs_hint = _ph.get("geographic_scope", "")
    if len(_gs) > 5 and not any(m in _gs.lower() for m in _LOC_MARKERS):
        _gs_example = re.sub(r"^e\.g\.,\s*", "", _gs_hint)
        st.warning(
            "Geographic scope should name specific districts, regions, or locations "
            f"(e.g., {_gs_example})."
        )


def _render_tab2_slot(slot: int):
    s, _ph = _tab_slot_setup(slot)
    _render_fix_notes(slot, 1)
    st.caption("Link this result to your approved logframe indicator. Fills up to +1.0 on your score.")

    # Show result statement as read-only reference so user can reconcile without scrolling back
    _rs_ref = st.session_state.get(f"result_statement{s}", "").strip()
    if _rs_ref:
        st.info(f"**Your result:** {_rs_ref}")

    _fill_later = st.checkbox(
        "I don't have my logframe to hand — I'll fill this later",
        key=f"logframe_fill_later{s}",
        help="Tick this to continue without logframe data. You can return to complete it before scoring.",
    )
    if _fill_later:
        st.caption("Note: unfilled logframe fields reduce your Clarity score.")
    if not _fill_later:
        _lib_email = st.session_state.get("user_email", "")
        if _lib_email:
            with st.expander("📁 Logframe Library — load or save indicators", expanded=False):
                _libs = list_logframe_libraries(_lib_email)
                _lib_names = {lib["name"]: lib["id"] for lib in _libs}
                if _libs:
                    _chosen_lib_name = st.selectbox(
                        "Load from a saved library",
                        options=["Choose a library..."] + list(_lib_names.keys()),
                        key=f"_lf_lib_pick{s}",
                    )
                    if _chosen_lib_name != "Choose a library...":
                        _chosen_lib_id = _lib_names[_chosen_lib_name]
                        _items = get_library_items(_chosen_lib_id, _lib_email)
                        if _items:
                            _item_labels = {
                                (it.get("indicator_name") or it.get("logframe_indicator") or f"Item {i + 1}"): it
                                for i, it in enumerate(_items)
                            }
                            _chosen_item_label = st.selectbox(
                                "Indicator to load",
                                options=["Choose an indicator..."] + list(_item_labels.keys()),
                                key=f"_lf_item_pick{s}",
                            )
                            if _chosen_item_label != "Choose an indicator..." and st.button(
                                "Load into this result", key=f"_lf_lib_load_btn{s}"
                            ):
                                _it = _item_labels[_chosen_item_label]
                                st.session_state[f"logframe_indicator{s}"] = _it.get("logframe_indicator", "")
                                st.session_state[f"logframe_baseline{s}"] = _it.get("logframe_baseline", "")
                                st.session_state[f"logframe_target{s}"] = _it.get("logframe_target", "")
                                st.session_state[f"logframe_achievement{s}"] = _it.get("logframe_achievement", "")
                                st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                                st.rerun()
                        else:
                            st.caption("This library has no saved indicators yet.")
                else:
                    st.caption("No saved libraries yet — save one below. Manage libraries from My Audits.")

                st.divider()
                st.caption("Save this result's indicator to a library, to reuse in future audits:")
                _save_lib_options = ["+ New library..."] + list(_lib_names.keys())
                _save_lib_choice = st.selectbox("Save to", options=_save_lib_options, key=f"_lf_save_lib_pick{s}")
                _new_lib_name = ""
                if _save_lib_choice == "+ New library...":
                    _new_lib_name = st.text_input(
                        "New library name", key=f"_lf_new_lib_name{s}",
                        placeholder="e.g., USAID WASH Program 2025",
                    )
                if st.button("Save this indicator to library", key=f"_lf_save_item_btn{s}"):
                    _cur_indicator = st.session_state.get(f"logframe_indicator{s}", "").strip()
                    if not _cur_indicator:
                        st.warning("Fill in a logframe indicator above before saving it to a library.")
                    else:
                        _target_lib_id = _lib_names.get(_save_lib_choice)
                        if _save_lib_choice == "+ New library...":
                            if _new_lib_name.strip():
                                _target_lib_id = create_logframe_library(_lib_email, _new_lib_name.strip())
                            else:
                                st.warning("Enter a name for the new library.")
                                _target_lib_id = None
                        if _target_lib_id and not _safe_rate_limit_ok(
                            _lib_email, "add_library_items", max_count=30, window_seconds=3600
                        ):
                            st.warning("You've saved a lot of indicators in the last hour — please wait a bit before adding more.")
                        elif _target_lib_id:
                            add_library_items(_target_lib_id, _lib_email, [{
                                "indicator_name": _cur_indicator,
                                "logframe_indicator": _cur_indicator,
                                "logframe_baseline": st.session_state.get(f"logframe_baseline{s}", ""),
                                "logframe_target": st.session_state.get(f"logframe_target{s}", ""),
                                "logframe_achievement": st.session_state.get(f"logframe_achievement{s}", ""),
                                "sector": st.session_state.get("sector", ""),
                            }])
                            _lib_save_err = last_audit_error()
                            if _lib_save_err:
                                st.warning(f"Could not save to your Logframe Library. ({_lib_save_err})")
                            else:
                                st.success("Saved to your Logframe Library.")

        _lf_api_key = (
            st.secrets.get("ANTHROPIC_API_KEY", "")
            if hasattr(st, "secrets") else
            os.environ.get("ANTHROPIC_API_KEY", "")
        )
        if _lf_api_key:
            with st.expander("🎯 AI Logframe Match — paste your indicators, get a suggested match", expanded=False):
                st.caption(
                    "Paste your approved logframe indicators (one per line). The AI suggests which "
                    "one this result reports against — it never forces a match, and it only quotes "
                    "words already in your result statement or your pasted indicators."
                )
                st.text_area(
                    "Your logframe indicators (one per line)",
                    key=f"_lf_paste{s}", height=100,
                    placeholder=(
                        "Indicator 1.2: Number of households with access to safe water\n"
                        "Indicator 2.1: % of farmers applying climate-smart practices"
                    ),
                )
                if st.button("Match my result to an indicator", key=f"lf_match_btn{s}"):
                    _lf_rs = st.session_state.get(f"result_statement{s}", "")
                    _lf_raw = st.session_state.get(f"_lf_paste{s}", "")
                    _lf_indicators = [ln.strip() for ln in _lf_raw.splitlines() if ln.strip()]
                    if not _lf_rs.strip():
                        st.warning("Enter a result statement in Tab 1 first.")
                    elif not _lf_indicators:
                        st.warning("Paste at least one logframe indicator above.")
                    else:
                        with st.spinner("Matching your result to an indicator…"):
                            from council import match_logframe_indicator
                            _lf_match = match_logframe_indicator(_lf_rs, _lf_indicators, _lf_api_key)
                        st.session_state[f"_lf_match_result{s}"] = _lf_match
                        if _lf_match.get("confidence_label") != "None" and _lf_match.get("best_match"):
                            st.session_state[f"logframe_indicator{s}"] = _lf_match["best_match"]
                            st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                        st.rerun()

                _lf_result = st.session_state.get(f"_lf_match_result{s}")
                if _lf_result:
                    _lf_cl = _lf_result.get("confidence_label", "None")
                    if _lf_cl == "None" or not _lf_result.get("best_match"):
                        st.info("No confident match found — enter the indicator manually below.")
                    else:
                        st.success(
                            f"AI-suggested match ({_lf_cl.lower()} confidence) — confirm against "
                            f"your approved logframe before submitting."
                        )
                        if _lf_result.get("justification"):
                            st.caption(_lf_result["justification"])

        _irc_widget(
            st.text_input,
            "Logframe indicator",
            f"logframe_indicator{s}", default="",
            placeholder=_ph.get("logframe_indicator", "e.g., Indicator 1.2: Number of [target group] achieving [outcome]"),
            help=(
                "Copy the exact indicator name and code from your approved Technical Proposal or logframe. "
                "If you cannot quote it, your donor cannot match your result to your commitment."
            ),
        )
        _irc_widget(
            st.text_input,
            "Pre-evaluation / baseline value",
            f"logframe_baseline{s}", default="",
            placeholder=_ph.get("logframe_baseline", "e.g., [baseline value] ([year] [data source] baseline)"),
            help=(
                "The value of this indicator before the programme began (or at the last measurement point). "
                "Required to compute % change from baseline and to validate the direction of change."
            ),
        )
        _irc_widget(
            st.text_input,
            "Approved target",
            f"logframe_target{s}", default="",
            placeholder=_ph.get("logframe_target", "e.g., 60% increase by Q4 2025"),
            help=(
                "The target as approved in the original Technical Proposal. Donors compare achievements "
                "against approved targets — not revised internal targets."
            ),
        )
        _data_forthcoming = st.checkbox(
            "Data not yet available for this indicator",
            key=f"logframe_data_forthcoming{s}",
            help=(
                "Tick if measurement has not yet been collected. The gap will be disclosed in your "
                "report rather than penalised — donors prefer honest disclosure to blank fields."
            ),
        )
        if not _data_forthcoming:
            _irc_widget(
                st.text_input,
                "Actual achievement",
                f"logframe_achievement{s}", default="",
                placeholder=_ph.get("logframe_achievement", "e.g., [achieved value] by [date] — [X]% of target"),
                help=(
                    "The actual delivered number, ideally with % achievement vs original target. "
                    "Must reconcile with your result statement above."
                ),
            )
        else:
            st.caption("✓ Noted — no achievement figure needed for now.")


_MONTH_NUM = {
    "january": 1, "february": 2, "march": 3, "april": 4, "may": 5, "june": 6,
    "july": 7, "august": 8, "september": 9, "october": 10, "november": 11, "december": 12,
    "jan": 1, "feb": 2, "mar": 3, "apr": 4, "jun": 6, "jul": 7, "aug": 8,
    "sep": 9, "sept": 9, "oct": 10, "nov": 11, "dec": 12,
}


def _derive_reporting_period(timeframe_text: str):
    """Parse the canonical timeframe formats _smart_extract_from_result() produces
    ("Month–Month YYYY", "QN YYYY", "Month YYYY", "YYYY–YYYY") into (start, end)
    dates, so the reporting-period pickers don't default to today for no reason
    when the timeframe the user already typed already answers the question.
    Returns (None, None) if the text doesn't match a known format."""
    import calendar
    if not timeframe_text:
        return None, None
    t = timeframe_text.strip()

    m = re.match(r"([A-Za-z]+)[–\-]([A-Za-z]+)\s+(\d{4})$", t)
    if m:
        m1, m2, yr = m.group(1).lower(), m.group(2).lower(), int(m.group(3))
        n1, n2 = _MONTH_NUM.get(m1), _MONTH_NUM.get(m2)
        if n1 and n2:
            return date(yr, n1, 1), date(yr, n2, calendar.monthrange(yr, n2)[1])

    m = re.match(r"[Qq]([1-4])\s+(\d{4})$", t)
    if m:
        q, yr = int(m.group(1)), int(m.group(2))
        n1 = (q - 1) * 3 + 1
        n2 = n1 + 2
        return date(yr, n1, 1), date(yr, n2, calendar.monthrange(yr, n2)[1])

    m = re.match(r"([A-Za-z]+)\s+(\d{4})$", t)
    if m:
        mo, yr = m.group(1).lower(), int(m.group(2))
        n = _MONTH_NUM.get(mo)
        if n:
            return date(yr, n, 1), date(yr, n, calendar.monthrange(yr, n)[1])

    m = re.match(r"(\d{4})[–\-](\d{4})$", t)
    if m:
        y1, y2 = int(m.group(1)), int(m.group(2))
        return date(y1, 1, 1), date(y2, 12, 31)

    return None, None


def _render_tab3_slot(slot: int):
    s, _ph = _tab_slot_setup(slot)
    _render_fix_notes(slot, 2)

    # ── CORE (4 fields — always visible) ────────────────────────────────────
    _irc_widget(
        st.text_area, "Describe your supporting evidence", f"evidence_description{s}", default="",
        placeholder=_ph["evidence_description"],
        height=120,
        help="Describe the actual document or data: who collected it, how, and what's in it.",
    )
    _ed_val = st.session_state.get(f"evidence_description{s}", "")
    if _ed_val and len(_ed_val.strip()) < 30:
        st.warning("Evidence description is brief. Specify: who collected it, how, and what it contains.")
    elif _ed_val:
        _smart_extract_ev_type(_ed_val, f"evidence_type{s}")  # auto-suggest type from keywords
    if _ed_val and _ed_val.strip():
        _wc = len(_ed_val.split())
        if 0 < _wc < 20:
            st.caption(f"📝 {_wc} words — add {20 - _wc}+ more for a stronger Confidence score.")

    st.caption("📊 **Affects Directness score** (max 2.0 — the system uses this to determine your evidence ceiling. Systematic Reviews and RCTs unlock the highest Directness score.)")
    _irc_widget(
        st.radio, "Primary evidence type (select your strongest source)", f"evidence_type{s}", default=EVIDENCE_TYPES[0],
        options=EVIDENCE_TYPES,  # placeholder included as a real option so "not yet chosen" is a detectable state
        help=EVIDENCE_TYPE_HELP,
    )
    ev_type = st.session_state.get(f"evidence_type{s}", EVIDENCE_TYPES[0])
    ev_desc = st.session_state.get(f"evidence_description{s}", "")
    _dl = _evaluator.get_directness_level(ev_type, ev_desc)
    _ds = round((_dl / 5) * 2.0, 1)

    if ev_type == "Other":
        st.text_input("Specify evidence type", key=f"evidence_type_other{s}")

    # Council XXIII — evidence type debate (re-evaluate button + reasoning expander)
    _ev_type_debate_key = f"_ev_type_debate{s}"
    _ev_debate_cols = st.columns([3, 1])
    with _ev_debate_cols[1]:
        if st.button("🏛 Re-evaluate type", key=f"reval_ev_type{s}",
                     help="Run a 5-member council debate to verify this is the closest-fit evidence type",
                     use_container_width=True):
            _rev_api_key = (
                st.secrets.get("ANTHROPIC_API_KEY", "")
                if hasattr(st, "secrets") else
                os.environ.get("ANTHROPIC_API_KEY", "")
            )
            if not _rev_api_key:
                st.warning("Evidence type debate is not available — API key not configured.")
            elif not ev_desc or len(ev_desc.strip()) < 10:
                st.warning("Add an evidence description first.")
            else:
                with st.spinner("Council debating evidence type…"):
                    from council import debate_evidence_type
                    _debate = debate_evidence_type(
                        description=ev_desc,
                        result_statement=st.session_state.get(f"result_statement{s}", ""),
                        evidence_types=EVIDENCE_TYPES[1:],
                        api_key=_rev_api_key,
                    )
                if _debate.get("recommended_type"):
                    st.session_state[f"evidence_type{s}"] = _debate["recommended_type"]
                    st.session_state[_ev_type_debate_key] = _debate
                    st.rerun()
                else:
                    st.warning("Council could not reach a confident recommendation — type left unchanged.")

    _ev_debate = st.session_state.get(_ev_type_debate_key) or st.session_state.get("_ev_type_debate")
    if _ev_debate and _ev_debate.get("recommended_type"):
        with st.expander("Why this type? — Council reasoning", expanded=False):
            st.caption(
                f"**Recommended:** {_ev_debate['recommended_type']} "
                f"({_ev_debate.get('confidence', 'medium')} confidence)"
            )
            if _ev_debate.get("reasoning"):
                st.caption(_ev_debate["reasoning"])
            if _ev_debate.get("donor_alignment"):
                st.caption(f"**Donor alignment:** {_ev_debate['donor_alignment']}")
            _votes = _ev_debate.get("member_votes", {})
            if _votes:
                st.markdown("**Council votes:**")
                for _mid, _v in _votes.items():
                    _mname = _mid.replace("_", " ").title()
                    st.markdown(f"- **{_mname}:** {_v.get('vote', '—')} — {_v.get('reasoning', '')}")

    _irc_widget(
        st.text_input, "Who verified this?", f"verifier{s}", default="",
        placeholder=_ph.get("verifier", "e.g., District Agriculture Officer, partner org M&E lead, external evaluator"),
        help="The person or organization that confirmed the data is accurate.",
    )

    st.caption("📊 **Affects Recency score** (max 1.0/1.0) — evidence >6 months from your reporting period end incurs a penalty")
    _irc_widget(st.date_input, "When was this evidence collected?", f"evidence_date{s}", default=date.today())
    _ed = st.session_state.get(f"evidence_date{s}")
    if _ed and hasattr(_evaluator, "get_recency_diagnostic"):
        _rec_diag = _evaluator.get_recency_diagnostic(_ed)
        if "0.4/1.0" in _rec_diag or "0.2/1.0" in _rec_diag:
            st.warning(_rec_diag)
        elif "0.6/1.0" in _rec_diag:
            st.info(_rec_diag)
        else:
            st.success(_rec_diag)

    # Determine whether to auto-open the detail expander (PII/safeguarding triggers)
    _ev_type_now = st.session_state.get(f"evidence_type{s}", "")
    _pii_triggered = _ev_type_now in PII_EVIDENCE_TYPES
    _safeguarding_triggered = _ev_type_now in SAFEGUARDING_EVIDENCE_TYPES
    _minors_triggered = _minors_possibly_involved(slot)
    _compliance_needed = _pii_triggered or _safeguarding_triggered or _minors_triggered

    # ── ADVANCED DETAILS (auto-opens for PII/safeguarding, or until Verification —
    # the single biggest lever in Confidence — has been answered at least once) ──
    _verification_unanswered = (
        st.session_state.get(f"internal_review{s}", "Choose an option...") == "Choose an option..."
        or st.session_state.get(f"external_review{s}", "Choose an option...") == "Choose an option..."
    )
    with st.expander(
        "⚙️ Improve your score — Advanced details (up to +2.0 points)" + (" — ⚠️ compliance required" if _compliance_needed else ""),
        expanded=(_compliance_needed or _verification_unanswered),
    ):
        st.caption("Filling these fields can improve your Confidence score by up to +2.0 points. Leave them blank to skip — you can always come back.")
        # Qualitative evidence checkbox + quality checks
        is_qualitative_evidence = (
            ev_type in QUALITATIVE_EVIDENCE_TYPES
            or st.session_state.get(f"qualitative_evidence{s}", False)
        )
        _irc_widget(
            st.checkbox,
            "Evidence is qualitative (case study, MSC, outcome harvesting)",
            f"qualitative_evidence{s}", default=False,
            help="Switches Definition and Measurement scoring to qualitative dimensions.",
        )
        is_qualitative_evidence = (
            ev_type in QUALITATIVE_EVIDENCE_TYPES
            or st.session_state.get(f"qualitative_evidence{s}", False)
        )
        _sub_checks = _EV_QUALITY_CHECKS.get(ev_type, [])
        if _sub_checks:
            st.markdown("**Quality checks**")
            for _ck_key, _ck_lbl in _sub_checks:
                st.checkbox(_ck_lbl, key=f"{_ck_key}{s}")
        elif is_qualitative_evidence:
            st.markdown("**Qualitative Rigor**")
            (
                _sourcing_lbl, _triangulated_lbl, _bias_lbl,
                _voice_lbl, _consent_lbl,
            ) = QUAL_RIGOR_CHECKLIST.get(ev_type, QUAL_RIGOR_CHECKLIST["Case study"])
            _irc_widget(st.checkbox, _sourcing_lbl,      f"qual_sourcing{s}",     default=False)
            _irc_widget(st.checkbox, _triangulated_lbl,  f"qual_triangulated{s}", default=False)
            _irc_widget(st.checkbox, _bias_lbl,          f"qual_bias{s}",         default=False)
            _irc_widget(st.checkbox, _voice_lbl,         f"qual_voice{s}",        default=False)
            _irc_widget(st.checkbox, _consent_lbl,       f"qual_consent{s}",      default=False)

        st.divider()
        st.markdown("**Internal & External Review**")
        st.caption("📊 **Affects Verification score** (max 2.0/2.0) — independent review is the biggest single lever in Confidence")
        # Dynamic review options based on org_type
        _rv_org = st.session_state.get("org_type", "International NGO (INGO)")
        _rv_community = "CBO" in _rv_org or "Government" in _rv_org
        _rv_national  = "National" in _rv_org
        if _rv_community or _rv_national:
            _INT_REVIEW_OPTS = [
                "Choose an option...",
                "Reviewed by Executive Director / Board",
                "Reviewed by community governance committee",
                "Reviewed by programme staff (no dedicated MEL)",
                "Not reviewed",
            ]
            _EXT_REVIEW_OPTS = [
                "Choose an option...",
                "Verified by ward / district committee",
                "Verified by community elder council",
                "Verified by peer organisation",
                "No external review",
            ]
        else:
            _INT_REVIEW_OPTS = INTERNAL_REVIEW_OPTIONS
            _EXT_REVIEW_OPTS = EXTERNAL_REVIEW_OPTIONS

        int_rev = st.session_state.get(f"internal_review{s}", _INT_REVIEW_OPTS[0])
        _irc_widget(
            st.selectbox, "Internal review", f"internal_review{s}", default=_INT_REVIEW_OPTS[0],
            options=_INT_REVIEW_OPTS,
            help="Did anyone in your organization review or cross-check this data?",
        )
        int_rev = st.session_state.get(f"internal_review{s}", _INT_REVIEW_OPTS[0])
        _int_vl = _evaluator.get_verification_level(int_rev, "No external review", "")
        _int_vs = round((_int_vl / 5) * 2.0, 1)
        if _int_vs == 0:
            st.warning("No internal review — adding a reviewer strengthens Verification.")
        if int_rev == "Other":
            st.text_input("Specify internal reviewer", key=f"internal_review_other{s}")

        _irc_widget(
            st.selectbox, "External review", f"external_review{s}", default=_EXT_REVIEW_OPTS[0],
            options=_EXT_REVIEW_OPTS,
            help="Did an outside party verify the data? Government, partner, auditor, or evaluator.",
        )
        ext_rev = st.session_state.get(f"external_review{s}", _EXT_REVIEW_OPTS[0])
        verifier_text = st.session_state.get(f"verifier{s}", "")
        _full_vl = _evaluator.get_verification_level(int_rev, ext_rev, verifier_text)
        _full_vs = round((_full_vl / 5) * 2.0, 1)
        if ext_rev == "No external review":
            st.warning("No external review — independent verification raises your score significantly.")
        if ext_rev == "Other":
            st.text_input("Specify external reviewer", key=f"external_review_other{s}")

        st.divider()
        st.markdown("**Reporting Period**")
        st.caption("The period this submission covers. Evidence outside this range is flagged.")
        _rp_derived_start, _rp_derived_end = _derive_reporting_period(st.session_state.get(f"timeframe{s}", ""))
        _rp_col1, _rp_col2 = st.columns(2)
        with _rp_col1:
            _irc_widget(st.date_input, "Period start", f"reporting_start{s}",
                       default=_rp_derived_start or date.today())
        with _rp_col2:
            _irc_widget(st.date_input, "Period end",   f"reporting_end{s}",
                       default=_rp_derived_end or date.today())
        _rp_s = st.session_state.get(f"reporting_start{s}")
        _rp_e = st.session_state.get(f"reporting_end{s}")
        if _ed and _rp_s and _rp_e and hasattr(_evaluator, "validate_reporting_period"):
            _, _rp_msg, _rp_sev = _evaluator.validate_reporting_period(_ed, _rp_s, _rp_e)
            if _rp_sev == "ERROR":
                st.error(_rp_msg)
            elif _rp_sev == "WARNING":
                st.warning(_rp_msg)
            elif _rp_msg:
                st.success(_rp_msg)

        st.divider()
        st.markdown("**Data Collection & Provenance**")
        st.caption("📊 **Affects Verification score** — each 'Yes' adds up to +0.1 on Confidence. 'Not applicable' is neutral where it honestly doesn't apply.")
        _ev_type_prov = st.session_state.get(f"evidence_type{s}", "")
        _prov_keys = _PROVENANCE_FOR_EV_TYPE.get(_ev_type_prov, _PROVENANCE_ALL)
        # Auto-set non-applicable provenance questions to "Not applicable" to avoid
        # silent -0.03 score penalties from placeholder "Choose an option..." defaults
        for _na_pk, _na_ss in {
            "sampling_documented":     f"provenance_sampling{s}",
            "double_counting_checked": f"provenance_dedup{s}",
            "collection_tool_named":   f"provenance_tool{s}",
            "collector_independent":   f"provenance_independence{s}",
            "recall_period_ok":        f"provenance_recall{s}",
        }.items():
            if _na_pk not in _prov_keys and st.session_state.get(_na_ss, PROVENANCE_YES_NO_NA_OPTIONS[0]) == PROVENANCE_YES_NO_NA_OPTIONS[0]:
                st.session_state[_na_ss] = "Not applicable"
        _prov_key_map = {
            "sampling_documented":     f"provenance_sampling{s}",
            "double_counting_checked": f"provenance_dedup{s}",
            "collection_tool_named":   f"provenance_tool{s}",
            "collector_independent":   f"provenance_independence{s}",
            "recall_period_ok":        f"provenance_recall{s}",
        }
        for _pk in _prov_keys:
            _pk_ss = _prov_key_map.get(_pk)
            if _pk_ss:
                _smart = _PROVENANCE_DEFAULTS.get((_ev_type_prov, _pk), PROVENANCE_YES_NO_NA_OPTIONS[0])
                _existing = st.session_state.get(_pk_ss, PROVENANCE_YES_NO_NA_OPTIONS[0])
                _prov_default = _existing if _existing != PROVENANCE_YES_NO_NA_OPTIONS[0] else _smart
                if _prov_default != PROVENANCE_YES_NO_NA_OPTIONS[0] and _existing == PROVENANCE_YES_NO_NA_OPTIONS[0]:
                    st.session_state[_pk_ss] = _prov_default
                _irc_widget(
                    st.selectbox, _PROVENANCE_LABELS[_pk],
                    _pk_ss, default=_prov_default,
                    options=PROVENANCE_YES_NO_NA_OPTIONS,
                )
        _irc_widget(
            st.selectbox, "Could an external auditor retrieve the original records?",
            f"provenance_traceability{s}", default=TRACEABILITY_OPTIONS[0],
            options=TRACEABILITY_OPTIONS,
            help="E.g. raw survey exports, signed registers, payment records.",
        )
        _prov_checklist = {
            "sampling_documented":     st.session_state.get(f"provenance_sampling{s}", PROVENANCE_YES_NO_NA_OPTIONS[0]),
            "double_counting_checked": st.session_state.get(f"provenance_dedup{s}", PROVENANCE_YES_NO_NA_OPTIONS[0]),
            "collection_tool_named":   st.session_state.get(f"provenance_tool{s}", PROVENANCE_YES_NO_NA_OPTIONS[0]),
            "collector_independent":   st.session_state.get(f"provenance_independence{s}", PROVENANCE_YES_NO_NA_OPTIONS[0]),
            "recall_period_ok":        st.session_state.get(f"provenance_recall{s}", PROVENANCE_YES_NO_NA_OPTIONS[0]),
            "auditor_traceable":       st.session_state.get(f"provenance_traceability{s}", TRACEABILITY_OPTIONS[0]),
        }

        st.divider()
        st.markdown("**Beneficiary Voice**")
        st.caption("Score by method — No voice: +0.0 · Anecdotal: +0.15 · Representatives/Systematic: +0.35 · Independent: +0.5")
        st.selectbox(
            "How were beneficiary voices captured?",
            key=f"beneficiary_voice{s}",
            options=_BV_OPTIONS,
            help="The strongest evidence includes beneficiary perspectives, not just provider reports.",
        )
        _bv_val = st.session_state.get(f"beneficiary_voice{s}", "")
        _BV_HIGH_2 = {
            "Direct beneficiary feedback collected (e.g., Lean Data survey, focus groups, NPS)",
            "Beneficiary representatives consulted (community leaders, beneficiary committees)",
        }
        if _bv_val in _BV_HIGH_2:
            st.text_input(
                "Briefly describe the method — when conducted and approximately how many people participated",
                key=f"bv_method_detail{s}",
                placeholder="e.g., Phone survey with 120 farmers, March 2025",
                help="Required to receive the full beneficiary voice bonus (≥20 characters).",
            )
        _bv_detail = st.session_state.get(f"bv_method_detail{s}", "")
        _bv_score = (_evaluator.compute_beneficiary_voice_bonus(_bv_val, _bv_detail)
                     if hasattr(_evaluator, "compute_beneficiary_voice_bonus") else 0.0)
        if _bv_val and _bv_val not in ("No beneficiary voice captured", "Choose an option..."):
            if _bv_val in _BV_HIGH_2:
                _bv_chars = len(_bv_detail.strip())
                if _bv_chars >= 20:
                    st.caption(f"✓ Beneficiary Voice bonus: **+{_bv_score}/0.5** — full bonus unlocked.")
                else:
                    st.caption(f"Beneficiary Voice bonus: **+{_bv_score}/0.5** · Add {20 - _bv_chars} more characters to unlock +0.5.")
            else:
                st.caption(f"Beneficiary Voice bonus: **+{_bv_score}/0.5**")

        st.divider()
        st.text_area(
            "Who owns this result, and what decision will it inform? (optional — improves Clarity)",
            key=f"additional_context{s}",
            placeholder="e.g., The MEL Lead owns this result. It will inform the Q3 budget reallocation.",
            help="Naming an owner and the decision this result informs strengthens your Governance sub-score.",
        )

        if _compliance_needed:
            st.divider()
            st.markdown("**🛡️ Compliance & Ethics** — your evidence type may involve personal data")

    # --- GOVERNANCE & COMPLIANCE LAYER — shown only when triggered ---
    if _compliance_needed:
        st.caption(
            "Required if your evidence involves personal data. "
            "Flags exposure under Ghana's Act 843 / Nigeria's NDPA."
        )
        if _pii_triggered:
            st.warning(
                "⚠️ **PII Alert:** One or more of your selected evidence types may "
                "contain Personally Identifiable Information (PII). Please answer the "
                "compliance checks below before proceeding."
            )
        if _safeguarding_triggered:
            st.warning(
                "⚠️ **Safeguarding Alert:** This evidence type may include beneficiary "
                "stories, photos, or testimony. Please confirm consent, a do-no-harm "
                "review, and secure handling below before proceeding."
            )
        if _minors_triggered:
            st.error(
                "🔴 **Child Safeguarding Alert:** This result statement or target group "
                "may involve minors. Please complete the child safeguarding check below "
                "(Core Humanitarian Standard, Commitment 4 — Keeping Children Safe)."
            )
        with st.expander(
            "📋 Data Governance Checklist",
            expanded=_compliance_needed,
        ):
            st.caption("⚠ Unanswered questions score as 'No' and reduce your Governance sub-score. Answer 'Not applicable' where it genuinely doesn't apply.")
            st.selectbox(
                "Do you have documented consent from beneficiaries for their data "
                "to be shared with the donor?",
                options=[
                    "Select consent status...",
                    "Yes — written consent forms on file",
                    "Yes — verbal consent documented",
                    "Partial — some beneficiaries consented",
                    "No — consent not obtained",
                    "Not applicable (no personal data)",
                ],
                key=f"gov_consent_status{s}",
            )
            st.selectbox(
                "Has this evidence been anonymized or de-identified where required?",
                options=[
                    "Select anonymization status...",
                    "Yes — fully anonymized",
                    "Partially anonymized",
                    "No — not anonymized",
                    "Not applicable",
                ],
                key=f"gov_anonymization_status{s}",
            )
            st.selectbox(
                "Does your evidence collection method comply with the data protection "
                "law in your project's country?",
                options=[
                    "Select compliance status...",
                    "Yes — compliant (e.g. Ghana Act 843, Nigeria NDPA, Kenya DPA)",
                    "Unsure — we haven't checked",
                    "No — we are not compliant",
                    "Not applicable",
                ],
                key=f"gov_compliance_law_status{s}",
            )
            st.selectbox(
                "If this evidence includes beneficiary stories, photos, or testimony, "
                "has it been reviewed for do-no-harm risks (identification, stigma, "
                "retraumatization, or safety) before sharing with the donor?",
                options=[
                    "Select safeguarding status...",
                    "Yes — reviewed, no concerns identified",
                    "Yes — reviewed, identifying details removed",
                    "Partial — some content not yet reviewed",
                    "No — not yet reviewed",
                    "Not applicable (no beneficiary stories/photos)",
                ],
                key=f"gov_safeguarding_status{s}",
            )
            st.selectbox(
                "If minors may be involved in this evidence, has a child safeguarding "
                "review been completed (organisational child safeguarding policy applied, "
                "and guardian consent obtained where applicable)?",
                options=[
                    "Select child safeguarding status...",
                    "Yes — child safeguarding policy applied, guardian consent obtained where applicable",
                    "Partial — some steps taken but not complete",
                    "No — not yet reviewed",
                    "Not applicable (no minors involved)",
                ],
                key=f"gov_child_safeguarding_status{s}",
            )
            st.selectbox(
                "If this evidence includes identifiable testimony, photos, or stories, "
                "is it stored securely with access limited to authorised staff "
                "(e.g. password-protected files, restricted folders)?",
                options=[
                    "Select secure handling status...",
                    "Yes — stored securely, access restricted to authorised staff",
                    "Partial — some identifiable material not yet secured",
                    "No — not yet secured",
                    "Not applicable (no identifiable testimony/photos)",
                ],
                key=f"gov_secure_handling_status{s}",
            )
        st.markdown(
            "#### 📁 Upload Organisational Data Protection Policy "
            "(optional — earns Governance Bonus)"
        )
        _dpp_file = st.file_uploader(
            "Upload your data protection policy (PDF or DOCX)",
            type=["pdf", "docx"],
            key=f"gov_dpp_upload{s}",
        )
        if _dpp_file is not None:
            st.session_state["gov_dpp_uploaded"] = True
            st.caption(
                "✅ Policy uploaded. **+5 Governance Bonus** applied to your "
                "Confidence Score for this session."
            )
        else:
            st.caption(
                "Uploading your policy grants a +5 Governance Bonus to your "
                "Confidence Score for this session."
            )
    # --- END GOVERNANCE & COMPLIANCE LAYER ---

    prev_files = st.session_state.get(f"draft_uploaded_filenames{s}", [])
    if prev_files:
        st.caption(
            f"For security, browsers don't let us keep re-uploaded files between sessions — "
            f"please re-attach: {', '.join(prev_files)}"
        )
    st.file_uploader(
        "Attach supporting documents (optional)", key=f"uploaded_files_widget{s}",
        accept_multiple_files=True,
        type=["pdf", "docx", "xlsx", "csv", "jpg", "jpeg", "png", "txt"],
        help="Attach raw evidence files — datasets, signed sheets, photos with metadata, partner letters.",
    )


# ---------------------------------------------------------------------------
# Screen 0 — Landing & Onboarding
# ---------------------------------------------------------------------------

# Pre-filled Ghana health NGO example — used by the "Try with a sample result" button.
# Values must be valid entries from the relevant option lists (EVIDENCE_TYPES,
# INTERNAL_REVIEW_OPTIONS, EXTERNAL_REVIEW_OPTIONS, _BV_OPTIONS, SECTOR_OPTIONS).
_DEMO_SUBMISSION = {
    # _BASE_FORM_KEYS text fields (slot 1, no suffix)
    "result_statement":    (
        "450 women of reproductive age in Ashanti Region received skilled antenatal care across "
        "3 health facilities, resulting in a 34% reduction in preventable maternal complications "
        "between January and December 2024 compared to the 2022 baseline."
    ),
    "target_group":        "Women of reproductive age (15–49) in peri-urban Kumasi",
    "timeframe":           "January – December 2024",
    "geographic_scope":    "Ashanti Region, Ghana (3 health facilities)",
    "evidence_description":(
        "Monthly DHIS2 data extracted from 3 facility registers by the district health "
        "information officer and verified against paper attendance registers by the project "
        "MEL officer."
    ),
    "logframe_indicator":  "% reduction in preventable maternal complications among ANC attendees",
    "logframe_target":     "30% reduction by December 2024",
    "logframe_achievement":"34% reduction (450 women served across 3 facilities)",
    "verifier":            "District Health Information Officer (independent of programme delivery)",
}
# Selectbox fields set separately (must match their option lists exactly)
_DEMO_SELECT_FIELDS = {
    "evidence_type":    "Raw datasets or survey exports",
    "internal_review":  "Reviewed by MEL Officer",
    "external_review":  "External partner review",
    "beneficiary_voice":"Anecdotal beneficiary quotes only (uncollected, not systematic)",
    "sector":           "Health & Nutrition",
    "donor_selected":   "USAID",
    "donor_framework":  "USAID",
}

# Second demo scenario — Agriculture & Livelihoods / GIZ (Council XXXII: sector diversity)
_DEMO_SUBMISSION_AGRIC = {
    "result_statement":    (
        "1,840 smallholder farmers in Upper East Region adopted improved maize varieties "
        "following ISFM training, achieving an average yield increase of 28% (from 1.2 to "
        "1.54 MT/ha) in the 2024 main season compared to the 2022 pre-programme baseline."
    ),
    "target_group":        "Smallholder maize farmers (≥50% women), Upper East Region",
    "timeframe":           "2024 main season (April – September 2024)",
    "geographic_scope":    "Upper East Region, Ghana (12 communities, 3 districts)",
    "evidence_description":(
        "Post-harvest yield assessment conducted by CSIR-SARI agronomists across 6 randomly "
        "selected communities. Yield measured from 5 sample plots per farmer household using "
        "the crop-cut method. Data cross-checked against GIZ results matrix records."
    ),
    "logframe_indicator":  "Average maize yield (MT/ha) among trained ISFM farmers",
    "logframe_target":     "25% yield increase by end of 2024 main season",
    "logframe_achievement":"28% yield increase — 1,840 farmers, avg 1.54 MT/ha vs. 1.2 MT/ha baseline",
    "verifier":            "CSIR-SARI agronomist (independent of GIZ programme delivery team)",
}
_DEMO_SELECT_FIELDS_AGRIC = {
    "evidence_type":    "Systematic observation or administrative data",
    "internal_review":  "Reviewed by MEL Officer",
    "external_review":  "External partner review",
    "beneficiary_voice":"Anecdotal beneficiary quotes only (uncollected, not systematic)",
    "sector":           "Agriculture & Livelihoods",
    "donor_selected":   "GIZ",
    "donor_framework":  "GIZ",
}


def _complete_email_login(email: str) -> None:
    st.session_state["user_email"] = email
    _is_new_user = get_user(email) is None  # check before upsert
    upsert_user(email)
    _u = get_user(email)
    if _u and is_still_paid(_u):
        st.session_state["is_paid"] = True
    # Send Day-1 welcome email for new users only
    if _is_new_user:
        try:
            from utils.email_otp import send_welcome_email
            send_welcome_email(email)
        except Exception:
            pass
    _pending_ref = st.session_state.pop("pending_paystack_ref", None)
    if _pending_ref:
        _pr = verify_payment(_pending_ref)
        if _pr.get("status") == "success":
            _pr_days = 365 if _pr.get("plan") == "annual" else (30 if _pr.get("plan") in ("monthly", "agency") else 1)
            mark_paid(email, days=_pr_days)
            st.session_state["is_paid"] = True
            metrics.log_event("payment_completed", email)
    for _k in ("_otp_email", "_otp_code", "_otp_sent_at", "_otp_attempts"):
        st.session_state.pop(_k, None)
    # Restore draft from Supabase if the user is returning after a refresh
    # and they haven't already started filling the form in this session
    _any_form_data = any(
        _ss_str(k).strip()
        for k in ("result_statement", "target_group", "evidence_description")
    )
    if not _any_form_data and not st.session_state.get("_form_is_resumption"):
        try:
            _saved_draft = load_user_draft(email)
            if _saved_draft:
                import json as _dj
                _load_from_inputs_json(_dj.loads(_saved_draft))
                st.session_state["_form_is_resumption"] = True
                st.session_state["_draft_restored_from_cloud"] = True
        except Exception:
            pass
    # Issue a durable session token so this browser is recognised on future
    # visits without retyping the email — mirrored into the URL like screen/tab.
    _session_tok = issue_session_token(email)
    if _session_tok:
        st.query_params["session"] = _session_tok
    st.rerun()


def _render_email_gate_inline(form_key_suffix: str = "") -> None:
    """Inline email-collection widget (OTP or simple). Stops rendering after display."""
    if otp_enabled() and st.session_state.get("_otp_email"):
        _otp_email = st.session_state["_otp_email"]
        if time.time() - st.session_state.get("_otp_sent_at", 0) > 600:
            st.warning("Your verification code expired. Please request a new one.")
            for _k in ("_otp_email", "_otp_code", "_otp_sent_at", "_otp_attempts"):
                st.session_state.pop(_k, None)
            st.rerun()
        st.caption(f"We sent a login link and a 6-digit code to **{_otp_email}**. Click the link, "
                   "or enter the code below (expires in 10 minutes).")
        with st.form(f"otp_verify_form{form_key_suffix}"):
            _otp_input = st.text_input("Verification code", max_chars=6, placeholder="123456")
            _verify_clicked = st.form_submit_button("Verify →", use_container_width=True)
        if _verify_clicked:
            if _otp_input.strip() == st.session_state.get("_otp_code"):
                _complete_email_login(_otp_email)
            else:
                st.session_state["_otp_attempts"] = st.session_state.get("_otp_attempts", 0) + 1
                if st.session_state["_otp_attempts"] >= 5:
                    st.error("Too many incorrect attempts. Please request a new code.")
                    for _k in ("_otp_email", "_otp_code", "_otp_sent_at", "_otp_attempts"):
                        st.session_state.pop(_k, None)
                    st.rerun()
                else:
                    st.error("Incorrect code. Please try again.")
        _otp_c1, _otp_c2 = st.columns(2)
        with _otp_c1:
            if st.button("Resend code", use_container_width=True, key=f"resend_otp{form_key_suffix}"):
                with st.spinner("Sending a new login link and code…"):
                    _ok, _err, _new_code = send_login_email(_otp_email, APP_URL)
                if _ok:
                    st.session_state["_otp_code"] = _new_code
                    st.session_state["_otp_sent_at"] = time.time()
                    st.session_state["_otp_attempts"] = 0
                    st.success("New link and code sent.")
                else:
                    st.error(f"Could not send code: {_err}")
        with _otp_c2:
            if st.button("Use a different email", use_container_width=True, key=f"diff_email{form_key_suffix}"):
                for _k in ("_otp_email", "_otp_code", "_otp_sent_at", "_otp_attempts"):
                    st.session_state.pop(_k, None)
                st.rerun()
    else:
        st.caption(
            "No password needed. We'll email you a login link and a 6-digit code."
            if otp_enabled() else
            "No password needed. We use your email to save your paid access."
        )
        with st.form(f"email_gate_form{form_key_suffix}"):
            # Privacy disclosure shown before email submission
            st.caption(
                "By continuing, you confirm that ImpactProof (a product of Impact-Receipts) "
                "may store your email and usage count to manage your account. "
                "Document content and result text are not stored unless you explicitly opt in "
                "to save an audit to your private history (you'll see that option after each "
                "check) — saved content is encrypted at rest, and you can permanently delete it "
                "anytime from My Audits. "
                "Processed by Supabase (Ireland) and Paystack (Nigeria). "
                "Contact: info@impact-receipts.com"
            )
            _gate_email = st.text_input("Email address", placeholder="you@organisation.org")
            _submit_label = "Email me a login link" if otp_enabled() else "Continue →"
            if st.form_submit_button(_submit_label, use_container_width=True):
                if "@" not in _gate_email or "." not in _gate_email.split("@")[-1]:
                    st.warning("Please enter a valid email address.")
                elif _is_disposable_email(_gate_email):
                    st.warning("Please use a permanent work or personal email — we don't accept temporary addresses.")
                else:
                    _e = _gate_email.strip().lower()
                    if otp_enabled():
                        with st.spinner("Sending your login link…"):
                            _ok, _err, _code = send_login_email(_e, APP_URL)
                        if _ok:
                            st.session_state["_otp_email"] = _e
                            st.session_state["_otp_code"] = _code
                            st.session_state["_otp_sent_at"] = time.time()
                            st.session_state["_otp_attempts"] = 0
                            st.rerun()
                        elif _err.startswith("DOMAIN_NOT_VERIFIED:"):
                            # Resend is in test mode — domain not yet verified.
                            # Fall back gracefully to simple email entry so
                            # users are not blocked. Verify a domain at
                            # resend.com/domains to enable OTP for all addresses.
                            _complete_email_login(_e)
                        else:
                            st.error(f"Could not send verification email: {_err}")
                    else:
                        _complete_email_login(_e)
    st.stop()


def _render_login_link_landing(raw_token: str) -> None:
    """Confirm-click landing for a magic-link login URL (?login_token=...).

    Deliberately requires an explicit click rather than logging in on the bare
    GET — some corporate email security scanners pre-fetch links, which would
    otherwise silently burn a single-use token before the real user opens it.
    """
    st.markdown("## Log in to ImpactProof")
    _email = verify_magic_link_token(raw_token)
    if not _email:
        st.error("This login link is invalid or has expired. Please request a new one.")
        if st.button("← Back to ImpactProof"):
            try:
                st.query_params.clear()
            except Exception:
                pass
            st.rerun()
        return
    st.markdown(f"Log in as **{_email}**?")
    if st.button("Yes, log me in →", type="primary", use_container_width=True):
        _confirmed_email = redeem_magic_link_token(raw_token)
        if _confirmed_email:
            try:
                st.query_params.clear()
            except Exception:
                pass
            _complete_email_login(_confirmed_email)
        else:
            st.error("This link was already used or has expired. Please request a new one.")


# ============================================================
# MATCH DAY — Patch 1: Scoreboard + Commentary ticker
# ============================================================

def inject_matchday_css():
    """Call ONCE per run, after st.set_page_config()."""
    st.markdown("""
    <style>
    .md-ticker { background:#2C2C2A; color:#fff; border-radius:8px; padding:10px 14px;
        font-size:14px; display:flex; align-items:center; gap:10px; margin:0 0 14px 0; }
    .md-ticker .mic { background:#A32D2D; color:#fff; font-size:9px; font-weight:600;
        padding:2px 6px; border-radius:4px; letter-spacing:.5px; flex-shrink:0; }
    .md-ticker .txt { line-height:1.4; }
    .md-sb { background:#fff; border:1px solid #d3d1c7; border-radius:8px;
        overflow:hidden; margin:0 0 16px 0; }
    .md-sb-head { background:#1B5E20; color:#fff; padding:8px 16px; font-size:11px;
        letter-spacing:1px; text-transform:uppercase; display:flex;
        justify-content:space-between; align-items:center; }
    .md-sb-live { font-size:10px; }
    .md-sb-body { display:grid; grid-template-columns:1fr 1px 1fr; }
    .md-stat { padding:18px 16px; text-align:center; }
    .md-stat .name { font-size:12px; color:#888780; margin-bottom:4px; }
    .md-stat .val { font-size:34px; font-weight:600; line-height:1; }
    .md-stat .sub { font-size:11px; color:#888780; margin-top:6px; }
    .md-div { background:#d3d1c7; }
    .md-green { color:#1B5E20; } .md-amber { color:#8A6500; } .md-red { color:#A32D2D; }
    .md-var { background:#1a1a18; color:#fff; border-radius:8px;
        padding:24px 20px; margin:16px 0; display:flex; align-items:center; gap:18px; }
    .md-var-badge { background:#A32D2D; color:#fff; font-size:11px; font-weight:700;
        padding:6px 10px; border-radius:6px; letter-spacing:1.5px; flex-shrink:0; }
    .md-var-text strong { font-size:1rem; }
    .md-var-text span { font-size:0.85rem; color:#aaa8a0; }
    .md-pitch { background:#1a1a18; border-radius:0; padding:10px 16px;
        margin:0; position:fixed; top:3.75rem; left:0; right:0; z-index:999999; }
    .md-pitch-stages { display:flex; align-items:flex-start; justify-content:space-between;
        position:relative; max-width:730px; margin:0 auto; width:100%; }
    .md-pitch-stages::before { content:""; position:absolute; top:16px; left:10%; right:10%;
        height:2px; background:#333; z-index:0; }
    .md-pstage { display:flex; flex-direction:column; align-items:center; gap:6px;
        flex:1; position:relative; z-index:1; cursor:help; }
    .md-pstage .dot { width:32px; height:32px; border-radius:50%; background:#333;
        color:#666; font-size:12px; font-weight:700; display:flex;
        align-items:center; justify-content:center; }
    .md-pstage .lbl { font-size:10px; color:#555; text-align:center; letter-spacing:.3px; }
    .md-pstage.done .dot { background:#1B5E20; color:#fff; }
    .md-pstage.done .lbl { color:#1B5E20; }
    .md-pstage.active .dot { background:#8A6500; color:#fff; }
    .md-pstage.active .lbl { color:#8A6500; font-weight:600; }
    .md-pstage:hover .lbl { text-decoration: underline dotted; }
    .md-fulltime { background:#2C2C2A; color:#fff; border-radius:8px;
        padding:20px 20px 16px; margin:12px 0; text-align:center; }
    .md-fulltime .whistle { font-size:10px; letter-spacing:1.5px; text-transform:uppercase;
        color:#888780; margin-bottom:8px; }
    .md-fulltime h3 { margin:0 0 8px; font-size:1.3rem; color:#8A6500; }
    .md-fulltime p { margin:0; font-size:0.85rem; color:#aaa8a0; line-height:1.5; }
    .md-card { background:#f9f8f5; border:1px solid #e0ded8; border-radius:8px;
        padding:14px 16px; margin:8px 0; }
    .stage-tip { font-size:10px; color:#aaa8a0; text-align:center;
        margin-top:4px; padding:0 2px; line-height:1.3; }
    .md-tab-header { margin:10px 0 14px; }
    .md-tab-header .tab-num { font-size:11px; color:#888780; letter-spacing:.5px;
        text-transform:uppercase; display:block; margin-bottom:2px; }
    .md-tab-header .tab-name { font-size:1.05rem; font-weight:700; color:#1B5E20; }
    </style>
    """, unsafe_allow_html=True)


MATCHDAY_COMMENTARY = {
    "enter":    "Step 1 — Define your result: who benefited, what changed, where and when.",
    "logframe": "Step 2 — Link your result to your approved logframe indicator and target.",
    "evidence": "Step 3 — Document and verify your supporting evidence.",
    "review":   "Step 4 — Review your scores and run the full diagnostic.",
    "report":   "Diagnostic complete. Here are your Confidence and Clarity scores.",
}


def render_commentary(stage_key: str):
    """Renders the one-line match commentary for the given stage."""
    line = MATCHDAY_COMMENTARY.get(stage_key, "")
    if not line:
        return
    st.markdown(
        f'<div class="md-ticker"><span class="mic">STEP</span>'
        f'<span class="txt">{line}</span></div>',
        unsafe_allow_html=True,
    )


def _score_class(value, verified):
    if value is None or not verified:
        return "md-amber"
    if value >= 70:
        return "md-green"
    if value >= 50:
        return "md-amber"
    return "md-red"


_SCORE_CLASS_LABEL = {"md-green": "strong", "md-amber": "acceptable", "md-red": "high risk"}


def render_scoreboard(confidence=None, clarity=None, verified=False):
    """Two-stat scoreboard for Confidence and Clarity (0-100 scale)."""
    conf_txt = "—" if confidence is None else str(int(round(confidence)))
    clar_txt = "—" if clarity is None else str(int(round(clarity)))
    conf_cls = _score_class(confidence, verified)
    clar_cls = _score_class(clarity, verified)
    # Accessibility D4: text labels alongside color — WCAG 1.4.1 (no color-only signal)
    conf_lbl = _SCORE_CLASS_LABEL.get(conf_cls, "")
    clar_lbl = _SCORE_CLASS_LABEL.get(clar_cls, "")
    status = "● scored" if verified else "● preview"
    status_color = "#A5D6A7" if verified else "#FFD54F"
    st.markdown(f"""
    <div class="md-sb">
      <div class="md-sb-head"><span>Score summary</span>
        <span class="md-sb-live" style="color:{status_color}">{status}</span></div>
      <div class="md-sb-body">
        <div class="md-stat"><div class="name">Confidence</div>
          <div class="val {conf_cls}" aria-label="{conf_txt} out of 100 — {conf_lbl}">{conf_txt}</div>
          <div class="sub">Will the donor trust the evidence?</div>
          {f'<div class="sub" style="font-size:11px;font-weight:600;">● {conf_lbl}</div>' if conf_lbl else ''}</div>
        <div class="md-div"></div>
        <div class="md-stat"><div class="name">Clarity</div>
          <div class="val {clar_cls}" aria-label="{clar_txt} out of 100 — {clar_lbl}">{clar_txt}</div>
          <div class="sub">Is the result clear enough to stand alone?</div>
          {f'<div class="sub" style="font-size:11px;font-weight:600;">● {clar_lbl}</div>' if clar_lbl else ''}</div>
      </div>
    </div>
    """, unsafe_allow_html=True)


def render_var_review():
    """Dark VAR panel shown while the confidence check is running."""
    st.markdown(
        '<div class="md-var">'
        '<div class="md-var-badge">CHECKING</div>'
        '<div class="md-var-text">'
        '<strong>Diagnostic in progress</strong><br>'
        '<span>Analysing your evidence and result — your scores are being computed.</span>'
        '</div></div>',
        unsafe_allow_html=True,
    )


MATCHDAY_STAGES = [
    ("enter",    "Result"),
    ("logframe", "Logframe"),
    ("evidence", "Evidence"),
    ("review",   "Review"),
    ("report",   "Report"),
]

MATCHDAY_TIPS = {
    "enter":    "Step 1 — Define your result: who benefited, what changed, where and when",
    "logframe": "Step 2 — Link your result to your logframe indicator and target",
    "evidence": "Step 3 — Document and verify your supporting evidence before scoring",
    "review":   "Step 4 — Review your live scores and run the full diagnostic",
    "report":   "Your Confidence &amp; Clarity scores and verdict are in",
}


def render_pitch_strip(current_stage: str):
    keys = [k for k, _ in MATCHDAY_STAGES]
    try:
        cur = keys.index(current_stage)
    except ValueError:
        cur = -1
    cells = ""
    for idx, (k, lbl) in enumerate(MATCHDAY_STAGES):
        cls = "done" if idx < cur else ("active" if idx == cur else "")
        mark = "✓" if idx < cur else str(idx + 1)
        tip = MATCHDAY_TIPS.get(k, "")
        tip_html = f'<div class="stage-tip">{tip}</div>' if cls == "active" else ""
        cells += (f'<div class="md-pstage {cls}" title="{tip}"><div class="dot">{mark}</div>'
                  f'<div class="lbl">{lbl}</div>{tip_html}</div>')
    st.markdown(
        f'<div class="md-pitch"><div class="md-pitch-stages">{cells}</div></div>'
        f'<div style="height:96px"></div>',
        unsafe_allow_html=True,
    )


def _esc(s: str) -> str:
    import html
    return html.escape(str(s))


def render_fulltime(confidence, clarity, summary):
    """Final-result card. Render directly above your st.download_button().

    confidence / clarity : final 0-100 scores (verified).
    summary : a short honest sentence about the result.
    """
    conf_txt = "—" if confidence is None else str(int(round(confidence)))
    clar_txt = "—" if clarity is None else str(int(round(clarity)))
    st.markdown(f"""
    <div class="md-fulltime">
      <div class="whistle">YOUR RESULT</div>
      <h3>Confidence {conf_txt} · Clarity {clar_txt}</h3>
      <p>{_esc(summary)}</p>
    </div>
    """, unsafe_allow_html=True)


def render_season_teaser():
    """One-line hook seeding the Phase 2 'Integrity League'. Copy only."""
    st.markdown(
        '<div class="md-card" style="text-align:center;">'
        '<h4 style="margin:0 0 4px 0;">Season record</h4>'
        '<p style="font-size:12px;color:#888780;margin:0;">This match is logged to your '
        'integrity record. Your rating rises only on verified evidence — never on volume. '
        '(Coming soon.)</p></div>',
        unsafe_allow_html=True,
    )


def render_pricing_page():
    """Public 3-tier pricing page. Shown when _show_pricing=True."""
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
    if st.button("← Back", key="pricing_back"):
        st.session_state.pop("_show_pricing", None)
        st.rerun()

    st.markdown("## Pricing")
    st.caption("Determine your evidence readiness. Prove your impact. First 3 checks always free.")

    # ROI micro-copy (Council XXVII — West Africa-specific framing)
    st.markdown(
        "<div style='background:#F1F8E9;border-left:4px solid #1B5E20;padding:10px 16px;"
        "border-radius:6px;font-size:0.9rem;margin-bottom:20px;'>"
        "💡 <strong>The ROI is immediate: GHS 50/month vs. GHS 12,000–17,000 in rework costs.</strong> "
        "DevEx MEL Salary Survey (2024): average Ghana consultant day rate ≈ GHS 1,200–1,800/day. "
        "One rejected USAID, Mastercard Foundation, or FCDO report = 40+ hours of rework. "
        "ImpactProof catches the gaps donors flag — before your report goes out. "
        "Score every KPI in 60 seconds. Download a citable Readiness Card with a reference ID."
        "</div>",
        unsafe_allow_html=True,
    )

    # Tier cards
    _t1, _t2, _t3 = st.columns(3)

    with _t1:
        st.markdown(
            f"<div style='border:2px solid #E0E0E0;border-radius:10px;padding:20px;height:100%;{_pca}'>"
            "<p style='font-size:0.75rem;color:#616161;text-transform:uppercase;letter-spacing:1px;margin:0 0 4px;'>For first-time users</p>"
            "<h3 style='color:#212121;margin:0 0 4px;'>Starter</h3>"
            "<p style='font-size:2rem;font-weight:700;color:#212121;margin:0;'>Free</p>"
            "<p style='font-size:0.8rem;color:#616161;margin:4px 0 16px;'>No card needed</p>"
            "<p style='font-size:0.85rem;color:#424242;margin-bottom:12px;font-style:italic;'>Try the tool on your next result</p>"
            "<hr style='border:none;border-top:1px solid #E0E0E0;margin:12px 0;'/>"
            "<ul style='padding-left:16px;font-size:0.85rem;color:#424242;margin:0;line-height:1.8;'>"
            "<li>3 free checks</li>"
            "<li>Quick Check — instant provisional score</li>"
            "<li>Full Confidence &amp; Clarity diagnostic</li>"
            "<li>Actionable fix list (top 3)</li>"
            "</ul>"
            "</div>",
            unsafe_allow_html=True,
        )
        if st.button("Start free →", key="pricing_starter", use_container_width=True):
            st.session_state.pop("_show_pricing", None)
            _go_to_screen(0)

    with _t2:
        st.markdown(
            f"<div style='border:3px solid #1B5E20;border-radius:10px;padding:20px;height:100%;background:#F9FFF9;{_pca}'>"
            "<p style='background:#1B5E20;color:white;font-size:0.7rem;font-weight:700;letter-spacing:1px;text-transform:uppercase;"
            "display:inline-block;padding:2px 8px;border-radius:20px;margin:0 0 8px;'>Most popular</p>"
            "<p style='font-size:0.75rem;color:#616161;text-transform:uppercase;letter-spacing:1px;margin:0 0 4px;'>For MEL practitioners reporting regularly</p>"
            "<h3 style='color:#1B5E20;margin:0 0 4px;'>Professional</h3>"
            "<p style='font-size:2rem;font-weight:700;color:#1B5E20;margin:0;'>GHS 50<span style='font-size:1rem;font-weight:400;'>/mo</span></p>"
            "<p style='font-size:0.8rem;color:#616161;margin:4px 0 16px;'>~£3.50 · or GHS 500/year (2 months free)</p>"
            "<p style='font-size:0.85rem;color:#424242;margin-bottom:12px;font-style:italic;'>Determine evidence readiness before every submission</p>"
            "<hr style='border:none;border-top:1px solid #C8E6C9;margin:12px 0;'/>"
            "<ul style='padding-left:16px;font-size:0.85rem;color:#424242;margin:0;line-height:1.8;'>"
            "<li><strong>Unlimited checks</strong></li>"
            "<li><strong>Audit My Report</strong> — upload a Word/PDF, get determinations for every result, download a decision audit Excel</li>"
            "<li>Instant Report Check — auto-fill form from uploaded document</li>"
            "<li>Readiness Card PDF (shareable with supervisor or donor)</li>"
            "<li>Donor-specific fixes (USAID, FCDO, GIZ, World Bank…)</li>"
            "<li>Full diagnostic report with logframe linkage scoring</li>"
            "<li>Score chat — 'Ask about your score' on every result</li>"
            "</ul>"
            "</div>",
            unsafe_allow_html=True,
        )
        if st.button("Check my results →", key="pricing_pro", type="primary", use_container_width=True):
            st.session_state.pop("_show_pricing", None)
            _go_to_screen(1, reset=True)

    with _t3:
        st.markdown(
            f"<div style='border:2px solid #8A6500;border-radius:10px;padding:20px;height:100%;{_pca}'>"
            "<p style='font-size:0.75rem;color:#616161;text-transform:uppercase;letter-spacing:1px;margin:0 0 4px;'>For MEL consultancies &amp; multi-donor programme teams</p>"
            "<h3 style='color:#8A6500;margin:0 0 4px;'>Agency</h3>"
            "<p style='font-size:2rem;font-weight:700;color:#8A6500;margin:0;'>GHS 200<span style='font-size:1rem;font-weight:400;'>/mo</span></p>"
            "<p style='font-size:0.8rem;color:#616161;margin:4px 0 16px;'>~£13 · multiple clients, multiple donors</p>"
            "<p style='font-size:0.85rem;color:#424242;margin-bottom:12px;font-style:italic;'>Score every client's evidence — USAID, MCF, GIZ Ghana, FCDO — one place</p>"
            "<hr style='border:none;border-top:1px solid #E8D5A3;margin:12px 0;'/>"
            "<ul style='padding-left:16px;font-size:0.85rem;color:#424242;margin:0;line-height:1.8;'>"
            "<li><strong>Multi-client, multi-donor</strong> — separate Readiness Cards per client, per donor</li>"
            "<li>Everything in Professional</li>"
            "<li><strong>Portfolio analysis</strong> — indicator heatmap, systemic gap report</li>"
            "<li><strong>Portfolio Q&amp;A chat</strong> — 'Which KPI needs the most work across all results?'</li>"
            "<li>Donor framework crosswalk reports (USAID DQA, Bond, OECD-DAC, World Bank)</li>"
            "<li>Up to 5 team seats</li>"
            "<li>Priority support</li>"
            "</ul>"
            "</div>",
            unsafe_allow_html=True,
        )
        _ag_email  = st.session_state.get("user_email", "")
        _ag_wa_key = "wa_agency_plan_clicked"
        if st.button("Talk to us →", key="wa_agency_btn", use_container_width=True):
            from utils.whatsapp import notify_founder
            notify_founder("agency_plan", user_email=_ag_email)
            st.session_state[_ag_wa_key] = True
        if st.session_state.get(_ag_wa_key):
            from utils.whatsapp import build_wa_url
            st.link_button("Open WhatsApp →",
                           build_wa_url("agency_plan", _ag_email),
                           use_container_width=True)
            st.success("✓ We've been notified — we'll reply within 24 hours.")

    st.divider()
    _pq_email  = st.session_state.get("user_email", "")
    _pq_wa_key = "wa_pricing_q_clicked"
    _pq_col1, _pq_col2 = st.columns([2, 1])
    with _pq_col1:
        st.caption("All prices in GHS (Ghana Cedis). Paid via Paystack (card, MoMo, bank). "
                   "Cancel anytime.")
    with _pq_col2:
        if st.button("Questions? WhatsApp →", key="wa_pricing_q_btn", use_container_width=True):
            from utils.whatsapp import notify_founder
            notify_founder("pricing_questions", user_email=_pq_email)
            st.session_state[_pq_wa_key] = True
    if st.session_state.get(_pq_wa_key):
        from utils.whatsapp import build_wa_url
        st.link_button("Open WhatsApp →",
                       build_wa_url("pricing_questions", _pq_email),
                       use_container_width=True)
        st.caption("✓ Notified — we'll reply within 24 hours.")


def render_billing_page():
    """Billing & account settings — current plan, invoice history,
    cancellation, and signed-in devices. Shown when _show_billing=True."""
    if st.button("← Back", key="billing_back"):
        st.session_state.pop("_show_billing", None)
        st.query_params.pop("billing", None)
        st.rerun()

    st.markdown("## Billing & account")

    email = st.session_state.get("user_email", "")
    if not email:
        st.info("Enter your email to view your billing details.")
        _render_email_gate_inline("_billing")
        return  # unreachable in practice -- the gate above calls st.stop()

    user = get_user(email) or {}
    is_paid = is_still_paid(user)
    plan = (user.get("plan") or "free").capitalize()
    paid_until = user.get("paid_until", "")
    subscription_status = user.get("subscription_status", "")

    st.markdown(f"**Account:** {email}")
    if is_paid:
        _renewal = f" · active until **{paid_until}**" if paid_until else ""
        st.success(f"**Plan:** {plan}{_renewal}")
        if subscription_status == "attention":
            st.warning("Your last renewal charge failed. Please update your payment method "
                       "to avoid losing access when your current period ends.")
    else:
        _checks_used = user.get("free_checks_used", 0)
        st.info("**Plan:** Free")
        st.caption(f"Free checks used: {_checks_used}/{FREE_CHECKS_LIMIT}")
    if st.button("View plans →", key="billing_view_plans"):
        st.session_state.pop("_show_billing", None)
        st.query_params.pop("billing", None)
        st.session_state["_show_pricing"] = True
        st.rerun()

    # --- Cancellation ---
    _sub_code = user.get("paystack_subscription_code", "")
    _sub_token = user.get("paystack_email_token", "")
    if is_paid and _sub_code:
        st.divider()
        st.markdown("#### Cancel subscription")
        if st.session_state.get("_confirm_cancel_sub"):
            st.warning(f"Cancel your subscription? You'll keep access until "
                       f"{paid_until or 'your current period ends'}.")
            _cc1, _cc2 = st.columns(2)
            with _cc1:
                if st.button("Yes, cancel", key="billing_cancel_confirm",
                             type="primary", use_container_width=True):
                    _ok, _msg = disable_subscription(_sub_code, _sub_token)
                    if _ok:
                        st.success(f"Subscription cancelled. {_msg}")
                    else:
                        st.error(f"Could not cancel automatically: {_msg}. "
                                 "Contact us and we'll cancel it manually within 24 hours.")
                        from utils.whatsapp import notify_founder
                        notify_founder("payment_support", user_email=email)
                    st.session_state.pop("_confirm_cancel_sub", None)
            with _cc2:
                if st.button("Keep subscription", key="billing_cancel_abort",
                             use_container_width=True):
                    st.session_state.pop("_confirm_cancel_sub", None)
                    st.rerun()
        else:
            if st.button("Cancel subscription", key="billing_cancel_start"):
                st.session_state["_confirm_cancel_sub"] = True
                st.rerun()

    # --- Invoice history ---
    st.divider()
    st.markdown("#### Invoice history")
    _history = get_payment_history(email)
    if not _history:
        st.caption("No payments yet.")
    else:
        _rows = [{
            "Date": (h.get("created_at") or "")[:10],
            "Plan": h.get("plan") or "—",
            "Amount (GHS)": round((h.get("amount_pesewas") or 0) / 100, 2),
            "Status": h.get("status", ""),
        } for h in _history]
        st.dataframe(_rows, use_container_width=True, hide_index=True)

    # --- Devices / sessions ---
    st.divider()
    st.markdown("#### Signed-in devices")
    _sessions = list_sessions(email)
    if not _sessions:
        st.caption("No active sessions found.")
    else:
        for _s in _sessions:
            _s_label = _s.get("user_agent") or "Unknown device"
            _s_last = (_s.get("last_seen_at") or "")[:16].replace("T", " ")
            _sc1, _sc2 = st.columns([3, 1])
            with _sc1:
                st.caption(f"{_s_label} · last used {_s_last}")
            with _sc2:
                if st.button("Sign out", key=f"revoke_{_s.get('token_hash', '')[:12]}",
                             use_container_width=True):
                    revoke_session(_s.get("token_hash", ""), email)
                    st.rerun()
        if len(_sessions) > 1:
            if st.button("Sign out of all devices", key="billing_revoke_all"):
                revoke_all_sessions(email)
                st.rerun()


def render_my_audits_page():
    """My Audits — saved audit history (opt-in only). Shown when
    _show_my_audits=True. Re-download regenerates the Readiness Card PDF from
    the stored submission+evaluation via the same builders render_screen_2
    uses; IRC field-source highlighting and Council Assessment content are
    session-only extras not captured at save time, so a re-download
    reproduces the core card faithfully but not those two enrichments."""
    if st.button("← Back", key="my_audits_back"):
        st.session_state.pop("_show_my_audits", None)
        st.query_params.pop("my_audits", None)
        st.rerun()

    st.markdown("## My Audits")
    st.caption("Audits you've explicitly saved. Nothing appears here unless you checked "
               "\"Save this audit to my private history\" after a check.")

    email = st.session_state.get("user_email", "")
    if not email:
        st.info("Enter your email to view your saved audits.")
        _render_email_gate_inline("_my_audits")
        return  # unreachable in practice -- the gate above calls st.stop()

    _audits = list_audits(email)
    if not _audits:
        st.caption("No saved audits yet.")

    for _a in _audits:
        _a_id = _a.get("id")
        _a_date = str(_a.get("created_at") or "")[:16].replace("T", " ")
        _a_conf = _a.get("primary_confidence_score")
        _a_clar = _a.get("primary_clarity_score")
        st.markdown(f"**{_a.get('donor') or 'No donor specified'} · {_a.get('sector') or 'No sector'}** — {_a_date}")
        if _a_conf is not None and _a_clar is not None:
            st.caption(f"Confidence {_a_conf:.1f} · Clarity {_a_clar:.1f} · {_a.get('primary_verdict', '')}")
        else:
            st.caption(_a.get("primary_verdict", ""))

        _mc1, _mc2 = st.columns(2)
        with _mc1:
            _pdf_key = f"_audit_pdf_{_a_id}"
            if st.session_state.get(_pdf_key):
                st.download_button(
                    "⬇️ Download PDF", data=st.session_state[_pdf_key],
                    file_name=f"readiness_card_{_a_id}.pdf", mime="application/pdf",
                    key=f"my_audits_dl_{_a_id}", use_container_width=True,
                )
            elif st.button("Re-download PDF", key=f"my_audits_redl_{_a_id}", use_container_width=True):
                _full = get_audit(email, _a_id)
                _subs = (_full or {}).get("submissions") or []
                _evs = (_full or {}).get("evaluations") or []
                if _full and _subs and _evs:
                    _ts = (_full.get("ref_id") or "").replace("IMP-", "")
                    _redl_html = _build_html_report_card(_subs[0], _evs[0], _ts,
                                                          field_sources=None, council_assessment=None)
                    _redl_pdf = _html_to_pdf_bytes(_redl_html)
                    if _redl_pdf:
                        st.session_state[_pdf_key] = _redl_pdf
                        st.rerun()
                    else:
                        st.warning("Could not regenerate this audit's PDF.")
                else:
                    st.warning("Could not load this audit's saved data.")
        with _mc2:
            if st.button("Delete", key=f"my_audits_delete_{_a_id}", use_container_width=True):
                delete_audit(email, _a_id)
                st.session_state.pop(f"_audit_pdf_{_a_id}", None)
                st.rerun()
        st.divider()

    # --- Logframe Library management ---
    st.markdown("#### Logframe Library")
    st.caption("Named indicator lists you've saved from Screen 1's Logframe tab, reusable across audits.")
    _libs = list_logframe_libraries(email)
    if not _libs:
        st.caption("No libraries saved yet — save one from the Logframe tab while filling out a result.")
    else:
        for _lib in _libs:
            _lib_id = _lib.get("id")
            _lib_items = get_library_items(_lib_id, email)
            _lc1, _lc2 = st.columns([3, 1])
            with _lc1:
                st.markdown(f"**{_lib.get('name', '')}** — {len(_lib_items)} indicator(s)")
            with _lc2:
                if st.button("Delete library", key=f"my_audits_del_lib_{_lib_id}", use_container_width=True):
                    delete_logframe_library(_lib_id, email)
                    st.rerun()
            if _lib_items:
                with st.expander("View indicators", expanded=False):
                    for _it in _lib_items:
                        st.caption(f"• {_it.get('indicator_name') or _it.get('logframe_indicator', '')} "
                                   f"— baseline {_it.get('logframe_baseline') or '—'}, "
                                   f"target {_it.get('logframe_target') or '—'}")

    # --- Danger zone: permanent data deletion (Ghana Data Protection Act 843) ---
    st.divider()
    st.markdown("#### Danger zone")
    st.caption(
        "Permanently delete every saved audit, Logframe Library, and pre-submission draft "
        "associated with your account — for Ghana Data Protection Act 843 consistency. "
        "This does not delete your account, active sign-ins, or payment/invoice history "
        "(kept for accounting purposes)."
    )
    if st.session_state.get("_confirm_purge_history"):
        st.warning(
            "This permanently deletes all your saved audits, Logframe Libraries, and any "
            "in-progress draft. This cannot be undone."
        )
        _pc1, _pc2 = st.columns(2)
        with _pc1:
            if st.button("Yes, permanently delete my history", key="purge_confirm",
                         type="primary", use_container_width=True):
                _safe_log_access(email, "account_purge", resource_type="account")  # logged BEFORE the deletes
                _counts = purge_account_audit_content(email)
                clear_user_draft(email)
                delete_wa_conversations(email)
                st.session_state.pop("_confirm_purge_history", None)
                st.success(
                    f"Deleted {_counts['audits_deleted']} audit(s) and "
                    f"{_counts['libraries_deleted']} Logframe Library/libraries. Your account, "
                    f"sign-in, and payment history are unaffected."
                )
        with _pc2:
            if st.button("Cancel", key="purge_cancel", use_container_width=True):
                st.session_state.pop("_confirm_purge_history", None)
                st.rerun()
    else:
        if st.button("🗑 Permanently erase my history", key="purge_start"):
            st.session_state["_confirm_purge_history"] = True
            st.rerun()


def _render_ph_landing():
    """Stripped-down landing for Product Hunt / referral traffic.
    Optimised for email capture + Quick Check conversion."""
    st.markdown(
        """
        <div style='text-align:center;padding:24px 0 8px;'>
          <p style='font-size:0.75rem;color:#8A6500;font-weight:700;letter-spacing:1px;text-transform:uppercase;margin:0 0 8px;'>
            Featured on Product Hunt 🎉
          </p>
          <h1 style='color:#1B5E20;font-size:1.8rem;margin:0 0 8px;line-height:1.2;'>
            About to submit a result to your donor?
          </h1>
          <p style='color:#424242;font-size:1rem;margin:0 0 4px;'>
            Run a 60-second evidence quality check first.
          </p>
          <p style='color:#616161;font-size:0.85rem;margin:0 0 24px;'>
            For MEL officers, programme leads, and consultants — the people who answer for evidence quality. Reporting to FCDO, GIZ, World Bank, EU, and 7 more donors.
          </p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    # Quick Check front-and-centre
    st.markdown("### ⚡ Try it now — no registration")
    with st.container(border=True):
        qc_result_ph  = st.text_area("Your result statement", key="qc_result_ph", height=80,
                                      placeholder="e.g., Trained 250 farmers in climate-smart practices in Ashanti Region, Jan–Jun 2025")
        qc_ev_type_ph = st.selectbox("Evidence type", key="qc_ev_type_ph",
                                      options=["(Select evidence type)"] + EVIDENCE_TYPES)
        if st.button("⚡ Check my result →", key="qc_ph_run", type="primary", use_container_width=True):
            if qc_result_ph and qc_ev_type_ph != "(Select evidence type)":
                try:
                    _s = {"result_statement": qc_result_ph, "target_group":"","timeframe":"","geographic_scope":"",
                          "logframe_indicator":"","logframe_target":"","logframe_achievement":"","additional_context":"",
                          "evidence":[{"type":qc_ev_type_ph,"description":"","verified_by":"",
                                       "internal_review":"Not reviewed","external_review":"No external review"}],
                          "beneficiary_voice":"","provenance_checklist":{}}
                    _e = _evaluator.evaluate_submission(_s)
                    _c  = _e.get("raw_confidence_score",0)
                    _cl = _e.get("clarity_score",0)
                    _cl_lbl, _ = _evaluator.interpret_score(_cl)
                    _c_lbl,  _ = _evaluator.interpret_score(_c)
                    st.session_state["_qc_ph_scores"] = {"c":_c,"cl":_cl,"c_lbl":_c_lbl,"cl_lbl":_cl_lbl,
                                                          "result":qc_result_ph,"ev_type":qc_ev_type_ph}
                    st.rerun()
                except Exception:
                    st.error("Something went wrong scoring this — please try again.")
            else:
                st.warning("Enter a result statement and select an evidence type.")

    _ph = st.session_state.get("_qc_ph_scores")
    if _ph:
        st.success(
            f"**Confidence: {_ph['c']}/5.0** ({_ph['c_lbl']}) · "
            f"**Clarity: {_ph['cl']}/5.0** ({_ph['cl_lbl']})"
        )
        st.caption("Provisional — full diagnosis needs evidence details, logframe linkage, and verification.")
        if st.button("Get the full diagnosis + Readiness Card PDF →", key="qc_ph_continue", type="primary", use_container_width=True):
            st.session_state["result_statement"] = _ph["result"]
            st.session_state["evidence_type"]    = _ph["ev_type"]
            st.session_state.pop("_qc_ph_scores", None)
            _go_to_screen(1, reset=False)

    st.divider()
    _ph_c1, _ph_c2, _ph_c3 = st.columns(3)
    with _ph_c1:
        st.metric("Reproducible score", "Always", help="Same inputs → same score, every time. No AI randomness — fully deterministic, rule-based scoring.")
    with _ph_c2:
        st.metric("Anchored to 4 frameworks", "USAID · FCDO · Bond · World Bank", help="Every sub-score traces to a named donor standard — USAID ADS 201, FCDO Evaluation Policy 2025, Bond Evidence Principles 2024, World Bank Results Framework.")
    with _ph_c3:
        st.metric("Portfolio in 60 seconds", "10+ results", help="Audit My Report: upload one document, every result extracted and determined. Download as a donor-ready Excel decision audit.")

    if st.button("← Back to full landing page", key="ph_back"):
        st.session_state.pop("_referral_source", None)
        st.rerun()


def render_screen_0():
    # PH / referral traffic gets a stripped-down Quick Check landing page
    _ref = st.session_state.get("_referral_source", "")
    if _ref in ("producthunt", "ph", "product_hunt"):
        _render_ph_landing()
        return

    _logo_path = pathlib.Path(__file__).parent / "logo.png.png"
    try:
        _logo_b64 = base64.b64encode(_logo_path.read_bytes()).decode()
        _logo_tag = (
            f'<div style="display:flex; align-items:center; gap:14px; margin-bottom:12px;">'
            f'<img src="data:image/png;base64,{_logo_b64}" alt="ImpactProof" style="height:56px;">'
            f'<span style="font-size:0.9rem; font-weight:600; line-height:1.2;">'
            f'<span style="color:#1B5E20;">ImpactProof</span><br>'
            f'<span style="color:#8A6500; font-weight:600;">Decide what to fix. Prove your impact.</span>'
            f'</span>'
            f'</div>'
        )
    except FileNotFoundError:
        _logo_tag = (
            '<div style="display:flex; align-items:center; gap:14px; margin-bottom:12px;">'
            '<span style="font-size:0.9rem; font-weight:600; line-height:1.2;">'
            '<span style="color:#1B5E20;">ImpactProof</span><br>'
            '<span style="color:#8A6500; font-weight:400;">Decide what to fix. Prove your impact.</span>'
            '</span>'
            '</div>'
        )
    # ── Hero ──────────────────────────────────────────────────────────────
    st.markdown(
        f"""
        <div class="hero-block">
          {_logo_tag}
          <h1 style="margin:8px 0 4px;">Know what to fix, in what order, before your donor does.</h1>
          <p style="font-size:0.85rem;color:#616161;margin:0 0 6px;">
            For MEL officers, programme leads, and consultants — the people who answer for evidence quality.
          </p>
          <p style="font-size:0.8rem;color:#8A6500;margin:0;font-style:italic;">
            Unlike a chatbot that generates suggestions, ImpactProof makes determinations — evidence quality decisions anchored to USAID ADS 201, Bond Evidence Principles 2024, FCDO, and World Bank standards. Traceable. Reproducible. Auditable.
          </p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    # ── Two-path choice cards ──────────────────────────────────────────────
    st.markdown("")
    _path_col1, _path_col2 = st.columns(2, gap="medium")

    with _path_col1:
        with st.container(border=True):
            st.markdown("#### 📄 Audit My Report")
            st.caption(
                "Upload a Word or PDF donor report. ImpactProof identifies every result, "
                "makes a determination for each, and ranks what to fix first — all in 60 seconds. "
                "Download a filled Excel with determinations and ranked priorities."
            )
            st.caption("First 3 uploads free · Just your email, no password")
            if st.button("Upload and Audit →", key="cta_score_report", type="primary",
                         use_container_width=True):
                _go_to_screen(3)

    with _path_col2:
        with st.container(border=True):
            st.markdown("#### ✏️ Check One Result")
            st.caption(
                "Fill a short form and get a determination: submission-ready, needs work, or high risk "
                "— with a ranked fix list ordered by score impact."
            )
            st.caption("Or check below for a 60-second instant read.")
            if st.button("Start Form →", key="cta_top", use_container_width=True):
                if not st.session_state.get("has_seen_tutorial"):
                    st.session_state["tutorial_step"] = 1
                _go_to_screen(1, reset=True)

    # ── Council XXVI — Decision Intelligence strip ─────────────────────────
    _di_c1, _di_c2, _di_c3 = st.columns(3)
    with _di_c1:
        st.metric("Decision type", "Rule-based",
                  help="No AI randomness — same inputs always produce the same determination.")
    with _di_c2:
        st.metric("Frameworks", "4 donor standards",
                  help="USAID ADS 201 · FCDO 2025 · Bond 2024 · World Bank RF")
    with _di_c3:
        st.metric("Fix routing", "Ranked by impact",
                  help="Each fix is ranked by score impact — the highest-leverage action is always first.")

    with st.expander("🤖 How the AI works", expanded=False):
        st.markdown(
            "- **AI reads & interrogates** — a panel of 5 AI reviewers (each modelled on a "
            "different reviewer type: evidence auditor, programme strategist, critical "
            "reviewer, implementation guide, and donor representative) asks the questions a "
            "real donor reviewer would ask, and AI Logframe Match suggests the closest-fit "
            "indicator from your own list.\n"
            "- **Rules score deterministically** — Confidence and Clarity always come from "
            "the same eight rule-based criteria above, never from the AI. Same inputs, same "
            "score, every time.\n"
            "- **Nothing is ever fabricated** — every AI-drafted rewrite is machine-checked "
            "against your own submission before it's shown; anything it can't verify is "
            "withheld, not guessed.\n\n"
            "Built with Ghana's Data Protection Act 843 and responsible-AI practice in mind — "
            "AI output is always labelled and editable, and your documents are never used to "
            "invent facts about your project."
        )

    # ── Quick Check (now secondary — inside expander) ──────────────────────
    with st.expander("⚡ Quick Check — instant provisional scores (60 seconds)", expanded=False):
        st.caption("Fill 3 fields to instantly see provisional Confidence and Clarity scores. No form, no email required.")
        qc_result   = st.text_area("Your result statement", key="qc_result", height=80,
                                    placeholder="e.g., Trained 250 farmers in climate-smart practices in Northern Region, Jan–Jun 2025")
        qc_ev_type  = st.selectbox("Evidence type", key="qc_ev_type",
                                    options=["(Select evidence type)"] + EVIDENCE_TYPES)
        qc_verifier = st.text_input("Who verified this?", key="qc_verifier",
                                     placeholder="e.g., District Agriculture Officer")
        if st.button("⚡ Quick Check →", key="qc_run", type="primary", use_container_width=True):
            if qc_result and qc_ev_type != "(Select evidence type)":
                try:
                    _qc_sub = {
                        "result_statement": qc_result,
                        "target_group": "", "timeframe": "", "geographic_scope": "",
                        "logframe_indicator": "", "logframe_target": "", "logframe_achievement": "",
                        "additional_context": "",
                        "evidence": [{"type": qc_ev_type, "description": "", "verified_by": qc_verifier,
                                      "internal_review": "Not reviewed", "external_review": "No external review"}],
                        "beneficiary_voice": "", "provenance_checklist": {},
                        "org_type": st.session_state.get("org_type", "International NGO (INGO)"),
                    }
                    _qc_ev = _evaluator.evaluate_submission(_qc_sub)
                    _qc_c  = _qc_ev.get("raw_confidence_score", 0)
                    _qc_cl = _qc_ev.get("clarity_score", 0)
                    _qc_cl_label, _ = _evaluator.interpret_score(_qc_cl)
                    _qc_c_label,  _ = _evaluator.interpret_score(_qc_c)
                    st.session_state["_qc_last_scores"] = {
                        "c": _qc_c, "cl": _qc_cl, "c_lbl": _qc_c_label, "cl_lbl": _qc_cl_label,
                        "result": qc_result, "ev_type": qc_ev_type, "verifier": qc_verifier,
                        "track_label": _qc_ev.get("track_label", "INGO standard"),
                        "threshold": _qc_ev.get("threshold_used", 4.0),
                    }
                    st.rerun()
                except Exception:
                    st.error("Something went wrong scoring this — please try again.")
            else:
                st.warning("Enter a result statement and select an evidence type.")

        _qc_scores = st.session_state.get("_qc_last_scores")
        if _qc_scores:
            st.success(
                f"**Provisional scores** — Confidence: **{_qc_scores['c']}/5.0** ({_qc_scores['c_lbl']}) · "
                f"Clarity: **{_qc_scores['cl']}/5.0** ({_qc_scores['cl_lbl']})"
            )
            st.caption(
                f"Provisional only — Verification, Recency, and Clarity improve with full form data. "
                f"Standard applied: {_qc_scores.get('track_label', 'INGO standard')} "
                f"(threshold {_qc_scores.get('threshold', 4.0)}/5.0)."
            )
            if st.button("Continue for full diagnosis →", key="qc_continue", use_container_width=True, type="primary"):
                st.session_state["result_statement"] = _qc_scores["result"]
                st.session_state["evidence_type"]    = _qc_scores["ev_type"]
                st.session_state["verifier"]         = _qc_scores["verifier"]
                st.session_state.pop("_qc_last_scores", None)
                st.session_state["_from_quick_check"] = True
                if not st.session_state.get("has_seen_tutorial"):
                    st.session_state["tutorial_step"] = 1
                _go_to_screen(1, reset=False)

    # ── Case study trust element (compact) ────────────────────────────────
    st.markdown(
        """
        <div style="border-left:4px solid #8A6500;padding:8px 12px;margin:12px 0;background:transparent;">
          <p style="margin:0;font-size:0.85rem;color:#212121;">
            <strong style="color:#1B5E20;">&#128204; Real case from 2024:</strong>
            An African consultancy&rsquo;s final donor report was rejected three times
            for missing M&amp;E data and logframe gaps. 40+ hours of senior staff rework.
            ImpactProof catches these issues before they reach your donor.
          </p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.caption(
        "Your scoring data stays in your browser — never stored on our servers."
    )
    _demo_scenario = st.radio(
        "Sample scenario:",
        options=["Health & Nutrition / USAID", "Agriculture & Livelihoods / GIZ"],
        key="_demo_scenario_choice",
        horizontal=True,
    )
    _footer_c1, _footer_c2 = st.columns(2)
    with _footer_c1:
        if st.button("Pricing →", key="cta_pricing", use_container_width=True,
                     help="View plans and pricing"):
            st.session_state["_show_pricing"] = True
            st.rerun()
    with _footer_c2:
        if st.button("🚀 Try with a sample →", key="cta_demo",
                     help="Loads a realistic example — runs in seconds",
                     use_container_width=True):
            _reset_all_slots()
            _use_agric = "Agric" in _demo_scenario
            _demo_sub = _DEMO_SUBMISSION_AGRIC if _use_agric else _DEMO_SUBMISSION
            _demo_sel = _DEMO_SELECT_FIELDS_AGRIC if _use_agric else _DEMO_SELECT_FIELDS
            for _k, _v in _demo_sub.items():
                st.session_state[_k] = _v
            for _k, _v in _demo_sel.items():
                st.session_state[_k] = _v
            st.session_state["_form_is_resumption"] = False
            if not st.session_state.get("has_seen_tutorial"):
                st.session_state["tutorial_step"] = 1
            st.query_params["demo"] = "1"
            _go_to_screen(1)

    with st.expander("More about this tool"):
        st.markdown(
            """
            <div style="border-radius:8px; background:#F1F8E9; border-left:4px solid #1B5E20;
                        padding:10px 16px; margin:10px 0; font-size:0.85rem; color:#374151;">
              <strong style="color:#1B5E20;">🔒 Your data stays in your browser.</strong>
              We never store your result statements, evidence, or uploaded documents on our servers.
              Your session ends when you close the tab &mdash; nothing is kept.<br>
              <strong style="color:#1B5E20;">&#10003; This is a self-check, not an audit.</strong>
              It helps you improve your result before submission.
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div style="border-left:4px solid #8A6500; border-radius:8px; padding:12px 16px; margin:10px 0; background:transparent;">
              <p style="margin:0; font-size:0.9rem; color:#212121;">
                <strong>Donors now ask:</strong> What changed? How do you know? How strong is the evidence?
                What did you learn? &mdash; <em>This check tells you before they do.</em>
              </p>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div style="border-left:4px solid #1565C0; border-radius:8px; padding:12px 16px; margin:10px 0; background:#F3F8FE;">
              <p style="margin:0 0 6px; font-size:0.9rem; color:#212121;">
                <strong style="color:#1565C0;">Why not just use ChatGPT?</strong>
              </p>
              <p style="margin:0 0 4px; font-size:0.85rem; color:#374151;">
                ChatGPT generates suggestions. ImpactProof makes determinations. The difference:
              </p>
              <ul style="margin:4px 0 0 16px; font-size:0.85rem; color:#374151; padding:0;">
                <li><strong>ImpactProof decides — ChatGPT suggests.</strong> Our scoring engine determines evidence quality against named donor standards and routes you to your highest-impact fix. A chatbot gives you a paragraph. We give you a determination and a ranked fix queue.</li>
                <li>Every determination traces to a <strong>named standard</strong> (USAID ADS 201.3.5.7, Bond 2024, FCDO)</li>
                <li>The same result <strong>always produces the same determination</strong> — no LLM randomness</li>
                <li>The output is a <strong>citable Readiness Card PDF with a reference ID</strong>, not a chat screenshot</li>
                <li>Audit My Report determines <strong>10+ results in 60 seconds</strong> against the same rubric — consistently</li>
              </ul>
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.markdown(
            """
            <div class="gtm-card">
              <p><strong>Want a deeper check?</strong></p>
              <p class="gtm-sub">I personally review results for MEL teams before their submission deadline.</p>
            </div>
            """,
            unsafe_allow_html=True,
        )
        # Landing review WhatsApp CTA — server-side notification (council XXIV)
        _lr_email  = st.session_state.get("user_email", "")
        _lr_wa_key = "wa_landing_review_clicked"
        if st.button("📱 Book a free first review with the founder",
                     key="wa_landing_review_btn", use_container_width=True):
            from utils.whatsapp import notify_founder
            notify_founder("landing_review", user_email=_lr_email)
            st.session_state[_lr_wa_key] = True
        if st.session_state.get(_lr_wa_key):
            from utils.whatsapp import build_wa_url
            st.link_button("Open WhatsApp →",
                           build_wa_url("landing_review", _lr_email),
                           use_container_width=True)
            st.success("✓ Notified — the founder will reach out within 24 hours.")

        _pq2_email  = st.session_state.get("user_email", "")
        _pq2_wa_key = "wa_s0_questions_clicked"
        st.caption("💬 Questions? MEL practitioner in Accra, built this to close a gap I kept hitting.")
        if st.button("Chat on WhatsApp →", key="wa_s0_questions_btn", use_container_width=True):
            from utils.whatsapp import notify_founder
            notify_founder("pricing_questions", user_email=_pq2_email)
            st.session_state[_pq2_wa_key] = True
        if st.session_state.get(_pq2_wa_key):
            from utils.whatsapp import build_wa_url
            st.link_button("Open WhatsApp +233 50 364 8195 →",
                           build_wa_url("pricing_questions", _pq2_email),
                           use_container_width=True)

        st.caption("📄 Already have a draft report? Use ⚡ Instant Report Check inside the form to upload it — AI pre-fills all fields. (Paid feature.)")

    _render_tagline_footer()


# ---------------------------------------------------------------------------
# Screen 1 — Submission Form
# ---------------------------------------------------------------------------

def render_screen_1():
    _cur_tab = st.session_state.get("current_tab", 0)
    render_pitch_strip(["enter", "logframe", "evidence", "review"][_cur_tab])

    if st.session_state.pop("_scroll_to_content", False):
        import streamlit.components.v1 as components
        components.html(
            """<script>
            (function() {
                var p = window.parent;
                var done = false;
                function doScroll() {
                    if (done) return;
                    try { p.scrollTo(0,0); } catch(e) {}
                    try { p.document.documentElement.scrollTop = 0; } catch(e) {}
                    try { p.document.body.scrollTop = 0; } catch(e) {}
                    try { var m=p.document.querySelector('[data-testid="stMain"]'); if(m) m.scrollTop=0; } catch(e) {}
                }
                var timer = null;
                function onMutation() {
                    clearTimeout(timer);
                    timer = setTimeout(function() {
                        doScroll();
                        done = true;
                        try { obs.disconnect(); } catch(e) {}
                    }, 350);
                }
                try {
                    var obs = new MutationObserver(onMutation);
                    obs.observe(p.document.body, {childList:true, subtree:true});
                    doScroll();
                    setTimeout(function() {
                        doScroll();
                        done = true;
                        try { obs.disconnect(); } catch(e) {}
                    }, 2500);
                } catch(e) {
                    var n=0, iv=setInterval(function(){ doScroll(); if(++n>=15) clearInterval(iv); },100);
                }
            })();
            </script>""",
            height=1,
        )

    _render_tutorial(1)

    # Your progress auto-saves every render (see _save_draft() call at the end of
    # this function) — no manual "Save" action needed. Download button for an
    # actual on-disk copy lives on Tab 4.
    if st.session_state.get("_last_saved_time"):
        st.caption(f"💾 Auto-saved at {st.session_state['_last_saved_time']} — download a copy on the last tab to keep it across sessions.")

    _has_prefill = any(
        _ss_str(k).strip()
        for k in ("result_statement", "target_group", "timeframe",
                   "geographic_scope", "evidence_description")
    )

    if st.session_state.pop("_payment_success", False):
        st.success("✅ Payment confirmed! Upload your document below to run the Instant Report Check — or fill in the form manually.")

    # Auto-save draft to Supabase so a page refresh doesn't lose work
    _email_for_draft = st.session_state.get("user_email", "")
    if _email_for_draft and _has_prefill:
        try:
            from datetime import datetime as _dt
            _draft_str = _build_inputs_json(_dt.now().strftime("%Y%m%d_%H%M%S"))
            save_user_draft(_email_for_draft, _draft_str)
        except Exception:
            pass

    _is_resumption = st.session_state.get("_form_is_resumption", False)
    _cloud_restored = st.session_state.pop("_draft_restored_from_cloud", False)
    if _has_prefill and _is_resumption:
        if _cloud_restored:
            st.success("✅ Your session was restored — we saved your form before the page refreshed. Pick up where you left off.")
        else:
            st.info("📂 Continuing from a previous session — your form fields are pre-filled.")
        if st.session_state.get("_confirm_clear_prefill"):
            st.warning("This will clear all pre-filled fields. Are you sure?")
            _cc1, _cc2 = st.columns(2)
            with _cc1:
                if st.button("Yes, clear fields", key="confirm_clear_prefill_yes", type="primary", use_container_width=True):
                    st.session_state.pop("_confirm_clear_prefill", None)
                    st.session_state["_form_is_resumption"] = False
                    _reset_all_slots()
                    st.rerun()
            with _cc2:
                if st.button("Cancel", key="confirm_clear_prefill_no", use_container_width=True):
                    st.session_state.pop("_confirm_clear_prefill", None)
                    st.rerun()
        else:
            if st.button("Clear and start fresh", key="clear_prefill"):
                st.session_state["_confirm_clear_prefill"] = True
                st.rerun()

    active = st.session_state.get("active_slots", 1)

    # --- UX: DYNAMIC SIDEBAR (v3.2) ---
    with st.sidebar:
        st.markdown("### Submission Summary")
        _sb_s = _slot_suffix(1)

        def _sb_field(label, key, trunc=None):
            val = str(st.session_state.get(key, "")).strip()
            if trunc and len(val) > trunc:
                val = val[:trunc] + "…"
            if val and val not in ("(No donor specified)", "(No sector selected)"):
                st.markdown(f"**{label}:** {val}")
            else:
                st.caption(f"{label}: —")

        _sb_field("Result",       "result_statement",  80)
        _sb_field("Target Group", "target_group")
        _sb_field("Geography",    "geographic_scope")
        _sb_field("Timeframe",    "timeframe")
        _sb_field("Donor",        "donor_selected")
        _logframe_skipped = st.session_state.get("logframe_fill_later", False)
        if _logframe_skipped:
            st.caption("Logframe: Skipped (fill before scoring)")
        else:
            _sb_field("Indicator", "logframe_indicator", 60)
        _t = st.session_state.get("logframe_target", "")
        _a = st.session_state.get("logframe_achievement", "")
        if _t or _a:
            _t_d = _t if _t else "—"; _a_d = _a if _a else "—"
            st.markdown(f"**Target → Actual:** {_t_d} → {_a_d}")
        else:
            st.caption("Target → Actual: —")
        _sb_field("Evidence Type", f"evidence_type{_sb_s}")
        _sb_field("Verifier",      f"verifier{_sb_s}")

        st.divider()
        try:
            _sb_sub = _build_submission_from_session(1)
            _sb_ev  = _evaluator.evaluate_submission(_sb_sub)
            _sb_c   = _sb_ev.get("raw_confidence_score", 0)
            _sb_cl  = _sb_ev.get("clarity_score", 0)
            _sb_min = min(_sb_c, _sb_cl)
            _sb_threshold = _sb_ev.get("threshold_used", 4.0)
            _sbe    = "🟢" if _sb_min >= _sb_threshold else "🟡" if _sb_min >= 3.0 else "🔴"
            st.markdown(f"{_sbe} Confidence **{_sb_c}** · Clarity **{_sb_cl}**")
        except Exception:
            st.caption("Fill in the form to see live scores")

        st.divider()
        if st.button("📊 Portfolio analysis →", key="sidebar_portfolio_cta", use_container_width=True):
            _go_to_screen(3)
    # --- END UX: DYNAMIC SIDEBAR (v3.2) ---


    st.markdown(
        f'<div class="md-tab-header">'
        f'<span class="tab-num">Step {_cur_tab + 1} of 4</span>'
        f'<span class="tab-name">{_UX_TAB_NAMES[_cur_tab]}</span>'
        f'</div>',
        unsafe_allow_html=True,
    )

    if _cur_tab == 0:

        # Org type gate — drives tiered threshold and community-appropriate review options
        def _on_org_type_change():
            _n_active = st.session_state.get("active_slots", 1)
            for _idx in range(1, _n_active + 1):
                _sf = "" if _idx == 1 else f"_{_idx}"
                st.session_state.pop(f"internal_review{_sf}", None)
                st.session_state.pop(f"external_review{_sf}", None)

        st.selectbox(
            "Organisation type",
            key="org_type",
            options=[
                "International NGO (INGO)",
                "National NGO",
                "Community-Based Organisation (CBO)",
                "Government department / local authority",
            ],
            help=(
                "Sets the evidence quality standard used for your determination. "
                "INGO standard (4.0 threshold) for bilateral donors such as USAID and FCDO. "
                "National / community standard (3.5–3.75) for national funders and district grants."
            ),
            on_change=_on_org_type_change,
        )

        # Sector selector always visible — gates placeholder quality for all fields below
        st.selectbox(
            "Sector (optional — tailors field examples)",
            key="sector",
            options=SECTOR_OPTIONS,
            help="Select your sector to see sector-specific example placeholders in the evidence description field.",
            on_change=lambda: st.session_state.pop("_sector_auto_inferred", None),
        )
        if st.session_state.get("_sector_auto_inferred"):
            st.caption("⚡ Sector auto-detected from your result text — change here if incorrect.")
        _sector_val = st.session_state.get("sector", SECTOR_OPTIONS[0])
        if _sector_val == "Other":
            st.text_input(
                "Specify your sector",
                key="sector_other",
                placeholder="e.g., Disaster Response, Gender Equality, Financial Inclusion",
            )

        with st.expander("Context — donor & project (optional)", expanded=not _has_prefill):
            st.selectbox(
                "Primary donor for this submission",
                key="donor_selected",
                options=["(No donor specified)", "World Bank", "USAID", "Global Fund", "Mastercard Foundation", "FCDO", "EU / EuropeAid", "AfDB", "GIZ", "SIDA", "RVO", "KOICA", "SDC", "Other"],
                index=0,
                help="Select your primary donor to receive tailored reporting tips and donor-specific diagnostic guidance.",
                on_change=lambda: st.session_state.pop("_donor_auto_inferred", None),
            )
            if st.session_state.get("_donor_auto_inferred"):
                st.caption("⚡ Donor auto-detected from your result text — change here if incorrect.")
            st.selectbox(
                "Donor reporting framework for the crosswalk table",
                key="donor_framework",
                options=list(DONOR_PROFILES.keys()),
                format_func=lambda k: DONOR_PROFILES[k]["label"],
                index=0,
                help="Choose which audit framework's standards are shown alongside each sub-score in the donor framework crosswalk.",
            )
            st.text_input(
                "Project name (optional — used to group/compare indicators in Trends over time)",
                key="project_name",
                placeholder="e.g., Northern Region WASH Programme",
            )
            _donor_val = st.session_state.get("donor_selected", "(No donor specified)")
            if _donor_val == "Other":
                st.text_input(
                    "Specify donor name",
                    key="donor_other",
                    placeholder="e.g., DFID, KfW, Bill & Melinda Gates Foundation",
                )
            _PARTIAL_DONORS = {"GIZ", "World Bank", "AfDB", "KOICA", "SDC"}
            if _donor_val in _PARTIAL_DONORS:
                st.caption(
                    f"Note: {_donor_val} diagnostic guidance covers 2 scoring dimensions — "
                    "full dimension coverage (Recency, Beneficiary Voice) coming in a future update."
                )
            if _donor_val in DONOR_GUIDANCE:
                _dg = DONOR_GUIDANCE[_donor_val]
                with st.expander(f"💡 {_donor_val} reporting tips (3 tips)", expanded=False):
                    st.markdown(f"**Key emphasis:** {_dg['key_emphasis']}")
                    st.markdown(f"**Most common rejection:** {_dg['common_rejection']}")
                    st.markdown(f"**Tip:** {_dg['tip']}")
            st.session_state["remembered_sector"] = st.session_state.get("sector", "")
            st.session_state["remembered_donor"]  = st.session_state.get("donor_selected", "")
        # --- IRC fill summary banner (shown once after extraction) ---
        _irc_summary = st.session_state.pop("_irc_summary", None)
        if _irc_summary:
            _filled   = _irc_summary.get("filled", 0)
            _skipped  = _irc_summary.get("skipped", "")
            _pages    = _irc_summary.get("pages", [])
            _conf_note = _irc_summary.get("confidence_note", "")
            _cgaps     = _irc_summary.get("compliance_gaps", "")
            _page_note = f" · extracted from {len(_pages)} page{'s' if len(_pages)!=1 else ''}" if _pages else ""

            # D2 & D5: extraction quality panel — always shown, adapts to fill count
            _TOTAL_IRC_FIELDS = 12  # approximate total extractable fields
            if _filled == 0:
                st.warning(
                    "⚠️ **IRC found 0 fields** — the document couldn't be read. "
                    "You have not been charged. A free retry has been added to your account.\n\n"
                    "**Try:** Export your report as Word (.docx) and re-upload, or fill the form below."
                )
            elif _filled < 4:
                st.warning(
                    f"⚡ **IRC found {_filled} of ~{_TOTAL_IRC_FIELDS} fields**{_page_note}. "
                    f"Several fields were left blank — review and fill the gaps below."
                )
                if _skipped:
                    st.caption(f"Not found in document: {_skipped}")
            else:
                st.success(
                    f"⚡ **IRC found {_filled} of ~{_TOTAL_IRC_FIELDS} fields**{_page_note}. "
                    "Review the pre-filled fields using the stage buttons at the top, then submit."
                )
                if _skipped:
                    st.info(f"ℹ️ Left blank (not found): {_skipped}")

            if _conf_note:
                st.caption(f"ℹ️ {_conf_note}")
            if _cgaps:
                st.warning(f"⚠️ Compliance fields not found: {_cgaps}")

            # D3: retry button — always offer it after IRC
            _n_active = st.session_state.get("active_slots", 1)
            if _n_active == 1:
                st.caption(
                    "📋 **More results in your report?** Click **＋ Add Another Result** below, "
                    "then re-run IRC targeting a different result."
                )
        # --- END IRC extraction quality panel ---

        col_h, col_add = st.columns([5, 1])
        with col_h:
            label = "Tell us about your result" if active == 1 else f"Tell us about your results ({active} added)"
            st.markdown(f"## {label}")
        with col_add:
            if active < 6:
                st.markdown("<div style='padding-top:22px'></div>", unsafe_allow_html=True)
                if st.button("＋ Add Another Result", use_container_width=True,
                             help="Add up to 6 results to this submission."):
                    st.session_state["active_slots"] = active + 1
                    st.rerun()

        # --- UX: INSTANT REPORT CHECK (v3.2) ---
        _entry_mode = st.session_state.get("entry_mode", "✍️ Fill in manually")
        st.markdown("**How would you like to fill in the form?**")
        _em_col1, _em_col2 = st.columns(2)
        with _em_col1:
            if st.button(
                "✍️ Fill in manually",
                use_container_width=True,
                type="primary" if _entry_mode == "✍️ Fill in manually" else "secondary",
                key="btn_fill_manual",
                help="Type your result details directly into each field.",
            ):
                st.session_state["entry_mode"] = "✍️ Fill in manually"
                st.rerun()
        with _em_col2:
            if st.button(
                "⚡ Instant Report Check (paid)",
                use_container_width=True,
                type="primary" if _entry_mode == "⚡ Instant Report Check" else "secondary",
                key="btn_irc",
                help="Upload your draft report — AI extracts and fills all fields automatically.",
            ):
                st.session_state["entry_mode"] = "⚡ Instant Report Check"
                st.rerun()
        if _entry_mode == "⚡ Instant Report Check":
            if not st.session_state.get("user_email"):
                st.markdown("#### 📧 Enter your email to unlock Instant Report Check")
                _render_email_gate_inline("_irc")
            with st.expander(
                "⚡ Instant Report Check — Upload your draft report to auto-fill this form",
                expanded=not st.session_state.get("_irc_used", False),
            ):
                st.caption(
                    "Upload your donor report to auto-fill all fields — AI extracts only what's "
                    "in the document, never invents. One result at a time."
                )
                _irc_paid_flag = (st.session_state.get("is_paid") or
                                  is_still_paid(get_user(st.session_state.get("user_email",""))))
                _irc_files = []
                if not _irc_paid_flag:
                    st.markdown("### Save 10+ minutes on every result:")
                    st.markdown(
                        "- **Upload your report** — AI reads it and pre-fills all form fields instantly\n"
                        "- **Fills every field your document contains** — skips only what isn't there, flags it clearly\n"
                        "- **Honest extraction** — only fills what's in your document, never invents\n\n"
                        f"*GHS {PRICE_PER_CHECK_GHS/100:.0f} per check · or GHS {PRICE_MONTHLY_GHS/100:.0f}/month for unlimited checks + IRC*\n\n"
                        "💡 *GHS 50/month vs. GHS 12,000–17,000 in rework costs from a donor-queried report.*"
                    )
                    _render_paywall(irc_context=True, prompt_context="irc_attempt")
                else:
                    # Multi-result extraction UI
                    _irc_n = st.radio(
                        "How many results to extract?",
                        options=[1, 2, 3],
                        index=0,
                        horizontal=True,
                        key="irc_n_results",
                        help="Extract 1–3 results from your uploaded document in a single run.",
                    )
                    _irc_hint1 = st.text_input(
                        "Result 1 hint (optional)",
                        key="irc_result_hint",
                        placeholder='e.g. "ANC attendance" or "Output 2.1"',
                        help="Name or describe the result you want extracted. Leave blank to use the most prominent one.",
                    )
                    if _irc_n >= 2:
                        _irc_hint2 = st.text_input(
                            "Result 2 hint (optional)",
                            key="irc_result_hint_2",
                            placeholder='e.g. "safe water access" or "Output 3.2"',
                        )
                    if _irc_n >= 3:
                        _irc_hint3 = st.text_input(
                            "Result 3 hint (optional)",
                            key="irc_result_hint_3",
                            placeholder='e.g. "youth employment" or "Output 4.1"',
                        )
                    _irc_files = st.file_uploader(
                        "Upload report file(s) (or a previously downloaded draft.json)",
                        type=["pdf", "docx", "txt", "csv", "pptx", "xlsx", "xls", "json"],
                        key="instant_report_upload",
                        accept_multiple_files=True,
                    )
                    # --- track and display when each file was added ---
                    _irc_upload_times = st.session_state.setdefault("_irc_upload_times", {})
                    for _f in _irc_files:
                        _fkey = f"{_f.name}_{_f.size}"
                        if _fkey not in _irc_upload_times:
                            _irc_upload_times[_fkey] = datetime.now().strftime("%H:%M:%S")
                    if _irc_files:
                        for _f in _irc_files:
                            _fkey = f"{_f.name}_{_f.size}"
                            st.caption(f"📄 {_f.name} — added {_irc_upload_times.get(_fkey, '')}")
                _irc_run_clicked = (_irc_paid_flag and st.button("🔍 Run Instant Check", key="run_instant_check")
                                     and bool(_irc_files))
                _irc_is_draft_json = (_irc_run_clicked and len(_irc_files) == 1
                                       and _irc_files[0].name.lower().endswith(".json"))

                if _irc_is_draft_json:
                    # --- v3.4: returning user re-upload of a previously downloaded draft ---
                    _irc_file = _irc_files[0]
                    try:
                        _irc_file.seek(0)
                        _draft_data = json.loads(_irc_file.read())
                        if not ("slots" in _draft_data or "active_slots" in _draft_data
                                or "result_statement" in _draft_data):
                            st.error("This JSON doesn't look like an ImpactProof draft. "
                                     "Please upload a file downloaded via 'Download Draft (JSON)' "
                                     "or 'Save Inputs (JSON)'.")
                        else:
                            _load_from_inputs_json(_draft_data)
                    except Exception as _draft_exc:
                        st.error(f"Could not read the draft file: {_draft_exc}")
                    # --- END v3.4 ---
                elif _irc_run_clicked and not _safe_rate_limit_ok(
                    st.session_state.get("user_email", ""), "irc_extraction", max_count=20, window_seconds=3600
                ):
                    st.warning("You've run a lot of Instant Report Checks in the last hour — please wait a bit before running more.")
                elif _irc_run_clicked:
                    _safe_log_access(st.session_state.get("user_email", ""), "irc_extraction")
                    _irc_should_rerun = False
                    with st.spinner("Reading your document(s) and pre-filling the form…"):
                        try:
                            # Step 1: extract and combine raw text from all uploaded documents
                            _doc_files = [f for f in _irc_files if not f.name.lower().endswith(".json")]
                            _full_text, _raw_fields, _ext_err = _irc_extract_combined(_doc_files)
                            if _ext_err:
                                st.warning(_ext_err)
                                st.stop()
                            else:
                                # Step 2: Claude API extraction
                                try:
                                    _irc_key = st.secrets.get("ANTHROPIC_API_KEY")
                                except Exception:
                                    _irc_key = None
                                _irc_key = _irc_key or os.environ.get("ANTHROPIC_API_KEY")
                                if not _irc_key:
                                    st.error("⚠️ ANTHROPIC_API_KEY not set. Rule-based extraction only.")
                                    # Fallback to rule-based
                                    _irc_filled = 0
                                    for _ef, _sf in _IRC_FIELD_MAP.items():
                                        _v = _raw_fields.get(_ef, "")
                                        if _v: st.session_state[_sf] = _v; _irc_filled += 1
                                    if _irc_filled:
                                        st.session_state["_irc_summary"] = {
                                            "filled": _irc_filled,
                                            "skipped": "",
                                            "confidence_note": "",
                                            "compliance_gaps": "",
                                        }
                                        st.session_state["_tab2_auto_advanced"] = True
                                        st.session_state["_irc_used"] = True
                                        st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                                        _irc_should_rerun = True
                                else:
                                    _irc_client = _anthropic.Anthropic(api_key=_irc_key)
                                    # Build few-shot block
                                    _fewshot = {}
                                    _irc_sector = st.session_state.get("sector", "Other")
                                    for _ff in ["result_statement","target_group","timeframe",
                                                "geographic_scope","logframe_indicator",
                                                "logframe_target","logframe_achievement","evidence_description"]:
                                        _fex = get_examples(_ff, _irc_sector, k=3)
                                        if _fex: _fewshot[_ff] = _fex
                                    import json as _ijsonfs
                                    _fewshot_str = _ijsonfs.dumps(_fewshot, indent=2) if _fewshot else ""
                                    _irc_n_res = st.session_state.get("irc_n_results", 1)
                                    _irc_hint1 = st.session_state.get("irc_result_hint", "").strip()
                                    _irc_hint2 = st.session_state.get("irc_result_hint_2", "").strip()
                                    _irc_hint3 = st.session_state.get("irc_result_hint_3", "").strip()
                                    _hints = [h for h in [_irc_hint1, _irc_hint2, _irc_hint3] if h]
                                    if _irc_n_res > 1:
                                        _hint_prefix = (
                                            f"MULTI-RESULT INSTRUCTION: Extract {_irc_n_res} distinct results from this document "
                                            f"and return them as a JSON array under a top-level 'results' key. "
                                            f"Each element in 'results' should follow the same schema as a single result_basics + logframe_linkage + evidence_verification object.\n"
                                        )
                                        if _hints:
                                            _hint_prefix += f"Focus on these results (in order): {'; '.join(f'{i+1}. {h}' for i,h in enumerate(_hints))}.\n"
                                        _hint_prefix += "\n"
                                    else:
                                        _hint_prefix = (
                                            f"FOCUS INSTRUCTION: The user wants to extract the result about: {_irc_hint1}\n"
                                            f"Prioritise this result statement over others in the document.\n\n"
                                        ) if _irc_hint1 else ""
                                    _irc_msgs = [{"role": "user", "content": [
                                        *([{"type":"text","text":f"Field examples for better extraction:\n{_fewshot_str}"}] if _fewshot_str else []),
                                        {"type":"text","text":f"{_hint_prefix}Extract all fields from this report:\n\n{_full_text[:60000]}"}
                                    ]}]
                                    # Run the AI call in a background thread so we can
                                    # show a live elapsed-time indicator while it works.
                                    _irc_timer_ph = st.empty()
                                    _irc_api_result = {}

                                    def _irc_call_api():
                                        try:
                                            _irc_api_result["resp"] = _irc_client.messages.create(
                                                model="claude-sonnet-4-6",
                                                max_tokens=4096,
                                                system=INSTANT_CHECK_SYSTEM_PROMPT,
                                                messages=_irc_msgs,
                                            )
                                        except Exception as _irc_api_exc:
                                            _irc_api_result["error"] = _irc_api_exc

                                    _irc_thread = threading.Thread(target=_irc_call_api, daemon=True)
                                    _irc_start_t = time.time()
                                    _irc_thread.start()
                                    while _irc_thread.is_alive():
                                        _irc_elapsed = time.time() - _irc_start_t
                                        _irc_timer_ph.info(f"🤖 Extracting fields with AI… {_irc_elapsed:.0f}s elapsed (usually 10-30s)")
                                        time.sleep(0.5)
                                    _irc_thread.join()
                                    _irc_timer_ph.empty()
                                    if "error" in _irc_api_result:
                                        raise _irc_api_result["error"]
                                    _irc_resp = _irc_api_result["resp"]
                                    import json as _ijson3
                                    _irc_raw = (_irc_resp.content[0].text if _irc_resp.content else "").strip()
                                    if _irc_raw.startswith("```"):
                                        _irc_parts = _irc_raw.split("```")
                                        if len(_irc_parts) >= 3:
                                            _irc_raw = _irc_parts[1].lstrip("json\n").strip()
                                    if not _irc_raw:
                                        raise ValueError("Model returned an empty response. Try a different document or fill the form manually.")
                                    _irc_data = _ijson3.loads(_irc_raw)

                                    # Multi-result: if response has a 'results' array, process each slot
                                    _irc_results_arr = _irc_data.get("results")
                                    if isinstance(_irc_results_arr, list) and len(_irc_results_arr) > 1:
                                        # Multi-result path: set active_slots, fill each slot
                                        _n_extracted = min(len(_irc_results_arr), 6)
                                        st.session_state["active_slots"] = _n_extracted
                                        _total_filled = 0
                                        for _ri, _ritem in enumerate(_irc_results_arr[:_n_extracted]):
                                            _rs = _slot_suffix(_ri + 1)
                                            _r_rb  = _ritem.get("result_basics", _ritem)
                                            _r_ll  = _ritem.get("logframe_linkage", {})
                                            _r_ev3 = _ritem.get("evidence_verification", {})
                                            for _fk, _fv in [
                                                (f"result_statement{_rs}", _r_rb.get("result_statement")),
                                                (f"target_group{_rs}",     _r_rb.get("target_group")),
                                                (f"timeframe{_rs}",        _r_rb.get("timeframe")),
                                                (f"geographic_scope{_rs}", _r_rb.get("geographic_scope")),
                                                (f"logframe_indicator{_rs}", _r_ll.get("indicator_name")),
                                                (f"logframe_target{_rs}",    _r_ll.get("original_target")),
                                                (f"logframe_achievement{_rs}", _r_ll.get("actual_achievement")),
                                                (f"evidence_description{_rs}", _r_ev3.get("evidence_description")),
                                            ]:
                                                _sval = ", ".join(str(v) for v in _fv) if isinstance(_fv, list) else str(_fv or "").strip()
                                                if _sval and _sval != "Not found":
                                                    st.session_state[_fk] = _sval
                                                    _total_filled += 1
                                        st.session_state["_irc_summary"] = {
                                            "filled": _total_filled,
                                            "skipped": "",
                                            "confidence_note": f"{_n_extracted} results extracted from your document.",
                                            "compliance_gaps": "",
                                            "pages": [],
                                        }
                                        st.session_state["_irc_used"] = True
                                        st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                                        _irc_should_rerun = True
                                        # Skip the single-result handler below
                                        _irc_data = {}

                                    _rb  = _irc_data.get("result_basics", {})
                                    _ll  = _irc_data.get("logframe_linkage", {})
                                    _ev3 = _irc_data.get("evidence_verification", {})
                                    _em  = _irc_data.get("extraction_metadata", {})
                                    _irc_filled = 0
                                    _skipped = []

                                    def _irc_to_str(val):
                                        """Coerce any extracted value to a plain string safe for text widgets."""
                                        if isinstance(val, list):
                                            return ", ".join(str(v) for v in val if v not in (None, "", "Not found"))
                                        if isinstance(val, dict):
                                            return ", ".join(f"{k}: {v}" for k, v in val.items())
                                        if val is None:
                                            return ""
                                        return str(val)

                                    def _irc_set(key, val):
                                        nonlocal _irc_filled
                                        try:
                                            sval = _irc_to_str(val)
                                            if sval and sval != "Not found":
                                                st.session_state[key] = sval; _irc_filled += 1
                                            else:
                                                _skipped.append(key)
                                        except Exception:
                                            _skipped.append(key)

                                    # --- Result Basics ---
                                    _irc_set("result_statement", _rb.get("result_statement"))
                                    _irc_set("target_group",     _rb.get("target_group"))
                                    _irc_set("timeframe",        _rb.get("timeframe"))
                                    _irc_set("geographic_scope", _rb.get("geographic_scope"))

                                    # --- Logframe Linkage ---
                                    _irc_set("logframe_indicator",   _ll.get("indicator_name"))
                                    _irc_set("logframe_target",      _ll.get("original_target"))
                                    _irc_set("logframe_achievement", _ll.get("actual_achievement"))

                                    # --- Evidence & Verification ---
                                    _irc_set("evidence_description", _ev3.get("evidence_description"))
                                    _vmt = None
                                    _ev_type_raw = _irc_to_str(_ev3.get("evidence_type",""))
                                    _ev_desc_for_debate = _irc_to_str(_ev3.get("evidence_description",""))
                                    _rs_for_debate = _irc_to_str(_rb.get("result_statement",""))
                                    _debate_result = None
                                    if _irc_key and _ev_desc_for_debate:
                                        # Council XXIII — 5-member debate on the closest-fit evidence type
                                        try:
                                            from council import debate_evidence_type
                                            _debate_result = debate_evidence_type(
                                                description=_ev_desc_for_debate,
                                                result_statement=_rs_for_debate,
                                                evidence_types=EVIDENCE_TYPES[1:],
                                                api_key=_irc_key,
                                            )
                                        except Exception:
                                            _debate_result = None
                                    if _debate_result and _debate_result.get("recommended_type"):
                                        st.session_state["evidence_type"] = _debate_result["recommended_type"]
                                        st.session_state["_ev_type_debate"] = _debate_result
                                        _irc_filled += 1
                                    else:
                                        # Fallback: simple fuzzy match (no API key, or debate failed)
                                        try:
                                            _vmt = _irc_match_option(_ev_type_raw, EVIDENCE_TYPES)
                                            if _vmt:
                                                st.session_state["evidence_type"] = _vmt; _irc_filled += 1
                                            elif _ev_type_raw and _ev_type_raw != "Not found":
                                                st.session_state["evidence_type"] = "Other"
                                                st.session_state["evidence_type_other"] = _ev_type_raw
                                                _irc_filled += 1
                                        except Exception:
                                            pass
                                    _ir_raw = _irc_to_str(_ev3.get("internal_review",""))
                                    _irmt = None
                                    try:
                                        _irmt = _irc_match_option(_ir_raw, INTERNAL_REVIEW_OPTIONS)
                                        if _irmt:
                                            st.session_state["internal_review"] = _irmt; _irc_filled += 1
                                        elif _ir_raw and _ir_raw != "Not found":
                                            st.session_state["internal_review"] = "Other"
                                            st.session_state["internal_review_other"] = _ir_raw
                                            _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _ermt = _irc_match_option(_irc_to_str(_ev3.get("external_review","")), EXTERNAL_REVIEW_OPTIONS)
                                        if _ermt: st.session_state["external_review"] = _ermt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    for _dkk, _skk in [("reporting_period_start","reporting_start"),
                                                        ("reporting_period_end","reporting_end"),
                                                        ("evidence_collection_date","evidence_date")]:
                                        try:
                                            _pdd = _irc_parse_date(_irc_to_str(_ev3.get(_dkk,"")))
                                            if _pdd: st.session_state[_skk] = _pdd; _irc_filled += 1
                                        except Exception:
                                            pass
                                    _ver3 = _irc_to_str(_ev3.get("independent_verifier",""))
                                    if not _ver3 or _ver3 == "Not found":
                                        _ver3 = ""
                                    _irc_set("verifier", _ver3)

                                    # --- Sector ---
                                    try:
                                        _sec_raw = _irc_to_str(_rb.get("sector",""))
                                        if _sec_raw and _sec_raw != "Not found":
                                            _sec_mt = _irc_match_option(_sec_raw, SECTOR_OPTIONS)
                                            if _sec_mt and _sec_mt != SECTOR_OPTIONS[0]:
                                                st.session_state["sector"] = _sec_mt; _irc_filled += 1
                                            else:
                                                st.session_state["sector"] = "Other"
                                                st.session_state["sector_other"] = _sec_raw
                                                _irc_filled += 1
                                    except Exception:
                                        pass

                                    # --- Primary donor ---
                                    try:
                                        _donor_raw = _irc_to_str(_rb.get("primary_donor",""))
                                        if _donor_raw and _donor_raw != "Not found":
                                            _don_mt = _irc_match_option(_donor_raw, ["USAID", "FCDO", "GIZ", "RVO", "World Bank", "AfDB", "EU / EuropeAid"])
                                            if _don_mt:
                                                st.session_state["donor_selected"] = _don_mt; _irc_filled += 1
                                                try:
                                                    _df_mt = _irc_match_option(_don_mt, list(DONOR_PROFILES.keys()))
                                                    if _df_mt: st.session_state["donor_framework"] = _df_mt
                                                except Exception:
                                                    pass
                                            else:
                                                st.session_state["donor_selected"] = "Other"
                                                st.session_state["donor_other"] = _donor_raw
                                                _irc_filled += 1
                                    except Exception:
                                        pass

                                    # --- Submission type ---
                                    _matched_sub_type = None
                                    try:
                                        _sub_raw = _irc_to_str(_rb.get("submission_type",""))
                                        if _sub_raw and _sub_raw != "Not found":
                                            _matched_sub_type = _irc_match_option(_sub_raw, list(SUBMISSION_CHECKLIST.keys()))
                                            if _matched_sub_type:
                                                st.session_state["submission_type"] = _matched_sub_type; _irc_filled += 1
                                    except Exception:
                                        pass

                                    # --- Documents referenced -> tick required checklist items ---
                                    try:
                                        _docs_ref = _irc_data.get("documents_referenced", [])
                                        if _matched_sub_type and isinstance(_docs_ref, list) and _docs_ref:
                                            for _ckey, _clabel in SUBMISSION_CHECKLIST.get(_matched_sub_type, []):
                                                for _docref in _docs_ref:
                                                    _docref_str = _irc_to_str(_docref)
                                                    if _docref_str and _irc_match_option(_docref_str, [_clabel]):
                                                        st.session_state[_ckey] = True
                                                        _irc_filled += 1
                                                        break
                                    except Exception:
                                        pass

                                    # --- Compliance fields ---
                                    try:
                                        _con_raw = _irc_to_str(_ev3.get("consent_documented",""))
                                        if _con_raw and _con_raw != "Not found":
                                            _con_mt = _irc_match_option(_con_raw, list(CONSENT_CHECKLIST_MAP.keys()))
                                            if _con_mt: st.session_state["gov_consent_status"] = _con_mt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _anon_raw = _irc_to_str(_ev3.get("data_anonymised",""))
                                        if _anon_raw and _anon_raw != "Not found":
                                            _anon_mt = _irc_match_option(_anon_raw, list(ANON_CHECKLIST_MAP.keys()))
                                            if _anon_mt: st.session_state["gov_anonymization_status"] = _anon_mt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _law_raw = _irc_to_str(_ev3.get("data_protection_compliant",""))
                                        if _law_raw and _law_raw != "Not found":
                                            _law_mt = _irc_match_option(_law_raw, list(LAW_CHECKLIST_MAP.keys()))
                                            if _law_mt: st.session_state["gov_compliance_law_status"] = _law_mt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _sfg_raw = _irc_to_str(_ev3.get("safeguarding_measures",""))
                                        if _sfg_raw and _sfg_raw != "Not found":
                                            _sfg_mt = _irc_match_option(_sfg_raw, list(SAFEGUARDING_CHECKLIST_MAP.keys()))
                                            if _sfg_mt: st.session_state["gov_safeguarding_status"] = _sfg_mt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _csg_raw = _irc_to_str(_ev3.get("child_safeguarding",""))
                                        if _csg_raw and _csg_raw != "Not found":
                                            _csg_mt = _irc_match_option(_csg_raw, list(CHILD_SAFEGUARDING_CHECKLIST_MAP.keys()))
                                            if _csg_mt: st.session_state["gov_child_safeguarding_status"] = _csg_mt; _irc_filled += 1
                                    except Exception:
                                        pass
                                    try:
                                        _sec_raw2 = _irc_to_str(_ev3.get("secure_data_handling",""))
                                        if _sec_raw2 and _sec_raw2 != "Not found":
                                            _sec_mt2 = _irc_match_option(_sec_raw2, list(SECURE_HANDLING_CHECKLIST_MAP.keys()))
                                            if _sec_mt2: st.session_state["gov_secure_handling_status"] = _sec_mt2; _irc_filled += 1
                                    except Exception:
                                        pass
                                    _irc_set("project_name", _em.get("project_name"))

                                    # --- Donor Readiness inputs (v3.5) ---
                                    _fri = _irc_data.get("funder_readiness_inputs", {})
                                    _irc_set("learning_notes",     _fri.get("learning_and_adaptation"))
                                    _irc_set("limitations_notes",  _fri.get("limitations"))
                                    _irc_set("additional_context", _fri.get("result_owner_and_decision"))
                                    # attribution_contribution and disaggregation_status are
                                    # advisory flags on the report page (Screen 2) — they must
                                    # only be set by explicit user selection, never by IRC
                                    # auto-extraction, to avoid stale values persisting across
                                    # sessions and appearing pre-filled without user confirmation.

                                    # --- Beneficiary voice ---
                                    try:
                                        _bv_raw = _irc_to_str(_rb.get("beneficiary_voice",""))
                                        if _bv_raw and _bv_raw != "Not found":
                                            _bv_mt = _irc_match_option(_bv_raw, _BV_OPTIONS)
                                            if _bv_mt and _bv_mt != _BV_OPTIONS[0]:
                                                st.session_state["beneficiary_voice"] = _bv_mt; _irc_filled += 1
                                    except Exception:
                                        pass

                                    # --- Strengthen this evidence -> tick verifiable checks ---
                                    try:
                                        _esc = _irc_data.get("evidence_strengthening_checks", [])
                                        if _vmt and isinstance(_esc, list) and _esc:
                                            for _ekey, _elabel in EVIDENCE_STRENGTHEN_CHECKLIST.get(_vmt, []):
                                                for _eitem in _esc:
                                                    _eitem_str = _irc_to_str(_eitem)
                                                    if _eitem_str and _irc_match_option(_eitem_str, [_elabel]):
                                                        st.session_state[_ekey] = True
                                                        _irc_filled += 1
                                                        break
                                    except Exception:
                                        pass

                                    # store summary for persistent banner; disable auto-advance
                                    _skip_str3 = ", ".join(_skipped[:6]) if _skipped else ""
                                    _conf3 = _irc_to_str(_em.get("confidence_note",""))
                                    _cgaps = [f for f in ["consent_documented","data_anonymised","data_protection_compliant"] if _ev3.get(f,"") == "Not found"]
                                    _glab = {"consent_documented":"Consent","data_anonymised":"Anonymisation","data_protection_compliant":"Data protection"}
                                    # Extract page references from field_sources
                                    _field_sources = _irc_data.get("field_sources", {})
                                    _page_nums = {k: v.get("page",0) for k,v in _field_sources.items() if isinstance(v, dict) and v.get("page",0) > 0}
                                    _unique_pages = sorted(set(_page_nums.values()))
                                    st.session_state["_irc_field_sources"] = _field_sources
                                    st.session_state["_irc_summary"] = {
                                        "filled": _irc_filled,
                                        "skipped": _skip_str3,
                                        "confidence_note": _conf3 if _conf3 and _conf3 != "Not found" else "",
                                        "compliance_gaps": ", ".join(_glab.get(g,g) for g in _cgaps),
                                        "pages": _unique_pages,
                                    }
                                    # prevent auto-advance from swallowing logframe/evidence tabs
                                    st.session_state["_tab2_auto_advanced"] = True
                                    st.session_state["_irc_used"] = True
                                    st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                                    # Auto-tick logframe_fill_later for slots where IRC found no logframe data
                                    _irc_active = st.session_state.get("active_slots", 1)
                                    for _irc_s in range(1, _irc_active + 1):
                                        _sf = _slot_suffix(_irc_s)
                                        _no_lf = all(
                                            not (st.session_state.get(f"{_k}{_sf}") or "").strip()
                                            for _k in ("logframe_indicator", "logframe_target", "logframe_achievement")
                                        )
                                        if _no_lf:
                                            st.session_state[f"logframe_fill_later{_sf}"] = True
                                    _irc_should_rerun = True
                                    # D1: zero-field extraction = no value delivered → grant free retry
                                    if _irc_filled == 0:
                                        st.session_state["_irc_retry_credit"] = True
                        except Exception as _irc_exc:
                            _exc_str = str(_irc_exc).lower()
                            # D4: classify failure type for targeted guidance
                            if "timeout" in _exc_str or "timed out" in _exc_str:
                                _fail_reason = "timeout"
                            elif "json" in _exc_str or "decode" in _exc_str or "parse" in _exc_str:
                                _fail_reason = "parse"
                            elif "rate" in _exc_str or "429" in _exc_str:
                                _fail_reason = "rate_limit"
                            else:
                                _fail_reason = "unknown"
                            _uploaded_name = (_irc_files[0].name if _irc_files else "").lower()
                            _is_pdf = _uploaded_name.endswith(".pdf")
                            if _fail_reason == "timeout":
                                _fail_msg = ("⚠️ **Extraction timed out** — your document may be too long. "
                                             "Try uploading just the results or outputs section.")
                            elif _fail_reason == "parse":
                                _fail_msg = ("⚠️ **IRC had trouble reading the document structure.** "
                                             "Try re-saving as a plain .docx and re-uploading.")
                            elif _fail_reason == "rate_limit":
                                _fail_msg = "⚠️ **Service busy — please wait 30 seconds and try again.**"
                            elif _is_pdf:
                                _fail_msg = ("⚠️ **IRC couldn't read this PDF.** "
                                             "If it's a scanned document, export as Word (.docx) first. "
                                             "Password-protected PDFs are not supported.")
                            else:
                                _fail_msg = "⚠️ **IRC couldn't extract your document.** Try a different file format or fill in manually below."
                            st.session_state["_irc_retry_credit"] = True
                            st.warning(
                                f"{_fail_msg}\n\n"
                                "**You have not been charged for this attempt.** "
                                "A free retry has been added to your account."
                            )
                    if _irc_should_rerun:
                        st.rerun()
                    # D3: show retry notice if retry credit is active
                    elif st.session_state.get("_irc_retry_credit"):
                        st.info("🔄 **Free retry available** — adjust your hint or upload a different file and run IRC again at no charge.")
        # --- END UX: INSTANT REPORT CHECK (v3.2) ---

        for slot in range(1, active + 1):
            if active > 1:
                st.markdown(f"---\n#### Result {slot}")
            _render_tab1_slot(slot)

        # --- v3.3: auto-advance to Logframe tab when tab1 complete ---
        _t1_done = all([
            _ss_str("result_statement").strip(),
            _ss_str("target_group").strip(),
            _ss_str("timeframe").strip(),
            _ss_str("geographic_scope").strip(),
        ])
        if _t1_done:
            try:
                _q_ev = _evaluator.evaluate_submission(_build_submission_from_session(1))
                _q_clar = _q_ev.get("clarity_score", 0)
                if _q_clar >= 4.0:
                    st.success(f"✓ Result statement looks strong ({_q_clar:.1f}/5.0 Clarity) — keep going.")
                elif _q_clar >= 2.5:
                    st.info(f"Result statement scores {_q_clar:.1f}/5.0 on Clarity — it will improve as you add evidence.")
                else:
                    st.warning(f"Result statement scores {_q_clar:.1f}/5.0 on Clarity — consider sharpening it before proceeding.")
            except Exception:
                pass
            # --- Quick Confidence signal (optional evidence fields) ---
            if active == 1:
                # Pre-fill from IRC-extracted evidence if available
                _irc_ev_desc = _ss_str("evidence_description").strip()
                _irc_ev_type = st.session_state.get("evidence_type", "")
                _quick_has_irc = bool(_irc_ev_desc and _irc_ev_type and _irc_ev_type != "—")
                _expander_label = (
                    "Quick Confidence signal — pre-filled from your uploaded report ✅"
                    if _quick_has_irc else
                    "Get a quick Confidence signal (optional — add your evidence here)"
                )
                with st.expander(_expander_label, expanded=_quick_has_irc):
                    _qe_desc = st.text_input(
                        "Brief evidence description",
                        key="quick_evidence_desc",
                        value=_irc_ev_desc,
                        placeholder="e.g. Monthly DHIS2 data from 3 facilities, reviewed by MEL officer",
                    )
                    _default_type_idx = (EVIDENCE_TYPES.index(_irc_ev_type) + 1
                                         if _irc_ev_type in EVIDENCE_TYPES else 0)
                    _qe_type = st.selectbox(
                        "Evidence type",
                        key="quick_evidence_type",
                        options=["—"] + EVIDENCE_TYPES,
                        index=_default_type_idx,
                    )
                    if _qe_desc.strip() and _qe_type != "—":
                        try:
                            _quick_sub = _build_submission_from_session(1)
                            _quick_sub["evidence"] = [{"type": _qe_type, "description": _qe_desc,
                                                       "recency": "", "verified_by": ""}]
                            _quick_ev2 = _evaluator.evaluate_submission(_quick_sub)
                            _qc = _quick_ev2.get("confidence_score", 0)
                            if _qc >= 3.5:
                                st.success(f"Provisional Confidence: {_qc:.1f}/5.0 ✅ — strong evidence signal.")
                            elif _qc >= 2.0:
                                st.info(f"Provisional Confidence: {_qc:.1f}/5.0 ⚠️ — strengthen your evidence in Tab 3.")
                            else:
                                st.warning(f"Provisional Confidence: {_qc:.1f}/5.0 🔴 — evidence needs significant work.")
                            st.caption("Provisional only — complete Tabs 2–4 for your full scored report.")
                            # Carry this through to Tab 3 instead of discarding it — only
                            # seeds the real fields if they're still untouched there, so it
                            # never overwrites something the user already typed in Tab 3.
                            if not _ss_str("evidence_description").strip():
                                st.session_state["evidence_description"] = _qe_desc
                            if st.session_state.get("evidence_type", EVIDENCE_TYPES[0]) == EVIDENCE_TYPES[0]:
                                st.session_state["evidence_type"] = _qe_type
                            st.session_state["_irc_fill_version"] = st.session_state.get("_irc_fill_version", 0) + 1
                        except Exception:
                            pass
            st.success("✓ Result defined. A reviewer can now check this against your logframe — donor question 1 answered.")
            # QC fast-path: skip Logframe (Tab 1) and land directly on Evidence (Tab 2)
            if st.session_state.pop("_from_quick_check", False):
                for _qc_sl in range(1, active + 1):
                    st.session_state[f"logframe_fill_later{_slot_suffix(_qc_sl)}"] = True
                st.session_state["current_tab"] = 2
                st.session_state["_tab2_auto_advanced"] = True
                st.session_state["_show_qc_tab2_hint"] = True
                st.session_state["_scroll_to_content"] = True
                st.query_params["tab"] = "2"
                st.rerun()
            _nb1, _pb1 = st.columns([3, 1])
            with _nb1:
                if st.button("Next: Logframe Linkage →", key="tab1_next_btn", type="primary", use_container_width=True):
                    st.session_state["current_tab"] = 1
                    st.session_state["_tab2_auto_advanced"] = False
                    st.session_state["_scroll_to_content"] = True
                    st.query_params["tab"] = "1"
                    st.rerun()
            with _pb1:
                if st.button("← Home", key="tab0_home_btn", use_container_width=True, help="Return to the landing page"):
                    _go_to_screen(0)
        else:
            st.caption("Fill in all four fields above to continue.")
        # --- END v3.3 ---

    elif _cur_tab == 1:
        for slot in range(1, active + 1):
            if active > 1:
                st.markdown(f"---\n#### Result {slot}")
            _render_tab2_slot(slot)

        # --- v3.3: next button / auto-advance to Evidence tab when logframe complete ---
        _t2_done = all([
            _ss_str("logframe_indicator").strip(),
            _ss_str("logframe_target").strip(),
            _ss_str("logframe_achievement").strip(),
        ])
        # IRC users OR users who ticked "fill later" may advance with blank logframe fields.
        _fill_later_any = any(
            st.session_state.get(f"logframe_fill_later{_slot_suffix(sl)}", False)
            for sl in range(1, active + 1)
        )
        _t2_can_advance = _t2_done or st.session_state.get("_irc_used", False) or _fill_later_any
        if _t2_can_advance:
            if _t2_done:
                try:
                    _q_ev2 = _evaluator.evaluate_submission(_build_submission_from_session(1))
                    _q_clar2 = _q_ev2.get("clarity_score", 0)
                    if _q_clar2 >= 4.0:
                        st.success(f"✓ Logframe linked. Your result is traceable to an approved commitment. ({_q_clar2:.1f}/5.0 Clarity so far)")
                    elif _q_clar2 >= 2.5:
                        st.info(f"Logframe linkage scores {_q_clar2:.1f}/5.0 so far — strong evidence on the next tab will raise your Confidence score.")
                    else:
                        st.warning(f"Logframe linkage scores {_q_clar2:.1f}/5.0 — check that your indicator and achievement match the result statement above.")
                except Exception:
                    pass
            _nb2, _pb2 = st.columns([3, 1])
            with _nb2:
                if st.button("Next: Evidence & Verification →", key="tab2_next_btn", type="primary", use_container_width=True):
                    st.session_state["current_tab"] = 2
                    st.session_state["_tab2_auto_advanced"] = True
                    st.session_state["_scroll_to_content"] = True
                    st.query_params["tab"] = "2"
                    st.rerun()
            with _pb2:
                if st.button("← Back", key="tab2_back_btn", use_container_width=True):
                    st.session_state["current_tab"] = 0
                    st.session_state["_scroll_to_content"] = True
                    st.query_params["tab"] = "0"
                    st.rerun()
            if not _t2_done:
                st.caption("Some logframe fields weren't in your uploaded report — you can fill them in now or continue and complete them later.")
        else:
            st.caption("Fill in all three logframe fields above to continue.")
        # --- END v3.3 ---

    elif _cur_tab == 2:
        st.caption("Describe your evidence. This is where the system makes its core determination: what your evidence is worth to a donor, and how to improve it.")
        if st.session_state.pop("_show_qc_tab2_hint", False):
            st.success("✓ **Evidence type and verifier pre-filled from Quick Check.** Add your evidence description below to complete the determination.")
        for slot in range(1, active + 1):
            if active > 1:
                st.markdown(f"---\n#### Result {slot}")
            _render_tab3_slot(slot)

        # --- v3.3: next button to Review & Submit ---
        try:
            _q_ev3 = _evaluator.evaluate_submission(_build_submission_from_session(1))
            _q_conf3 = _q_ev3.get("confidence_score", 0)
            if _q_conf3 >= 4.0:
                st.success(f"✓ Evidence documented. Your donor will find this verifiable. ({_q_conf3:.1f}/5.0 Confidence) — moving to final review.")
            elif _q_conf3 >= 2.5:
                st.info(f"Evidence scores {_q_conf3:.1f}/5.0 on Confidence — complete the verification fields to strengthen before review.")
            else:
                st.warning(f"Evidence scores {_q_conf3:.1f}/5.0 on Confidence — strengthen your evidence description or add a verifier before submitting.")
        except Exception:
            pass
        _nb3, _pb3 = st.columns([3, 1])
        with _nb3:
            if st.button("Next: Review & Submit →", key="tab3_next_btn", type="primary", use_container_width=True):
                st.session_state["current_tab"] = 3
                st.session_state["_scroll_to_content"] = True
                st.query_params["tab"] = "3"
                st.rerun()
        with _pb3:
            if st.button("← Back", key="tab3_back_btn", use_container_width=True):
                st.session_state["current_tab"] = 1
                st.session_state["_scroll_to_content"] = True
                st.query_params["tab"] = "1"
                st.rerun()
        # --- END v3.3 ---

    elif _cur_tab == 3:
        st.caption("Final check: See your submission the way your donor will see it — and close any gaps before they do.")

        _REQUIRED_FIELDS_B = [
            ("result_statement",     "Result statement (Tab 1)"),
            ("target_group",         "Target group (Tab 1)"),
            ("timeframe",            "Timeframe (Tab 1)"),
            ("geographic_scope",     "Geographic scope (Tab 1)"),
            ("evidence_description", "Evidence description (Tab 3)"),
            ("evidence_type",        "Evidence type (Tab 3)"),
        ]
        _TAB_IDX_B = {
            "result_statement": 0, "target_group": 0, "timeframe": 0, "geographic_scope": 0,
            "evidence_description": 2, "evidence_type": 2,
        }
        # The Child Safeguarding Alert (Tab 3) implies a hard block -- make it one,
        # rather than a red banner a user can submit straight past.
        if _minors_possibly_involved(1):
            _REQUIRED_FIELDS_B.append(
                ("gov_child_safeguarding_status", "Child safeguarding check (Tab 3 — Data Governance Checklist)")
            )
            _TAB_IDX_B["gov_child_safeguarding_status"] = 2
        _missing_b = [
            (key, lbl) for key, lbl in _REQUIRED_FIELDS_B
            if not str(st.session_state.get(key, "")).strip()
            or st.session_state.get(key, "") in (
                EVIDENCE_TYPES[0], "Choose an option...", "Select child safeguarding status...", ""
            )
        ]
        _completed_b = len(_REQUIRED_FIELDS_B) - len(_missing_b)

        # Auto-save status + draft download at TOP of Tab 3
        _save_draft()
        _t3_col_save, _t3_col_ts = st.columns([1, 2])
        with _t3_col_ts:
            st.caption(f"💾 Auto-saved · Last saved {st.session_state.get('_last_saved_time', '--:--')}")
        with _t3_col_save:
            _draft_bytes_top = st.session_state.get("_draft_bytes", b"")
            if _draft_bytes_top:
                st.download_button(
                    "💾 Save Draft",
                    data=_draft_bytes_top,
                    file_name="impact_receipts_draft.json",
                    mime="application/json",
                    use_container_width=True,
                    key="tab3_save_draft_top",
                )

        st.progress(_completed_b / len(_REQUIRED_FIELDS_B),
                    text=f"Form completion: {_completed_b}/{len(_REQUIRED_FIELDS_B)} required fields")
        if _missing_b:
            with st.expander(f"⚠ {len(_missing_b)} required field(s) incomplete", expanded=True):
                for _fk, _fl in _missing_b:
                    st.markdown(f"- {_fl}")
                if st.button("→ Fix: Jump to First Missing Field", key="jump_missing_b", type="primary"):
                    _first_b = _TAB_IDX_B[_missing_b[0][0]]
                    st.session_state["current_tab"] = _first_b
                    st.session_state["_scroll_to_content"] = True
                    st.query_params["tab"] = str(_first_b)
                    st.rerun()

        # --- Score preview (unconditional — no email needed to see your own scores) ---
        # Compute evaluation ONCE and cache; _render_live_score_preview reuses it (lean: no triple eval)
        try:
            _tab3_sub = _build_submission_from_session(1)
            _tab3_ev  = _evaluator.evaluate_submission(_tab3_sub)
            st.session_state["_tab3_ev_cache"] = {"sub": _tab3_sub, "ev": _tab3_ev}
            _banner_c  = round(_tab3_ev.get("raw_confidence_score", 0) * 20, 1)
            _banner_cl = round(_tab3_ev.get("clarity_score", 0) * 20, 1)
            if _banner_c >= 75 and _banner_cl >= 75:
                st.success("✅ Strong Submission — Your result meets quality thresholds for donor submission.")
                st.caption("Decision: Submission-ready — proceed to generate your Pre-Submission Readiness Card.")
            elif _banner_c >= 50 or _banner_cl >= 50:
                st.warning("⚠️ Submission Needs Work — Address the items below before submitting.")
                st.caption("Decision: Needs work — address the top fix before generating your Readiness Card.")
            else:
                st.error("🔴 High Risk — Your donor will likely query or reject this result. Fix critical issues first.")
                st.caption("Decision: High risk — act on all critical fixes before submission.")
        except Exception:
            st.session_state.pop("_tab3_ev_cache", None)

        # Live score breakdown — uses cached evaluation for slot 1
        for _lsp_slot in range(1, active + 1):
            _render_live_score_preview(_lsp_slot)

        st.divider()

        # Email gate — only gates the Submit action, not score visibility.
        # _render_email_gate_inline calls st.stop() so nothing below renders until email is set.
        _has_email = bool(st.session_state.get("user_email"))
        if not _has_email:
            st.warning("📧 **Enter your email to run your check.** We use it to track your free checks — no password needed.")
            _render_email_gate_inline("_check")
            # st.stop() inside _render_email_gate_inline halts rendering here until email is set

        # Reached only after email is set
        _has_email = True
        _email_now = st.session_state.get("user_email", "")
        _report_allowed = check_access(_email_now)["allowed"]  # DB-authoritative; session state is display-only
        if not _report_allowed:
            st.info(
                "You've used your 3 free checks — scoring this result is still free. "
                "Downloading the Readiness Card PDF or using Instant Report Check will need an upgrade."
            )

        # GDPR consent — shown after email confirmed
        st.checkbox(
            "📚 Allow my anonymised entries to improve extraction quality for other MEL officers. "
            "(Act 843 / NDPA compliant — no names or organisations are stored.)",
            key="consent_examples",
        )
        # Cache for Screen 2 download gate — scoring itself is always free
        st.session_state["_report_allowed"] = _report_allowed
        _sb4, _bb4 = st.columns([3, 1])
        with _bb4:
            if st.button("← Back", key="tab4_back_btn", use_container_width=True):
                st.session_state["current_tab"] = 2
                st.session_state["_scroll_to_content"] = True
                st.query_params["tab"] = "2"
                st.rerun()
        with _sb4:
            pass
        # Submit is always enabled — paywall moves to the Download button on Screen 2
        if st.button(
            "Get Determination & Fix Queue →",
            type="primary",
            use_container_width=True,
        ):
            mandatory = [
                st.session_state.get("result_statement", ""),
                st.session_state.get("target_group", ""),
                st.session_state.get("timeframe", ""),
                st.session_state.get("geographic_scope", ""),
                st.session_state.get("evidence_description", ""),
            ]
            ev_type = st.session_state.get("evidence_type", "")
            ev_other = _ss_str("evidence_type_other").strip()
            if ev_type == "Other" and not ev_other:
                st.warning("Please specify your evidence type in Tab 3 — Evidence & Verification.")
            elif not all(mandatory) or _missing_b:
                _missing_labels = ", ".join(lbl for _, lbl in _missing_b) or "required fields"
                st.warning(f"Please complete the following before running the check: {_missing_labels}.")
            else:
                if not st.session_state.get("has_seen_tutorial"):
                    st.session_state["tutorial_step"] = 2
                for slot in range(1, active + 1):
                    s = _slot_suffix(slot)
                    raw = st.session_state.get(f"uploaded_files_widget{s}") or []
                    st.session_state[f"uploaded_files{s}"] = [f.name for f in raw]
                st.session_state["active_slots_run"] = active
                st.session_state["evaluations"]       = None
                st.session_state["submissions_snapshot"] = None
                # --- Save anonymised examples on consent ---
                if st.session_state.get("consent_examples") and _email_now:
                    _ex_sector = st.session_state.get("sector", "Other")
                    for _ex_field in ["result_statement", "target_group", "timeframe",
                                      "geographic_scope", "logframe_indicator",
                                      "logframe_target", "logframe_achievement",
                                      "evidence_description"]:
                        _ex_val = st.session_state.get(_ex_field, "")
                        _ex_clean = _anonymize_value(_ex_val)
                        if _ex_clean:
                            save_example(_ex_field, _ex_sector, _ex_clean)
                    # Phase B: log provenance answers for future adaptive learning
                    _ev_type_log = st.session_state.get("evidence_type", "")
                    _prov_log_map = {
                        "sampling_documented":     "provenance_sampling",
                        "double_counting_checked": "provenance_dedup",
                        "collection_tool_named":   "provenance_tool",
                        "collector_independent":   "provenance_independence",
                        "recall_period_ok":        "provenance_recall",
                    }
                    for _plk, _plss in _prov_log_map.items():
                        _plv = st.session_state.get(_plss, "")
                        if _plv and _plv not in ("Choose an option...", ""):
                            save_example(f"provenance__{_plk}", _ev_type_log, _plv)
                # --- Track usage ---
                # D1: if user had a failed IRC attempt they're retrying, don't charge again
                _using_retry_credit = st.session_state.pop("_irc_retry_credit", False)
                # D2: if user is re-scoring the same result_statement in this session, don't charge again
                _stmt_key = (st.session_state.get("result_statement", "") or "").strip()[:120]
                _scored_stmts = st.session_state.setdefault("_scored_stmts", set())
                _is_rescore = bool(_stmt_key and _stmt_key in _scored_stmts)
                if _stmt_key:
                    _scored_stmts.add(_stmt_key)
                if _email_now and not _using_retry_credit and not _is_rescore:
                    record_check(_email_now)
                # Clear saved draft — user has submitted, fresh start next time
                if _email_now:
                    try:
                        clear_user_draft(_email_now)
                    except Exception:
                        pass
                # --- End tracking ---
                # --- Metrics: privacy-safe usage instrumentation (no PII, no result text) ---
                try:
                    _m_ev = _evaluator.evaluate_submission(_build_submission_from_session(1))
                    _m_session_id = _metrics_session_id()
                    metrics.log_event(
                        "check_completed", _m_session_id,
                        score_band=_m_ev.get("confidence_label", ""),
                    )
                    _m_prev_conf = st.session_state.get("_last_logged_confidence")
                    if _is_rescore and _m_prev_conf is not None:
                        metrics.log_event(
                            "score_uplift", _m_session_id,
                            score_uplift=round(_m_ev.get("confidence_score", 0) - _m_prev_conf, 2),
                        )
                    st.session_state["_last_logged_confidence"] = _m_ev.get("confidence_score", 0)
                except Exception:
                    pass
                # --- End metrics ---
                st.session_state["screen"] = 2
                st.rerun()

        # Portfolio CTA — contextual: user is about to score a result, may want to do the whole logframe
        if st.button(
            "📊 Run portfolio analysis →",
            key="tab3_portfolio_cta",
            use_container_width=False,
            help="Score all your indicators at once by uploading a logframe CSV.",
        ):
            _go_to_screen(3)

        st.divider()

        # Hibernate option — preserves data, pauses billing intent. Acts immediately:
        # it only saves (already happens automatically every render) and navigates
        # home, so a confirm step protected nothing and just cost an extra click.
        if st.session_state.get("confirm_reset"):
            # Exit survey before clearing
            st.warning("Clear all inputs and start over?")
            _exit_reason = st.selectbox(
                "Quick question: why are you leaving? (helps us improve)",
                ["(Select reason — optional)", "Finished my reporting cycle",
                 "Result doesn't need checking", "Too many fields / too complex",
                 "Switching to a different tool", "Organisation budget cut", "Other"],
                key="exit_survey_reason",
            )
            cf1, cf2 = st.columns(2)
            with cf1:
                if st.button("Yes, clear everything", type="primary", use_container_width=True):
                    if _exit_reason and _exit_reason != "(Select reason — optional)":
                        try:
                            from utils.db import save_example
                            save_example("exit_reason", st.session_state.get("sector","Other"), _exit_reason)
                        except Exception:
                            pass
                    st.session_state["confirm_reset"] = False
                    _clear_draft()
                    _go_to_screen(1, reset=True)
            with cf2:
                if st.button("Cancel", use_container_width=True):
                    st.session_state["confirm_reset"] = False
                    st.rerun()
        else:
            _rst_c1, _rst_c2 = st.columns(2)
            with _rst_c1:
                if st.button("💤 Hibernate (save & pause)", use_container_width=True, help="Saves your form. Come back next reporting cycle."):
                    _save_draft()
                    _go_to_screen(0)
            with _rst_c2:
                if st.button("🗑 Clear all & restart", use_container_width=True):
                    st.session_state["confirm_reset"] = True
                    st.rerun()

    _save_draft()
    _render_tagline_footer()


# ---------------------------------------------------------------------------
# Screen 2 — Confidence Snapshot & Next Steps
# ---------------------------------------------------------------------------

def render_how_scoring_works_panel() -> None:
    """Render a reference panel explaining all eight scoring criteria."""
    with st.expander("How scoring works", expanded=False):
        st.caption(
            "Confidence and Clarity are each built from rule-based sub-criteria. "
            "No AI scoring is used — every score below follows a fixed rubric."
        )
        for axis, axis_title in (("confidence", "Confidence"), ("clarity", "Clarity")):
            st.markdown(f"**{axis_title} criteria**")
            for key, guide in _SCORING_GUIDE.items():
                if guide["axis"] != axis:
                    continue
                label = guide["label"]
                definition = guide["definition"]
                weak = guide["weak_example"]
                strong = guide["strong_example"]
                improve = guide["improve_actions"]
                if key in ("measurement", "definition"):
                    st.markdown(f"**{label} / {guide['qualitative_label']}** (max {guide['max_score']})")
                    st.markdown(f"{definition}")
                    st.caption(guide["qualitative_definition"])
                else:
                    st.markdown(f"**{label}** (max {guide['max_score']})")
                    st.markdown(definition)
                st.caption(f"Why it matters: {guide['why_it_matters']}")
                col_weak, col_strong = st.columns(2)
                with col_weak:
                    st.markdown("*Weak example*")
                    st.caption(weak)
                with col_strong:
                    st.markdown("*Strong example*")
                    st.caption(strong)
                st.markdown("How to improve:")
                for action in improve:
                    st.markdown(f"- {action}")
                if key in ("measurement", "definition"):
                    st.markdown(f"How to improve ({guide['qualitative_label']} — for case studies, "
                                 f"outcome harvesting, beneficiary narratives):")
                    for action in guide["qualitative_improve_actions"]:
                        st.markdown(f"- {action}")
                st.divider()


def render_personalized_weakness_panel(
    direct_score: float, verify_score: float, recency_score: float,
    definition_score: float, measurement_score: float,
    integrity_score: float, scope_score: float, governance_score: float,
    is_qualitative: bool = False,
    n: int = 3,
) -> None:
    """Show the n weakest scoring criteria for this submission with improvement actions."""
    scores = {
        "directness": direct_score,
        "verification": verify_score,
        "recency": recency_score,
        "definition": definition_score,
        "measurement": measurement_score,
        "integrity": integrity_score,
        "scope": scope_score,
        "governance": governance_score,
    }
    axis_order = {"confidence": 0, "clarity": 1}
    ranked = sorted(
        scores.items(),
        key=lambda item: (
            item[1] / _SCORING_GUIDE[item[0]]["max_score"],
            axis_order[_SCORING_GUIDE[item[0]]["axis"]],
        ),
    )
    st.markdown("#### Where to focus first")
    for key, score in ranked[:n]:
        guide = _SCORING_GUIDE[key]
        max_score = guide["max_score"]
        if key in ("measurement", "definition") and is_qualitative:
            label = guide["qualitative_label"]
            improve = guide["qualitative_improve_actions"]
        else:
            label = guide["label"]
            improve = guide["improve_actions"]
        st.warning(f"**{label}** is one of your weakest areas ({score}/{max_score})")
        for action in improve:
            st.markdown(f"- {action}")


def _build_reviewer_info(slot: int) -> dict:
    """Read the current review-handoff fields for one submission slot from
    session state. Read live (not from the evaluation snapshot) so exports
    reflect reviewer input entered after the snapshot was taken."""
    s = _slot_suffix(slot)
    return {
        "review_status":     st.session_state.get(f"review_status{s}", _SUBMISSION_STATUS_OPTIONS[0]),
        "reviewer_name":     st.session_state.get(f"reviewer_name{s}", ""),
        "reviewer_role":     st.session_state.get(f"reviewer_role{s}", ""),
        "reviewer_date":     st.session_state.get(f"reviewer_date{s}", ""),
        "reviewer_decision": st.session_state.get(f"reviewer_decision{s}", ""),
        "reviewer_notes":    st.session_state.get(f"reviewer_notes{s}", ""),
    }


def _render_review_handoff(submission: dict, ev: dict, card_idx: int):
    """Lightweight review-handoff layer (v3.6): per-submission status,
    reviewer name/role/date/decision/notes, and an exportable review summary
    that reuses the existing HTML/PDF/DOCX report builders.

    Known limitation: no user accounts, authentication, or real-time
    multi-user support — reviewer identity is free text within this session,
    consistent with the rest of the app."""
    slot = card_idx + 1
    s = _slot_suffix(slot)

    with st.expander("📋 Team review & sign-off (agency feature)", expanded=False, key=f"team_review_{card_idx}"):
        st.caption(
            "Optional — record a reviewer's decision before handing this result back "
            "to the field officer or passing it up the chain. No accounts or logins: "
            "these fields live only in this browser session unless you click "
            "**Save Draft** below, which writes them (with the rest of this submission) "
            "to inputs/draft.json so they survive a reload."
        )
        rc1, rc2 = st.columns(2)
        with rc1:
            st.selectbox("Status", _SUBMISSION_STATUS_OPTIONS, key=f"review_status{s}")
            st.text_input("Reviewer name", key=f"reviewer_name{s}")
            st.text_input("Reviewer role", key=f"reviewer_role{s}")
        with rc2:
            st.selectbox(
                "Decision", _REVIEW_DECISION_OPTIONS, key=f"reviewer_decision{s}",
                format_func=lambda d: d or "Select decision...",
            )
            st.text_input("Review date", key=f"reviewer_date{s}", placeholder="YYYY-MM-DD")
        st.text_area(
            "Reviewer notes (overall or per criterion)",
            key=f"reviewer_notes{s}", height=100,
        )

        if st.button("💾 Save Draft (incl. review notes)", key=f"review_save_draft_{card_idx}"):
            _save_draft()
            st.toast("Draft saved, including review notes!", icon="💾")

        review_info = _build_reviewer_info(slot)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        summary_html = _build_review_summary_html(submission, ev, review_info, timestamp, chart_id=str(card_idx))

        dl1, dl2, dl3 = st.columns(3)
        with dl1:
            st.download_button(
                "⬇️ Review Summary (HTML)",
                data=summary_html,
                file_name=f"review_summary_{slot}_{timestamp}.html",
                mime="text/html",
                key=f"review_summary_html_btn_{card_idx}",
            )
        with dl2:
            summary_pdf = _html_to_pdf_bytes(summary_html)
            if summary_pdf:
                st.download_button(
                    "📄 Review Summary (PDF)",
                    data=summary_pdf,
                    file_name=f"review_summary_{slot}_{timestamp}.pdf",
                    mime="application/pdf",
                    key=f"review_summary_pdf_btn_{card_idx}",
                )
            else:
                st.caption("PDF: install xhtml2pdf to enable PDF export.")
        with dl3:
            if _HAS_DOCX:
                st.download_button(
                    "📄 Review Summary (DOCX)",
                    data=_build_review_summary_docx(submission, ev, review_info, timestamp),
                    file_name=f"review_summary_{slot}_{timestamp}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    key=f"review_summary_docx_btn_{card_idx}",
                )
            else:
                st.caption("DOCX: install python-docx to enable DOCX export.")


_PLAIN_ENGLISH_VERDICT = {
    "STRONG":             "This result is well-positioned for submission. The fixes below are refinements, not blockers.",
    "MISLEADING":         "Your evidence is solid, but the result definition needs sharpening before a donor will accept it.",
    "UNDEREVIDENCED":     "The result is clearly defined, but the evidence won't survive a donor's scrutiny yet.",
    "NEEDS REFINEMENT":   "This result is on track — address the top fix below and it will be submission-ready.",
    "FUNDAMENTALLY WEAK": "Both the definition and the evidence need rework. Start with the Confidence fixes below.",
}


# ---------------------------------------------------------------------------
# Council XXII — Score gap chart
# ---------------------------------------------------------------------------

def _council_score_gap_chart(conf_score: float, clar_score: float,
                              proj_conf: float, proj_clar: float):
    """Horizontal grouped bar chart: Current vs Projected scores for each axis."""
    import pandas as pd
    import altair as alt

    def _band(score: float) -> str:
        if score >= 4.0:
            return "Strong"
        if score >= 3.0:
            return "Acceptable"
        return "Below threshold"

    rows = [
        {"Axis": "Confidence", "State": "Current",   "Score": conf_score,  "Band": _band(conf_score)},
        {"Axis": "Confidence", "State": "Projected", "Score": proj_conf,   "Band": _band(proj_conf)},
        {"Axis": "Clarity",    "State": "Current",   "Score": clar_score,  "Band": _band(clar_score)},
        {"Axis": "Clarity",    "State": "Projected", "Score": proj_clar,   "Band": _band(proj_clar)},
    ]
    df = pd.DataFrame(rows)

    # Delta annotation text
    delta_conf = round(proj_conf - conf_score, 2)
    delta_clar = round(proj_clar - clar_score, 2)
    ann_rows = [
        {"Axis": "Confidence", "State": "Projected", "Score": proj_conf,
         "Label": f"{proj_conf}/5  (+{delta_conf})"},
        {"Axis": "Clarity",    "State": "Projected", "Score": proj_clar,
         "Label": f"{proj_clar}/5  (+{delta_clar})"},
    ]
    ann_df = pd.DataFrame(ann_rows)

    bars = (
        alt.Chart(df)
        .mark_bar(cornerRadiusTopRight=3, cornerRadiusBottomRight=3)
        .encode(
            x=alt.X("Score:Q", scale=alt.Scale(domain=[0, 5]), title="Score (0–5)"),
            y=alt.Y("Axis:N", sort=None, title=None),
            yOffset=alt.YOffset("State:N", sort=["Current", "Projected"]),
            color=alt.Color(
                "Band:N",
                scale=alt.Scale(
                    domain=["Strong", "Acceptable", "Below threshold"],
                    range=["#1B5E20", "#F57F17", "#C62828"],
                ),
                legend=alt.Legend(title="Band", orient="bottom"),
            ),
            opacity=alt.condition(
                alt.datum["State"] == "Current",
                alt.value(0.55),
                alt.value(1.0),
            ),
            tooltip=[
                alt.Tooltip("Axis:N"),
                alt.Tooltip("State:N"),
                alt.Tooltip("Score:Q", format=".2f"),
                alt.Tooltip("Band:N"),
            ],
        )
    )

    text = (
        alt.Chart(ann_df)
        .mark_text(align="left", dx=4, fontSize=11, fontWeight="bold")
        .encode(
            x=alt.X("Score:Q"),
            y=alt.Y("Axis:N", sort=None),
            yOffset=alt.YOffset("State:N", sort=["Current", "Projected"]),
            text=alt.Text("Label:N"),
            color=alt.value("#1B5E20"),
        )
    )

    return (bars + text).properties(
        width="container",
        height=100,
        title=alt.TitleParams(
            "Projected score if all priority fixes are implemented",
            fontSize=11,
            color="#616161",
        ),
    )


# ---------------------------------------------------------------------------
# Council XXII — Assessment renderer (council XXII)
# ---------------------------------------------------------------------------

def _render_council_assessment(submission: dict, ev: dict, card_idx: int, api_key: str):
    """Render the 5-member Council Assessment section inside its expander."""
    from council import run_council_assessment, _calculate_projected_scores, COUNCIL_MEMBERS

    st.caption(
        "🤖 Generated by Claude AI (Anthropic — Fable 5 + Haiku models). "
        "Verify specific donor requirements directly with the donor before submission."
    )
    email = st.session_state.get("user_email", "")
    has_access = check_access(email)["allowed"]

    proj_conf, proj_clar = _calculate_projected_scores(ev)
    conf_score = ev.get("confidence_score", 0)
    clar_score = ev.get("clarity_score", 0)

    def _score_label(s: float) -> str:
        if s >= 4.0:
            return "Strong"
        if s >= 3.0:
            return "Acceptable"
        return "High Risk"

    if not has_access:
        metrics.log_event("upgrade_prompt_shown", _metrics_session_id(), context="council_attempt")
        st.markdown(
            "**5-Member Council Assessment** reviews your result from 5 expert lenses "
            "and produces a plain-English upgrade brief your reporting team can act on."
        )
        col_a, col_b = st.columns(2)
        with col_a:
            st.metric("Confidence (current)", f"{conf_score}/5.0", label_visibility="visible")
        with col_b:
            st.metric("Confidence (projected)", f"{proj_conf}/5.0",
                      delta=f"+{round(proj_conf - conf_score, 2)}")
        col_c, col_d = st.columns(2)
        with col_c:
            st.metric("Clarity (current)", f"{clar_score}/5.0")
        with col_d:
            st.metric("Clarity (projected)", f"{proj_clar}/5.0",
                      delta=f"+{round(proj_clar - clar_score, 2)}")
        st.caption("Upgrade to Professional to run the full council assessment — "
                   "GHS 50/month vs. GHS 12,000–17,000 in rework costs from a donor-queried report.")
        if st.button("Upgrade to Professional →", key=f"council_upgrade_{card_idx}", type="primary"):
            metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context="council_attempt")
            st.session_state["_show_pricing"] = True
            st.rerun()
        return

    cache_key = f"council_xxii_{card_idx}"
    assessment = st.session_state.get(cache_key)

    if assessment is None:
        st.markdown(
            "Run a 5-member council review. Each member assesses the result from their assigned "
            "lens — evidence quality, strategy, critical review, implementation steps, and donor "
            "acceptance — then the council produces a plain-English upgrade brief for your reporting team."
        )
        if not api_key:
            st.warning("Council Assessment is not available — API key not configured.")
            return
        if st.button("🏛 Run Council Assessment", key=f"run_council_{card_idx}", type="primary",
                     use_container_width=True):
            with st.spinner("5 council members reviewing your result…"):
                assessment = run_council_assessment(submission, ev, api_key)
                st.session_state[cache_key] = assessment
            try:
                _m_session_id = _metrics_session_id()
                metrics.log_event("ai_questions_generated", _m_session_id)
                _m_withheld = assessment.get("withheld", {})
                if _m_withheld.get("upgraded_result_statement") or _m_withheld.get("upgraded_evidence_statement"):
                    metrics.log_event("draft_withheld_fabrication", _m_session_id)
            except Exception:
                pass
            st.rerun()
        return

    # --- Render cached assessment ---

    if assessment.get("error"):
        st.caption(f"Note: some council members were unavailable ({assessment['error']}).")

    # 1. Council member cards — 2-column grid, Donor Lens full-width last
    verdicts = assessment.get("verdicts", {})
    member_ids = [m["id"] for m in COUNCIL_MEMBERS]
    paired = [(member_ids[i], member_ids[i + 1]) for i in range(0, 4, 2)]

    for left_id, right_id in paired:
        col_l, col_r = st.columns(2)
        for col, mid in ((col_l, left_id), (col_r, right_id)):
            v = verdicts.get(mid, {})
            with col:
                st.markdown(
                    f"<div style='border:1px solid {v.get('color','#ccc')};border-radius:6px;"
                    f"padding:10px 12px;margin-bottom:6px;'>"
                    f"<div style='font-size:13px;font-weight:700;color:{v.get('color','#333')};'>"
                    f"{v.get('icon','')} {v.get('name','')}</div>"
                    f"<div style='font-size:11px;color:#757575;margin-bottom:6px;'>"
                    f"{v.get('archetype','')}</div>"
                    f"<div style='font-size:13px;line-height:1.6;'>"
                    f"{v.get('verdict_text','').replace(chr(10),'<br>')}</div>"
                    f"</div>",
                    unsafe_allow_html=True,
                )

    # Donor Lens — full width
    donor_v = verdicts.get("donor_rep", {})
    st.markdown(
        f"<div style='border:2px solid {donor_v.get('color','#B71C1C')};border-radius:6px;"
        f"padding:12px 14px;margin-bottom:10px;'>"
        f"<div style='font-size:12px;font-weight:700;color:{donor_v.get('color','#B71C1C')};'>"
        f"{donor_v.get('icon','')} {donor_v.get('name','')} — {donor_v.get('archetype','')}</div>"
        f"<div style='font-size:12px;line-height:1.6;margin-top:6px;'>"
        f"{donor_v.get('verdict_text','').replace(chr(10),'<br>')}</div>"
        f"</div>",
        unsafe_allow_html=True,
    )

    st.divider()

    # 2. Score gap chart
    if not st.session_state.get("lite_mode", False):
        proj_c  = assessment.get("projected_conf", proj_conf)
        proj_cl = assessment.get("projected_clar", proj_clar)
        try:
            st.altair_chart(
                _council_score_gap_chart(conf_score, clar_score, proj_c, proj_cl),
                use_container_width=True,
                key=f"council_gap_chart_{card_idx}",
            )
        except Exception:
            pass
    else:
        st.markdown(
            f"**Score uplift:** Confidence {conf_score} → {proj_conf} "
            f"(+{round(proj_conf - conf_score, 2)})  |  "
            f"Clarity {clar_score} → {proj_clar} (+{round(proj_clar - clar_score, 2)})"
        )

    # 3. Reporting Team Brief
    brief = assessment.get("reporting_team_brief", {})
    if brief:
        st.markdown("#### For Your Reporting Team")
        with st.container(border=True):
            wm = brief.get("what_score_means", "")
            if wm:
                st.markdown(f"**What the score means right now:**\n{wm}")
            changes = brief.get("what_to_change", [])
            if changes:
                st.markdown("**What needs to change:**")
                for item in changes:
                    st.markdown(f"- {item}")
            hl = brief.get("how_long", "")
            if hl:
                st.markdown(f"**How long this takes:** {hl}")
            ps = brief.get("projected_status", "")
            if ps:
                st.markdown(
                    f"<div style='background:#E8F5E9;border-left:3px solid #1B5E20;"
                    f"padding:8px 12px;border-radius:4px;font-size:12px;margin-top:8px;'>"
                    f"<strong>After fixes:</strong> {ps}</div>",
                    unsafe_allow_html=True,
                )

    # 4. Upgraded Statements
    upg_rs   = assessment.get("upgraded_result_statement", "")
    upg_ev   = assessment.get("upgraded_evidence_statement", "")
    withheld = assessment.get("withheld", {})
    rs_withheld = withheld.get("upgraded_result_statement", False)
    ev_withheld = withheld.get("upgraded_evidence_statement", False)
    if upg_rs or upg_ev or rs_withheld or ev_withheld:
        with st.expander("📝 Upgraded statements (council draft — review before use)", expanded=False):
            if rs_withheld:
                st.markdown("**Upgraded result statement:**")
                st.warning("AI draft withheld — it introduced content not in your evidence.")
            elif upg_rs:
                st.markdown("**Upgraded result statement:**")
                st.markdown(
                    f"<blockquote style='border-left:3px solid #1565C0;padding:8px 12px;"
                    f"margin:0;font-size:12px;color:#1a1a1a;background:#F3F8FE;'>"
                    f"{upg_rs}</blockquote>",
                    unsafe_allow_html=True,
                )
            if ev_withheld:
                st.markdown("**Upgraded evidence statement:**")
                st.warning("AI draft withheld — it introduced content not in your evidence.")
            elif upg_ev:
                st.markdown("**Upgraded evidence statement:**")
                st.markdown(
                    f"<blockquote style='border-left:3px solid #E65100;padding:8px 12px;"
                    f"margin:0;font-size:12px;color:#1a1a1a;background:#FFF3E0;'>"
                    f"{upg_ev}</blockquote>",
                    unsafe_allow_html=True,
                )
            st.caption(
                "These are council-drafted suggestions. Review all values carefully "
                "before using in a submission — replace any [placeholder] tokens with "
                "your actual data."
            )


def _render_help_chat(submission: dict, ev: dict, donor: str = "", card_idx: int = 0):
    """Score-explanation chat assistant (council XV).

    Scoped strictly to rubric Q&A using the user's actual scores as context.
    Gated behind paid tier. Chat history lives in st.session_state for the session.
    """
    from diagnostics import build_chat_system_prompt

    email = st.session_state.get("user_email", "")
    has_access = check_access(email)["allowed"]

    if not has_access:
        metrics.log_event("upgrade_prompt_shown", _metrics_session_id(), context="chat_attempt")
        st.info(
            "Score chat is available on the **Professional plan**. "
            "Upgrade to ask questions about your score and get rubric-based guidance — "
            "GHS 50/month vs. GHS 12,000–17,000 in rework costs from a donor-queried report."
        )
        if st.button("Upgrade to Professional →", key=f"chat_upgrade_{card_idx}", type="primary"):
            metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context="chat_attempt")
            st.session_state["_show_pricing"] = True
            st.rerun()
        return

    # Chat history key — per card_idx to support multi-result submissions
    hist_key = f"chat_messages_{card_idx}" if card_idx else "chat_messages"
    if hist_key not in st.session_state:
        st.session_state[hist_key] = []

    msgs = st.session_state[hist_key]

    # Scope notice (shown once at top)
    st.caption(
        "Ask about your score, any of the 8 scoring criteria, or your donor's requirements. "
        "This assistant answers from the scoring rubric only — not general MEL advice."
    )

    # Render existing conversation
    for msg in msgs:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # Chat input
    prompt = st.chat_input("Ask about your score…", key=f"chat_input_{card_idx}")
    if not prompt:
        return

    # Append user message and show it immediately
    msgs.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    # Build system prompt with live scores + rubric
    system_prompt = build_chat_system_prompt(ev, submission, donor)

    # Call Claude Haiku (fast, cheap for Q&A)
    _api_key = (
        st.secrets.get("ANTHROPIC_API_KEY", "")
        if hasattr(st, "secrets") else
        __import__("os").environ.get("ANTHROPIC_API_KEY", "")
    )
    if not _api_key:
        reply = "Score chat is not available right now — API key not configured."
    else:
        try:
            import anthropic as _anthropic_mod
            _client = _anthropic_mod.Anthropic(api_key=_api_key)
            _resp = _client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=512,
                system=system_prompt,
                messages=[{"role": m["role"], "content": m["content"]} for m in msgs],
            )
            reply = _resp.content[0].text if _resp.content else "No response received."
        except Exception as exc:
            reply = f"Could not reach the scoring assistant right now. ({type(exc).__name__})"

    msgs.append({"role": "assistant", "content": reply})
    with st.chat_message("assistant"):
        st.markdown(reply)


def _render_result_card(submission: dict, ev: dict, card_idx: int = 0, donor: str = ""):
    conf_score   = ev.get("confidence_score", 0)
    clar_score   = ev.get("clarity_score", 0)
    conf_label   = ev.get("confidence_label", "High Risk")
    clar_label   = ev.get("clarity_label",   "High Risk")
    conf_meaning = ev.get("confidence_meaning", "")
    clar_meaning = ev.get("clarity_meaning",    "")
    verdict      = ev.get("verdict", "")
    conf_comp    = ev.get("confidence_components", {})
    clar_comp    = ev.get("clarity_components", {})
    fixes        = ev.get("fixes", [])

    snippet = submission.get("result_statement", "")
    if len(snippet) > 120:
        snippet = snippet[:120] + "..."
    st.markdown(f"**{snippet}**")
    st.divider()

    # Diagnostic state badge
    content_issues    = ev.get("content_issues", [])
    bv_voice_field    = submission.get("beneficiary_voice", "")
    diag_state, diag_sub = get_diagnostic_state(conf_score, clar_score, content_issues, bv_voice_field)

    # Single status signal — diagnostic badge carries state + description
    diag_cfg = _DIAGNOSTIC_BADGE.get(diag_state, {"bg": "#9E9E9E", "text": "#FFFFFF", "subtitle": ""})
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
    if diag_state != "INVALID INPUT":
        st.markdown(
            f"<div class='diagnostic-badge' style='background:{diag_cfg['bg']};color:{diag_cfg['text']};{_pca}'>"
            f"{diag_state} &nbsp;·&nbsp; {diag_sub}"
            f"</div>",
            unsafe_allow_html=True,
        )

    # Scores follow immediately — no ticker or redundant banners
    render_scoreboard(
        confidence=round(conf_score * 20),
        clarity=round(clar_score * 20),
        verified=True,
    )
    st.caption(
        f"Shown out of 100 for a quick read — the same score is **{conf_score}/5.0** Confidence "
        f"and **{clar_score}/5.0** Clarity in the breakdown below."
    )

    # "How you compare" — anonymized percentile benchmark, derived from other
    # users' opted-in saved audits. Viewing this doesn't require this user to
    # have opted in themselves; the underlying data is score-only, no content.
    _bm_donor, _bm_sector, _bm_org_type = (
        submission.get("donor", ""), submission.get("sector", ""), submission.get("org_type", ""))
    if _bm_donor and _bm_sector and _bm_org_type:
        _benchmark = get_benchmark(_bm_donor, _bm_sector, _bm_org_type, conf_score, clar_score)
        if _benchmark:
            st.caption(
                f"📊 **How you compare:** your Confidence score is higher than "
                f"{_benchmark['confidence_percentile']}% and your Clarity score is higher than "
                f"{_benchmark['clarity_percentile']}% of {_benchmark['sample_size']} saved {_bm_donor} "
                f"audits in {_bm_sector} at your organisation type."
            )
        else:
            st.caption(
                f"📊 **How you compare:** not enough saved audits yet in {_bm_sector} for {_bm_donor} "
                f"at your organisation type to show a comparison — check back as more MEL teams use ImpactProof."
            )

    _pev = _PLAIN_ENGLISH_VERDICT.get(diag_state, "")
    if _pev:
        st.info(_pev)

    # Methodology disclaimer — shown for every verdict state, immediately after the badge
    _used_threshold = ev.get("threshold_used", 4.0)
    st.caption(
        f"This is a heuristic pre-submission check against your evidence *description*, not your "
        f"underlying data — not an expert audit. A passing score (≥{_used_threshold}) does not guarantee "
        f"donor acceptance if the described evidence is incomplete or inaccurate. Your donor reviewer "
        f"makes the final determination."
    )

    # Near-boundary notice — shown when either axis is within 0.1 of the org-type threshold
    from evaluator import SUBMISSION_THRESHOLD, NEAR_THRESHOLD_BAND
    _conf_near = abs(conf_score - _used_threshold) <= NEAR_THRESHOLD_BAND
    _clar_near = abs(clar_score - _used_threshold) <= NEAR_THRESHOLD_BAND
    if (_conf_near or _clar_near) and diag_state not in ("INVALID INPUT", "INCOMPLETE"):
        st.caption(
            "Note: your score is within 0.1 of the submission threshold. "
            "A single evidence improvement could change the verdict either way — "
            "review the breakdown below before submitting."
        )

    # Single biggest-impact fix surfaced up-front (all non-STRONG states)
    if fixes and diag_state not in ("STRONG", "INVALID INPUT"):
        _top = fixes[0]
        _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
        st.markdown(
            f"<div style='background:#FFF9C4;border-left:4px solid #F57F17;border-radius:8px;"
            f"padding:10px 14px;margin:8px 0;font-size:0.9rem;{_pca}'>"
            f"🎯 <strong>Biggest single fix:</strong> {_top['message']} "
            f"<em>({_top['score_impact']})</em></div>",
            unsafe_allow_html=True,
        )


    # Evidence statement — inline, most actionable output after the top fix
    _ev_stmt = _generate_evidence_statement(submission) if callable(globals().get("_generate_evidence_statement")) else None
    if _ev_stmt:
        st.markdown("**Evidence statement for your report** — copy and edit before pasting into your narrative:")
        st.code(_ev_stmt, language=None)

    # INVALID INPUT early exit
    if diag_state == "INVALID INPUT":
        st.error("Input Quality Issue Detected")
        st.markdown(
            "Your responses appear to be placeholder text. ImpactProof scores **real** "
            "reported results. Please return to Screen 1 and provide genuine content."
        )
        for issue in content_issues:
            st.markdown(f"- {issue}")
        raw_conf = ev.get("raw_confidence_score")
        mult = ev.get("content_quality_multiplier", 1.0)
        if raw_conf is not None:
            st.caption(
                f"Raw score before quality adjustment: {raw_conf}/5.0 — multiplier applied: ×{mult}"
            )
        if st.button("← Edit Submission", key=f"invalid_back_{card_idx}"):
            st.session_state["screen"] = 1
            st.session_state["evaluations"] = None
            st.rerun()
        st.divider()
        return

    # --- Full breakdown expander — all detail panels (council XVIII) ---
    # The verdict, top fixes, and download are already shown above.
    # Everything else lives in this single expander to reduce default noise.
    _s2_breakdown_label = (
        "✅ Reference breakdown — no action required" if diag_state == "STRONG"
        else "🔍 See full breakdown — scores, evidence ladder, donor readiness"
    )
    # Open by default for anything that isn't already STRONG -- those are exactly
    # the users who need to see the diagnosis, not the ones who least need it.
    _s2_breakdown_default = diag_state != "STRONG"
    _s2_breakdown_open = st.session_state.get(f"_s2_breakdown_open_{card_idx}", _s2_breakdown_default)
    if st.button(
        ("▼ Hide breakdown" if _s2_breakdown_open else "▶ " + _s2_breakdown_label),
        key=f"s2_breakdown_toggle_{card_idx}",
        use_container_width=True,
    ):
        st.session_state[f"_s2_breakdown_open_{card_idx}"] = not _s2_breakdown_open
        st.rerun()

    if not st.session_state.get(f"_s2_breakdown_open_{card_idx}", _s2_breakdown_default):
        return  # User hasn't opened breakdown — show nothing more

    # --- Four Funder Questions summary (top of report) ---
    ev_top      = (submission.get("evidence") or [{}])[0]
    ev_type_top = ev_top.get("type", "") or "Not specified"
    ladder_top  = ev.get("evidence_ladder", {})
    fr_top      = ev.get("funder_readiness", {})
    direct_level = conf_comp.get("direct_level", 0)
    verify_level = conf_comp.get("verify_level", 0)
    def_score    = clar_comp.get("definition_score", 0)

    if diag_state == "STRONG":
        st.success("✅ Your result is submission-ready. The full breakdown is below for reference — no action required.")
    st.markdown("### What Donors Want to Know")
    st.caption("ImpactProof checked your evidence against these 4 donor questions. Each maps to a scored sub-criterion.")
    fq_col1, fq_col2 = st.columns(2)
    with fq_col1:
        st.markdown("**1. What has changed?**")
        st.markdown(snippet if snippet else "_Not yet described._")
        st.caption(f"Directness: Level {direct_level}/5 · Definition: {def_score}/1.25")

        st.markdown("**3. How strong is the evidence?**")
        st.markdown(f"Confidence: **{conf_score}/5.0** ({conf_label})")
        st.caption(_VERIFICATION_TIPS.get(verify_level, ""))

    with fq_col2:
        st.markdown("**2. How do you know?**")
        dominant = ladder_top.get("dominant_tier")
        if dominant:
            st.markdown(f"Evidence type: **{ev_type_top}** — evidence base is mainly **{dominant}**-tier.")
        else:
            st.markdown(f"Evidence type: **{ev_type_top}**")
        st.caption(_DIRECTNESS_TIPS.get(direct_level, ""))

        st.markdown("**4. What did you learn?**")
        learn = fr_top.get("learning", {})
        if learn.get("detected"):
            st.markdown("Yes — the report describes what your team learned and how the programme adapted.")
        else:
            st.markdown("_Not yet stated._ Add a sentence on what you learned and changed as a result.")

    st.divider()

    # Logframe linkage panel (guarded for backward-compat with stale evaluator deploys)
    linkage = ev.get("logframe_linkage", {})
    if linkage:
        lk_state  = linkage.get("state", "MISSING")
        lk_rat    = linkage.get("rationale", "")
        lk_issues = linkage.get("issues", [])
        pct_of_target     = linkage.get("pct_of_target")
        direction_mismatch = linkage.get("direction_mismatch", False)
        st.markdown("#### Logframe Linkage")

        # % of target badge — shown when computable
        if pct_of_target is not None:
            _pct_color = "#1B5E20" if pct_of_target >= 100 else ("#8A6500" if pct_of_target >= 80 else "#B71C1C")
            st.markdown(
                f"<span style='background:{_pct_color};color:#fff;padding:3px 10px;"
                f"border-radius:4px;font-size:0.85rem;font-weight:700;'>"
                f"{pct_of_target:.0f}% of target reached</span>",
                unsafe_allow_html=True,
            )

        if lk_state == "STRONG":
            st.success(f"✓ {lk_rat}")
        elif lk_state == "DATA_FORTHCOMING":
            st.info("Data gap disclosed — not penalised. State this in your submission narrative.")
            for iss in lk_issues:
                st.markdown(f"- {iss}")
        elif lk_state in ("WEAK", "DIRECTION_MISMATCH"):
            st.warning(f"⚠️ {lk_rat}")
            if direction_mismatch:
                st.error(
                    "Direction mismatch detected: the indicator implies a change in one direction "
                    "but the baseline-to-achievement comparison shows the opposite. "
                    "Donors will flag this — review the framing before submission."
                )
            for iss in lk_issues:
                st.markdown(f"- {iss}")
        else:
            st.error(f"❌ {lk_rat}")
            st.markdown(
                "**Highest-impact missing piece:** Donors will reject results that cannot "
                "be traced to an approved indicator from your Technical Proposal."
            )
            for iss in lk_issues:
                st.markdown(f"- {iss}")

    # Score-explanation chat assistant (council XV)
    with st.expander("💬 Ask about your score", expanded=False):
        _render_help_chat(submission, ev, donor=donor, card_idx=card_idx)

    # Council Assessment — upgrade recommendations (council XXII)
    _ca_api_key = (
        st.secrets.get("ANTHROPIC_API_KEY", "")
        if hasattr(st, "secrets") else
        __import__("os").environ.get("ANTHROPIC_API_KEY", "")
    )
    with st.expander("🏛 Council Assessment — Upgrade Recommendations", expanded=(diag_state != "STRONG")):
        _render_council_assessment(submission, ev, card_idx, _ca_api_key)

    # Reporting period validation on Screen 2
    rp_start_str = submission.get("reporting_start", "")
    rp_end_str   = submission.get("reporting_end", "")
    ev_date_str  = (submission.get("evidence") or [{}])[0].get("recency", "")
    if rp_start_str and rp_end_str and ev_date_str and hasattr(_evaluator, "validate_reporting_period"):
        try:
            from datetime import date as _d
            _rp_s2 = _d.fromisoformat(rp_start_str)
            _rp_e2 = _d.fromisoformat(rp_end_str)
            _ev2   = _d.fromisoformat(ev_date_str)
            _, _rp2_msg, _rp2_sev = _evaluator.validate_reporting_period(_ev2, _rp_s2, _rp_e2)
            if _rp2_sev == "WARNING":
                st.warning(f"📅 Reporting Period Issue: {_rp2_msg}")
                st.caption("This is a common cause of donor flags. Address before submission.")
            elif _rp2_sev == "ERROR":
                st.error(f"📅 Reporting Period Error: {_rp2_msg}")
        except (ValueError, TypeError, AttributeError):
            pass

    # Dual-axis columns
    col_conf, col_clar = st.columns(2)

    with col_conf:
        st.markdown("#### Confidence Score")
        st.markdown(_axis_badge_html(conf_label, conf_score, 5.0), unsafe_allow_html=True)
        st.caption(conf_meaning)
        dl = conf_comp.get("direct_level", 0)
        vl = conf_comp.get("verify_level", 0)
        rl = conf_comp.get("recency_level", 0)
        ds = conf_comp.get("direct_score", 0)
        vs = conf_comp.get("verify_score", 0)
        rs = conf_comp.get("recency_score", 0)
        _dir_rationale = conf_comp.get("direct_rationale") or _evaluator.get_score_rationale("directness", dl, ds, 2.0)
        _verify_rationale = conf_comp.get("verify_rationale") or _evaluator.get_score_rationale("verification", vl, vs, 2.0)
        st.metric("Directness (max 2.0)", f"{ds:.1f}", help=_dir_rationale)
        st.metric("Verification (max 2.0)", f"{vs:.1f}", help=_verify_rationale)
        st.metric("Recency (max 1.0)", f"{rs:.1f}",
                  help=_evaluator.get_score_rationale("recency", rl, rs, 1.0))
        bv_bonus = conf_comp.get("bv_bonus", 0.0)
        st.metric("Beneficiary Voice Bonus", f"+{bv_bonus:.1f}",
                  help=BENEFICIARY_VOICE_TOOLTIP)
        _render_subscore_chart([
            ("Directness", ds, 2.0, _dir_rationale),
            ("Verification", vs, 2.0, _verify_rationale),
            ("Recency", rs, 1.0, _evaluator.get_score_rationale("recency", rl, rs, 1.0)),
        ], key=f"snapshot_conf_chart_{card_idx}")

    with col_clar:
        st.markdown("#### Clarity Score")
        st.markdown(_axis_badge_html(clar_label, clar_score, 5.0), unsafe_allow_html=True)
        st.caption(clar_meaning)
        def_s  = clar_comp.get("definition_score",  0)
        meas_s = clar_comp.get("measurement_score", 0)
        integ  = clar_comp.get("integrity_score",   0)
        scope  = clar_comp.get("scope_score",       0)
        gov    = clar_comp.get("governance_score",  0)
        is_qual    = clar_comp.get("is_qualitative", False)
        def_label  = "Narrative Definition" if is_qual else "Definition"
        def_tip    = _CLARITY_TIPS["definition_qualitative"] if is_qual else _CLARITY_TIPS["definition"]
        meas_label = "Sourcing & Triangulation" if is_qual else "Measurement"
        meas_tip   = _CLARITY_TIPS["measurement_qualitative"] if is_qual else _CLARITY_TIPS["measurement"]
        st.metric(f"{def_label} (max 1.25)", f"{def_s:.2f}",
                  help=f"{def_tip}\n\n{TOOLTIP_DEFINITION}")
        st.metric(f"{meas_label} (max 1.25)", f"{meas_s:.2f}",
                  help=f"{meas_tip}\n\n{TOOLTIP_MEASUREMENT}")
        st.metric("Integrity (max 1.0)", f"{integ:.2f}",
                  help=f"{_CLARITY_TIPS['integrity']}\n\n{TOOLTIP_INTEGRITY}")
        st.metric("Scope (max 0.75)", f"{scope:.2f}",
                  help=f"{_CLARITY_TIPS['scope']}\n\n{TOOLTIP_SCOPE}")
        st.metric("Governance (max 0.75)", f"{gov:.2f}",
                  help=f"{_CLARITY_TIPS['governance']}\n\n{TOOLTIP_GOVERNANCE}")
        _render_subscore_chart([
            (def_label, def_s, 1.25, def_tip),
            (meas_label, meas_s, 1.25, meas_tip),
            ("Integrity", integ, 1.0, _CLARITY_TIPS["integrity"]),
            ("Scope", scope, 0.75, _CLARITY_TIPS["scope"]),
            ("Governance", gov, 0.75, _CLARITY_TIPS["governance"]),
        ], key=f"snapshot_clar_chart_{card_idx}")

    # Governance gap callout — surfaced here so it's visible outside the Clarity column
    if gov < 0.6:
        st.warning(
            "⚠️ **Governance gap flagged.** Your data-compliance sub-score is low. "
            "Return to the **Governance tab** on the form to complete the checklist — "
            "this affects your Clarity score and Act 843 / NDPA compliance (Ghana)."
        )

    # Scoring explanation + weakness panel — skip for STRONG (nothing to fix)
    if diag_state != "STRONG":
        render_how_scoring_works_panel()
        render_personalized_weakness_panel(
            direct_score=ds, verify_score=vs, recency_score=rs,
            definition_score=def_s, measurement_score=meas_s,
            integrity_score=integ, scope_score=scope, governance_score=gov,
            is_qualitative=is_qual,
        )

    # Evidence Ladder (rule-based, no score impact)
    ladder = ev.get("evidence_ladder", {})
    if ladder:
        st.markdown("#### Evidence Ladder")
        st.caption(
            "Rule-based check of the evidence sources you described — does this "
            "report rely mainly on Basic, Moderate, or Stronger evidence? Hover "
            "over a rung below for details."
        )
        _render_evidence_ladder_chart(ladder, key=f"ladder_chart_{card_idx}")
        st.info(ladder.get("suggestion", ""))

    # Indicator Maturity (rule-based, count-only indicator detection)
    maturity = ev.get("indicator_maturity", {})
    if maturity.get("flagged"):
        st.markdown("#### Indicator Maturity")
        st.caption(
            "This indicator is written as a raw count. Donors increasingly expect "
            "indicators that show whether the result was sustained or verified."
        )
        _your_indicator = (submission.get("logframe_indicator", "") or "").strip() or "(not specified)"
        st.markdown(f"**Your indicator (as written):** {_your_indicator}")
        _maturity_rows = [{"Level": level, "Example wording": wording} for level, wording in maturity["rows"]]
        _maturity_rows[0]["Level"] = "👈 " + _maturity_rows[0]["Level"] + " — what you wrote"
        st.table(_maturity_rows)
        st.caption(
            f"Measurement score adjusted by **{maturity['adjustment']}** for this "
            "count-only indicator framing."
        )

    # Donor Readiness flags (informational only — no score impact)
    st.markdown("#### Donor Readiness")
    st.caption(
        "Two quick checks donors increasingly look for — these do not affect "
        "your Confidence or Clarity scores."
    )
    fr = ev.get("funder_readiness", {})
    lim = fr.get("limitations", {})
    learn = fr.get("learning", {})

    if lim.get("detected"):
        st.success("Limitations disclosed — the report states what the data can't confidently say.")
    else:
        st.warning(
            "No limitations disclosure detected. Consider adding a sentence on what "
            "this data cannot confirm or cannot be generalized to."
        )

    if learn.get("detected"):
        st.success("Learning & adaptation stated — the report describes what your team learned and changed.")
    else:
        st.warning(
            "No learning/adaptation statement detected. Consider adding what your "
            "organization learned and how the program adapted as a result."
        )

    # Additional advisory flags (v3.4, score-neutral)
    attrib = submission.get("attribution_contribution", "Not specified")
    disagg = submission.get("disaggregation_status", "Not specified")
    if attrib != "Not specified" or disagg != "Not specified":
        st.markdown("#### Additional Advisory Flags")
        st.caption("Optional checklist answers — advisory only, no effect on your score.")
        if attrib != "Not specified":
            st.markdown(f"- **Attribution vs. contribution distinguished:** {attrib}")
        if disagg != "Not specified":
            st.markdown(f"- **Beneficiary data disaggregated (women, youth, PWD, rural):** {disagg}")

    # Verdict banner
    css_class = _VERDICT_CSS.get(verdict, "")
    st.markdown(
        f"<div class='verdict-banner {css_class}'>{verdict}</div>",
        unsafe_allow_html=True,
    )

    # What To Fix — simple bullet list (no checkboxes: users thought ticking = fixed)
    conf_fixes = [f for f in fixes if f.get("dimension") == "confidence"]
    clar_fixes = [f for f in fixes if f.get("dimension") == "clarity"]

    # IRC page references — annotate fixes with source page when available
    _irc_sources = st.session_state.get("_irc_field_sources", {}) if st.session_state.get("_irc_used") else {}

    def _page_note_for_fix(fix):
        if not _irc_sources:
            return ""
        dim = fix.get("dimension", "")
        fields = _FIX_FIELD_SOURCE_MAP.get(dim, [])
        page = next((
            _irc_sources[f].get("page", 0)
            for f in fields
            if f in _irc_sources and isinstance(_irc_sources[f], dict) and _irc_sources[f].get("page", 0) > 0
        ), 0)
        return f" · *see page {page} of your document*" if page else ""

    def _render_fix_bullets(fix_list, label):
        if not fix_list:
            return
        st.markdown(f"**{label}**")
        for fix in fix_list:
            st.markdown(f"- {fix['message']} *({fix['score_impact']})*{_page_note_for_fix(fix)}")

    if diag_state == "MISLEADING":
        _render_fix_bullets(clar_fixes, "Sharpen your definition (Clarity)")
        _render_fix_bullets(conf_fixes, "Confidence fixes")

    elif diag_state == "UNDEREVIDENCED":
        _render_fix_bullets(conf_fixes, "Strengthen your evidence (Confidence)")
        _render_fix_bullets(clar_fixes, "Clarity fixes")

    elif diag_state == "FUNDAMENTALLY WEAK":
        st.error("This result requires fundamental rework on both axes.")
        _render_fix_bullets(conf_fixes, "Strengthen your evidence (Confidence)")
        _render_fix_bullets(clar_fixes, "Sharpen your definition (Clarity)")

    elif fixes:
        _render_fix_bullets(fixes[:3], "Top fixes before submission")

    # Fear-of-rejection ROI hook (Council XXVII) — surface rework cost for highest-risk results
    _rr_paid = st.session_state.get("is_paid", False)
    if diag_state in ("FUNDAMENTALLY WEAK", "UNDEREVIDENCED", "MISLEADING") and not _rr_paid:
        _rr_context = f"high_risk_{diag_state.lower().replace(' ', '_')}"
        metrics.log_event("upgrade_prompt_shown", _metrics_session_id(), context=_rr_context)
        st.markdown(
            "<div style='background:#FFFBF2;border:1px solid #FFE0B2;border-radius:8px;"
            "padding:14px 18px;margin:12px 0;'>"
            "<p style='font-weight:700;color:#8A6500;margin:0 0 6px;font-size:0.95rem;'>"
            "💡 Worth fixing before you submit</p>"
            "<p style='font-size:0.85rem;color:#374151;margin:0 0 8px;'>"
            "USAID Learning Lab (2024): 3 of 5 DQA failures are predictable from evidence "
            "quality gaps like these — catching them now is far cheaper than reworking a "
            "rejected report later. At Ghana MEL consultant rates (GHS 1,200–1,800/day), "
            "40+ hours of rework runs <strong>GHS 12,000–17,000</strong> — this fix list is free."
            "</p>"
            "<p style='font-size:0.85rem;color:#374151;margin:0;'>"
            "<strong>ImpactProof Professional catches this before your donor does "
            "— GHS 50/month, unlimited checks.</strong>"
            "</p></div>",
            unsafe_allow_html=True,
        )
        if st.button(
            "Upgrade to catch this before your donor does →",
            key=f"fear_rejection_upgrade_{card_idx}",
            type="primary",
            use_container_width=True,
        ):
            metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context=_rr_context)
            st.session_state["_show_pricing"] = True
            st.rerun()

    filenames = submission.get("attached_filenames", [])
    if filenames:
        st.caption(f"Attached documents: {', '.join(filenames)}")

    bv_bonus = conf_comp.get("bv_bonus", 0.0)
    if bv_bonus < 0.5:
        bv_fix = BENEFICIARY_VOICE_WHATTOFIX.get(bv_bonus, BENEFICIARY_VOICE_WHATTOFIX[0.0])
        with st.expander("Improve Beneficiary Voice →", key=f"bv_improve_{card_idx}"):
            st.caption(bv_fix)
            # Sector-specific HOW-TO guidance (Council XXVII)
            _bv_sector = st.session_state.get("sector", "")
            if _bv_sector and _bv_sector in _BV_SECTOR_GUIDANCE:
                st.markdown(f"**How to collect beneficiary voice in {_bv_sector}:**")
                st.info(_BV_SECTOR_GUIDANCE[_bv_sector])

    if donor and donor != "Other/Not specified" and donor in DONOR_DIAGNOSTICS:
        st.subheader(f"{donor}-Specific Guidance")
        st.caption(f"Diagnostic language adapted for {donor} submission standards")
        donor_map = DONOR_DIAGNOSTICS[donor]
        checks = [
            ("Directness",       conf_comp.get("direct_score", 0),  2.0),
            ("Verification",     conf_comp.get("verify_score", 0),  2.0),
            ("Recency",          conf_comp.get("recency_score", 0), 1.0),
            ("BeneficiaryVoice", conf_comp.get("bv_bonus", 0.0),    0.5),
        ]
        for dim, raw_score, max_val in checks:
            if dim not in donor_map:
                continue
            level = "low" if (raw_score / max_val) < 0.6 else "high"
            st.markdown(f"**{dim}:** {donor_map[dim][level]}")


    with st.expander("Share this result", key=f"share_result_{card_idx}"):
        def _share_icon(s):
            return "✅" if s >= 4.0 else "⚠️" if s >= 3.0 else "🔴"
        _tf = fixes[0]["message"] if fixes else "No major gaps — ready to refine."
        _wa_text = (
            f"📊 ImpactProof — Evidence Quality Check\n"
            f"Confidence: {conf_score}/5.0 {_share_icon(conf_score)}  ·  "
            f"Clarity: {clar_score}/5.0 {_share_icon(clar_score)}\n"
            f"Top fix: {_tf}\n"
            f"Verdict: {verdict}\n"
            f"Checked with: ImpactProof ({APP_URL}/)"
        )
        _wa_url = "https://wa.me/?text=" + urllib.parse.quote(_wa_text)
        st.markdown(
            f'<a href="{_wa_url}" target="_blank" style="display:inline-block;'
            f'background:#25D366;color:white;padding:8px 18px;border-radius:8px;'
            f'text-decoration:none;font-weight:700;font-size:0.9rem;margin:4px 0;">'
            f'📱 Send to WhatsApp</a>',
            unsafe_allow_html=True,
        )
        st.caption("Opens WhatsApp — choose who to send to from your contacts.")
        st.code(_wa_text, language=None)

    _render_review_handoff(submission, ev, card_idx)

    st.divider()


def render_screen_2():
    render_pitch_strip("report")
    # Run evaluations once, cache results
    if not st.session_state.get("evaluations"):
        active = st.session_state.get("active_slots_run", st.session_state.get("active_slots", 1))
        subs, evs = [], []
        try:
            render_var_review()
            with st.spinner("Running diagnostic…"):
                for slot in range(1, active + 1):
                    sub = _build_submission_from_session(slot)
                    ev  = _evaluator.evaluate_submission(sub)
                    save_all_files(sub, ev)
                    subs.append(sub)
                    evs.append(ev)
            st.session_state["evaluations"]        = evs
            st.session_state["submissions_snapshot"] = subs
            st.rerun()
        except Exception as exc:
            import logging as _logging
            _logging.error("Evaluation failed on Screen 2", exc_info=True)
            st.session_state["error_message"] = (
                "Something went wrong while scoring your result. Please go back and try again. "
                "If the problem persists, contact us: info@impact-receipts.com"
            )

    if st.session_state.get("error_message"):
        st.error(st.session_state["error_message"])
        if st.button("← Edit Submission"):
            st.session_state["screen"] = 1
            st.session_state["evaluations"]  = None
            st.session_state["error_message"] = None
            st.rerun()
        return

    evs  = st.session_state.get("evaluations") or []
    subs = st.session_state.get("submissions_snapshot") or []

    if evs and not st.session_state.get("_results_email_sent"):
        _re_email = st.session_state.get("user_email", "")
        if _re_email:
            try:
                from utils.email_otp import send_results_email as _sre
                _re_ev = evs[0]
                _ok_re, _err_re = _sre(
                    to_email=_re_email,
                    conf_score=_re_ev.get("confidence_score", 0),
                    clar_score=_re_ev.get("clarity_score", 0),
                    top_fixes=_re_ev.get("fixes", [])[:3],
                    result_snippet=st.session_state.get("result_statement", ""),
                    verdict=_re_ev.get("verdict", ""),
                )
                if _ok_re:
                    st.toast(f"Results emailed to {_re_email}", icon="📧")
                    st.session_state["_results_email_sent"] = True
                else:
                    st.caption(f"📧 Results email could not send: {_err_re}")
            except Exception as _re_exc:
                st.caption(f"📧 Results email error: {_re_exc}")
        else:
            st.session_state["_results_email_sent"] = True

    if not evs:
        st.warning("No evaluation results found. Please go back and try again.")
        if st.button("← Edit Submission"):
            _go_to_screen(1)
        return

    _render_tutorial(2)

    n = len(evs)
    st.markdown(
        "<h2 style='color:#1B5E20;margin-bottom:4px;'>Your Confidence Snapshot</h2>"
        "<p style='color:#8A6500;font-style:italic;font-size:0.95rem;margin-bottom:8px;'>"
        "Evidence quality determination — what the evidence is worth to your donor, and what to fix first.</p>",
        unsafe_allow_html=True,
    )
    # Council XXIV — DRCA reproducibility badge (always visible, not in expander)
    from diagnostics import REPRODUCIBILITY_STATEMENT, STANDARDS_ANCHOR
    st.markdown(
        f'<div style="background:#F3F8FE;border-left:3px solid #1565C0;border-radius:6px;'
        f'padding:8px 12px;margin:0 0 12px 0;font-size:0.78rem;color:#374151;">'
        f'<strong style="color:#1565C0;">Scored against:</strong> {STANDARDS_ANCHOR}<br>'
        f'<span style="color:#616161;">Deterministic rule-based scoring — no AI judgement on scores. '
        f'Same inputs always produce the same result.</span>'
        f'</div>',
        unsafe_allow_html=True,
    )

    if evs:
        _bv_conf = evs[0].get("confidence_score", 0)
        _bv_clar = evs[0].get("clarity_score", 0)
        _bv_fixes = evs[0].get("fixes", [])
        if _bv_conf >= 3.5 and _bv_clar >= 3.5:
            _bv_sym, _bv_msg, _bv_bg, _bv_border = (
                "✓", "This result is ready to submit.",
                "#EDF7F1", "#1B5E20"
            )
        elif _bv_conf >= 2.5 or _bv_clar >= 2.5:
            _bv_sym, _bv_msg, _bv_bg, _bv_border = (
                "⚠", "Almost ready — fix the top issue before submitting.",
                "#FFFEF7", "#8A6500"
            )
        else:
            _top_fix_msg = (
                _bv_fixes[0].get("message", "review the flagged issues")
                if _bv_fixes else "review the flagged issues"
            )
            _bv_sym, _bv_msg, _bv_bg, _bv_border = (
                "✗", f"Not ready — {_top_fix_msg.rstrip('.')} before submitting.",
                "#FEF3F2", "#B71C1C"
            )
        _bv_threshold   = evs[0].get("threshold_used", 4.0) if evs else 4.0
        _bv_track_label = evs[0].get("track_label", "INGO standard") if evs else "INGO standard"
        _bv_scope_note = f" — Result 1 of {n}, see below for the rest" if n > 1 else ""
        st.caption(f"📋 Determination ({_bv_track_label}, threshold {_bv_threshold}){_bv_scope_note}:")
        st.markdown(
            f'<div style="background:{_bv_bg};border-left:4px solid {_bv_border};'
            f'border-radius:8px;padding:14px 20px;margin:0 0 16px 0;'
            f'font-size:1.05rem;font-weight:700;color:{_bv_border};">'
            f'{_bv_sym} {_bv_msg}</div>',
            unsafe_allow_html=True,
        )
        # Journey indicator — show INGO gap for CBO/National NGO users
        if _bv_track_label != "INGO standard":
            _ingo_gap = max(0.0, 4.0 - min(_bv_conf, _bv_clar))
            if _ingo_gap == 0:
                st.caption("📈 Capacity journey: ✓ also meets INGO standard (4.0)")
            else:
                st.caption(f"📈 Capacity journey: INGO-equivalent gap — {_ingo_gap:.1f} more needed on the weaker axis to meet the 4.0 bilateral-donor threshold.")

    # (Unfair advantage message removed — council XVIII: marketing copy at result stage is noise)

    # Primary download — 2–3 page Readiness Card (shareable with MEL lead / donor)
    timestamp   = datetime.now().strftime("%Y%m%d_%H%M%S")
    _report_allowed = st.session_state.get("_report_allowed", True)
    # Include council assessment (Page 4) if one has been run for this result
    _council_data = st.session_state.get("council_xxii_0")
    # Build the card (primary) and the full analysis (secondary / Advanced exports)
    _card_html   = _build_html_report_card(
        subs[0], evs[0], timestamp,
        field_sources=st.session_state.get("_irc_field_sources") if st.session_state.get("_irc_used") else None,
        council_assessment=_council_data,
    )
    _card_pdf    = _html_to_pdf_bytes(_card_html)
    html_report  = _build_html_report(subs[0], evs[0], timestamp) if n == 1 else \
                   _build_combined_html_report(subs, evs, timestamp)
    _pdf_primary = _html_to_pdf_bytes(html_report)

    # Generate a reproducible reference ID for the card — based on timestamp
    _ref_id = f"IMP-{timestamp}"

    if _report_allowed:
        if _card_pdf:
            st.download_button(
                f"📄 Download Pre-Submission Readiness Card  [{_ref_id}]",
                data=_card_pdf,
                file_name=f"readiness_card_{timestamp}.pdf",
                mime="application/pdf",
                use_container_width=True,
                type="primary",
                key="pdf_card_btn",
                help="Cite this card in your donor submission — includes scores, methodology anchor, and priority fixes. Reference ID: " + _ref_id,
            )
            st.caption(f"Ref: {_ref_id} · Decision record anchored to USAID ADS 201 · Bond 2024 · FCDO 2025 · Same inputs → same determination, always.")
        elif _card_html:
            st.download_button(
                f"⬇️ Download Pre-Submission Readiness Card (HTML)  [{_ref_id}]",
                data=_card_html.encode("utf-8"),
                file_name=f"readiness_card_{timestamp}.html",
                mime="text/html",
                use_container_width=True,
                type="primary",
                key="html_card_btn",
            )
    else:
        st.info("📄 **Your score is above.** Upgrade to download the Readiness Card.")
        _render_paywall(prompt_context="limit_hit")

    # Opt-in saved audit history — consent toggle right after the download,
    # the natural "also save this" moment. Off by default; a user who never
    # checks this sees no other change to this screen.
    _audit_email = st.session_state.get("user_email", "")
    if _audit_email:
        st.checkbox(
            "💾 Save this audit to my private history (encrypted at rest)",
            key="save_audit_consent",
            help="Stored in your account only. View, re-download, or delete past "
                 "audits anytime from My Audits in the sidebar.",
        )
        _audit_saved_key = f"_audit_saved_{_ref_id}"
        if st.session_state.get("save_audit_consent") and not st.session_state.get(_audit_saved_key):
            if not _safe_rate_limit_ok(_audit_email, "save_audit", max_count=10, window_seconds=3600):
                st.warning("You've saved a lot of audits in the last hour — please wait a bit before saving more.")
            elif save_audit(_audit_email, subs, evs, _ref_id):
                st.session_state[_audit_saved_key] = True
                st.success("✓ Saved to your private history.")
            else:
                st.warning(f"Could not save this audit right now — your download above still works. "
                           f"({last_audit_error()})")

    # Post-download CTA — primary next step after getting the determination
    _portfolio_cta_label = (
        f"📊 You've scored {n} results — run a portfolio analysis to see which is weakest →"
        if n >= 2 else
        "📊 Portfolio analysis — score all your indicators at once →"
    )
    if st.button(
        _portfolio_cta_label,
        key="s2_portfolio_cta",
        use_container_width=True,
        help="Upload a CSV of your full logframe to see which indicators are weakest.",
    ):
        _go_to_screen(3)

    # "Same programme" shortcut: clears result-specific fields, keeps programme context
    _SAME_PROG_CLEAR = [
        "result_statement", "target_group",
        "logframe_indicator", "logframe_baseline", "logframe_target", "logframe_achievement",
        "logframe_fill_later", "logframe_data_forthcoming",
        "evidence_description", "evidence_type", "verifier", "evidence_date",
        "internal_review", "external_review", "reporting_start", "reporting_end",
        "beneficiary_voice", "qual_evidence_flag", "qual_rigor_checklist",
        "additional_context", "limitations_notes", "learning_notes",
        "quick_evidence_desc", "quick_evidence_type",
    ]
    if st.button(
        "+ Score another result — same programme  (keeps sector, donor & timeframe)",
        key="s2_same_prog_cta",
        use_container_width=True,
        help="Clears only result-specific fields. Sector, donor, project name, timeframe, and geographic scope stay filled.",
    ):
        for _spk in _SAME_PROG_CLEAR:
            for _si in range(1, (st.session_state.get("active_slots", 1) + 1)):
                st.session_state.pop(f"{_spk}{_slot_suffix(_si)}", None)
            st.session_state.pop(_spk, None)
        for _ck in [k for k in list(st.session_state) if k.startswith("council_xxii_")]:
            del st.session_state[_ck]
        st.session_state["evaluations"] = None
        st.session_state["current_tab"] = 0
        st.query_params["tab"] = "0"
        _go_to_screen(1, reset=False)

    # Return hook — plant the seed for next quarter
    st.caption(
        "💡 **Save your Readiness Card.** The next time you check a result, your "
        "📈 Improvement over time tracker will show how your evidence quality has grown."
    )

    st.divider()

    # Ranked fix queue — visible before detail cards so users act, not just read (Council XXXI)
    _all_fixes = []
    for _ev in evs:
        _all_fixes.extend(_ev.get("fixes", []))
    if _all_fixes:
        st.markdown("### Ranked fix queue — act on these to improve your score")
        for _fi, _fix in enumerate(_all_fixes[:5], 1):
            _fix_msg    = _fix.get("message", "")
            _fix_impact = _fix.get("score_impact", "")
            _fix_val    = _fix.get("score_impact_value", 0)
            _fix_dim    = _fix.get("dimension", "")
            _fix_msg_lc = (_fix_msg or "").lower()
            if _fix_dim == "confidence":
                _fix_tab = 2
            elif any(_kw in _fix_msg_lc for _kw in ["indicator", "logframe", "baseline", "target", "achievement"]):
                _fix_tab = 1
            else:
                _fix_tab = 0
            _fq_c1, _fq_c2 = st.columns([5, 1])
            with _fq_c1:
                if _fix_impact and _fix_val:
                    st.markdown(f"**{_fi}.** {_fix_msg} — *adds {_fix_val:.2f} pts · {_fix_impact}*")
                elif _fix_impact:
                    st.markdown(f"**{_fi}.** {_fix_msg} *({_fix_impact})*")
                else:
                    st.markdown(f"**{_fi}.** {_fix_msg}")
            with _fq_c2:
                if st.button("↗ Fix", key=f"jump_fix_{_fi}", use_container_width=True, help="Go to this form field"):
                    st.session_state["evaluations"] = None
                    st.session_state["current_tab"] = _fix_tab
                    st.query_params["tab"] = str(_fix_tab)
                    _go_to_screen(1, reset=False)
        if len(evs) == 1:
            try:
                from council import _calculate_projected_scores
                _proj_c, _proj_cl = _calculate_projected_scores(evs[0])
                _cur_c  = evs[0].get("confidence_score", 0)
                _cur_cl = evs[0].get("clarity_score", 0)
                if _proj_c > _cur_c or _proj_cl > _cur_cl:
                    st.info(
                        f"**If you act on all {len(_all_fixes[:5])} fixes:** "
                        f"Projected Confidence → {_proj_c}/5.0 · Projected Clarity → {_proj_cl}/5.0"
                    )
            except Exception:
                pass

    # Elevated WhatsApp share CTA (Council XXVII) — moved here, after the fix queue
    # (when there is one), so sharing happens once the user has seen what to fix.
    # Unconditional on evs so a STRONG result with no fixes can still be shared.
    if evs:
        import urllib.parse as _urlparse
        _wa_ev0 = evs[0]
        _wa_cf = _wa_ev0.get("confidence_score", 0)
        _wa_cl = _wa_ev0.get("clarity_score", 0)
        _wa_fixes = _wa_ev0.get("fixes", [])
        _wa_state = _wa_ev0.get("diagnostic_state", "")
        _wa_tf = _wa_fixes[0]["message"] if _wa_fixes else "Review complete."
        def _wa_icon(s): return "✅" if s >= 4.0 else "⚠️" if s >= 3.0 else "🔴"
        _wa_text = (
            f"📊 ImpactProof Evidence Check\n"
            f"Confidence: {_wa_cf}/5.0 {_wa_icon(_wa_cf)}  ·  Clarity: {_wa_cl}/5.0 {_wa_icon(_wa_cl)}\n"
            f"Verdict: {_wa_state}\n"
            f"Top action: {_wa_tf}\n"
            f"Check your report: https://impact-proof.streamlit.app"
        )
        _wa_url = "https://wa.me/?text=" + _urlparse.quote(_wa_text)
        st.markdown(
            f'<a href="{_wa_url}" target="_blank" '
            f'style="display:block;text-align:center;background:#25D366;color:white;'
            f'padding:13px 18px;border-radius:8px;text-decoration:none;font-weight:700;'
            f'font-size:1rem;margin:8px 0;">📱 Share result on WhatsApp — tell your team what to fix</a>',
            unsafe_allow_html=True,
        )
        st.caption("Opens WhatsApp with your score, verdict, and top action. Works on phone and desktop.")
        st.divider()

    for i, (sub, ev) in enumerate(zip(subs, evs)):
        if n > 1:
            st.markdown(f"### Result {i + 1}")
        _render_result_card(sub, ev, card_idx=i,
                            donor=st.session_state.get("donor_selected", ""))

    # Optional reporting flags — moved here from Tab 2 so users can fill after seeing scores
    with st.expander("📝 Optional reporting flags (no score impact)", expanded=False):
        st.caption("These appear as advisory flags in your report — they do not affect your Confidence or Clarity scores.")
        _s2_slot_suffix = _slot_suffix(1)
        st.selectbox(
            "Does your report distinguish attribution from contribution?",
            options=["Not specified", "Yes", "No", "Not sure"],
            key=f"attribution_contribution{_s2_slot_suffix}",
            help="Attribution claims your program caused the change on its own. Contribution acknowledges multiple factors.",
        )
        st.selectbox(
            "Is beneficiary data disaggregated (women, youth, PWD, rural)?",
            options=["Not specified", "Yes — fully disaggregated", "Partially disaggregated", "No"],
            key=f"disaggregation_status{_s2_slot_suffix}",
            help="Many donors now expect results broken down by sex, age, disability, and location.",
        )
        st.text_area(
            "What did you learn from this result, and how did your program adapt?",
            key=f"learning_notes{_s2_slot_suffix}",
            placeholder="e.g., We learned that follow-up calls increased survey response rates, so we adjusted our M&E plan.",
        )
        st.text_area(
            "What can this data NOT confirm or be generalized to?",
            key=f"limitations_notes{_s2_slot_suffix}",
            placeholder="e.g., This sample covers only urban participants and cannot be generalized to rural areas.",
        )

    with st.expander("⚖️ Fairness & limitations", expanded=False):
        st.markdown(
            "**This tool now offers two evidence quality tracks:**\n\n"
            "- **INGO standard (threshold 4.0)** — calibrated to USAID, FCDO, GIZ, and "
            "Mastercard Foundation bilateral reporting requirements.\n"
            "- **Community / National standard (threshold 3.5–3.75)** — calibrated to "
            "local NGO, CBO, and national-funder contexts including STAR-Ghana, District "
            "Assembly grants, and national government reporting.\n\n"
            "Select your **Organisation type** on the evidence form (Tab 1) to activate "
            "the appropriate track. Your determination banner shows which standard was applied.\n\n"
            "**Community track recognises:**\n"
            "- Community elder council and ward committee verification\n"
            "- Community registers, village books, community scorecards, PRA outputs\n"
            "- Executive Director / board review in place of a dedicated MEL Officer\n"
            "- Oral testimony and participatory assessments in the qualitative evidence track\n\n"
            "**Capacity journey indicator:** CBO and National NGO users also see their "
            "INGO-equivalent score gap so the tool doubles as a capacity development roadmap — "
            "showing both where you stand for current funders and what is needed for bilateral donors.\n\n"
            "**Residual limitation:** Provenance bonuses (up to +0.60 on Confidence) still reward "
            "formal sampling documentation, independent enumerators, and auditor-retrievable records. "
            "Community organisations with informal but legitimate data collection may reach 3.5 "
            "without maximising provenance. Describe your collection method in **Additional Context** "
            "to document this in the Readiness Card.\n\n"
            "**On AI-assisted features:** The Council Assessment and Audit My Report functions "
            "use Anthropic's Claude API. These are advisory — not authoritative determinations. "
            "The evidence scoring engine itself makes no AI calls."
        )

    with st.expander("📚 Methodology & Citations", expanded=False):
        st.markdown("""
**ImpactProof scoring methodology is anchored in:**

- **Data Quality Standards** — adapted from USAID ADS 201.3.5.7, OECD-DAC 2019 evaluation criteria, and FCDO DQA guidance. Used for all Confidence and Clarity sub-scores.

- **Bond Evidence Principles 2024 (refresh)** — particularly Voice & Inclusion (operationalised as the Beneficiary Voice dimension) and Triangulation.

- **60 Decibels Lean Data Methodology** — informs the Beneficiary Voice scoring rubric.

- **Audit Logic** — classical audit independence principle (auditor independence from preparer) operationalised in Verification scoring.

**Why this matters for your donor:** Every sub-score traces to a named, citable standard. Hover over any score to see its specific anchor.

*ImpactProof is a pre-submission verification tool, not an audit. It does not replace formal evaluation but identifies gaps that would weaken external review.*
""")
        st.markdown("**Donor framework crosswalk** — how each sub-score maps to the standards your donor audits against:")
        st.markdown(
            build_donor_crosswalk_html(st.session_state.get("donor_framework", "Generic")),
            unsafe_allow_html=True,
        )

    st.divider()

    # Navigation buttons — simplified to 2 (council XVIII: remove LinkedIn, reduce noise)
    _nav_c1, _nav_c2 = st.columns(2)
    with _nav_c1:
        if st.button("← Edit this result", key="back_to_form", use_container_width=True):
            st.session_state["evaluations"] = None
            _go_to_screen(1)
    with _nav_c2:
        if st.button("✓ Check another result →", key="check_another", type="primary", use_container_width=True):
            _reset_all_slots()
            st.session_state["evaluations"] = None
            st.session_state["current_tab"] = 0
            st.query_params["tab"] = "0"
            _go_to_screen(1, reset=True)

    # Stage 2 Engagement Card — shown only for weak results
    if evs and evs[0].get("diagnostic_state", "") in ("FUNDAMENTALLY WEAK", "UNDEREVIDENCED", "MISLEADING"):
        st.divider()
        with st.container(border=True):
            st.markdown("#### Want someone to look at this with you?")
            _s2_ev0       = evs[0]
            _s2_rd        = {
                "conf":    _s2_ev0.get("confidence_score", "?"),
                "clar":    _s2_ev0.get("clarity_score", "?"),
                "verdict": _s2_ev0.get("diagnostic_state", ""),
            }
            _s2_email     = st.session_state.get("user_email", "")
            _s2_wa_key    = "wa_s2_weak_clicked"
            if st.button("📱 Book a free first review with the founder",
                         key="wa_weak_review_btn", type="primary",
                         use_container_width=True):
                from utils.whatsapp import notify_founder
                notify_founder("weak_result_review", user_email=_s2_email, result_data=_s2_rd)
                st.session_state[_s2_wa_key] = True
            if st.session_state.get(_s2_wa_key):
                from utils.whatsapp import build_wa_url
                _s2_wa_url = build_wa_url("weak_result_review", _s2_email, _s2_rd)
                st.link_button("Open WhatsApp to send →", _s2_wa_url, use_container_width=True)
                st.success("✓ The founder has been notified — you'll hear back within 24 hours.")
            st.caption("Download your report above and share it before the call.")

    _render_tutorial(3)

    # Secondary / additional exports — simplified 2-button row + advanced expander
    _sec_col1, _sec_col2 = st.columns(2)
    with _sec_col1:
        st.download_button(
            label="💾 Save Inputs (JSON)",
            data=_build_inputs_json(timestamp),
            file_name=f"impactproof-inputs-{timestamp}.json",
            mime="application/json",
            use_container_width=True,
            help="Save your form inputs as JSON so you can reload them later for iteration.",
        )
    with _sec_col2:
        if n < 3:
            if st.button("＋ Add Another Result", use_container_width=True):
                st.session_state["active_slots"] = n + 1
                st.session_state["evaluations"]  = None
                st.session_state["submissions_snapshot"] = None
                st.session_state["screen"] = 1
                st.rerun()

    with st.expander("Advanced exports (full analysis, verification summary, donor templates)"):
        # Full analysis — the verbose 20-section diagnostic report
        if _pdf_primary:
            st.download_button(
                label="📊 Download Full Analysis (PDF) — detailed 20-section diagnostic",
                data=_pdf_primary,
                file_name=f"full_analysis_{timestamp}.pdf",
                mime="application/pdf",
                use_container_width=True,
                key="pdf_full_analysis_btn",
            )
        elif html_report:
            st.download_button(
                label="📊 Download Full Analysis (HTML)",
                data=html_report,
                file_name=f"full_analysis_{timestamp}.html",
                mime="text/html",
                use_container_width=True,
                key="html_full_analysis_btn",
            )
        st.divider()
        if _pdf_primary:
            with st.expander("Report handoff details"):
                st.text_input("Prepared by (your name)", key="report_prepared_by")
                st.selectbox("Status", _REPORT_STATUS_OPTIONS, key="report_status")
                st.text_area("Notes for reviewer (optional)", key="report_notes", height=80)
            _summary_html = _build_verification_summary_html(subs, evs, timestamp)
            _summary_pdf  = _html_to_pdf_bytes(_summary_html)
            if _summary_pdf:
                st.download_button(
                    label="📋 Download Verification Summary (PDF)",
                    data=_summary_pdf,
                    file_name=f"verification_summary_{timestamp}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key="verification_summary_pdf_btn",
                )
            if _HAS_DOCX:
                _donor_key = st.selectbox(
                    _DONOR_TEMPLATE_COPY["selector_label"],
                    options=list(_DONOR_TEMPLATES.keys()),
                    format_func=lambda k: _DONOR_TEMPLATES[k]["label"],
                    key="donor_template_select",
                )
                _template = _DONOR_TEMPLATES[_donor_key]
                _field_rows = _donor_template_field_rows(_template, subs[0], evs[0])
                _missing_required = [r["Template field"] for r in _field_rows if r["required"] and not r["provided"]]
                if _missing_required:
                    st.warning(
                        _DONOR_TEMPLATE_COPY["missing_required_intro"] + "\n"
                        + "\n".join(f"- {m}" for m in _missing_required)
                    )
                st.download_button(
                    label=f"📄 Download {_template['label']} Template (DOCX)",
                    data=_build_donor_template_docx(_template, subs[0], evs[0], timestamp),
                    file_name=f"{_donor_key}_template_{timestamp}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key="donor_template_docx_btn",
                )
            else:
                st.caption(_DONOR_TEMPLATE_COPY["no_docx"])
        else:
            st.caption("PDF: install xhtml2pdf to enable one-click PDF download.")

    _render_tagline_footer()


# ---------------------------------------------------------------------------
# File persistence
# ---------------------------------------------------------------------------

def save_all_files(submission: dict, evaluation: dict) -> dict:
    os.makedirs("inputs", exist_ok=True)
    os.makedirs("evaluations", exist_ok=True)
    os.makedirs("outputs", exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    slug      = _make_slug(submission.get("result_statement", "result"))
    base      = f"{timestamp}_{slug}"

    paths = {
        "input":      os.path.join("inputs",      f"{base}_input.json"),
        "evaluation": os.path.join("evaluations", f"{base}_evaluation.json"),
        "output":     os.path.join("outputs",     f"{base}_report.md"),
    }

    with open(paths["input"], "w", encoding="utf-8") as f:
        json.dump(submission, f, indent=2, ensure_ascii=False, default=str)

    save_eval = {k: v for k, v in evaluation.items() if not k.startswith("_")}

    # --- Trends over time: additive fields, no scoring math changed ---
    _conf_comp = evaluation.get("confidence_components", {}) or {}
    _clar_comp = evaluation.get("clarity_components", {}) or {}
    _indicator_label = (submission.get("logframe_indicator") or "").strip()
    save_eval.update({
        "indicator_id":    _indicator_label.lower(),
        "indicator_label": _indicator_label,
        "project_name":    submission.get("project_name", ""),
        "submission_date": datetime.now().strftime("%Y-%m-%d"),
        "import_source":   submission.get("import_source", "manual"),
        "sub_scores": {
            "Directness":   _conf_comp.get("direct_score"),
            "Verification": _conf_comp.get("verify_score"),
            "Recency":      _conf_comp.get("recency_score"),
            "Definition":   _clar_comp.get("definition_score"),
            "Measurement":  _clar_comp.get("measurement_score"),
            "Integrity":    _clar_comp.get("integrity_score"),
            "Scope":        _clar_comp.get("scope_score"),
            "Governance":   _clar_comp.get("governance_score"),
        },
    })
    # --- END Trends over time ---

    with open(paths["evaluation"], "w", encoding="utf-8") as f:
        json.dump(save_eval, f, indent=2, ensure_ascii=False)

    report = _build_html_report(submission, evaluation, timestamp)
    with open(paths["output"], "w", encoding="utf-8") as f:
        f.write(report)

    return paths


def _make_slug(text: str, max_len: int = 45) -> str:
    slug = re.sub(r"[^\w\s-]", "", text.lower())
    slug = re.sub(r"[\s_-]+", "-", slug).strip("-")
    return slug[:max_len]


# ---------------------------------------------------------------------------
# Trends over time
# ---------------------------------------------------------------------------

_TREND_COPY = {
    "header": "📈 Your improvement over time",
    "intro": "See how each indicator's Confidence and Clarity scores improve across reporting cycles. Run a new portfolio check each quarter to build your history.",
    "no_data": "No saved submissions yet. Run a check on Screen 2 to start building history for this view.",
    "one_point": "Not enough history yet to show a trend — only one submission recorded so far for this indicator.",
    "select_primary": "Select an indicator/result to view its trend",
    "select_compare": "Compare with other indicators (optional)",
    "delta_caption": "Change since previous submission ({prev_date} → {curr_date}):",
}

_TREND_SUB_SCORE_DIMS = [
    "Directness", "Verification", "Recency", "Definition",
    "Measurement", "Integrity", "Scope", "Governance",
]


def _load_trend_history():
    """Scan evaluations/*_evaluation.json and build a tidy history of saved
    submissions: one row per file with indicator_id, indicator_label,
    project_name, date, confidence_score, clarity_score, and the 8 sub-scores.

    Pre-existing evaluation files (saved before this feature) are backfilled
    from their paired inputs/*_input.json and filename timestamp. Rows with
    no indicator label are skipped — never fabricated.
    """
    import glob
    import pandas as pd

    rows = []
    for path in sorted(glob.glob(os.path.join("evaluations", "*_evaluation.json"))):
        try:
            with open(path, encoding="utf-8") as f:
                ev = json.load(f)
        except Exception:
            continue

        base = os.path.basename(path)[: -len("_evaluation.json")]

        indicator_label = ev.get("indicator_label")
        project_name    = ev.get("project_name")
        date_str        = ev.get("submission_date")
        sub_scores      = ev.get("sub_scores")

        if indicator_label is None or project_name is None:
            sub = {}
            input_path = os.path.join("inputs", f"{base}_input.json")
            if os.path.exists(input_path):
                try:
                    with open(input_path, encoding="utf-8") as f:
                        sub = json.load(f)
                except Exception:
                    sub = {}
            if indicator_label is None:
                indicator_label = (sub.get("logframe_indicator") or "").strip()
            if project_name is None:
                project_name = sub.get("project_name", "")

        if not date_str:
            try:
                date_str = datetime.strptime(base[:15], "%Y%m%d_%H%M%S").strftime("%Y-%m-%d")
            except ValueError:
                date_str = ""

        if not indicator_label:
            continue  # can't place on a per-indicator trend without a label

        if not sub_scores:
            _conf_comp = ev.get("confidence_components", {}) or {}
            _clar_comp = ev.get("clarity_components", {}) or {}
            sub_scores = {
                "Directness":   _conf_comp.get("direct_score"),
                "Verification": _conf_comp.get("verify_score"),
                "Recency":      _conf_comp.get("recency_score"),
                "Definition":   _clar_comp.get("definition_score"),
                "Measurement":  _clar_comp.get("measurement_score"),
                "Integrity":    _clar_comp.get("integrity_score"),
                "Scope":        _clar_comp.get("scope_score"),
                "Governance":   _clar_comp.get("governance_score"),
            }

        row = {
            "indicator_id":     indicator_label.strip().lower(),
            "indicator_label":  indicator_label,
            "project_name":     project_name or "",
            "date":             date_str,
            "confidence_score": ev.get("confidence_score"),
            "clarity_score":    ev.get("clarity_score"),
        }
        for dim in _TREND_SUB_SCORE_DIMS:
            row[dim] = sub_scores.get(dim)
        rows.append(row)

    return pd.DataFrame(rows)


def _trend_indicator_options(history_df):
    """Return {indicator_id: display_label} for the selectors, in first-seen order."""
    options = {}
    for _, row in history_df.iterrows():
        iid = row["indicator_id"]
        if iid in options:
            continue
        label = row["indicator_label"]
        if row.get("project_name"):
            label = f"{label} — {row['project_name']}"
        options[iid] = label
    return options


def _render_trend_chart(plot_df, value_col, y_title, color_col=None):
    """Render a date-vs-score line chart, Altair normally, st.line_chart in lite mode."""
    if st.session_state.get("lite_mode", False):
        if color_col:
            pivoted = plot_df.pivot(index="date", columns=color_col, values=value_col)
        else:
            pivoted = plot_df.set_index("date")[[value_col]]
        st.line_chart(pivoted)
        return

    import altair as alt
    enc = {
        "x": alt.X("date:T", title="Date"),
        "y": alt.Y(f"{value_col}:Q", title=y_title, scale=alt.Scale(domain=[0, 5])),
        "tooltip": ["date:T", f"{value_col}:Q"] + ([f"{color_col}:N"] if color_col else []),
    }
    if color_col:
        enc["color"] = alt.Color(f"{color_col}:N", title=None)
        enc["tooltip"] = ["date:T", f"{color_col}:N", f"{value_col}:Q"]
    chart = alt.Chart(plot_df).mark_line(point=True).encode(**enc).properties(width="container")
    st.altair_chart(chart, use_container_width=True)


def render_trends_view(history_df) -> None:
    """Render the 'Trends over time' section for a selected indicator,
    with optional side-by-side comparison against other indicators."""
    import pandas as pd

    if history_df.empty:
        st.info(_TREND_COPY["no_data"])
        return

    options = _trend_indicator_options(history_df)
    ids = list(options.keys())

    primary_id = st.selectbox(
        _TREND_COPY["select_primary"],
        options=ids,
        format_func=lambda i: options[i],
        key="trend_primary_indicator",
    )
    compare_ids = st.multiselect(
        _TREND_COPY["select_compare"],
        options=[i for i in ids if i != primary_id],
        format_func=lambda i: options[i],
        key="trend_compare_indicators",
    )

    primary_df = (history_df[history_df["indicator_id"] == primary_id]
                   .sort_values("date").reset_index(drop=True))

    if len(primary_df) == 1:
        row = primary_df.iloc[0]
        st.info(_TREND_COPY["one_point"])
        cols = st.columns(2)
        cols[0].metric("Confidence", f"{row['confidence_score']}/5.0")
        cols[1].metric("Clarity", f"{row['clarity_score']}/5.0")
        sub_cols = st.columns(4)
        for i, dim in enumerate(_TREND_SUB_SCORE_DIMS):
            val = row.get(dim)
            if pd.notna(val):
                sub_cols[i % 4].metric(dim, val)
    else:
        long_rows = []
        for _, row in primary_df.iterrows():
            long_rows.append({"date": row["date"], "series": "Confidence", "score": row["confidence_score"]})
            long_rows.append({"date": row["date"], "series": "Clarity", "score": row["clarity_score"]})
        _render_trend_chart(pd.DataFrame(long_rows), "score", "Score (0-5)", color_col="series")

        prev, curr = primary_df.iloc[-2], primary_df.iloc[-1]
        st.caption(_TREND_COPY["delta_caption"].format(prev_date=prev["date"], curr_date=curr["date"]))
        cols = st.columns(2)
        cols[0].metric("Confidence", f"{curr['confidence_score']}/5.0",
                        delta=round(curr["confidence_score"] - prev["confidence_score"], 2))
        cols[1].metric("Clarity", f"{curr['clarity_score']}/5.0",
                        delta=round(curr["clarity_score"] - prev["clarity_score"], 2))
        sub_cols = st.columns(4)
        for i, dim in enumerate(_TREND_SUB_SCORE_DIMS):
            cv, pv = curr.get(dim), prev.get(dim)
            if pd.isna(cv) or pd.isna(pv):
                continue
            delta = round(cv - pv, 2)
            if delta != 0:
                sub_cols[i % 4].metric(dim, cv, delta=delta)

    if compare_ids:
        compare_df = history_df[history_df["indicator_id"].isin([primary_id] + compare_ids)].copy()
        compare_df["label"] = compare_df["indicator_id"].map(options)
        st.markdown("##### Confidence comparison")
        _render_trend_chart(compare_df.sort_values("date"), "confidence_score", "Confidence (0-5)", color_col="label")
        st.markdown("##### Clarity comparison")
        _render_trend_chart(compare_df.sort_values("date"), "clarity_score", "Clarity (0-5)", color_col="label")


# ---------------------------------------------------------------------------
# Donor template export
# ---------------------------------------------------------------------------

_DONOR_TEMPLATE_COPY = {
    "section_header": "Donor Template Export",
    "selector_label": "Export for:",
    "preview_label": "Preview field mapping",
    "missing_required_intro": "Before exporting, note: these required fields have no verified data yet and will be marked \"Not provided\" in the export:",
    "no_docx": "Donor template export needs python-docx. Install it to enable this.",
}

_DONOR_TEMPLATES = {
    "generic": {
        "label": "Generic",
        "sections": [
            {
                "heading": "Result Overview",
                "fields": [
                    {"label": "Project / Programme Name", "source": "submission.project_name"},
                    {"label": "Result Statement",          "source": "submission.result_statement"},
                    {"label": "Logframe Indicator",         "source": "submission.logframe_indicator"},
                    {"label": "Reporting Period",           "source": "submission.timeframe"},
                    {"label": "Target Group",               "source": "submission.target_group"},
                    {"label": "Geographic Scope",           "source": "submission.geographic_scope"},
                ],
                "required": ["Project / Programme Name", "Result Statement", "Logframe Indicator"],
            },
            {
                "heading": "Evidence Quality Scores",
                "fields": [
                    {"label": "Confidence Score (0-5)", "source": "evaluation.confidence_score"},
                    {"label": "Confidence Rating",       "source": "evaluation.confidence_label"},
                    {"label": "Clarity Score (0-5)",     "source": "evaluation.clarity_score"},
                    {"label": "Clarity Rating",          "source": "evaluation.clarity_label"},
                ],
                "required": [],
            },
            {
                "heading": "Verification & Sign-off",
                "fields": [
                    {"label": "Internal Review Level", "source": "submission.internal_review"},
                    {"label": "External Review Level", "source": "submission.external_review"},
                    {"label": "Prepared By",            "source": "session.report_prepared_by"},
                    {"label": "Status",                 "source": "session.report_status"},
                ],
                "required": [],
            },
        ],
    },
    "usaid": {
        "label": "USAID",
        "sections": [
            {
                "heading": "Activity / Indicator Information",
                "fields": [
                    {"label": "Activity Name",          "source": "submission.project_name"},
                    {"label": "Performance Indicator",  "source": "submission.logframe_indicator"},
                    {"label": "Reporting Period",        "source": "submission.timeframe"},
                    {"label": "Target",                  "source": "submission.logframe_target"},
                    {"label": "Achieved Result",         "source": "submission.logframe_achievement"},
                ],
                "required": ["Activity Name", "Performance Indicator", "Target", "Achieved Result"],
            },
            {
                "heading": "Data Quality Assessment (DQA)",
                "fields": [
                    {"label": "Validity (Directness)",                  "source": "subscore.Directness"},
                    {"label": "Reliability / Integrity (Verification)", "source": "subscore.Verification"},
                    {"label": "Timeliness (Recency)",                   "source": "subscore.Recency"},
                    {"label": "Precision (Definition)",                 "source": "subscore.Definition"},
                    {"label": "Precision (Measurement)",                "source": "subscore.Measurement"},
                    {"label": "Integrity (Ethics)",                     "source": "subscore.Integrity"},
                    {"label": "Confidence Rating",                      "source": "evaluation.confidence_label"},
                    {"label": "Clarity Rating",                         "source": "evaluation.clarity_label"},
                ],
                "required": ["Validity (Directness)", "Reliability / Integrity (Verification)"],
            },
            {
                "heading": "DQA Limitations & Sign-off",
                "fields": [
                    {"label": "Known Data Limitations", "source": "submission.limitations_notes"},
                    {"label": "Prepared By",             "source": "session.report_prepared_by"},
                    {"label": "Status",                  "source": "session.report_status"},
                ],
                "required": [],
            },
        ],
    },
    "fcdo": {
        "label": "FCDO",
        "sections": [
            {
                "heading": "Output / Outcome Information",
                "fields": [
                    {"label": "Programme Name",         "source": "session.project_name"},
                    {"label": "Output / Outcome Title", "source": "submission.result_statement"},
                    {"label": "Logframe Indicator",     "source": "submission.logframe_indicator"},
                    {"label": "Reporting Period",       "source": "submission.timeframe"},
                    {"label": "Target",                 "source": "submission.logframe_target"},
                    {"label": "Achievement",            "source": "submission.logframe_achievement"},
                    {"label": "Geographic Scope",       "source": "submission.geographic_scope"},
                ],
                "required": ["Programme Name", "Output / Outcome Title", "Logframe Indicator", "Reporting Period"],
            },
            {
                "heading": "Evidence Quality (Bond Evidence Principles 2024)",
                "fields": [
                    {"label": "Confidence Score (0–5)",     "source": "evaluation.confidence_score"},
                    {"label": "Confidence Rating",          "source": "evaluation.confidence_label"},
                    {"label": "Clarity Score (0–5)",        "source": "evaluation.clarity_score"},
                    {"label": "Clarity Rating",             "source": "evaluation.clarity_label"},
                    {"label": "Beneficiary Voice",          "source": "submission.beneficiary_voice"},
                    {"label": "Attribution / Contribution", "source": "submission.attribution_contribution"},
                ],
                "required": [],
            },
            {
                "heading": "VfM & Learning",
                "fields": [
                    {"label": "What the team learned",    "source": "submission.learning_notes"},
                    {"label": "Limitations of this data", "source": "submission.limitations_notes"},
                    {"label": "Internal Review Level",    "source": "submission.internal_review"},
                    {"label": "External Review Level",    "source": "submission.external_review"},
                    {"label": "Prepared By",              "source": "session.report_prepared_by"},
                    {"label": "Status",                   "source": "session.report_status"},
                ],
                "required": [],
            },
        ],
    },
}


def _resolve_donor_field_value(source: str, submission: dict, evaluation: dict):
    """Resolve a 'namespace.key' source string from a donor template config
    against the current submission/evaluation/session data.
    Returns (value, provided) — provided=False means no real data exists,
    caller must render 'Not provided' and never fabricate a value."""
    if not source:
        return None, False
    namespace, _, key = source.partition(".")
    if namespace == "submission":
        val = submission.get(key)
    elif namespace == "evaluation":
        val = evaluation.get(key)
    elif namespace == "subscore":
        conf = evaluation.get("confidence_components", {}) or {}
        clar = evaluation.get("clarity_components", {}) or {}
        val = {
            "Directness":   conf.get("direct_score"),
            "Verification": conf.get("verify_score"),
            "Recency":      conf.get("recency_score"),
            "Definition":   clar.get("definition_score"),
            "Measurement":  clar.get("measurement_score"),
            "Integrity":    clar.get("integrity_score"),
            "Scope":        clar.get("scope_score"),
            "Governance":   clar.get("governance_score"),
        }.get(key)
    elif namespace == "session":
        val = st.session_state.get(key)
    else:
        val = None
    provided = val is not None and str(val).strip() != ""
    return val, provided


def _donor_template_field_rows(template: dict, submission: dict, evaluation: dict) -> list:
    """One row per template field: section, label, source, value (or
    'Not provided'), provided flag, required flag. Used for both the
    mapping-preview table and the pre-export checklist."""
    rows = []
    for section in template["sections"]:
        required = set(section.get("required", []))
        for field in section["fields"]:
            value, provided = _resolve_donor_field_value(field["source"], submission, evaluation)
            rows.append({
                "Section": section["heading"],
                "Template field": field["label"],
                "Diagnostic source": field["source"],
                "Value": value if provided else "Not provided",
                "provided": provided,
                "required": field["label"] in required,
            })
    return rows


def _build_donor_template_docx(template: dict, submission: dict, evaluation: dict, timestamp: str) -> bytes:
    """Render submission/evaluation into the selected donor template as an
    editable DOCX. Fields with no verified data are written literally as
    'Not provided' — never inferred or fabricated."""
    import io

    doc = _docx.Document()
    doc.add_heading(f"{template['label']} Report", level=0)
    doc.add_paragraph(f"Generated by ImpactProof: {timestamp}")

    for section in template["sections"]:
        doc.add_heading(section["heading"], level=1)
        table = doc.add_table(rows=0, cols=2)
        table.style = "Light Grid Accent 1"
        for field in section["fields"]:
            value, provided = _resolve_donor_field_value(field["source"], submission, evaluation)
            cells = table.add_row().cells
            cells[0].text = field["label"]
            cells[1].text = str(value) if provided else "Not provided"

    doc.add_heading("About this export", level=1)
    doc.add_paragraph(
        f"Evaluated using {METHODOLOGY_STACK}. {_LIMITS_DISCLAIMER} "
        "Fields marked 'Not provided' had no verified data entered for this "
        "submission and were intentionally left blank."
    )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _html_to_pdf_bytes(html: str) -> bytes | None:
    """Convert an HTML report string to PDF bytes via xhtml2pdf, or None if unavailable."""
    try:
        from xhtml2pdf import pisa as _pisa
        import io as _io2
        buf = _io2.BytesIO()
        _pisa.CreatePDF(src=html, dest=buf)
        if buf.tell() > 0:
            buf.seek(0)
            return buf.getvalue()
    except ImportError:
        pass
    return None


def _generate_evidence_statement(submission: dict) -> str:
    """Generate a copy-pasteable evidence statement for the narrative report
    from the submission data already entered. No AI — template only."""
    ev      = (submission.get("evidence") or [{}])[0]
    ev_type = ev.get("type", "")
    recency = ev.get("recency", "")
    verifier= ev.get("verified_by", "")
    tg      = submission.get("target_group", "")
    tf      = submission.get("timeframe", "")
    bv      = submission.get("beneficiary_voice", "")
    parts   = [f"Evidence for this result was collected via {ev_type}" if ev_type else "Evidence was collected"]
    if tg:  parts[0] += f" among {tg}"
    if tf:  parts[0] += f" during {tf}"
    parts[0] += "."
    if verifier:
        parts.append(f"Data were independently reviewed by {verifier}.")
    if bv and not any(skip in bv for skip in ("No beneficiary", "Not applicable", "Choose")):
        parts.append(f"Beneficiary perspective: {bv}.")
    if recency:
        parts.append(f"Evidence date: {recency}.")
    return " ".join(parts)


# ---------------------------------------------------------------------------
# Score My Report — batch document extraction + scoring pipeline (council XVII)
# ---------------------------------------------------------------------------

def _recover_partial_json_results(raw: str) -> list[dict]:
    """Walk character-by-character through a potentially truncated JSON response
    and return every complete result object found inside the "results" array.

    This handles the case where max_tokens cuts the response mid-string,
    producing an unterminated JSON document.
    """
    import json as _json

    # Locate the results array
    results_key = raw.find('"results"')
    if results_key == -1:
        return []
    array_start = raw.find('[', results_key)
    if array_start == -1:
        return []

    # Walk through characters tracking brace depth to find complete objects
    depth = 0
    in_string = False
    escape_next = False
    obj_start = -1
    complete_results = []

    for i in range(array_start + 1, len(raw)):
        ch = raw[i]
        if escape_next:
            escape_next = False
            continue
        if ch == '\\' and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == '{':
            if depth == 0:
                obj_start = i
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0 and obj_start != -1:
                try:
                    obj = _json.loads(raw[obj_start:i + 1])
                    complete_results.append(obj)
                except _json.JSONDecodeError:
                    pass
                obj_start = -1

    return complete_results


def _extract_all_results_from_document(document_text: str, api_key: str,
                                        max_chars: int = 60000) -> tuple[list[dict], str]:
    """Call Claude with BATCH_EXTRACTION_SYSTEM_PROMPT to extract all results.

    Returns (results_list, error_message).  error_message is "" on success.
    Each item in results_list is a raw dict from the JSON response.

    If the response is truncated (hits max_tokens), _recover_partial_json_results()
    salvages every complete result object from the partial output.
    """
    import anthropic as _anthr
    import json as _json

    text = document_text[:max_chars]
    try:
        client = _anthr.Anthropic(api_key=api_key)
        resp = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=8192,
            system=BATCH_EXTRACTION_SYSTEM_PROMPT,
            messages=[{"role": "user", "content": f"Document text:\n\n{text}"}],
        )
        raw = resp.content[0].text.strip() if resp.content else ""
        # Strip markdown fences if present
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

        # Try clean parse first
        try:
            data = _json.loads(raw)
            results = data.get("results", [])
            if isinstance(results, list) and results:
                return results, ""
        except _json.JSONDecodeError:
            pass

        # Response was likely truncated — recover complete objects
        recovered = _recover_partial_json_results(raw)
        if recovered:
            return recovered, ""

        return [], "Could not extract any complete results from the document. The response may have been too large. Try uploading a shorter document or a document with fewer indicators."
    except Exception as e:
        import logging as _logging
        _logging.error("Batch extraction failed", exc_info=True)
        return [], "Extraction failed unexpectedly. Please try again, or try a shorter document."


def _batch_results_to_portfolio_df(raw_results: list[dict]) -> tuple["pd.DataFrame", list[dict]]:
    """Convert raw extraction dicts to a portfolio DataFrame and field status tracker.

    Returns (df, statuses) where statuses is a parallel list of {field: STATUS_*} dicts.
    """
    import pandas as pd
    from excel_report import STATUS_CONFIRMED, STATUS_AUTO_POPULATED, STATUS_NOT_FOUND

    _NOT_FOUND_VALUES = {"Not found", "not found", "NOT FOUND", "", None}

    rows = []
    statuses = []
    for r in raw_results:
        confidence = r.get("field_confidence", {})
        row = {}
        status = {}
        for col_name, _, _ in _PORTFOLIO_COLUMNS:
            # Map extraction field names to portfolio column names
            val = r.get(col_name, "Not found")
            if val is None or str(val).strip() in _NOT_FOUND_VALUES:
                val = ""
                status[col_name] = STATUS_NOT_FOUND
            else:
                conf = confidence.get(col_name, "medium")
                status[col_name] = STATUS_CONFIRMED if conf == "high" else STATUS_AUTO_POPULATED
            # geographic_scope may be a list
            if isinstance(val, list):
                val = "; ".join(str(v) for v in val)
            row[col_name] = str(val)
        # Overall review status — NOT_FOUND if any critical field is missing
        critical = ["result_statement", "evidence_description", "evidence_type"]
        overall = STATUS_AUTO_POPULATED
        if any(status.get(k) == STATUS_NOT_FOUND for k in critical):
            overall = STATUS_NOT_FOUND
        status["_overall"] = overall
        rows.append(row)
        statuses.append(status)

    col_names = [c[0] for c in _PORTFOLIO_COLUMNS]
    df = pd.DataFrame(rows, columns=col_names) if rows else pd.DataFrame(columns=col_names)
    return df, statuses


def _score_report_from_document(document_text: str, api_key: str) -> tuple[
    "pd.DataFrame", list[dict], list[dict], str
]:
    """Full pipeline: document text → scored portfolio DataFrame.

    Returns (input_df, evaluations, field_statuses, error_message).
    evaluations[i] is the result of evaluator.evaluate_submission(row_i).
    """
    raw_results, err = _extract_all_results_from_document(document_text, api_key)
    if err:
        return None, [], [], err
    if not raw_results:
        return None, [], [], "No results found in document."

    input_df, statuses = _batch_results_to_portfolio_df(raw_results)
    _, warnings = _evaluate_portfolio(input_df)

    # Run evaluation per row
    evaluations = []
    for _, row in input_df.iterrows():
        sub = _portfolio_row_to_submission(row.to_dict())
        try:
            ev = _evaluator.evaluate_submission(sub)
        except Exception:
            ev = {"confidence_score": 0, "clarity_score": 0,
                  "verdict": "Could not evaluate", "fixes": []}
        evaluations.append(ev)

    return input_df, evaluations, statuses, ""


def _build_markdown_report(submission: dict, evaluation: dict, timestamp: str) -> str:
    conf_score = evaluation.get("confidence_score", 0)
    clar_score = evaluation.get("clarity_score", 0)
    conf_label = evaluation.get("confidence_label", "")
    clar_label = evaluation.get("clarity_label", "")
    verdict    = evaluation.get("verdict", "")
    fixes      = evaluation.get("fixes", [])
    conf_comp  = evaluation.get("confidence_components", {})
    clar_comp  = evaluation.get("clarity_components", {})
    filenames  = submission.get("attached_filenames", [])

    lines = [
        "# ImpactProof — Evidence Quality Report",
        f"**Generated:** {timestamp}",
        f"**Confidence Score:** {conf_score}/5.0 ({conf_label})",
        f"**Clarity Score:** {clar_score}/5.0 ({clar_label})",
        f"**Verdict:** {verdict}",
        "",
        "---",
        "",
        "## Result Statement",
        f"{submission.get('result_statement', '-')}",
        "",
        f"- **Target Group:** {submission.get('target_group', '-')}",
        f"- **Timeframe:** {submission.get('timeframe', '-')}",
        f"- **Geographic Scope:** {submission.get('geographic_scope', '-')}",
        "",
    ]

    if filenames:
        lines += [f"- **Attached documents:** {', '.join(filenames)}", ""]

    # Logframe comparison table — shown when baseline + target + achievement are all present
    _lf_baseline    = submission.get("logframe_baseline", "").strip()
    _lf_indicator   = submission.get("logframe_indicator", "").strip()
    _lf_target      = submission.get("logframe_target", "").strip()
    _lf_achievement = submission.get("logframe_achievement", "").strip()
    _lf_pct         = (evaluation.get("logframe_linkage") or {}).get("pct_of_target")
    _lf_dir_miss    = (evaluation.get("logframe_linkage") or {}).get("direction_mismatch", False)
    _lf_forthcoming = submission.get("logframe_data_forthcoming", False)

    if _lf_indicator:
        lines += ["---", "", "## Logframe Linkage", "", f"**Indicator:** {_lf_indicator}", ""]
        if _lf_baseline or _lf_target or _lf_achievement:
            lines += [
                "| | Value |",
                "|---|---|",
                f"| **Baseline / Pre-evaluation** | {_lf_baseline or '—'} |",
                f"| **Approved target** | {_lf_target or '—'} |",
                f"| **Actual achievement** | {_lf_achievement or '—'} |",
                f"| **% of target reached** | {f'{_lf_pct:.0f}%' if _lf_pct is not None else '—'} |",
                "",
            ]
        if _lf_dir_miss:
            lines += ["> **Direction mismatch:** The indicator implies a change in one direction but the "
                      "baseline-to-achievement comparison shows the opposite. Review framing before submission.", ""]
        if _lf_forthcoming:
            lines += ["> **Data gap disclosed:** Measurement not yet collected — state this in the submission narrative.", ""]

    lines += [
        "---",
        "",
        "## Confidence Score Breakdown",
        "",
        "| Component | Level | Score |",
        "|---|---|---|",
        f"| Directness   | {conf_comp.get('direct_level', '-')}/5  | {conf_comp.get('direct_score', '-')}/2.0 |",
        f"| Verification | {conf_comp.get('verify_level', '-')}/5  | {conf_comp.get('verify_score', '-')}/2.0 |",
        f"| Recency      | {conf_comp.get('recency_level', '-')}/5 | {conf_comp.get('recency_score', '-')}/1.0 |",
        f"| **Total**    |                                          | **{conf_score}/5.0** |",
        "",
        "## Clarity Score Breakdown",
        "",
        "| Component   | Score | Max  |",
        "|---|---|---|",
        f"| Definition  | {clar_comp.get('definition_score', '-')}  | 1.25 |",
        f"| Measurement | {clar_comp.get('measurement_score', '-')} | 1.25 |",
        f"| Integrity   | {clar_comp.get('integrity_score', '-')}   | 1.0  |",
        f"| Scope       | {clar_comp.get('scope_score', '-')}       | 0.75 |",
        f"| Governance  | {clar_comp.get('governance_score', '-')}  | 0.75 |",
        f"| **Total**   | **{clar_score}**                          | **5.0** |",
        "",
        "---",
        "",
        "## What To Fix",
        "",
    ]

    conf_fixes = [f for f in fixes if f.get("dimension") == "confidence"]
    clar_fixes = [f for f in fixes if f.get("dimension") == "clarity"]

    if conf_fixes:
        lines.append("### Strengthen your evidence (Confidence)")
        for fix in conf_fixes:
            lines.append(f"- [ ] {fix['message']}  *({fix['score_impact']})*")
        lines.append("")

    if clar_fixes:
        lines.append("### Sharpen your definition (Clarity)")
        for fix in clar_fixes:
            lines.append(f"- [ ] {fix['message']}  *({fix['score_impact']})*")
        lines.append("")

    if not fixes:
        lines += ["This check found no further fixes to address.", ""]

    # Append last 3 chat exchanges as a due-diligence record
    _card_idx = submission.get("_card_idx", 0)
    _hist_key = f"chat_messages_{_card_idx}" if _card_idx else "chat_messages"
    _chat_msgs = st.session_state.get(_hist_key, []) or st.session_state.get("chat_messages", [])
    if _chat_msgs:
        lines += ["---", "", "## Score Discussion (session)", ""]
        for _m in _chat_msgs[-6:]:  # last 3 exchanges (user + assistant × 3)
            _role = "You" if _m["role"] == "user" else "Scoring Assistant"
            lines += [f"**{_role}:** {_m['content']}", ""]

    lines += [
        "---",
        "",
        f"*Evaluated using: {METHODOLOGY_STACK}*",
    ]

    return "\n".join(lines)



# ---------------------------------------------------------------------------
# Portfolio / Framework Dashboard — helpers
# ---------------------------------------------------------------------------

_PORTFOLIO_SUBSCORE_DIMENSIONS = [
    ("Directness",   2.0),
    ("Verification", 2.0),
    ("Recency",      1.0),
    ("Definition",   1.25),
    ("Measurement",  1.25),
    ("Integrity",    1.0),
    ("Scope",        0.75),
    ("Governance",   0.75),
]


def _portfolio_template_csv() -> bytes:
    """Build a one-row example CSV for the Portfolio Dashboard upload (full 29-column version)."""
    import pandas as pd

    scoring_hints = (
        "# SCORING IMPACT: internal_review + external_review + verifier → Verification (Confidence axis) | "
        "provenance_* fields → Verification score | beneficiary_voice → +0 to +0.5 Confidence bonus | "
        "additional_context → Governance sub-score (Clarity) | "
        "qual_* fields → only apply when evidence_type is qualitative"
    )
    headers = [c[0] for c in _PORTFOLIO_COLUMNS]
    example = {c[0]: c[2] for c in _PORTFOLIO_COLUMNS}
    df = pd.DataFrame([example], columns=headers)
    csv_body = df.to_csv(index=False)
    return (scoring_hints + "\n" + csv_body).encode("utf-8")


def _portfolio_minimal_template_csv() -> bytes:
    """Build a minimal 7-column template for first-time Portfolio users."""
    import pandas as pd
    required = [(c[0], c[2]) for c in _PORTFOLIO_COLUMNS if c[1]]
    headers = [c[0] for c in required]
    example = {c[0]: c[1] for c in required}
    note = "# Fill one row per indicator/result. All 7 columns are required. Download the full template for optional scoring-boost columns."
    df = pd.DataFrame([example], columns=headers)
    return (note + "\n" + df.to_csv(index=False)).encode("utf-8")


def _portfolio_row_to_submission(row: dict) -> dict:
    """Map one CSV/Excel row to an evaluator-compatible submission dict."""
    def _get(key, default=""):
        val = row.get(key, default)
        if val is None:
            return default
        val = str(val).strip()
        if val.lower() == "nan":
            return default
        return val or default

    internal_review = _get("internal_review", "Not reviewed") or "Not reviewed"
    external_review = _get("external_review", "No external review") or "No external review"

    def _get_bool(key):
        return _get(key, "FALSE").strip().lower() in ("true", "yes", "1", "x")

    return {
        "result_statement":   _get("result_statement"),
        "target_group":       _get("target_group"),
        "timeframe":          _get("timeframe"),
        "geographic_scope":   _get("geographic_scope"),
        "additional_context": _get("additional_context"),
        "learning_notes":     _get("learning_notes"),
        "limitations_notes":  _get("limitations_notes"),
        "internal_review":    internal_review,
        "external_review":    external_review,
        "logframe_indicator":   _get("logframe_indicator"),
        "logframe_target":      _get("logframe_target"),
        "logframe_achievement": _get("logframe_achievement"),
        "beneficiary_voice":    _get("beneficiary_voice", "No beneficiary voice captured"),
        "evidence": [{
            "type":        _get("evidence_type"),
            "description": _get("evidence_description"),
            "recency":     _get("evidence_date"),
            "verified_by": _get("verifier"),
        }],
        "qualitative_rigor_checklist": {
            "sourcing_documented":           _get_bool("qual_sourcing_documented"),
            "triangulated":                  _get_bool("qual_triangulated"),
            "bias_considered":               _get_bool("qual_bias_considered"),
            "beneficiary_voice_represented": _get_bool("qual_beneficiary_voice_represented"),
            "consent_ethics_addressed":      _get_bool("qual_consent_ethics_addressed"),
        },
        "provenance_checklist": {
            "sampling_documented":     _get("provenance_sampling",    "Not applicable"),
            "double_counting_checked": _get("provenance_dedup",       "Not applicable"),
            "collection_tool_named":   _get("provenance_tool",        "Not applicable"),
            "collector_independent":   _get("provenance_independent", "Not applicable"),
            "recall_period_ok":        _get("provenance_recall",      "Not applicable"),
            "auditor_traceable":       _get("provenance_traceable",   "Choose an option..."),
        },
    }


def _evaluate_portfolio(df):
    """Evaluate each row of a portfolio DataFrame.

    Returns (results_df, warnings) where results_df is sorted by
    confidence_score ascending (weakest indicators first).
    """
    import pandas as pd

    required = [c[0] for c in _PORTFOLIO_COLUMNS if c[1]]
    missing_cols = [c for c in required if c not in df.columns]
    if missing_cols:
        return None, [f"Missing required column(s): {', '.join(missing_cols)}"]

    rows = []
    warnings = []
    for i, row in df.iterrows():
        row_dict = row.to_dict()
        label = str(row_dict.get("indicator_name", "")).strip() or f"Row {i + 2}"
        try:
            sub = _portfolio_row_to_submission(row_dict)
            ev = _evaluator.evaluate_submission(sub)
            conf_comp = ev.get("confidence_components", {})
            clar_comp = ev.get("clarity_components", {})
            sub_scores = {
                "Directness":   conf_comp.get("direct_score", 0),
                "Verification": conf_comp.get("verify_score", 0),
                "Recency":      conf_comp.get("recency_score", 0),
                "Definition":   clar_comp.get("definition_score", 0),
                "Measurement":  clar_comp.get("measurement_score", 0),
                "Integrity":    clar_comp.get("integrity_score", 0),
                "Scope":        clar_comp.get("scope_score", 0),
                "Governance":   clar_comp.get("governance_score", 0),
            }
            _linkage = ev.get("logframe_linkage", {})
            result_row = {
                "indicator_name":    label,
                "logframe_indicator": sub.get("logframe_indicator", ""),
                "confidence_score":  ev.get("confidence_score", 0),
                "clarity_score":     ev.get("clarity_score", 0),
                "verdict":           ev.get("verdict", ""),
                "top_fix":           (ev.get("fixes") or [{}])[0].get("message", "") if ev.get("fixes") else "",
                "pct_of_target":     _linkage.get("pct_of_target"),        # float or None
                "direction_mismatch": _linkage.get("direction_mismatch", False),
            }
            for dim, max_val in _PORTFOLIO_SUBSCORE_DIMENSIONS:
                pct = round(min(sub_scores[dim] / max_val, 1.0) * 100, 1) if max_val else 0.0
                result_row[dim] = pct
            rows.append(result_row)
        except Exception as exc:
            warnings.append(f"{label}: could not be evaluated ({exc})")

    if not rows:
        return None, warnings or ["No rows could be evaluated."]

    results_df = pd.DataFrame(rows).sort_values("confidence_score", ascending=True).reset_index(drop=True)
    return results_df, warnings


def _portfolio_heatmap_chart(results_df):
    """Build an Altair heatmap: rows = indicators, columns = sub-score dimensions."""
    import pandas as pd
    import altair as alt

    dims = [d for d, _ in _PORTFOLIO_SUBSCORE_DIMENSIONS]
    long_df = results_df.melt(
        id_vars=["indicator_name"], value_vars=dims,
        var_name="Dimension", value_name="% of target",
    )
    order = list(results_df["indicator_name"])

    return (
        alt.Chart(long_df)
        .mark_rect()
        .encode(
            x=alt.X("Dimension:N", sort=dims, title=None),
            y=alt.Y("indicator_name:N", sort=order, title=None),
            color=alt.Color(
                "% of target:Q",
                scale=alt.Scale(domain=[0, 100], scheme="redyellowgreen"),
                legend=alt.Legend(title="% of target"),
            ),
            tooltip=[
                alt.Tooltip("indicator_name:N", title="Indicator"),
                alt.Tooltip("Dimension:N", title="Dimension"),
                alt.Tooltip("% of target:Q", title="% of target"),
            ],
        )
        .properties(width="container", height=alt.Step(28))
    )


# ---------------------------------------------------------------------------
# CSV Import — bulk submissions from external M&E exports
# ---------------------------------------------------------------------------

_CSV_IMPORT_COPY = {
    "header": "📥 Import Submissions from CSV",
    "intro": (
        "Upload a CSV export from your M&E system (KoboToolbox, DHIS2, Excel, "
        "etc.), map its columns to diagnostic fields, then preview and confirm "
        "before importing. Imported submissions are scored by the same checks "
        "as manual entries and flagged \"Imported — unverified\" until reviewed."
    ),
    "not_mapped": "— Not mapped —",
    "new_mapping": "— New mapping —",
    "profile_label": "Mapping profile",
    "save_profile_label": "Save this mapping as a profile",
    "save_profile_button": "💾 Save profile",
    "profile_saved": "Mapping profile saved.",
    "preview_header": "Preview",
    "summary_template": "{n} row(s) • {m} with missing required fields • {d} duplicate(s)",
    "skip_incomplete": "Skip rows with missing required fields",
    "confirm_checkbox": "I've reviewed this preview and want to import these rows",
    "import_button": "Import & Score",
    "import_done": (
        "Imported {k} submission(s), scored, and flagged \"Imported — unverified\". "
        "Review the Internal/External Review fields and verify the underlying "
        "evidence before relying on these results in a donor report."
    ),
    "no_rows": "No rows left to import after applying your settings.",
    "parse_error_prefix": "Could not read this CSV",
}

# (target_key, label, required, value_type) — value_type in {"text", "bool", "date"}
_CSV_IMPORT_FIELDS = [
    ("result_statement",         "Result Statement",        True,  "text"),
    ("target_group",              "Target Group",            True,  "text"),
    ("timeframe",                 "Timeframe",               True,  "text"),
    ("geographic_scope",          "Geographic Scope",        True,  "text"),
    ("evidence_type",             "Evidence Type",           True,  "text"),
    ("evidence_description",      "Evidence Description",    True,  "text"),
    ("evidence_date",             "Evidence Date / Recency", False, "date"),
    ("verified_by",               "Evidence Verified By",    False, "text"),
    ("internal_review",           "Internal Review Level",   False, "text"),
    ("external_review",           "External Review Level",   False, "text"),
    ("logframe_indicator",        "Logframe Indicator",      False, "text"),
    ("logframe_target",           "Logframe Target",         False, "text"),
    ("logframe_achievement",      "Logframe Achievement",    False, "text"),
    ("learning_notes",            "Learning Notes",          False, "text"),
    ("limitations_notes",         "Limitations Notes",       False, "text"),
    ("beneficiary_voice",         "Beneficiary Voice",       False, "text"),
    ("project_name",              "Project Name",            False, "text"),
    ("donor",                     "Donor",                    False, "text"),
    ("sector",                    "Sector",                   False, "text"),
    ("qual_sourcing_documented",  "Qual: Sourcing Documented (true/false)", False, "bool"),
    ("qual_triangulated",         "Qual: Triangulated (true/false)",       False, "bool"),
    ("qual_bias_considered",      "Qual: Bias Considered (true/false)",    False, "bool"),
    ("qual_beneficiary_voice_represented", "Qual: Beneficiary Voice Represented (true/false)", False, "bool"),
    ("qual_consent_ethics_addressed",      "Qual: Consent/Ethics Addressed (true/false)",      False, "bool"),
]

_CSV_IMPORT_PROFILES_PATH = "csv_import_profiles.json"


def _load_import_profiles() -> dict:
    """Load saved CSV column-mapping profiles. Returns {} if none exist yet."""
    if not os.path.exists(_CSV_IMPORT_PROFILES_PATH):
        return {}
    try:
        with open(_CSV_IMPORT_PROFILES_PATH, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_import_profiles(profiles: dict) -> None:
    """Persist CSV column-mapping profiles to disk."""
    with open(_CSV_IMPORT_PROFILES_PATH, "w", encoding="utf-8") as f:
        json.dump(profiles, f, indent=2, ensure_ascii=False)


def _parse_import_csv(uploaded_file):
    """Read an uploaded CSV robustly: tries utf-8-sig/latin-1, sniffs the
    delimiter, detects an extra title row before the real header, strips
    whitespace, and drops fully-blank rows.

    Returns (df, None) on success or (None, error_message) on failure.
    """
    import csv
    import io
    import pandas as pd

    raw = uploaded_file.getvalue()
    text = None
    for enc in ("utf-8-sig", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        return None, f"{_CSV_IMPORT_COPY['parse_error_prefix']}: unrecognised text encoding."

    lines = [ln for ln in text.splitlines() if ln.strip() != ""]
    if not lines:
        return None, f"{_CSV_IMPORT_COPY['parse_error_prefix']}: the file appears to be empty."

    sample = "\n".join(lines[:10])
    try:
        delimiter = csv.Sniffer().sniff(sample, delimiters=",;\t").delimiter
    except csv.Error:
        delimiter = ","

    # Header-row detection: find the first line whose field count (>= 2)
    # matches the next line's field count — handles an extra title row.
    header_idx = 0
    for i in range(min(len(lines) - 1, 5)):
        this_count = len(lines[i].split(delimiter))
        next_count = len(lines[i + 1].split(delimiter))
        if this_count >= 2 and this_count == next_count:
            header_idx = i
            break

    try:
        df = pd.read_csv(
            io.StringIO("\n".join(lines)),
            sep=delimiter,
            skiprows=header_idx,
            skip_blank_lines=True,
            dtype=str,
        )
    except Exception as exc:
        return None, f"{_CSV_IMPORT_COPY['parse_error_prefix']}: {exc}"

    df.columns = [str(c).strip() for c in df.columns]
    for col in df.columns:
        df[col] = df[col].apply(lambda v: v.strip() if isinstance(v, str) else v)
    df = df.dropna(how="all")

    if df.empty:
        return None, f"{_CSV_IMPORT_COPY['parse_error_prefix']}: no data rows found."

    return df.reset_index(drop=True), None


def _csv_row_to_submission(row: dict, column_map: dict, profile_name: str) -> dict:
    """Build a full submission dict from one CSV row using column_map
    (target_key -> CSV column name). Unmapped fields stay empty/False —
    never inferred or fabricated. Marks the result as an unverified import."""

    def _val(target_key, default=""):
        col = column_map.get(target_key)
        if not col:
            return default
        v = row.get(col)
        if v is None:
            return default
        v = str(v).strip()
        if not v or v.lower() == "nan":
            return default
        return v

    def _bool(target_key):
        return _val(target_key).strip().lower() in ("true", "yes", "1", "x")

    return {
        "result_statement":   _val("result_statement"),
        "target_group":       _val("target_group"),
        "timeframe":          _val("timeframe"),
        "geographic_scope":   _val("geographic_scope"),
        "additional_context": "",
        "learning_notes":     _val("learning_notes"),
        "limitations_notes":  _val("limitations_notes"),
        "internal_review":    _val("internal_review", "Not reviewed") or "Not reviewed",
        "external_review":    _val("external_review", "No external review") or "No external review",
        "attached_filenames": [],
        "beneficiary_voice":  _val("beneficiary_voice"),
        "logframe_indicator":   _val("logframe_indicator"),
        "logframe_target":      _val("logframe_target"),
        "logframe_achievement": _val("logframe_achievement"),
        "reporting_start": "",
        "reporting_end":   "",
        "provenance_checklist": {
            "sampling_documented":     "Choose an option...",
            "double_counting_checked": "Choose an option...",
            "collection_tool_named":   "Choose an option...",
            "collector_independent":   "Choose an option...",
            "recall_period_ok":        "Choose an option...",
            "auditor_traceable":       "Choose an option...",
        },
        "qualitative_evidence": False,
        "qualitative_rigor_checklist": {
            "sourcing_documented":           _bool("qual_sourcing_documented"),
            "triangulated":                  _bool("qual_triangulated"),
            "bias_considered":               _bool("qual_bias_considered"),
            "beneficiary_voice_represented": _bool("qual_beneficiary_voice_represented"),
            "consent_ethics_addressed":      _bool("qual_consent_ethics_addressed"),
        },
        "attribution_contribution": "",
        "disaggregation_status": "",
        "review_status":     _SUBMISSION_STATUS_OPTIONS[0],
        "reviewer_name":     "",
        "reviewer_role":     "",
        "reviewer_date":     "",
        "reviewer_decision": "",
        "reviewer_notes":    "",
        "donor":          _val("donor"),
        "sector":         _val("sector"),
        "project_name":   _val("project_name"),
        "submission_type": "",
        "evidence": [{
            "type":        _val("evidence_type"),
            "description": _val("evidence_description"),
            "recency":     _val("evidence_date"),
            "verified_by": _val("verified_by"),
        }],
        "verification_status": "Imported — unverified",
        "import_source":    "csv",
        "import_profile":   profile_name,
        "import_timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }


def _validate_import_rows(submissions: list) -> list:
    """Return one issue-record per submission: missing required fields,
    type mismatches, and within-import duplicates. Annotation only — never
    blocks or alters the mapped values."""
    import pandas as pd

    required_targets = [(t, label) for t, label, req, _ in _CSV_IMPORT_FIELDS if req]
    bool_targets = [t for t, _, _, vt in _CSV_IMPORT_FIELDS if vt == "bool"]
    date_targets = [t for t, _, _, vt in _CSV_IMPORT_FIELDS if vt == "date"]

    issues = []
    seen = set()
    for i, sub in enumerate(submissions):
        row_issues = {"row": i + 1, "missing_required": [], "type_issues": [], "duplicate": False}

        ev0 = (sub.get("evidence") or [{}])[0]
        field_values = {
            "result_statement":    sub.get("result_statement", ""),
            "target_group":        sub.get("target_group", ""),
            "timeframe":           sub.get("timeframe", ""),
            "geographic_scope":    sub.get("geographic_scope", ""),
            "evidence_type":       ev0.get("type", ""),
            "evidence_description": ev0.get("description", ""),
            "evidence_date":       ev0.get("recency", ""),
        }
        for target_key, label in required_targets:
            if not str(field_values.get(target_key, "")).strip():
                row_issues["missing_required"].append(label)

        for target_key in bool_targets:
            qual = sub.get("qualitative_rigor_checklist", {})
            raw = ""  # already coerced to bool in _csv_row_to_submission; nothing to flag here
        for target_key in date_targets:
            raw = ev0.get("recency", "")
            if raw and pd.to_datetime(raw, errors="coerce") is pd.NaT:
                row_issues["type_issues"].append(f"{target_key}: '{raw}' is not a recognisable date")

        dup_key = (
            sub.get("result_statement", "").strip().lower(),
            sub.get("logframe_indicator", "").strip().lower(),
        )
        if dup_key in seen:
            row_issues["duplicate"] = True
        seen.add(dup_key)

        issues.append(row_issues)

    return issues


# ---------------------------------------------------------------------------
# Screen 3 — Portfolio / Framework Dashboard
# ---------------------------------------------------------------------------

def _render_portfolio_chat(input_df, evaluations: list, statuses: list) -> None:
    """Portfolio-level chat assistant for the Score My Report tab (council XIX).

    Knows all N results' scores simultaneously — enables cross-result questions
    like "which KPI needs the most work?" and "what's my systemic gap?"
    Reuses the same Haiku + session-state pattern as _render_help_chat().
    """
    from diagnostics import build_portfolio_chat_system_prompt

    hist_key = "smr_chat_messages"
    if hist_key not in st.session_state:
        st.session_state[hist_key] = []
    msgs = st.session_state[hist_key]

    st.caption(
        "Ask about any of your scored results or your portfolio as a whole. "
        "Examples: 'Which result needs the most work?' · 'What's my systemic gap?' · "
        "'Which results can I submit now?' — answers are based solely on the scores shown above."
    )

    for msg in msgs:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    prompt = st.chat_input("Ask about your results…", key="smr_chat_input")
    if not prompt:
        return

    msgs.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    _api_key = (
        st.secrets.get("ANTHROPIC_API_KEY", "")
        if hasattr(st, "secrets") else
        __import__("os").environ.get("ANTHROPIC_API_KEY", "")
    )
    if not _api_key:
        reply = "Portfolio chat is not available — API key not configured."
    else:
        try:
            import anthropic as _anthr
            _client = _anthr.Anthropic(api_key=_api_key)
            system_prompt = build_portfolio_chat_system_prompt(input_df, evaluations, statuses)
            _resp = _client.messages.create(
                model="claude-haiku-4-5-20251001",
                max_tokens=512,
                system=system_prompt,
                messages=[{"role": m["role"], "content": m["content"]} for m in msgs],
            )
            reply = _resp.content[0].text if _resp.content else "No response received."
        except Exception as exc:
            reply = f"Could not reach the assistant right now. ({type(exc).__name__})"

    msgs.append({"role": "assistant", "content": reply})
    with st.chat_message("assistant"):
        st.markdown(reply)


def _render_score_my_report_tab():
    """Score My Report — document-in, scored-Excel-out pipeline (council XVII)."""
    import pandas as pd
    from excel_report import build_scored_excel, build_rescore_excel, STATUS_AUTO_POPULATED, STATUS_NOT_FOUND

    st.markdown("### Audit My Report")
    st.markdown(
        "Upload your donor report. ImpactProof extracts every result, makes a determination for each "
        "(submission-ready / needs work / high risk), and delivers a colour-coded "
        "**Excel decision audit** — one row per result, traceable to named donor standards. "
        "Use it to know exactly where to focus before a DQA, before submission, or before a partner review meeting. "
        "**For MEL consultancies managing multiple clients:** audit each client's report separately "
        "and compare evidence quality across USAID, Mastercard Foundation, GIZ Ghana, and FCDO "
        "in a single session."
    )
    st.caption(
        "Anchored to: USAID ADS 201 · FCDO 2025 · Bond Evidence Principles 2024 · World Bank RF  "
        "— Deterministic assessment: same document always produces the same determination."
    )

    _api_key = (
        st.secrets.get("ANTHROPIC_API_KEY", "")
        if hasattr(st, "secrets") else
        __import__("os").environ.get("ANTHROPIC_API_KEY", "")
    )
    if not _api_key:
        st.error("Audit My Report requires an Anthropic API key. Configure ANTHROPIC_API_KEY in secrets.")
        return

    # Council XXIII — SMR paywall gate (mirrors single-result check counter)
    _smr_email = st.session_state.get("user_email", "")
    if not _smr_email:
        st.warning("📧 **Enter your email to use Audit My Report.** We use it to track your free checks — no password needed.")
        _render_email_gate_inline("_smr")
        # st.stop() inside _render_email_gate_inline halts rendering here until email is set
    _smr_email = st.session_state.get("user_email", "")
    _smr_access    = check_access(_smr_email)  # DB-authoritative
    _smr_is_paid   = _smr_access["is_paid"]
    _smr_allowed   = _smr_access["allowed"]

    if not _smr_is_paid:
        st.caption(f"Free checks remaining: **{_smr_access['checks_remaining']}/{FREE_CHECKS_LIMIT}** — "
                   "each report upload uses 1 check.")

    st.info(
        "🔒 **Data processing notice:** Your document is sent to Anthropic's Claude API "
        "(claude-sonnet-4-6) for result extraction. Up to 60,000 characters of text are "
        "transmitted. ImpactProof does not store your document after your session ends. "
        "Anthropic's privacy policy applies to API processing."
    )
    uploaded_doc = st.file_uploader(
        "Upload your donor report (Word or PDF)",
        type=["pdf", "docx", "txt"],
        key="smr_upload",
    )

    if not uploaded_doc:
        st.info("Upload a Word or PDF progress report to extract and determine all results automatically.")
        return

    org_name = st.text_input(
        "Organisation name (optional — appears in the Excel header)",
        key="smr_org_name",
        placeholder="e.g., Action Aid Ghana",
    )

    if not _smr_allowed:
        st.warning(
            "You've used your 3 free checks. Upgrade to Professional to process more reports."
        )
        _render_paywall(prompt_context="audit_attempt", custom_message=(
            "### Upgrade to keep using Audit My Report\n\n"
            "You've used your 3 free document uploads. Upgrade to Professional to keep going:\n\n"
            "- **Unlimited document uploads** — extract and determine every result automatically\n"
            "- **Colour-coded Excel decision audit** — one row per result, traceable to named donor standards\n"
            "- **Unlimited single-result checks and re-scores** too\n\n"
            f"*GHS {PRICE_PER_CHECK_GHS/100:.0f} per check · or GHS {PRICE_MONTHLY_GHS/100:.0f}/month for unlimited*"
        ))
        return

    run_btn = st.button("Extract & Get Determinations", type="primary", key="smr_run")
    _smr_state = st.session_state.get("smr_results")

    if run_btn:
        # Extract document text
        with st.spinner("Reading document..."):
            try:
                _raw_bytes = uploaded_doc.read()
                _fname_lower = uploaded_doc.name.lower()
                doc_text, _doc_err = _extract_text_from_file(_fname_lower, _raw_bytes)
                if _doc_err:
                    st.error(f"Could not read the document: {_doc_err}")
                    return
            except Exception as exc:
                import logging as _logging
                _logging.error("Could not read uploaded document", exc_info=True)
                st.error("Could not read the document. Try re-saving it and uploading again.")
                return
        if not doc_text or len(doc_text.strip()) < 100:
            st.error("The document appears to be empty or image-based. Export as a text-based DOCX and try again.")
            return

        # Live second reader — threads extraction so the timer can tick in the main thread
        import threading as _threading, time as _time
        _smr_result_box: dict = {"data": None, "exc": None}
        _smr_done = _threading.Event()

        def _run_smr():
            try:
                _smr_result_box["data"] = _score_report_from_document(doc_text, _api_key)
            except Exception as _e:
                import logging as _logging
                _logging.error("Score My Report extraction failed", exc_info=True)
                _smr_result_box["exc"] = True
            finally:
                _smr_done.set()

        _threading.Thread(target=_run_smr, daemon=True).start()
        _smr_ph = st.empty()
        _smr_t0 = _time.time()
        while not _smr_done.wait(timeout=1.0):
            _smr_s = int(_time.time() - _smr_t0)
            _smr_ph.info(
                f"🔍 **Second reader active** — extracting and scoring all results…  "
                f"**{_smr_s}s** elapsed *(15–60+ seconds)*"
            )
        _smr_ph.empty()

        if _smr_result_box["exc"]:
            st.error("Extraction failed unexpectedly. Please try again, or try a shorter document.")
            return
        input_df, evaluations, statuses, err = _smr_result_box["data"]

        if err:
            st.error(f"Extraction failed: {err}")
            return
        if input_df is None or input_df.empty:
            st.warning("No reportable results were found in the document. Try a progress report with a logframe or KPI table.")
            return

        st.session_state["smr_results"] = {
            "input_df": input_df,
            "evaluations": evaluations,
            "statuses": statuses,
            "doc_name": uploaded_doc.name,
            "org_name": org_name,
        }
        _smr_state = st.session_state["smr_results"]

        # Council XXIII — increment shared free-check counter after a successful run
        if _smr_email:
            record_check(_smr_email)

    if not _smr_state:
        return

    input_df   = _smr_state["input_df"]
    evaluations = _smr_state["evaluations"]
    statuses   = _smr_state["statuses"]
    doc_name   = _smr_state.get("doc_name", "")
    org_name   = _smr_state.get("org_name", "") or org_name

    n = len(evaluations)
    strong  = sum(1 for ev in evaluations if ev.get("confidence_score", 0) >= 4.0 and ev.get("clarity_score", 0) >= 4.0)
    medium  = sum(1 for ev in evaluations if 2.5 <= ev.get("confidence_score", 0) < 4.0 or 2.5 <= ev.get("clarity_score", 0) < 4.0)
    weak    = n - strong - medium

    st.success(f"Extracted and scored **{n} result{'s' if n != 1 else ''}** from {doc_name}.")
    c1, c2, c3 = st.columns(3)
    c1.metric("Submission-ready", strong, help="Both axes ≥ 4.0 (INGO standard — portfolio uses the bilateral-donor threshold)")
    c2.metric("Needs improvement", medium, help="At least one axis 2.5–3.9")
    c3.metric("High risk", weak, help="One or both axes < 2.5")

    # Council XXIV — Portfolio narrative synthesis (deterministic, no API call)
    if n > 0:
        # Find weakest sub-score dimension across all results
        _dim_sums: dict[str, float] = {}
        _dim_maxs: dict[str, float] = {
            "Directness": 2.0, "Verification": 2.0, "Recency": 1.0,
            "Definition": 1.25, "Measurement": 1.25, "Integrity": 1.0,
            "Scope": 0.75, "Governance": 0.75,
        }
        for _ev in evaluations:
            _cc = _ev.get("confidence_components", {})
            _cl = _ev.get("clarity_components", {})
            for _dim, _key in [
                ("Directness",   "direct_score"),  ("Verification", "verify_score"),
                ("Recency",      "recency_score"),  ("Definition",   "definition_score"),
                ("Measurement",  "measurement_score"), ("Integrity", "integrity_score"),
                ("Scope",        "scope_score"),    ("Governance",   "governance_score"),
            ]:
                _src = _cc if _dim in ("Directness", "Verification", "Recency") else _cl
                _dim_sums[_dim] = _dim_sums.get(_dim, 0) + _src.get(_key, 0)

        _weakest_dim = min(_dim_sums, key=lambda d: _dim_sums[d] / (_dim_maxs[d] * n))
        _weakest_pct = round((_dim_sums[_weakest_dim] / (_dim_maxs[_weakest_dim] * n)) * 100)
        _strongest_dim = max(_dim_sums, key=lambda d: _dim_sums[d] / (_dim_maxs[d] * n))
        _strongest_pct = round((_dim_sums[_strongest_dim] / (_dim_maxs[_strongest_dim] * n)) * 100)

        _conf_scores = [_ev.get("confidence_score", 0) for _ev in evaluations]
        _avg_conf = round(sum(_conf_scores) / n, 1)
        _recency_weak = sum(1 for _ev in evaluations if _ev.get("confidence_components", {}).get("recency_score", 1) <= 0.4)

        _narr_lines = [f"**Portfolio Snapshot — {n} result{'s' if n != 1 else ''} from {doc_name}**"]
        _narr_lines.append(f"• {strong} submission-ready · {medium} need work · {weak} high risk · Average Confidence: {_avg_conf}/5.0")
        if _weakest_pct < 60:
            _narr_lines.append(f"• **Gap: {_weakest_dim}** is your weakest sub-score ({_weakest_pct}% of target) — address this first for the biggest improvement across results")
        if _strongest_pct >= 75:
            _narr_lines.append(f"• **Strength: {_strongest_dim}** scored well ({_strongest_pct}% of target) — this is a consistent strong point in your portfolio")
        if _recency_weak > 0:
            _narr_lines.append(f"• {_recency_weak} result{'s' if _recency_weak > 1 else ''} ha{'ve' if _recency_weak > 1 else 's'} evidence date issues — check dates against reporting period end")

        st.info("\n".join(_narr_lines))
        if _weakest_pct < 80:
            st.caption(
                f"**System decision:** Fix {_weakest_dim} first — it is your portfolio's highest-leverage action "
                f"({_weakest_pct}% of target across all results). Improving it will lift your average Confidence score "
                f"more than any other single action."
            )
        st.caption(
            "Each score is traceable to named sub-criteria (USAID ADS 201.3.5.7 for Validity, Integrity, Precision, "
            "Reliability, Timeliness). Re-running this check with the same document produces identical scores — "
            "unlike AI chatbot feedback."
        )

    st.caption(
        "⚠️ **Review required:** Auto-populated fields are shown in amber in the Excel. "
        "Confirm extracted values match your documentation before treating scores as final."
    )

    # Score interpretation guide
    with st.expander("📖 How to read your determination", expanded=False):
        st.markdown(
            "Each result receives a determination, not just a score. The determination tells you whether "
            "a donor is likely to accept, query, or reject the result based on evidence quality alone.\n\n"
            "**Confidence (0–5)** — How much should we trust the evidence?\n\n"
            "- ≥ 4.0 🟢 Submission-ready &nbsp;·&nbsp; 2.5–3.9 🟡 Needs work &nbsp;·&nbsp; < 2.5 🔴 High risk\n\n"
            "**Clarity (0–5)** — Can someone else interpret this result the same way?\n\n"
            "- Same thresholds apply.\n\n"
            "**Auto-populated fields** — Extracted by AI from your document. May contain errors. "
            "Open the Excel and review every amber cell before sharing.\n\n"
            "**Not found fields** — Missing from your document. Fill them manually in the Excel, "
            "or upload a more complete version and re-score.\n\n"
            "**What to do next:** Download the Excel → review amber cells → "
            "address the top-3 results with lowest scores → re-upload the CSV to re-score if needed."
        )

    # Preview table
    with st.expander("Preview extracted results", expanded=True):
        preview_rows = []
        for i, (_, row) in enumerate(input_df.iterrows()):
            ev = evaluations[i]
            stat = statuses[i]
            auto_count = sum(1 for v in stat.values() if v == STATUS_AUTO_POPULATED)
            not_found  = sum(1 for v in stat.values() if v == STATUS_NOT_FOUND)
            preview_rows.append({
                "Indicator": row.get("indicator_name", "") or row.get("result_statement", "")[:60],
                "Confidence": ev.get("confidence_score", 0),
                "Clarity": ev.get("clarity_score", 0),
                "Verdict": ev.get("verdict", ""),
                "Auto-populated fields": auto_count,
                "Not found fields": not_found,
            })
        st.dataframe(pd.DataFrame(preview_rows), use_container_width=True, hide_index=True)

    # Downloads
    st.divider()
    _ts = __import__("datetime").datetime.now().strftime("%Y%m%d_%H%M")
    dl_c1, dl_c2 = st.columns(2)

    with dl_c1:
        try:
            rows_list = [row.to_dict() for _, row in input_df.iterrows()]
            excel_bytes = build_scored_excel(
                rows_list, evaluations, statuses,
                org_name=org_name, document_name=doc_name,
            )
            st.download_button(
                f"📊 Download Portfolio Readiness Workbook — {n} result{'s' if n != 1 else ''} · {_ts}",
                data=excel_bytes,
                file_name=f"portfolio_readiness_{_ts}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="smr_excel_dl",
                type="primary",
                use_container_width=True,
                help="Anchored to USAID ADS 201 + Bond 2024. Sheet 1: scored results. Sheet 2: portfolio gap summary. Share with MEL lead or donor as a pre-submission evidence quality check.",
            )
            st.caption("Share with your MEL lead or donor as a pre-submission evidence quality check — each score cites the standard it was measured against.")
        except Exception as exc:
            st.error(f"Could not generate Excel: {exc}")

    with dl_c2:
        try:
            _port_rows = []
            for _, row in input_df.iterrows():
                r = {k: row.get(k, "") for k in [c[0] for c in _PORTFOLIO_COLUMNS]}
                _port_rows.append(r)
            _rescore_xls = build_rescore_excel(
                _port_rows, evaluations, statuses,
                org_name=org_name, doc_name=doc_name,
            )
            st.download_button(
                "📥 Download re-score workbook",
                data=_rescore_xls,
                file_name=f"portfolio_rescore_{_ts}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="smr_csv_dl",
                use_container_width=True,
                help=(
                    "Two-sheet Excel: 'Re-score Data' (edit & save as CSV to re-upload) + "
                    "'Determinations' (colour-coded results view). "
                    "Fix amber/red fields in Re-score Data sheet, save as CSV, "
                    "then re-upload to the CSV Portfolio tab."
                ),
            )
            st.caption(
                "✏️ **To re-score after fixes:** Open 'Re-score Data' sheet → correct fields → "
                "File > Save As > CSV → re-upload to **CSV Portfolio tab**."
            )
        except Exception as exc:
            st.error(f"Could not generate re-score workbook: {exc}")

    # Portfolio chat — payment-gated, context-aware of all N results (council XIX)
    st.divider()
    _portfolio_chat_allowed = check_access(_smr_email)["allowed"]
    with st.expander("💬 Ask about your results", expanded=False):
        if not _portfolio_chat_allowed:
            metrics.log_event("upgrade_prompt_shown", _metrics_session_id(), context="portfolio_chat_attempt")
            st.info(
                "Portfolio Q&A is available on the **Professional plan**. "
                "Ask the system direct decision questions: 'Which KPI needs the most work?', "
                "'What is my systemic gap?', 'Which results are at risk of a donor query?' "
                "— the system routes you to the highest-leverage actions across your entire portfolio. "
                "GHS 50/month vs. GHS 12,000–17,000 in rework costs from a donor-queried report."
            )
            if st.button("Upgrade to Professional →", key="smr_chat_upgrade", type="primary"):
                metrics.log_event("upgrade_prompt_clicked", _metrics_session_id(), context="portfolio_chat_attempt")
                st.session_state["_show_pricing"] = True
                st.rerun()
        else:
            _render_portfolio_chat(input_df, evaluations, statuses)

    if st.button("Clear results and upload a different report", key="smr_clear"):
        st.session_state.pop("smr_results", None)
        st.session_state.pop("smr_chat_messages", None)
        st.rerun()


def render_screen_3():
    import pandas as pd

    if st.button("← Back to Home", key="portfolio_back"):
        _go_to_screen(0)

    st.markdown("## 📊 Portfolio Decision Review")

    _s3_tab_smr, _s3_tab_csv = st.tabs(["📄 Audit My Report", "📊 CSV Portfolio"])

    with _s3_tab_smr:
        _render_score_my_report_tab()

    with _s3_tab_csv:
        st.info(
            "**Two ways to use this tab:**\n\n"
            "**1. Re-assess after fixes** — After running Audit My Report, download the "
            "**re-score CSV**, fix amber/red fields in a spreadsheet editor, and re-upload here "
            "to see revised determinations without re-uploading the original document.\n\n"
            "**2. Assess your full logframe** — Upload your logframe as a CSV or Excel file "
            "to determine all indicators at once and see the portfolio heatmap."
        )

        _tmpl_c1, _tmpl_c2 = st.columns(2)
        with _tmpl_c1:
            st.download_button(
                "📥 Minimal template (7 required columns)",
                data=_portfolio_minimal_template_csv(),
                file_name="portfolio_template_minimal.csv",
                mime="text/csv",
                key="portfolio_minimal_template_dl",
                help="Start here — fill the 7 required columns, upload, and see your heatmap.",
            )
        with _tmpl_c2:
            st.download_button(
                "📥 Full template (29 columns)",
                data=_portfolio_template_csv(),
                file_name="impact_receipts_portfolio_template.csv",
                mime="text/csv",
                key="portfolio_template_dl",
                help="Includes provenance, beneficiary voice, and governance fields for higher scores.",
            )
        with st.expander("Column reference & accepted values"):
            st.caption(_PORTFOLIO_REVIEW_HINT)

        _csvpf_email = st.session_state.get("user_email", "")
        if not _csvpf_email:
            st.warning("📧 **Enter your email to use CSV Portfolio.** We use it to track your free checks — no password needed.")
            _render_email_gate_inline("_csvpf")
        _csvpf_email  = st.session_state.get("user_email", "")
        _csvpf_access  = check_access(_csvpf_email)
        _csvpf_paid    = _csvpf_access["is_paid"]
        _csvpf_allowed = _csvpf_access["allowed"]

        if not _csvpf_paid:
            st.caption(f"Free checks remaining: **{_csvpf_access['checks_remaining']}/{FREE_CHECKS_LIMIT}** — "
                       "each upload uses 1 check.")

        uploaded = st.file_uploader(
            "Upload your completed logframe (CSV or Excel)",
            type=["csv", "xlsx", "xls"],
            key="portfolio_upload",
        )

        if uploaded is not None and not _csvpf_allowed:
            st.warning("You've used your 3 free checks. Upgrade to Professional to keep using CSV Portfolio.")
            _render_paywall(prompt_context="csv_portfolio_attempt", custom_message=(
                "### Upgrade to keep using CSV Portfolio\n\n"
                "You've used your 3 free uploads. Upgrade to Professional to keep going:\n\n"
                "- **Unlimited portfolio uploads** — assess your full logframe or re-score after fixes\n"
                "- **Portfolio heatmap** — see systemic gaps across every indicator\n"
                "- **Unlimited single-result checks and Audit My Report uploads** too\n\n"
                f"*GHS {PRICE_PER_CHECK_GHS/100:.0f} per check · or GHS {PRICE_MONTHLY_GHS/100:.0f}/month for unlimited*"
            ))
        elif uploaded is not None and not _safe_rate_limit_ok(
            _csvpf_email, "portfolio_upload", max_count=15, window_seconds=3600
        ):
            st.warning("You've uploaded a lot of portfolios in the last hour — please wait a bit before uploading more.")
        elif uploaded is not None:
            _safe_log_access(_csvpf_email, "portfolio_upload")
            try:
                if uploaded.name.lower().endswith((".xlsx", ".xls")):
                    df = pd.read_excel(uploaded)
                else:
                    df = pd.read_csv(uploaded)
            except Exception as exc:
                import logging as _logging
                _logging.error("Could not read uploaded CSV/Excel file", exc_info=True)
                st.error("Could not read the file. Confirm it's a valid CSV or Excel file and try again.")
                df = None

            if df is not None:
                results_df, warnings = _evaluate_portfolio(df)
                st.session_state["portfolio_results"] = results_df
                st.session_state["portfolio_warnings"] = warnings
                if _csvpf_email:
                    record_check(_csvpf_email)

        results_df = st.session_state.get("portfolio_results")
        warnings = st.session_state.get("portfolio_warnings") or []

        for w in warnings:
            st.warning(w)

        if results_df is not None and not results_df.empty:
            dims = [d for d, _ in _PORTFOLIO_SUBSCORE_DIMENSIONS]
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Indicators evaluated", len(results_df))
            with col2:
                st.metric("Avg Confidence", f"{results_df['confidence_score'].mean():.1f}/5.0")
            with col3:
                st.metric("Avg Clarity", f"{results_df['clarity_score'].mean():.1f}/5.0")

            weakest_dim = results_df[dims].mean().idxmin()
            weakest_dim_pct = results_df[dims].mean().min()
            st.markdown(
                f"**System determination — highest-leverage action:** *{weakest_dim}* is your portfolio's "
                f"most critical gap ({weakest_dim_pct:.0f}% of target on average). Prioritise this "
                f"sub-criterion across all indicators for the greatest combined score uplift."
            )

            st.markdown("#### Heatmap — weakest indicators & sub-scores")
            st.caption(
                "For rows using a qualitative evidence type (case study, outcome harvesting, "
                "beneficiary narrative or testimony), the \"Measurement\" column reflects "
                "Sourcing & Triangulation (case/respondent selection, triangulation, and bias "
                "mitigation) instead of measurement precision."
            )
            if st.session_state.get("lite_mode", False):
                st.caption(
                    "Heatmap hidden in low-bandwidth mode — see the results table "
                    "below for the same data."
                )
            else:
                st.altair_chart(_portfolio_heatmap_chart(results_df), use_container_width=True, key="portfolio_heatmap")

            st.markdown("#### Results table")
            st.dataframe(
                results_df[["indicator_name", "confidence_score", "clarity_score", "verdict", "top_fix"]],
                use_container_width=True,
            )

            st.download_button(
                "Download results (CSV)",
                data=results_df.to_csv(index=False).encode("utf-8"),
                file_name="impact_receipts_portfolio_results.csv",
                mime="text/csv",
                key="portfolio_results_dl",
            )

            with st.expander("Report handoff details"):
                st.text_input("Prepared by (your name)", key="report_prepared_by")
                st.selectbox("Status", _REPORT_STATUS_OPTIONS, key="report_status")
                st.text_area("Notes for reviewer (optional)", key="report_notes", height=80)

            _timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            _portfolio_summary_html = _build_portfolio_verification_summary_html(results_df, warnings, _timestamp)
            _portfolio_summary_pdf  = _html_to_pdf_bytes(_portfolio_summary_html)
            if _portfolio_summary_pdf:
                st.download_button(
                    "📋 Download Verification Summary (PDF)",
                    data=_portfolio_summary_pdf,
                    file_name=f"portfolio_verification_summary_{_timestamp}.pdf",
                    mime="application/pdf",
                    key="portfolio_verification_summary_pdf_btn",
                )
            else:
                st.caption("PDF: install xhtml2pdf to enable one-click PDF download.")

    st.divider()
    st.markdown(f"### {_TREND_COPY['header']}")
    st.caption(_TREND_COPY["intro"])
    render_trends_view(_load_trend_history())


def _build_council_page_html(council_assessment: dict, conf_score: float,
                              clar_score: float, timestamp: str, P: str) -> str:
    """Build the council summary Page 4 HTML block (xhtml2pdf-safe, plain English).

    Appended to the Readiness Card when a council assessment has been run.
    """
    if not council_assessment:
        return ""

    from council import COUNCIL_MEMBERS

    brief       = council_assessment.get("reporting_team_brief", {})
    proj_conf   = council_assessment.get("projected_conf", conf_score)
    proj_clar   = council_assessment.get("projected_clar", clar_score)
    upg_rs      = council_assessment.get("upgraded_result_statement", "")
    upg_ev      = council_assessment.get("upgraded_evidence_statement", "")
    verdicts    = council_assessment.get("verdicts", {})

    # Unicode block bars (10-unit scale, PDF-safe)
    def _bar(score: float, max_val: float = 5.0) -> str:
        filled = round(min(score / max_val, 1.0) * 10)
        return "&#9632;" * filled + "&#9633;" * (10 - filled)

    def _score_color(s: float) -> str:
        if s >= 4.0:
            return "#1B5E20"
        if s >= 3.0:
            return "#8A6500"
        return "#B71C1C"

    # Score bars block
    conf_bar_html = (
        f"<tr><td style='font-size:11px;padding:3px 0;width:100px;'>Confidence</td>"
        f"<td style='font-size:10px;letter-spacing:1px;color:{_score_color(conf_score)};'>"
        f"{_bar(conf_score)}</td>"
        f"<td style='font-size:11px;padding-left:8px;color:{_score_color(conf_score)};font-weight:700;'>"
        f"{conf_score}/5.0</td>"
        f"<td style='font-size:10px;letter-spacing:1px;color:{_score_color(proj_conf)};padding-left:16px;'>"
        f"{_bar(proj_conf)}</td>"
        f"<td style='font-size:11px;padding-left:8px;color:{_score_color(proj_conf)};font-weight:700;'>"
        f"{proj_conf}/5.0 (+{round(proj_conf - conf_score, 2)})</td></tr>"
    )
    clar_bar_html = (
        f"<tr><td style='font-size:11px;padding:3px 0;'>Clarity</td>"
        f"<td style='font-size:10px;letter-spacing:1px;color:{_score_color(clar_score)};'>"
        f"{_bar(clar_score)}</td>"
        f"<td style='font-size:11px;padding-left:8px;color:{_score_color(clar_score)};font-weight:700;'>"
        f"{clar_score}/5.0</td>"
        f"<td style='font-size:10px;letter-spacing:1px;color:{_score_color(proj_clar)};padding-left:16px;'>"
        f"{_bar(proj_clar)}</td>"
        f"<td style='font-size:11px;padding-left:8px;color:{_score_color(proj_clar)};font-weight:700;'>"
        f"{proj_clar}/5.0 (+{round(proj_clar - clar_score, 2)})</td></tr>"
    )

    # Plain-English brief section
    wm      = brief.get("what_score_means", "")
    changes = brief.get("what_to_change", [])
    hl      = brief.get("how_long", "")
    ps      = brief.get("projected_status", "")

    changes_html = "".join(
        f"<li style='font-size:11px;margin-bottom:3px;'>{c}</li>" for c in changes
    ) if changes else ""

    # Verdict cards (Evidence Auditor + Programme Strategist shown; others summarised)
    def _member_cell(mid: str) -> str:
        v = verdicts.get(mid, {})
        color = v.get("color", "#333333")
        txt   = (v.get("verdict_text") or "")[:280].replace("<", "&lt;").replace(">", "&gt;")
        return (
            f"<td style='vertical-align:top;padding:6px 8px;border:1px solid #E0E0E0;width:50%;'>"
            f"<p style='font-size:10px;font-weight:700;color:{color};margin:0 0 3px;'>"
            f"{v.get('icon','')} {v.get('name','')}</p>"
            f"<p style='font-size:9px;color:#757575;margin:0 0 4px;'>{v.get('archetype','')}</p>"
            f"<p style='font-size:10px;line-height:1.5;margin:0;'>{txt}</p>"
            f"</td>"
        )

    council_table = "".join(
        f"<tr>{_member_cell(m['id'])}{_member_cell(COUNCIL_MEMBERS[i+1]['id'])}</tr>"
        for i, m in enumerate(COUNCIL_MEMBERS[:4:2])
    )

    donor_v  = verdicts.get("donor_rep", {})
    donor_txt = (donor_v.get("verdict_text") or "")[:400].replace("<", "&lt;").replace(">", "&gt;")

    upgraded_block = ""
    if upg_rs:
        upgraded_block += (
            f"<p style='font-size:11px;font-weight:700;margin:8px 0 2px;color:#1565C0;'>Upgraded result statement:</p>"
            f"<p style='font-size:11px;color:#1a1a1a;background:#F3F8FE;padding:6px 10px;"
            f"border-left:3px solid #1565C0;margin:0 0 6px;{P}'>{upg_rs}</p>"
        )
    if upg_ev:
        upgraded_block += (
            f"<p style='font-size:11px;font-weight:700;margin:6px 0 2px;color:#E65100;'>Upgraded evidence statement:</p>"
            f"<p style='font-size:11px;color:#1a1a1a;background:#FFF3E0;padding:6px 10px;"
            f"border-left:3px solid #E65100;margin:0;{P}'>{upg_ev}</p>"
        )

    return f"""
<!-- ═══ PAGE 4: COUNCIL ASSESSMENT ═══ -->
<p style="page-break-before:always;"></p>

<h2 style="color:#1B5E20;border-bottom:2px solid #1B5E20;padding-bottom:4px;margin-bottom:12px;">
Council Assessment Summary</h2>

<p style="font-size:11px;color:#424242;margin-bottom:12px;">
Five council reviewers assessed this result from distinct lenses. Their assessments
and the upgrade recommendations below are AI-generated guidance — review all content
before use. Score projections are calculated deterministically from the priority fixes above.
</p>

<h3 style="color:#1B5E20;font-size:12px;margin:0 0 6px;">Score Uplift Projection</h3>
<table style="width:100%;border-collapse:collapse;margin-bottom:14px;">
<tr>
  <td style="font-size:10px;color:#616161;padding:2px 0;width:100px;"></td>
  <td colspan="2" style="font-size:10px;color:#616161;font-weight:700;padding:2px 0;">Current</td>
  <td colspan="2" style="font-size:10px;color:#1B5E20;font-weight:700;padding:2px 0 2px 16px;">Projected (after all fixes)</td>
</tr>
{conf_bar_html}
{clar_bar_html}
</table>

<h3 style="color:#1B5E20;font-size:12px;margin:0 0 8px;">For Your Reporting Team</h3>
<div style="background:#F9FBE7;border:1px solid #C5E1A5;padding:10px 12px;border-radius:4px;margin-bottom:12px;{P}">
{"<p style='font-size:11px;margin:0 0 6px;'>" + wm + "</p>" if wm else ""}
{"<p style='font-size:11px;font-weight:700;margin:6px 0 2px;'>What needs to change:</p><ul style='margin:0 0 6px;padding-left:16px;'>" + changes_html + "</ul>" if changes_html else ""}
{"<p style='font-size:11px;margin:0 0 4px;'><strong>How long:</strong> " + hl + "</p>" if hl else ""}
{"<p style='font-size:11px;margin:0;color:#1B5E20;font-weight:700;'>After fixes: " + ps + "</p>" if ps else ""}
</div>

<h3 style="color:#1B5E20;font-size:12px;margin:0 0 8px;">Council Verdicts</h3>
<table style="width:100%;border-collapse:collapse;margin-bottom:8px;">
{council_table}
</table>

<div style="border:1px solid #B71C1C;padding:8px 10px;border-radius:4px;margin-bottom:12px;{P}">
<p style="font-size:10px;font-weight:700;color:#B71C1C;margin:0 0 3px;">
{donor_v.get('icon','')} {donor_v.get('name','')} — {donor_v.get('archetype','')}</p>
<p style="font-size:10px;line-height:1.5;margin:0;">{donor_txt}</p>
</div>

{"<h3 style='color:#1B5E20;font-size:12px;margin:0 0 6px;'>Upgraded Statements (council draft — review before use)</h3>" + upgraded_block if upgraded_block else ""}

<p style="color:#616161;font-style:italic;font-size:9px;border-top:1px solid #E0E0E0;margin-top:12px;padding-top:6px;">
Council XXII assessment &middot; Generated: {timestamp} &middot;
Council output is AI-generated. Verify all content before submitting to a donor.
</p>"""


def _build_html_report_card(submission: dict, evaluation: dict, timestamp: str,
                             field_sources: dict | None = None,
                             council_assessment: dict | None = None) -> str:
    """Lean 2–3 page Submission Readiness Card.
    Flat table layout only — no nested percentage-width tables (xhtml2pdf constraint).
    If council_assessment is provided, a Page 4 council summary is appended."""
    P = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"

    conf_score = round(evaluation.get("confidence_score", 0), 1)
    clar_score = round(evaluation.get("clarity_score", 0), 2)
    conf_label = evaluation.get("confidence_label", "")
    clar_label = evaluation.get("clarity_label", "")
    verdict    = evaluation.get("verdict", "")
    fixes      = evaluation.get("fixes", [])
    conf_comp  = evaluation.get("confidence_components", {})
    clar_comp  = evaluation.get("clarity_components", {})
    diag_state = evaluation.get("diagnostic_state", "")
    ev_stmt    = _generate_evidence_statement(submission) if callable(globals().get("_generate_evidence_statement")) else ""

    # "How you compare" — same anonymized benchmark shown on-screen in
    # _render_result_card(), included here so it appears on the exported
    # Readiness Card PDF too, not just the live view.
    _bm_donor, _bm_sector, _bm_org_type = (
        submission.get("donor", ""), submission.get("sector", ""), submission.get("org_type", ""))
    _benchmark_html = ""
    if _bm_donor and _bm_sector and _bm_org_type:
        _bm = get_benchmark(_bm_donor, _bm_sector, _bm_org_type, conf_score, clar_score)
        if _bm:
            _benchmark_html = (
                f"<p style='font-size:11px;color:#424242;margin:0 0 12px;{P}'>"
                f"📊 How you compare: your Confidence is higher than {_bm['confidence_percentile']}% and your "
                f"Clarity is higher than {_bm['clarity_percentile']}% of {_bm['sample_size']} saved {_bm_donor} "
                f"audits in {_bm_sector} at your organisation type.</p>"
            )

    rs       = submission.get("result_statement", "—")
    li       = submission.get("logframe_indicator", "")
    lt       = submission.get("logframe_target", "")
    la       = submission.get("logframe_achievement", "")
    ev       = (submission.get("evidence") or [{}])[0]
    ev_type  = ev.get("type", "")
    ev_desc  = ev.get("description", "")
    verifier = ev.get("verified_by", "")
    ev_date  = str(ev.get("recency", "") or "")

    # Score colours (used for boxes, bars, text)
    _score_palette = {
        "Strong":     ("#C8E6C9", "#1B5E20"),
        "Acceptable": ("#FFF9C4", "#F57F17"),
        "Weak":       ("#FFE0B2", "#E65100"),
        "High Risk":  ("#FFCDD2", "#B71C1C"),
    }
    cbg, cfg = _score_palette.get(conf_label, ("#F5F5F5","#212121"))
    lbg, lfg = _score_palette.get(clar_label, ("#F5F5F5","#212121"))

    # Verdict colour
    _verdict_map = {
        "Strong KPI":          ("#C8E6C9","#1B5E20"),
        "Misleading KPI":      ("#FFE0B2","#E65100"),
        "Well-defined but":    ("#FFF9C4","#F57F17"),
        "High risk":           ("#FFCDD2","#B71C1C"),
        "STRONG":              ("#C8E6C9","#1B5E20"),
        "MISLEADING":          ("#FFE0B2","#E65100"),
        "NEEDS REFINEMENT":    ("#FFF9C4","#F57F17"),
        "FUNDAMENTALLY WEAK":  ("#FFCDD2","#B71C1C"),
        "UNDEREVIDENCED":      ("#FFE0B2","#E65100"),
    }
    vbg, vfg = ("#F5F5F5","#424242")
    for key, (bg, fg) in _verdict_map.items():
        if key.lower() in (verdict or diag_state).lower():
            vbg, vfg = bg, fg; break

    def bar_row(val, max_v, label):
        """Score row with label + coloured score value. No bar height tricks — fully xhtml2pdf-safe."""
        pct = min(int((val / max_v) * 100), 100) if max_v else 0
        bar_color = "#1B5E20" if pct >= 70 else ("#F57F17" if pct >= 50 else "#B71C1C")
        score_str = f"{val} / {max_v}"
        # Unicode block bar: ■ repeated proportionally (renders as text, never crashes)
        filled = max(0, round(pct / 10))
        empty  = 10 - filled
        bar_txt = f"<font color='{bar_color}'>{'&#9632;' * filled}</font><font color='#E0E0E0'>{'&#9632;' * empty}</font>"
        return (
            f"<tr>"
            f"<td width='130' style='font-size:11px;color:#424242;padding:3px 0;'>{label}</td>"
            f"<td width='100' style='font-size:10px;padding:3px 4px;letter-spacing:1px;'>{bar_txt}</td>"
            f"<td width='55' style='font-size:11px;font-weight:700;color:{bar_color};padding:3px 0;{P}'>{score_str}</td>"
            f"</tr>"
        )

    _fs = field_sources or {}

    def _card_page_note(fix):
        if not _fs:
            return ""
        dim = fix.get("dimension", "")
        fields = _FIX_FIELD_SOURCE_MAP.get(dim, [])
        page = next((
            _fs[f].get("page", 0)
            for f in fields
            if f in _fs and isinstance(_fs[f], dict) and _fs[f].get("page", 0) > 0
        ), 0)
        return f" <i style='color:#8A6500;font-size:10px;'>(p.{page})</i>" if page else ""

    top_fixes = fixes[:3]
    fixes_rows = "".join(
        f"<tr><td style='vertical-align:top;padding:4px 8px 4px 0;font-size:12px;color:#1B5E20;font-weight:700;'>{i+1}.</td>"
        f"<td style='padding:4px 0;font-size:12px;color:#212121;'>{f.get('message','')} "
        f"<span style='color:#8A6500;font-size:11px;'>({f.get('score_impact','')})</span>"
        f"{_card_page_note(f)}</td></tr>"
        for i, f in enumerate(top_fixes)
    ) if top_fixes else f"<tr><td colspan='2' style='font-size:12px;color:#1B5E20;padding:4px 0;'>&#10003; No critical fixes — result is ready to submit.</td></tr>"

    ev_stmt_block = ""
    if ev_stmt:
        ev_stmt_block = (
            f"<table border='0' cellspacing='0' cellpadding='0' width='100%' style='margin:12px 0;{P}'><tr>"
            f"<td width='4' style='background:#1B5E20;{P}'></td>"
            f"<td style='background:#F1F8E9;padding:10px 14px;font-size:11px;color:#212121;{P}'>"
            f"<strong>Evidence statement (ready to paste into your report):</strong><br/>{ev_stmt}"
            f"</td></tr></table>"
        )

    is_qual    = clar_comp.get("is_qualitative", False)
    def_label  = "Narrative Definition" if is_qual else "Definition"
    meas_label = "Sourcing &amp; Triangulation" if is_qual else "Measurement"

    user_email = st.session_state.get("user_email", "")
    donor = st.session_state.get("donor_selected", "")
    donor_note = f" &middot; Prepared for {donor}" if donor and donor not in ("(No donor specified)","") else ""

    lf_rows = ""
    if li:
        lf_rows += f"<tr><td style='color:#616161;font-size:11px;padding:3px 10px 3px 0;width:90px;'>Indicator</td><td style='font-size:11px;'>{li}</td></tr>"
    if lt:
        lf_rows += f"<tr><td style='color:#616161;font-size:11px;padding:3px 10px 3px 0;'>Target</td><td style='font-size:11px;'>{lt}</td></tr>"
    if la:
        lf_rows += f"<tr><td style='color:#616161;font-size:11px;padding:3px 10px 3px 0;'>Achievement</td><td style='font-size:11px;color:#1B5E20;font-weight:700;'>{la}</td></tr>"

    _irc_note_html = (
        f'<p style="font-size:11px;color:#6D4C41;background:#FFF3E0;padding:8px 12px;'
        f'border-left:3px solid #FF6F00;margin-bottom:10px;{P}">'
        f'<strong>Instant Report Check used:</strong> Some form fields were pre-filled '
        f'using AI extraction from an uploaded document. Review all values to confirm '
        f'they accurately reflect your documentation before treating this report as final.</p>'
    ) if st.session_state.get("_irc_used") else ""

    # ── Change A: Verdict rationale (council XXI) ─────────────────────────────
    _conf_subs = [
        ("Directness",   round(conf_comp.get("direct_score", 0), 1),  2.0),
        ("Verification", round(conf_comp.get("verify_score", 0), 1),  2.0),
        ("Recency",      round(conf_comp.get("recency_score", 0), 1), 1.0),
    ]
    _sorted_conf = sorted(_conf_subs, key=lambda x: x[1] / x[2])
    _weakest_c   = _sorted_conf[0]
    _verdict_rationale_html = ""
    if conf_score < 3.0:
        _verdict_rationale_html = (
            f'<p style="font-size:11px;color:#616161;font-style:italic;margin:0 0 12px;">'
            f'Your Confidence of {conf_score}/5.0 is driven primarily by '
            f'<strong>{_weakest_c[0]}</strong> at {_weakest_c[1]}/{_weakest_c[2]:.1f}. '
            f'Addressing the priority fixes below is the fastest path to re-scoring.'
            f'</p>'
        )
    elif clar_score < 3.0:
        _clar_subs = [
            ("Definition",  round(clar_comp.get("definition_score", 0), 2),  1.25),
            ("Measurement", round(clar_comp.get("measurement_score", 0), 2), 1.25),
            ("Integrity",   round(clar_comp.get("integrity_score", 0), 2),   1.0),
            ("Scope",       round(clar_comp.get("scope_score", 0), 2),       0.75),
            ("Governance",  round(clar_comp.get("governance_score", 0), 2),  0.75),
        ]
        _weakest_l = min(_clar_subs, key=lambda x: x[1] / x[2])
        _verdict_rationale_html = (
            f'<p style="font-size:11px;color:#616161;font-style:italic;margin:0 0 12px;">'
            f'Your Clarity of {clar_score}/5.0 is driven primarily by '
            f'<strong>{_weakest_l[0]}</strong> at {_weakest_l[1]}/{_weakest_l[2]:.1f}. '
            f'Addressing the priority fixes below is the fastest path to re-scoring.'
            f'</p>'
        )

    # ── Change B: Recency anomaly warning (council XXI) ───────────────────────
    _recency_level = conf_comp.get("recency_level", 5)
    _recency_score = round(conf_comp.get("recency_score", 1.0), 1)
    _recency_warning_html = ""
    if _recency_level <= 1 and _recency_score <= 0.2 and ev_date:
        _recency_warning_html = (
            f'<p style="font-size:11px;color:#6D4C41;background:#FFF3E0;padding:8px 12px;'
            f'border-left:3px solid #FF6F00;margin:8px 0;{P}">'
            f'<strong>⚠️ Evidence date may be incorrect:</strong> Your evidence date '
            f'({ev_date}) is scoring as very old or outside your reporting period, '
            f'which is penalising your Recency score ({_recency_score}/1.0). '
            f'This may be a data entry error — correct the evidence date and re-score.'
            f'</p>'
        )
        # Suppress the evidence statement when date is clearly anomalous
        ev_stmt_block = (
            f'<p style="font-size:11px;color:#6D4C41;font-style:italic;margin:12px 0;">'
            f'Evidence statement not generated — evidence date appears incorrect. '
            f'Correct the date and re-score to generate this.'
            f'</p>'
        )

    # ── Change D: Per-sub-score rationales for bar rows (council XXI) ─────────
    from diagnostics import _DIRECTNESS_TIPS, _VERIFICATION_TIPS  # noqa: E402
    _direct_level  = conf_comp.get("direct_level", 0)
    _verify_level  = conf_comp.get("verify_level", 0)

    def _recency_rationale(level: int) -> str:
        return {
            0: "Evidence date could not be parsed — check the date format",
            1: "Evidence >12 months from reporting period — refresh or flag explicitly",
            2: "Evidence 7–12 months from reporting period — within acceptable range for some donors",
            3: "Evidence 4–6 months from reporting period — acceptable",
            4: "Evidence 1–3 months from reporting period — good recency",
            5: "Evidence within reporting month — maximum recency",
        }.get(level, "")

    def _clarity_rationale(score: float, max_v: float, label: str) -> str:
        pct = (score / max_v * 100) if max_v else 0
        if pct >= 90:
            return f"{label} at or near maximum"
        if pct >= 60:
            return f"{label} acceptable — minor gaps present"
        return f"{label} below threshold — primary Clarity gap"

    def bar_row_with_note(val, max_v, label, note: str = ""):
        """Score row with optional rationale note below the bar."""
        pct = min(int((val / max_v) * 100), 100) if max_v else 0
        bar_color = "#1B5E20" if pct >= 70 else ("#F57F17" if pct >= 50 else "#B71C1C")
        score_str = f"{val} / {max_v}"
        filled = max(0, round(pct / 10))
        empty  = 10 - filled
        bar_txt = f"<font color='{bar_color}'>{'&#9632;' * filled}</font><font color='#E0E0E0'>{'&#9632;' * empty}</font>"
        note_cell = (
            f"<td colspan='3' style='font-size:9px;color:#757575;font-style:italic;"
            f"padding:0 0 5px 0;'>{note}</td>"
        ) if note else ""
        row = (
            f"<tr>"
            f"<td width='130' style='font-size:11px;color:#424242;padding:3px 0;'>{label}</td>"
            f"<td width='100' style='font-size:10px;padding:3px 4px;letter-spacing:1px;'>{bar_txt}</td>"
            f"<td width='55' style='font-size:11px;font-weight:700;color:{bar_color};padding:3px 0;{P}'>{score_str}</td>"
            f"</tr>"
        )
        if note_cell:
            row += f"<tr>{note_cell}</tr>"
        return row

    _card_threshold   = evaluation.get("threshold_used", 4.0)
    _card_track_label = evaluation.get("track_label", "INGO standard")
    _card_org_type    = submission.get("org_type", "International NGO (INGO)")

    return f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<title>Submission Readiness Card</title>
<style>
body{{font-family:Arial,Helvetica,sans-serif;color:#212121;margin:24px 32px;font-size:12px;line-height:1.5;}}
h1{{color:#1B5E20;font-size:18px;margin:0 0 2px;border-bottom:2px solid #1B5E20;padding-bottom:6px;}}
h2{{color:#1B5E20;font-size:13px;font-weight:700;border-bottom:1px solid #8A6500;padding-bottom:3px;margin:16px 0 8px;}}
</style>
</head><body>

<!-- ═══ HEADER ═══ -->
<h1>Pre-Submission Readiness Card</h1>
<p style="color:#616161;font-size:10px;margin:2px 0 4px;">ImpactProof &middot; Ref: IMP-{timestamp}{donor_note}{(' &middot; ' + user_email) if user_email else ''}</p>
<table border="0" cellspacing="0" cellpadding="0" width="100%" style="margin-bottom:10px;{P}"><tr>
<td style="font-size:10px;color:#424242;padding:2px 0;">Submitted to: ___________________________________</td>
<td style="font-size:10px;color:#424242;padding:2px 0;text-align:right;">Date submitted: _______________</td>
</tr></table>

<!-- Reproducibility statement on Page 1 (council XXIV DRCA) -->
<p style="font-size:9px;color:#374151;background:#F3F8FE;border-left:3px solid #1565C0;padding:5px 8px;margin:0 0 6px;{P}">
<strong style="color:#1565C0;">Deterministic score</strong> — computed by rule-based criteria anchored to
<strong>USAID ADS 201</strong> &middot; <strong>FCDO Evaluation Policy 2025</strong> &middot;
<strong>Bond Evidence Principles 2024</strong> &middot; <strong>World Bank Results Framework</strong>.
No AI judgement applied to scores. Same inputs always produce the same result.
</p>
<!-- Evidence standard track badge (council XXXVI) -->
<p style="font-size:9px;color:{vfg};background:{vbg};padding:3px 8px;display:inline-block;border-radius:3px;margin:0 0 10px;{P}">
Evidence standard: <strong>{_card_track_label}</strong> &middot; threshold {_card_threshold}/5.0 &middot; {_card_org_type}
</p>

<!-- Result -->
<table border="0" cellspacing="0" cellpadding="0" width="100%" style="margin-bottom:10px;{P}"><tr>
<td style="background:#F5F5F5;padding:10px 14px;border-radius:6px;font-size:12px;{P}">
<strong>Result statement:</strong> {rs}
</td></tr></table>

<!-- Verdict badge -->
<table border="0" cellspacing="0" cellpadding="0" width="100%" style="margin-bottom:10px;{P}"><tr>
<td style="background:{vbg};color:{vfg};padding:10px 16px;border-radius:6px;font-weight:700;font-size:13px;{P}">
{verdict or diag_state}
</td></tr></table>

<!-- Score boxes (fixed-px widths to avoid xhtml2pdf negative-width bug) -->
<table border="0" cellspacing="8" cellpadding="0" style="margin-bottom:12px;{P}"><tr>
<td width="240" bgcolor="{cbg}" style="padding:14px;text-align:center;{P}">
  <p style="font-size:28px;font-weight:700;color:{cfg};margin:0;font-family:Courier,monospace;{P}">{conf_score}/5.0</p>
  <p style="font-size:10px;font-weight:700;color:{cfg};margin:4px 0 0;{P}">CONFIDENCE &mdash; {conf_label.upper()}</p>
</td>
<td width="240" bgcolor="{lbg}" style="padding:14px;text-align:center;{P}">
  <p style="font-size:28px;font-weight:700;color:{lfg};margin:0;font-family:Courier,monospace;{P}">{clar_score}/5.0</p>
  <p style="font-size:10px;font-weight:700;color:{lfg};margin:4px 0 0;{P}">CLARITY &mdash; {clar_label.upper()}</p>
</td>
</tr></table>

{_benchmark_html}

<!-- Verdict rationale (council XXI) -->
{_verdict_rationale_html}

<!-- Priority fixes -->
<h2>Priority fixes before submission</h2>
<table border="0" cellspacing="0" cellpadding="0" width="100%" style="margin-bottom:12px;">
{fixes_rows}
</table>

{ev_stmt_block}

<!-- ═══ PAGE 2 ═══ -->
<p style="page-break-before:always;"></p>

<h2>Score Breakdown</h2>

<p style="font-size:11px;font-weight:700;color:#424242;margin:8px 0 4px;">CONFIDENCE</p>
<table border="0" cellspacing="0" cellpadding="0" style="margin-bottom:10px;">
{bar_row_with_note(round(conf_comp.get('direct_score',0),1), 2.0, 'Directness', _DIRECTNESS_TIPS.get(_direct_level,'')[:90])}
{bar_row_with_note(round(conf_comp.get('verify_score',0),1), 2.0, 'Verification', _VERIFICATION_TIPS.get(_verify_level,'')[:90])}
{bar_row_with_note(round(conf_comp.get('recency_score',0),1), 1.0, 'Recency', _recency_rationale(_recency_level))}
</table>
{_recency_warning_html}

<p style="font-size:11px;font-weight:700;color:#424242;margin:8px 0 4px;">CLARITY</p>
<table border="0" cellspacing="0" cellpadding="0" style="margin-bottom:10px;">
{bar_row_with_note(round(clar_comp.get('definition_score',0),2), 1.25, def_label, _clarity_rationale(clar_comp.get('definition_score',0),1.25,'Definition'))}
{bar_row_with_note(round(clar_comp.get('measurement_score',0),2), 1.25, meas_label, _clarity_rationale(clar_comp.get('measurement_score',0),1.25,'Measurement'))}
{bar_row_with_note(round(clar_comp.get('integrity_score',0),2), 1.0, 'Integrity', _clarity_rationale(clar_comp.get('integrity_score',0),1.0,'Integrity'))}
{bar_row_with_note(round(clar_comp.get('scope_score',0),2), 0.75, 'Scope', _clarity_rationale(clar_comp.get('scope_score',0),0.75,'Scope'))}
{bar_row_with_note(round(clar_comp.get('governance_score',0),2), 0.75, 'Governance', _clarity_rationale(clar_comp.get('governance_score',0),0.75,'Governance'))}
</table>

{'<h2>Logframe Linkage</h2><table border="0" cellspacing="0" cellpadding="0" style="margin-bottom:10px;">' + lf_rows + '</table>' if lf_rows else ''}

<h2>Evidence Details</h2>
<table border="0" cellspacing="0" cellpadding="0" width="100%" style="margin-bottom:6px;">
<tr><td style="color:#616161;font-size:11px;padding:2px 10px 2px 0;width:90px;">Type</td><td style="font-size:11px;">{ev_type}</td></tr>
{'<tr><td style="color:#616161;font-size:11px;padding:2px 10px 2px 0;">Verifier</td><td style="font-size:11px;">' + verifier + '</td></tr>' if verifier else ''}
{'<tr><td style="color:#616161;font-size:11px;padding:2px 10px 2px 0;">Date</td><td style="font-size:11px;">' + ev_date + '</td></tr>' if ev_date else ''}
</table>
<p style="font-size:11px;color:#616161;margin:0;">{ev_desc[:350] + ('...' if len(ev_desc)>350 else '')}</p>

<!-- ═══ PAGE 3 ═══ -->
<p style="page-break-before:always;"></p>

<h2>Methodology &amp; Standards</h2>
<p style="font-size:11px;color:#424242;margin-bottom:8px;">
This check scored 8 evidence-quality dimensions anchored in:
<strong>USAID ADS 201</strong> (Validity, Integrity, Precision, Reliability, Timeliness) &middot;
<strong>FCDO Evaluation Policy January 2025</strong> &middot;
<strong>Bond Evidence Principles 2024</strong> &middot;
<strong>World Bank Results Framework</strong>.
Scoring is fully deterministic — no AI judgement was applied; all decisions are rule-based and reproducible.
Scores reflect patterns in submitted form fields and do not constitute expert review, audit, or guarantee of donor acceptance.
</p>
<p style="font-size:11px;color:#424242;background:#FFF9C4;padding:8px 12px;border-left:3px solid #8A6500;margin-bottom:10px;{P}">
<strong>Important:</strong> Guidance only — your donor makes the final call, not this tool.
Score generated: {timestamp}.
</p>
{_irc_note_html}
<p style="color:#616161;font-style:italic;font-size:10px;border-top:1px solid #E0E0E0;margin-top:20px;padding-top:8px;">
ImpactProof &middot; Built in Accra for MEL teams across West Africa &middot; {APP_URL.replace('https://','').rstrip('/')}
</p>
{_build_council_page_html(council_assessment, conf_score, clar_score, timestamp, P) if council_assessment else ""}
</body></html>"""


def _build_html_report(submission: dict, evaluation: dict, timestamp: str, chart_id: str = "0") -> str:
    lite = st.session_state.get("lite_mode", False)
    conf_score = evaluation.get("confidence_score", 0)
    clar_score = evaluation.get("clarity_score", 0)
    conf_label = evaluation.get("confidence_label", "")
    clar_label = evaluation.get("clarity_label", "")
    verdict    = evaluation.get("verdict", "")
    fixes      = evaluation.get("fixes", [])
    conf_comp  = evaluation.get("confidence_components", {})
    clar_comp  = evaluation.get("clarity_components", {})
    filenames  = submission.get("attached_filenames", [])

    badge_colors = {
        "Strong":     ("#C8E6C9", "#1B5E20"),
        "Acceptable": ("#FFF9C4", "#F57F17"),
        "Weak":       ("#FFE0B2", "#E65100"),
        "High Risk":  ("#FFCDD2", "#B71C1C"),
    }
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"

    _SCORE_GLOSS = {
        "Strong":     "Low risk — suitable for donor reporting.",
        "Acceptable": "Meets basic thresholds — room to strengthen further.",
        "Weak":       "Significant gaps — strengthen before submitting.",
        "High Risk":  "Both axes need work — address all fixes before sharing with donors.",
    }

    def badge(label, score, max_s):
        bg, fg = badge_colors.get(label, ("#F5F5F5", "#212121"))
        gloss = _SCORE_GLOSS.get(label, "")
        gloss_html = (f"<div style='font-size:0.8rem;font-weight:400;color:#616161;"
                      f"margin-top:4px;'>{gloss}</div>") if gloss else ""
        return (f"<div style='background:{bg};color:{fg};padding:10px 14px;"
                f"border-radius:8px;font-weight:700;font-size:0.9rem;margin-bottom:6px;{_pca}'>"
                f"{score}/{max_s} &nbsp; {label.upper()}{gloss_html}</div>")

    def bar(value, max_v):
        pct = min(value / max_v * 100, 100) if max_v else 0
        fill_px = round(pct / 100 * 120)
        empty_px = 120 - fill_px
        return (
            f"<table style='border-collapse:collapse;width:120px;height:10px;margin:4px 0 6px 0;{_pca}'>"
            f"<tr>"
            f"<td style='background:#1B5E20;width:{fill_px}px;height:10px;padding:0;{_pca}'></td>"
            f"<td style='background:#E0E0E0;width:{empty_px}px;height:10px;padding:0;'></td>"
            f"</tr></table>"
        )

    def row(label, score, max_v, tip=""):
        title = f' title="{tip}"' if tip else ""
        return (f"<tr><td{title} style='padding:6px 8px;font-size:0.88rem;'>{label}</td>"
                f"<td style='padding:6px 8px;font-family:monospace;'>{score}/{max_v}</td>"
                f"<td style='padding:6px 8px;width:120px;'>{bar(score, max_v)}</td></tr>")

    conf_fixes = [f for f in fixes if f.get("dimension") == "confidence"]
    clar_fixes = [f for f in fixes if f.get("dimension") == "clarity"]

    def fix_items(items):
        if not items:
            return ""
        return "".join(
            f"<li style='margin-bottom:6px;'>{f['message']} "
            f"<em style='color:#616161;'>({f['score_impact']})</em></li>"
            for f in items
        )

    verdict_colors = {
        "Strong KPI — well-positioned for submission": "#1B5E20",
        "Misleading KPI — sharpen the definition before submission": "#E65100",
        "Well-defined but weak evidence — strengthen the verification chain": "#F57F17",
        "High risk — strengthen both axes before relying on this result": "#B71C1C",
    }
    verdict_bg = verdict_colors.get(verdict, "#1B5E20")

    # Headline: is this good enough to submit?
    _diag_state, _ = get_diagnostic_state(
        conf_score, clar_score,
        evaluation.get("content_issues", []),
        submission.get("beneficiary_voice", ""),
    )
    _readiness_html = _readiness_banner_html(_diag_state)

    files_row = (f"<p><strong>Attached documents:</strong> {', '.join(filenames)}</p>"
                 if filenames else "")

    dl = conf_comp.get("direct_level", 0)
    vl = conf_comp.get("verify_level", 0)
    rl = conf_comp.get("recency_level", 0)
    ds = conf_comp.get("direct_score", 0)
    vs = conf_comp.get("verify_score", 0)
    rs = conf_comp.get("recency_score", 0)
    def_s  = clar_comp.get("definition_score", 0)
    meas_s = clar_comp.get("measurement_score", 0)
    integ  = clar_comp.get("integrity_score", 0)
    scope  = clar_comp.get("scope_score", 0)
    gov    = clar_comp.get("governance_score", 0)
    is_qual    = clar_comp.get("is_qualitative", False)
    def_label  = "Narrative Definition" if is_qual else "Definition"
    def_tip    = _CLARITY_TIPS["definition_qualitative"] if is_qual else _CLARITY_TIPS["definition"]
    meas_label = "Sourcing & Triangulation" if is_qual else "Measurement"
    meas_tip   = _CLARITY_TIPS["measurement_qualitative"] if is_qual else _CLARITY_TIPS["measurement"]

    # --- Overview chart for report ---
    _r_conf, _r_clar, _r_eth, _r_comp = _overview_score_values(evaluation)
    _overview_b64 = "" if lite else _build_overview_chart_b64(_r_conf, _r_clar, _r_eth, _r_comp)
    _overview_alt = (
        f"Overview chart of diagnostic scores out of 100: Confidence {_r_conf:.0f}, "
        f"Clarity {_r_clar:.0f}, Ethics {_r_eth:.0f}, Compliance {_r_comp:.0f}."
    )
    _overview_img = (f'<img src="data:image/png;base64,{_overview_b64}" '
                     f'alt="{_overview_alt}" style="max-width:420px;width:100%;display:block;margin:0 auto 12px;" />'
                     ) if _overview_b64 else ""
    _overview_legend = (
        "<ul style='font-size:0.8rem;color:#616161;max-width:560px;margin:0 auto 16px;padding-left:20px;'>"
        "<li><strong>Confidence</strong> — how directly the evidence supports the result, "
        "how independently it was verified, and how recent it is.</li>"
        "<li><strong>Clarity</strong> — how precisely the result is defined, measured, "
        "and bounded in scope, with a documented audit trail.</li>"
        "<li><strong>Ethics</strong> — completeness and integrity of the underlying data "
        "(missing data, audit trail, sample adequacy). Feeds into Clarity.</li>"
        "<li><strong>Compliance</strong> — consent, anonymisation, and data-protection-law status "
        "for any beneficiary data used as evidence. Feeds into Clarity.</li>"
        "<li style='margin-top:6px;'>Bar color shows the score band: "
        "<span style='color:#1B5E20;font-weight:700;'>green = 75&ndash;100 (strong)</span>, "
        "<span style='color:#8A6500;font-weight:700;'>gold = 50&ndash;74 (acceptable)</span>, "
        "<span style='color:#C62828;font-weight:700;'>red = below 50 (needs work)</span>.</li>"
        "</ul>"
    ) if _overview_b64 else ""

    # --- Interactive (vega-embed) sub-score charts for the downloadable report ---
    if lite:
        _conf_chart_div = ""
        _clar_chart_div = ""
        _charts_script = ""
    else:
        _conf_spec_json = _subscore_chart([
            (f"Directness (Level {dl}/5)", ds, 2.0, _DIRECTNESS_TIPS.get(dl, "")),
            (f"Verification (Level {vl}/5)", vs, 2.0, _VERIFICATION_TIPS.get(vl, "")),
            (f"Recency (Level {rl}/5)", rs, 1.0, _RECENCY_TIPS.get(rl, "")),
        ]).to_json()
        _clar_spec_json = _subscore_chart([
            (def_label, def_s, 1.25, def_tip),
            (meas_label, meas_s, 1.25, meas_tip),
            ("Integrity", integ, 1.0, _CLARITY_TIPS["integrity"]),
            ("Scope", scope, 0.75, _CLARITY_TIPS["scope"]),
            ("Governance", gov, 0.75, _CLARITY_TIPS["governance"]),
        ]).to_json()
        _conf_chart_div = f'<div id="conf-chart-{chart_id}" class="no-print" style="margin-top:10px;"></div>'
        _clar_chart_div = f'<div id="clar-chart-{chart_id}" class="no-print" style="margin-top:10px;"></div>'
        _charts_script = f"""
<script type="application/json" id="conf-spec-{chart_id}">{_conf_spec_json}</script>
<script type="application/json" id="clar-spec-{chart_id}">{_clar_spec_json}</script>
<script>
(function(){{
  function renderImpactReceiptsCharts_{chart_id}(){{
    if (!window.vegaEmbed) {{ setTimeout(renderImpactReceiptsCharts_{chart_id}, 200); return; }}
    var confSpec = JSON.parse(document.getElementById('conf-spec-{chart_id}').textContent);
    var clarSpec = JSON.parse(document.getElementById('clar-spec-{chart_id}').textContent);
    vegaEmbed('#conf-chart-{chart_id}', confSpec, {{actions: false}});
    vegaEmbed('#clar-chart-{chart_id}', clarSpec, {{actions: false}});
  }}
  renderImpactReceiptsCharts_{chart_id}();
}})();
</script>
"""

    _methodology_table = (
        "<div style='max-width:620px;margin:0 auto 24px;'>"
        "<p style='font-size:0.85rem;font-weight:700;color:#1B5E20;margin-bottom:6px;'>What this report checks:</p>"
        "<table style='width:100%;font-size:0.8rem;'>"
        "<tr><th>Check</th><th>Reflected in</th></tr>"
        "<tr><td>Logframe linkage — does your result tie to an approved indicator?</td>"
        "<td>Clarity &rarr; Measurement</td></tr>"
        "<tr><td>Evidence quality — direct, verified, recent, defensible?</td>"
        "<td>Confidence</td></tr>"
        "<tr><td>Beneficiary voice — were they part of the evidence?</td>"
        "<td>Confidence (bonus)</td></tr>"
        "<tr><td>Definition clarity — would two readers interpret it the same way?</td>"
        "<td>Clarity &rarr; Definition</td></tr>"
        "<tr><td>Submission completeness — is your package donor-ready?</td>"
        "<td>Advisory only — not scored</td></tr>"
        "</table>"
        "<p style='font-size:0.875rem;color:#424242;margin-top:6px;'>"
        "Confidence and Clarity are your two top-line scores. Ethics and Compliance are the "
        "Integrity and Governance sub-scores within Clarity, shown separately above so you can "
        "see what's driving your Clarity score — see definitions above.</p>"
        "<p style='font-size:0.85rem;font-weight:700;color:#1B5E20;margin:16px 0 6px;'>"
        "Donor framework crosswalk &mdash; how each sub-score maps to the standards your donor audits against:</p>"
        f"{build_donor_crosswalk_html(st.session_state.get('donor_framework', 'Generic'))}"
        "</div>"
    )
    _meta_donor  = submission.get("donor") or st.session_state.get("donor_selected", "")
    if not _meta_donor or _meta_donor in ("(No donor specified)", "Not specified"):
        _meta_donor = "—"

    _meta_sector = submission.get("sector") or st.session_state.get("sector", "")
    if not _meta_sector or _meta_sector in ("(No sector selected)", "Not specified"):
        _meta_sector = "—"

    _meta_rtype = submission.get("submission_type") or st.session_state.get("submission_type", "")
    if not _meta_rtype or _meta_rtype in ("Select submission type...", "Not specified"):
        _meta_rtype = "—"
    _meta_html = (
        "<table style='margin-bottom:20px;'><tbody>"
        f"<tr><td style='padding:5px 12px;font-weight:700;'>Donor</td><td style='padding:5px 12px;'>{_meta_donor}</td></tr>"
        f"<tr><td style='padding:5px 12px;font-weight:700;'>Sector</td><td style='padding:5px 12px;'>{_meta_sector}</td></tr>"
        f"<tr><td style='padding:5px 12px;font-weight:700;'>Report Type</td><td style='padding:5px 12px;'>{_meta_rtype}</td></tr>"
        "</tbody></table>"
    )
    # --- End radar ---

    fixes_html = ""
    if conf_fixes:
        fixes_html += ("<h3 style='color:#1B5E20;'>Strengthen your evidence (Confidence)</h3>"
                       f"<ul>{fix_items(conf_fixes)}</ul>")
    if clar_fixes:
        fixes_html += ("<h3 style='color:#1B5E20;'>Sharpen your definition (Clarity)</h3>"
                       f"<ul>{fix_items(clar_fixes)}</ul>")
    if not fixes:
        fixes_html = "<p style='color:#1B5E20;font-weight:700;'>This check found no further fixes to address.</p>"

    # --- What Donors Want to Know (four-question summary) ---
    ladder   = evaluation.get("evidence_ladder", {})
    maturity = evaluation.get("indicator_maturity", {})
    fr       = evaluation.get("funder_readiness", {})
    ev_top      = (submission.get("evidence") or [{}])[0]
    ev_type_top = ev_top.get("type", "") or "Not specified"
    dominant    = ladder.get("dominant_tier")
    learn       = fr.get("learning", {})
    lim         = fr.get("limitations", {})

    q2_answer = f"Evidence type: <strong>{ev_type_top}</strong>"
    if dominant:
        q2_answer += f" &mdash; evidence base is mainly <strong>{dominant}</strong>-tier."
    q4_answer = ("Yes &mdash; the report describes what your team learned and how the programme adapted."
                  if learn.get("detected") else
                  "Not yet stated. Add a sentence on what you learned and changed as a result.")

    four_questions_html = f"""
<h2>What Donors Want to Know</h2>
<table>
  <tr><th>Question</th><th>Answer</th></tr>
  <tr><td><strong>1. What has changed?</strong></td>
      <td>{submission.get('result_statement', '-')}<br/>
          <span style="color:#616161;font-size:0.85rem;">Directness: Level {dl}/5 &middot; Definition: {def_s}/1.25</span></td></tr>
  <tr><td><strong>2. How do you know?</strong></td><td>{q2_answer}</td></tr>
  <tr><td><strong>3. How strong is the evidence?</strong></td>
      <td>Confidence: <strong>{conf_score}/5.0</strong> ({conf_label})</td></tr>
  <tr><td><strong>4. What did you learn?</strong></td><td>{q4_answer}</td></tr>
</table>
"""

    # --- Evidence Ladder ---
    ladder_html = ""
    if ladder:
        tier_descriptions = {
            "Basic": "Attendance, registration, logs, photos",
            "Moderate": "Follow-up surveys, testimonials",
            "Stronger": "Business/regulatory records, mentor verification, "
                        "baseline/endline, external evaluation, comparison groups",
        }
        counts = ladder.get("tier_counts", {})
        ladder_rows = "".join(
            f"<tr><td>{tier}{' &#128072;' if tier == dominant else ''}</td>"
            f"<td>{tier_descriptions[tier]}</td>"
            f"<td style='font-family:monospace;'>{counts.get(tier, 0)}</td></tr>"
            for tier in _evaluator.EVIDENCE_LADDER_TIERS
        )
        ladder_html = f"""
<h2>Evidence Ladder</h2>
<p style="color:#616161;font-size:0.85rem;">Rule-based check of the evidence sources you described — does this report rely mainly on Basic, Moderate, or Stronger evidence?</p>
<table><tr><th>Tier</th><th>Description</th><th>Sources detected</th></tr>{ladder_rows}</table>
<p>{ladder.get('suggestion', '')}</p>
"""

    # --- Indicator Maturity ---
    maturity_html = ""
    if maturity.get("flagged"):
        _maturity_row_parts = []
        for i, (level, wording) in enumerate(maturity["rows"]):
            tr_style = ' style="background:#FFF9E0;"' if i == 0 else ""
            level_label = f"\U0001F448 {level} — what you wrote" if i == 0 else level
            _maturity_row_parts.append(
                f"<tr{tr_style}><td>{level_label}</td><td>{wording}</td></tr>"
            )
        maturity_rows = "".join(_maturity_row_parts)
        _your_indicator = (submission.get("logframe_indicator", "") or "").strip() or "(not specified)"
        maturity_html = f"""
<h2>Indicator Maturity</h2>
<p style="color:#616161;font-size:0.85rem;">This indicator is written as a raw count. Donors increasingly expect indicators that show whether the result was sustained or verified.</p>
<p><strong>Your indicator (as written):</strong> {_your_indicator}</p>
<table><tr><th>Level</th><th>Example wording</th></tr>{maturity_rows}</table>
<p>Measurement score adjusted by <strong>{maturity['adjustment']}</strong> for this count-only indicator framing.</p>
"""

    # --- Donor Readiness flags ---
    lim_text = ("Limitations disclosed &mdash; the report states what the data can't confidently say."
                 if lim.get("detected") else
                 "No limitations disclosure detected. Consider adding a sentence on what this data "
                 "cannot confirm or cannot be generalized to.")
    learn_text = ("Learning &amp; adaptation stated &mdash; the report describes what your team learned and changed."
                   if learn.get("detected") else
                   "No learning/adaptation statement detected. Consider adding what your organization "
                   "learned and how the program adapted as a result.")
    funder_readiness_html = f"""
<h2>Donor Readiness</h2>
<p style="color:#616161;font-size:0.85rem;">These checks do not affect your Confidence or Clarity scores.</p>
<ul>
  <li>{lim_text}</li>
  <li>{learn_text}</li>
</ul>
"""

    # --- Additional advisory flags (v3.4, score-neutral) ---
    attrib = submission.get("attribution_contribution", "Not specified")
    disagg = submission.get("disaggregation_status", "Not specified")
    advisory_html = ""
    if attrib != "Not specified" or disagg != "Not specified":
        advisory_items = ""
        if attrib != "Not specified":
            advisory_items += f"<li><strong>Attribution vs. contribution distinguished:</strong> {attrib}</li>"
        if disagg != "Not specified":
            advisory_items += f"<li><strong>Beneficiary data disaggregated (women, youth, PWD, rural):</strong> {disagg}</li>"
        advisory_html = f"""
<h2>Additional Advisory Flags</h2>
<p style="color:#616161;font-size:0.85rem;">Optional checklist answers — advisory only, no effect on your score.</p>
<ul>{advisory_items}</ul>
"""

    _vega_scripts = "" if lite else (
        '<script src="https://cdn.jsdelivr.net/npm/vega@5"></script>\n'
        '<script src="https://cdn.jsdelivr.net/npm/vega-lite@6"></script>\n'
        '<script src="https://cdn.jsdelivr.net/npm/vega-embed@6"></script>'
    )
    _lite_note = (
        '<p style="color:#8A6500;font-size:0.85rem;">'
        'Low-bandwidth mode &mdash; simplified report (no embedded charts).</p>'
    ) if lite else ""

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>ImpactProof — Report</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;700&display=swap" rel="stylesheet"/>
{_vega_scripts}
<style>
  body{{font-family:'Inter',sans-serif;color:#212121;max-width:860px;margin:40px auto;padding:0 24px;}}
  h1,h2,h3{{color:#1B5E20;}} h1{{font-size:1.6rem;}} h2{{font-size:1.2rem;border-bottom:1px solid #8A6500;padding-bottom:4px;margin-top:28px;}}
  table{{width:100%;border-collapse:collapse;margin-bottom:16px;}}
  td,th{{border:1px solid #E0E0E0;text-align:left;}}
  th{{background:#F5F5F5;padding:7px 8px;font-size:0.85rem;}}
  .grid{{display:grid;grid-template-columns:1fr 1fr;gap:24px;}}
  .footer{{color:#424242;font-style:italic;font-size:0.875rem;border-top:1px solid #E0E0E0;margin-top:32px;padding-top:12px;}}
  @media print{{
    body{{margin:20px;}}
    .no-print{{display:none;}}
    *{{-webkit-print-color-adjust:exact !important;print-color-adjust:exact !important;}}
  }}
</style>
</head>
<body>
<h1>ImpactProof — Evaluation Report</h1>
<p style="color:#616161;font-size:0.88rem;">Generated: {timestamp}</p>
{_lite_note}
{_meta_html}
{_overview_img}
{_overview_legend}
{_methodology_table}
<h2>Result Statement</h2>
<p>{submission.get('result_statement', '-')}</p>
<p><strong>Target Group:</strong> {submission.get('target_group', '-')}<br/>
   <strong>Timeframe:</strong> {submission.get('timeframe', '-')}<br/>
   <strong>Geographic Scope:</strong> {submission.get('geographic_scope', '-')}</p>
{files_row}
{_readiness_html}
{four_questions_html}
<div style="background:{verdict_bg};color:white;border-radius:10px;padding:14px 20px;
     font-weight:700;text-align:center;margin:20px 0;font-size:1rem;
     -webkit-print-color-adjust:exact;print-color-adjust:exact;">
  {verdict}
</div>

<h2>Score Summary</h2>
<div class="grid">
  <div>
    <strong>Confidence Score</strong><br/>
    {badge(conf_label, conf_score, 5.0)}
    {bar(conf_score, 5.0)}
    <p style="color:#616161;font-size:0.82rem;margin:2px 0 6px 0;">
      Your Confidence score of {round(conf_score/5*100)}/100 reflects a sub-component total of {conf_score:.2f}/5.0.
    </p>
    <p style="color:#616161;font-size:0.85rem;">{evaluation.get('confidence_meaning','')}</p>
    <table>
      <tr><th>Component</th><th>Score</th><th>Bar</th></tr>
      {row(f"Directness (Level {dl}/5)", ds, 2.0, _DIRECTNESS_TIPS.get(dl,''))}
      {row(f"Verification (Level {vl}/5)", vs, 2.0, _VERIFICATION_TIPS.get(vl,''))}
      {row(f"Recency (Level {rl}/5)", rs, 1.0, _RECENCY_TIPS.get(rl,''))}
    </table>
    {_conf_chart_div}
  </div>
  <div>
    <strong>Clarity Score</strong><br/>
    {badge(clar_label, clar_score, 5.0)}
    {bar(clar_score, 5.0)}
    <p style="color:#616161;font-size:0.82rem;margin:2px 0 6px 0;">
      Your Clarity score of {round(clar_score/5*100)}/100 reflects a sub-component total of {clar_score:.2f}/5.0.
    </p>
    <p style="color:#616161;font-size:0.85rem;">{evaluation.get('clarity_meaning','')}</p>
    <table>
      <tr><th>Component</th><th>Score</th><th>Bar</th></tr>
      {row(def_label, def_s, 1.25, def_tip)}
      {row(meas_label, meas_s, 1.25, meas_tip)}
      {row("Integrity", integ, 1.0, _CLARITY_TIPS['integrity'])}
      {row("Scope", scope, 0.75, _CLARITY_TIPS['scope'])}
      {row("Governance", gov, 0.75, _CLARITY_TIPS['governance'])}
    </table>
    {_clar_chart_div}
  </div>
</div>
<p class="no-print" style="color:#424242;font-size:0.875rem;">Hover over the charts above for details (requires an internet connection to load the chart library).</p>
{_charts_script}
{ladder_html}
{maturity_html}
{funder_readiness_html}
{advisory_html}
<div style="page-break-inside:avoid;">
<h2>What To Fix</h2>
{fixes_html}
</div>

<div class="footer">
  Evaluated using {METHODOLOGY_STACK}.<br/>
  Contact: <a href="https://wa.me/233503648195">WhatsApp +233 50 364 8195</a> &nbsp;&middot;&nbsp; <a href="{APP_URL}/">ImpactProof</a>
</div>
</body>
</html>"""


def _build_reviewer_signoff_section_html(review_info: dict) -> str:
    """Render the review-handoff fields (v3.6) as an HTML section to be
    appended to an existing report. Nothing is auto-filled — blanks are
    shown literally as 'Not provided'."""
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
    status   = review_info.get("review_status", "") or _SUBMISSION_STATUS_OPTIONS[0]
    decision = review_info.get("reviewer_decision", "") or "Not yet decided"
    name     = review_info.get("reviewer_name", "") or "Not provided"
    role     = review_info.get("reviewer_role", "") or "Not provided"
    rdate    = review_info.get("reviewer_date", "") or "Not provided"
    notes    = review_info.get("reviewer_notes", "")

    bg, fg = _SUBMISSION_STATUS_COLORS.get(status, ("#F5F5F5", "#616161"))
    notes_html = (
        f"<div style='background:#FFF9C4;border-radius:8px;padding:10px 14px;margin:10px 0;font-size:0.85rem;'>"
        f"<strong>Reviewer notes:</strong> {notes}</div>"
    ) if notes else ""

    return f"""
<h2>Review &amp; Sign-off</h2>
<div style="background:{bg};color:{fg};padding:6px 12px;border-radius:8px;font-weight:700;
     font-size:0.85rem;display:inline-block;margin-bottom:10px;{_pca}">
  STATUS: {status.upper()}
</div>
<p><strong>Decision:</strong> {decision}</p>
{notes_html}
<p style="margin-top:14px;">
  <strong>Reviewer:</strong> {name} &nbsp;&nbsp;
  <strong>Role:</strong> {role} &nbsp;&nbsp;
  <strong>Date:</strong> {rdate}
</p>
<p style="color:#424242;font-size:0.875rem;">
  No accounts or authentication are used — this reflects the reviewer information
  entered for this submission in the current session.
</p>
"""


def _build_review_summary_html(submission: dict, evaluation: dict, review_info: dict,
                                timestamp: str, chart_id: str = "0") -> str:
    """Review summary export (v3.6): reuses _build_html_report unchanged and
    appends a Review & Sign-off section capturing the submission, scores,
    gaps (via the existing report), and the reviewer's notes/decision."""
    base = _build_html_report(submission, evaluation, timestamp, chart_id=chart_id)
    section = _build_reviewer_signoff_section_html(review_info)
    return base.replace("</body>", section + "</body>")


def _build_review_summary_docx(submission: dict, evaluation: dict, review_info: dict, timestamp: str) -> bytes:
    """DOCX review summary (v3.6): reuses the existing donor-template DOCX
    table style for the submission/score/gaps content, then appends a
    Review & Sign-off table. Blanks are written literally as 'Not provided'."""
    import io

    conf_score = evaluation.get("confidence_score", 0)
    clar_score = evaluation.get("clarity_score", 0)
    conf_comp  = evaluation.get("confidence_components", {})
    clar_comp  = evaluation.get("clarity_components", {})
    fixes      = evaluation.get("fixes", [])

    doc = _docx.Document()
    doc.add_heading("ImpactProof — Review Summary", level=0)
    doc.add_paragraph(f"Generated: {timestamp}")

    doc.add_heading("Result Statement", level=1)
    doc.add_paragraph(submission.get("result_statement", "") or "Not provided")
    table = doc.add_table(rows=0, cols=2)
    table.style = "Light Grid Accent 1"
    for label, value in [
        ("Target Group",     submission.get("target_group", "")),
        ("Timeframe",        submission.get("timeframe", "")),
        ("Geographic Scope", submission.get("geographic_scope", "")),
        ("Verdict",          evaluation.get("verdict", "")),
    ]:
        cells = table.add_row().cells
        cells[0].text = label
        cells[1].text = str(value) if value else "Not provided"

    doc.add_heading("Score Summary", level=1)
    score_table = doc.add_table(rows=0, cols=3)
    score_table.style = "Light Grid Accent 1"
    header = score_table.add_row().cells
    header[0].text, header[1].text, header[2].text = "Component", "Score", "Max"
    for label, value, max_val in [
        ("Confidence (overall)", conf_score, 5.0),
        ("Directness",   conf_comp.get("direct_score", "-"), 2.0),
        ("Verification", conf_comp.get("verify_score", "-"), 2.0),
        ("Recency",       conf_comp.get("recency_score", "-"), 1.0),
        ("Clarity (overall)", clar_score, 5.0),
        ("Definition",   clar_comp.get("definition_score", "-"), 1.25),
        ("Measurement",  clar_comp.get("measurement_score", "-"), 1.25),
        ("Integrity",    clar_comp.get("integrity_score", "-"), 1.0),
        ("Scope",        clar_comp.get("scope_score", "-"), 0.75),
        ("Governance",   clar_comp.get("governance_score", "-"), 0.75),
    ]:
        cells = score_table.add_row().cells
        cells[0].text = label
        cells[1].text = str(value)
        cells[2].text = str(max_val)

    doc.add_heading("Gaps / What To Fix", level=1)
    if fixes:
        for fix in fixes:
            doc.add_paragraph(f"{fix.get('message', '')} ({fix.get('score_impact', '')})", style="List Bullet")
    else:
        doc.add_paragraph("This check found no further fixes to address.")

    doc.add_heading("Review & Sign-off", level=1)
    review_table = doc.add_table(rows=0, cols=2)
    review_table.style = "Light Grid Accent 1"
    for label, value in [
        ("Status",        review_info.get("review_status", "") or _SUBMISSION_STATUS_OPTIONS[0]),
        ("Decision",      review_info.get("reviewer_decision", "") or "Not yet decided"),
        ("Reviewer name", review_info.get("reviewer_name", "")),
        ("Reviewer role", review_info.get("reviewer_role", "")),
        ("Review date",   review_info.get("reviewer_date", "")),
    ]:
        cells = review_table.add_row().cells
        cells[0].text = label
        cells[1].text = str(value) if value else "Not provided"

    notes = review_info.get("reviewer_notes", "")
    if notes:
        doc.add_paragraph("Reviewer notes:")
        doc.add_paragraph(notes)

    doc.add_heading("About this export", level=1)
    doc.add_paragraph(
        f"Evaluated using {METHODOLOGY_STACK}. {_LIMITS_DISCLAIMER} "
        "This review summary has no user accounts or authentication — reviewer "
        "identity is free text entered for the current session."
    )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_combined_html_report(submissions: list, evaluations: list, timestamp: str) -> str:
    parts = []
    for i, (sub, ev) in enumerate(zip(submissions, evaluations)):
        section = _build_html_report(sub, ev, timestamp, chart_id=str(i))
        # Strip outer HTML wrapper from all but the first, append just the body content
        if i == 0:
            parts.append(section)
        else:
            start = section.find("<h2>Result Statement</h2>")
            if start == -1:
                start = section.find("<h2 ")
            end   = section.rfind("</body>")
            insert_at = parts[0].rfind("</body>")
            divider = f"<hr style='margin:40px 0;border:2px solid #8A6500;'/><h2 style='color:#1B5E20;'>Result {i+1}</h2>"
            parts[0] = parts[0][:insert_at] + divider + section[start:end] + parts[0][insert_at:]
    return parts[0]


_VERIFICATION_SUMMARY_CSS = """
  body{font-family:'Inter',sans-serif;color:#212121;max-width:760px;margin:40px auto;padding:0 24px;}
  h1,h2,h3{color:#1B5E20;} h1{font-size:1.5rem;} h2{font-size:1.05rem;border-bottom:1px solid #8A6500;padding-bottom:4px;margin-top:22px;}
  table{width:100%;border-collapse:collapse;margin-bottom:14px;}
  td,th{border:1px solid #E0E0E0;text-align:left;padding:5px 8px;font-size:0.85rem;}
  th{background:#F5F5F5;}
  .signoff{margin-top:28px;border-top:1px solid #E0E0E0;padding-top:18px;font-size:0.9rem;}
  .signoff p{margin:6px 0;}
  .signoff .line{display:inline-block;min-width:220px;border-bottom:1px solid #212121;margin:0 6px;}
  .footer{color:#616161;font-style:italic;font-size:0.8rem;border-top:1px solid #E0E0E0;margin-top:28px;padding-top:12px;}
  @media print{ body{margin:20px;} *{-webkit-print-color-adjust:exact !important;print-color-adjust:exact !important;} }
"""

_VERIFICATION_SUMMARY_BADGE_COLORS = {
    "Strong":     ("#C8E6C9", "#1B5E20"),
    "Acceptable": ("#FFF9C4", "#F57F17"),
    "Weak":       ("#FFE0B2", "#E65100"),
    "High Risk":  ("#FFCDD2", "#B71C1C"),
}

_VERIFICATION_SUMMARY_VERDICT_COLORS = {
    "Strong KPI — well-positioned for submission": "#1B5E20",
    "Misleading KPI — sharpen the definition before submission": "#E65100",
    "Well-defined but weak evidence — strengthen the verification chain": "#F57F17",
    "High risk — strengthen both axes before relying on this result": "#B71C1C",
}

_REPORT_STATUS_OPTIONS = [
    "Draft – pending review",
    "Submitted for review",
    "Reviewed – changes requested",
    "Approved – ready to submit",
]

_REPORT_STATUS_COLORS = {
    "Draft – pending review":        ("#F5F5F5", "#616161"),
    "Submitted for review":           ("#FFF9C4", "#F57F17"),
    "Reviewed – changes requested":   ("#FFE0B2", "#E65100"),
    "Approved – ready to submit":     ("#C8E6C9", "#1B5E20"),
}


def _verification_summary_badge(label: str, score: float, max_s: float) -> str:
    bg, fg = _VERIFICATION_SUMMARY_BADGE_COLORS.get(label, ("#F5F5F5", "#212121"))
    _pca = "-webkit-print-color-adjust:exact;print-color-adjust:exact;"
    return (f"<div style='background:{bg};color:{fg};padding:8px 12px;"
            f"border-radius:8px;font-weight:700;font-size:0.85rem;display:inline-block;{_pca}'>"
            f"{score}/{max_s} &nbsp; {label.upper()}</div>")


def _verification_summary_signoff(prepared_by: str = "", status: str = "", notes: str = "", timestamp: str = "") -> str:
    bg, fg = _REPORT_STATUS_COLORS.get(status, ("#F5F5F5", "#212121"))
    status_html = (
        f"<div style='background:{bg};color:{fg};padding:6px 12px;border-radius:8px;"
        f"font-weight:700;font-size:0.85rem;display:inline-block;margin-bottom:10px;"
        f"-webkit-print-color-adjust:exact;print-color-adjust:exact;'>STATUS: {status.upper()}</div>"
    ) if status else ""
    notes_html = (
        f"<div style='background:#FFF9C4;border-radius:8px;padding:10px 14px;margin:10px 0;font-size:0.85rem;'>"
        f"<strong>Notes for reviewer:</strong> {notes}</div>"
    ) if notes else ""
    prepared_line = prepared_by or "&nbsp;"
    return (
        f"{status_html}{notes_html}"
        "<div class='signoff'>"
        f"<p><strong>Prepared by:</strong> {prepared_line} &nbsp;&nbsp; <strong>Date:</strong> {timestamp}</p>"
        "<p><strong>Reviewed by:</strong> <span class='line'>&nbsp;</span> &nbsp;&nbsp; "
        "<strong>Date:</strong> <span class='line'>&nbsp;</span></p>"
        "<p><strong>Comments:</strong> <span class='line' style='min-width:400px;'>&nbsp;</span></p>"
        "<p><strong>Approved by:</strong> <span class='line'>&nbsp;</span> &nbsp;&nbsp; "
        "<strong>Date:</strong> <span class='line'>&nbsp;</span></p>"
        "</div>"
    )


def _build_verification_summary_html(submissions: list, evaluations: list, timestamp: str) -> str:
    """Build a short, printable per-result gap report for attaching to a submission."""
    _prepared_by = st.session_state.get("report_prepared_by", "")
    _status      = st.session_state.get("report_status", "")
    _notes       = st.session_state.get("report_notes", "")
    sections = []
    for i, (submission, evaluation) in enumerate(zip(submissions, evaluations)):
        conf_score = evaluation.get("confidence_score", 0)
        clar_score = evaluation.get("clarity_score", 0)
        conf_label = evaluation.get("confidence_label", "")
        clar_label = evaluation.get("clarity_label", "")
        verdict    = evaluation.get("verdict", "")
        fixes      = evaluation.get("fixes", [])
        conf_comp  = evaluation.get("confidence_components", {})
        clar_comp  = evaluation.get("clarity_components", {})
        verdict_bg = _VERIFICATION_SUMMARY_VERDICT_COLORS.get(verdict, "#1B5E20")
        is_qual    = clar_comp.get("is_qualitative", False)
        def_label  = "Narrative Definition" if is_qual else "Definition"
        meas_label = "Sourcing & Triangulation" if is_qual else "Measurement"

        sub_rows = "".join(
            f"<tr><td>{dim}</td><td style='font-family:monospace;'>{score}/{max_v}</td></tr>"
            for dim, score, max_v in [
                ("Directness",   conf_comp.get("direct_score", 0), 2.0),
                ("Verification", conf_comp.get("verify_score", 0), 2.0),
                ("Recency",      conf_comp.get("recency_score", 0), 1.0),
                (def_label,      clar_comp.get("definition_score", 0), 1.25),
                (meas_label,     clar_comp.get("measurement_score", 0), 1.25),
                ("Integrity",    clar_comp.get("integrity_score", 0), 1.0),
                ("Scope",        clar_comp.get("scope_score", 0), 0.75),
                ("Governance",   clar_comp.get("governance_score", 0), 0.75),
            ]
        )

        if fixes:
            fixes_html = "<ul>" + "".join(
                f"<li>{f['message']} <em style='color:#616161;'>({f['score_impact']})</em></li>"
                for f in fixes
            ) + "</ul>"
        else:
            fixes_html = "<p style='color:#1B5E20;font-weight:700;'>No further gaps flagged by this tool's checks.</p>"

        heading = "<h1>Verification Summary</h1>" if i == 0 else \
                  f"<hr style='margin:36px 0;border:2px solid #8A6500;'/><h1>Result {i + 1}</h1>"

        sections.append(f"""
{heading}
<p style="color:#616161;font-size:0.85rem;">Generated: {timestamp}</p>
<h2>Result Statement</h2>
<p>{submission.get('result_statement', '-')}</p>
<p><strong>Logframe Indicator:</strong> {submission.get('logframe_indicator', '-') or '-'}<br/>
   <strong>Target Group:</strong> {submission.get('target_group', '-')}<br/>
   <strong>Timeframe:</strong> {submission.get('timeframe', '-')}<br/>
   <strong>Geographic Scope:</strong> {submission.get('geographic_scope', '-')}</p>

<div style="background:{verdict_bg};color:white;border-radius:10px;padding:12px 18px;
     font-weight:700;text-align:center;margin:16px 0;font-size:0.95rem;
     -webkit-print-color-adjust:exact;print-color-adjust:exact;">
  {verdict}
</div>

<h2>Scores</h2>
<p>
  {_verification_summary_badge(conf_label, conf_score, 5.0)}
  &nbsp; {_verification_summary_badge(clar_label, clar_score, 5.0)}
</p>
<table><tr><th>Sub-score</th><th>Score</th></tr>{sub_rows}</table>

<h2>Gaps to Address</h2>
{fixes_html}
""")

    sections.append(f"""
<h2>Donor Framework Crosswalk</h2>
{build_donor_crosswalk_html(st.session_state.get('donor_framework', 'Generic'))}
<p style="color:#616161;font-size:0.8rem;">{_LIMITS_DISCLAIMER}</p>
{_verification_summary_signoff(_prepared_by, _status, _notes, timestamp)}
<div class="footer">
  Evaluated using {METHODOLOGY_STACK}.<br/>
  Contact: <a href="mailto:info@impact-receipts.com">info@impact-receipts.com</a>
</div>
""")

    body = "".join(sections)
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Verification Summary</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;700&display=swap" rel="stylesheet"/>
<style>{_VERIFICATION_SUMMARY_CSS}</style>
</head>
<body>
{body}
</body>
</html>"""


def _build_portfolio_verification_summary_html(results_df, warnings: list, timestamp: str) -> str:
    """Build a short, printable portfolio-level gap report for a reporting cycle."""
    _prepared_by = st.session_state.get("report_prepared_by", "")
    _status      = st.session_state.get("report_status", "")
    _notes       = st.session_state.get("report_notes", "")
    dims = [d for d, _ in _PORTFOLIO_SUBSCORE_DIMENSIONS]
    n = len(results_df)
    avg_conf = results_df["confidence_score"].mean() if n else 0
    avg_clar = results_df["clarity_score"].mean() if n else 0

    if n and dims:
        weakest_dim = results_df[dims].mean().idxmin()
        weakest_pct = results_df[dims].mean().min()
        systemic_gap_html = (
            f"<p><strong>Systemic gap:</strong> <em>{weakest_dim}</em> is this portfolio's "
            f"weakest sub-score on average ({weakest_pct:.0f}% of target) — start here for "
            f"the biggest improvement across multiple indicators.</p>"
        )
    else:
        systemic_gap_html = ""

    table_rows = "".join(
        f"<tr><td>{r['indicator_name']}</td><td>{r['logframe_indicator'] or '-'}</td>"
        f"<td style='font-family:monospace;'>{r['confidence_score']}/5.0</td>"
        f"<td style='font-family:monospace;'>{r['clarity_score']}/5.0</td>"
        f"<td>{r['verdict']}</td><td>{r['top_fix'] or '-'}</td></tr>"
        for _, r in results_df.iterrows()
    )

    warnings_html = ""
    if warnings:
        warnings_html = (
            "<h2>Rows Skipped</h2><ul>"
            + "".join(f"<li>{w}</li>" for w in warnings)
            + "</ul>"
        )

    body = f"""
<h1>Portfolio Verification Summary</h1>
<p style="color:#616161;font-size:0.85rem;">Generated: {timestamp}</p>
<p><strong>Indicators evaluated:</strong> {n}<br/>
   <strong>Average Confidence:</strong> {avg_conf:.1f}/5.0<br/>
   <strong>Average Clarity:</strong> {avg_clar:.1f}/5.0</p>
{systemic_gap_html}

<h2>Results by Indicator</h2>
<table>
  <tr><th>Indicator</th><th>Logframe Reference</th><th>Confidence</th><th>Clarity</th><th>Verdict</th><th>Top Fix</th></tr>
  {table_rows}
</table>
{warnings_html}

<h2>Donor Framework Crosswalk</h2>
{build_donor_crosswalk_html(st.session_state.get('donor_framework', 'Generic'))}
<p style="color:#616161;font-size:0.8rem;">{_LIMITS_DISCLAIMER}</p>
{_verification_summary_signoff(_prepared_by, _status, _notes, timestamp)}
<div class="footer">
  Evaluated using {METHODOLOGY_STACK}.<br/>
  Contact: <a href="mailto:info@impact-receipts.com">info@impact-receipts.com</a>
</div>
"""
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Portfolio Verification Summary</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;700&display=swap" rel="stylesheet"/>
<style>{_VERIFICATION_SUMMARY_CSS}</style>
</head>
<body>
{body}
</body>
</html>"""


# ---------------------------------------------------------------------------
# Admin view — hidden usage-metrics dashboard, ?admin=1 + passphrase-gated
# ---------------------------------------------------------------------------

_ADMIN_EVENT_LABELS = [
    ("demo_viewed", "Demo views"),
    ("check_completed", "Checks completed"),
    ("ai_questions_generated", "AI reviews run"),
    ("draft_withheld_fabrication", "Drafts withheld"),
    ("payment_initiated", "Payments started"),
    ("payment_completed", "Payments completed"),
]


def _admin_totals_chart(totals: dict):
    """Horizontal bar chart of event totals — single brand hue (nominal
    categories, no status/identity dimension to encode), sorted by magnitude,
    value labelled at the bar tip."""
    import pandas as pd
    import altair as alt

    df = pd.DataFrame([
        {"Event": label, "Count": totals.get(key, 0)} for key, label in _ADMIN_EVENT_LABELS
    ])
    bar = alt.Chart(df).mark_bar(size=20, cornerRadiusEnd=4, color="#1B5E20").encode(
        x=alt.X("Count:Q", title="Count"),
        y=alt.Y("Event:N", sort="-x", title=None),
        tooltip=[alt.Tooltip("Event:N", title="Event"), alt.Tooltip("Count:Q", title="Count")],
    )
    text = bar.mark_text(align="left", dx=4, color="#212121").encode(text="Count:Q")
    return (bar + text).properties(width="container", height=alt.Step(28))


def _admin_funnel_chart(funnel: dict):
    """Funnel-stage bar chart — ordinal (stage order carries meaning), one
    hue with a monotone opacity ramp so darker = further down the funnel."""
    import pandas as pd
    import altair as alt

    stages = [
        ("Demo viewed", funnel.get("demo_viewed", 0)),
        ("Check completed", funnel.get("check_completed", 0)),
        ("Payment completed", funnel.get("payment_completed", 0)),
    ]
    base = stages[0][1] or 1
    df = pd.DataFrame([
        {"Stage": s, "Sessions": v, "Order": i, "% of demo views": round(v / base * 100)}
        for i, (s, v) in enumerate(stages)
    ])
    bar = alt.Chart(df).mark_bar(size=28, cornerRadiusEnd=4, color="#1B5E20").encode(
        x=alt.X("Sessions:Q", title="Distinct sessions"),
        y=alt.Y("Stage:N", sort=[s for s, _ in stages], title=None),
        opacity=alt.Opacity("Order:O", scale=alt.Scale(range=[0.45, 1.0]), legend=None),
        tooltip=[
            alt.Tooltip("Stage:N", title="Stage"),
            alt.Tooltip("Sessions:Q", title="Sessions"),
            alt.Tooltip("% of demo views:Q", title="% of demo views"),
        ],
    )
    text = bar.mark_text(align="left", dx=4, color="#212121").encode(text="Sessions:Q")
    return (bar + text).properties(width="container", height=alt.Step(32))


def _render_admin_charts(totals: dict, funnel: dict, daily: list):
    """Render the three admin charts, Altair normally, Streamlit-native
    charts in lite mode (same fallback convention as _render_trend_chart)."""
    import pandas as pd

    lite = st.session_state.get("lite_mode", False)

    st.markdown("##### Totals by event")
    if lite:
        st.bar_chart(pd.DataFrame(
            {"Count": [totals.get(k, 0) for k, _ in _ADMIN_EVENT_LABELS]},
            index=[label for _, label in _ADMIN_EVENT_LABELS],
        ))
    else:
        st.altair_chart(_admin_totals_chart(totals), use_container_width=True)

    st.markdown("##### Conversion funnel")
    if lite:
        st.bar_chart(pd.DataFrame(
            {"Sessions": [funnel.get("demo_viewed", 0), funnel.get("check_completed", 0),
                          funnel.get("payment_completed", 0)]},
            index=["Demo viewed", "Check completed", "Payment completed"],
        ))
    else:
        st.altair_chart(_admin_funnel_chart(funnel), use_container_width=True)

    st.markdown("##### Daily activity")
    if not daily:
        st.caption("No events logged yet.")
        return
    df = pd.DataFrame(daily)
    if lite:
        st.line_chart(df.set_index("date")[["count"]])
    else:
        import altair as alt
        df["date"] = pd.to_datetime(df["date"])
        line = (
            alt.Chart(df)
            .mark_line(point=alt.OverlayMarkDef(size=50, filled=True, color="#1B5E20"),
                       color="#1B5E20", strokeWidth=2)
            .encode(
                x=alt.X("date:T", title="Date"),
                y=alt.Y("count:Q", title="Events per day"),
                tooltip=[alt.Tooltip("date:T", title="Date"), alt.Tooltip("count:Q", title="Events")],
            )
            .properties(width="container", height=200)
        )
        st.altair_chart(line, use_container_width=True)


def _render_admin_view():
    st.markdown("### Admin — usage metrics")
    st.caption("Anonymous usage counts only — no result text or documents are ever logged.")

    _admin_secret = (
        st.secrets.get("ADMIN_PASSPHRASE", "")
        if hasattr(st, "secrets") else
        os.environ.get("ADMIN_PASSPHRASE", "")
    )
    if not _admin_secret:
        st.error("Admin view is not configured — set ADMIN_PASSPHRASE in secrets.")
        return

    _entered = st.text_input("Passphrase", type="password", key="_admin_passphrase_input")
    if _entered != _admin_secret:
        if _entered:
            st.warning("Incorrect passphrase.")
        return

    summary = metrics.summarize()
    totals = summary["totals"]
    funnel = summary["funnel"]

    st.markdown("#### Totals by event")
    _cols = st.columns(3)
    for i, (key, label) in enumerate(_ADMIN_EVENT_LABELS):
        with _cols[i % 3]:
            st.metric(label, totals.get(key, 0))

    st.markdown("#### Score uplift")
    st.metric("Average uplift per re-score", summary["average_uplift"])

    st.markdown("#### Conversion funnel (distinct sessions)")
    _fcols = st.columns(3)
    with _fcols[0]:
        st.metric("Demo viewed", funnel["demo_viewed"])
    with _fcols[1]:
        st.metric("Check completed", funnel["check_completed"])
    with _fcols[2]:
        st.metric("Payment completed", funnel["payment_completed"])

    st.divider()
    st.markdown("#### Dashboard")
    _render_admin_charts(totals, funnel, metrics.daily_counts())


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def main():
    st.set_page_config(
        page_title="ImpactProof",
        page_icon="✅",
        layout="centered",
        initial_sidebar_state="collapsed",
    )

    st.markdown(CSS, unsafe_allow_html=True)
    inject_matchday_css()
    _init_session_state()
    _restore_session_from_query_param()

    with st.sidebar:
        st.toggle(
            "⚡ Low-bandwidth mode",
            key="lite_mode",
            help="Fewer interactive charts and a smaller, self-contained "
                 "downloadable report — for offices with slow or unreliable "
                 "internet.",
        )
        if st.button("💳 Billing & account", key="sidebar_billing_btn", use_container_width=True):
            st.session_state["_show_billing"] = True
            st.query_params["billing"] = "1"
            st.rerun()
        if st.button("📁 My Audits", key="sidebar_my_audits_btn", use_container_width=True):
            st.session_state["_show_my_audits"] = True
            st.query_params["my_audits"] = "1"
            st.rerun()

    # --- Paystack payment callback handler ---
    # Paystack always appends "?trxref=...&reference=..." to the callback_url
    # on redirect, even if it already contains a "?" — so the callback_url is
    # kept bare (see utils/paystack.initialize_payment) and the paying user's
    # email is recovered from Paystack's own verify response below, rather
    # than round-tripped through our own query-string params.
    _paystack_ref = (
        st.query_params.get("paystack_ref", "")
        or st.query_params.get("reference", "")
        or st.query_params.get("trxref", "")
    )
    if _paystack_ref:
        _pay_result = verify_payment(_paystack_ref)
        if _pay_result.get("status") == "success":
            _pay_email = (_pay_result.get("email") or st.session_state.get("user_email") or "").strip().lower()
            if _pay_email:
                _days = 365 if _pay_result.get("plan") == "annual" else (30 if _pay_result.get("plan") in ("monthly", "agency") else 1)
                # Track agency-tier users
                if _pay_result.get("plan") == "agency":
                    st.session_state["_is_agency"] = True
                upsert_user(_pay_email)
                mark_paid(_pay_email, days=_days)
                st.session_state["user_email"] = _pay_email
                st.session_state["is_paid"] = True
                metrics.log_event("payment_completed", _pay_email)
                st.session_state.pop("_pay_once_url", None)
                st.session_state.pop("_pay_monthly_url", None)
                st.session_state.pop("_pay_agency_url", None)
                st.session_state.pop("_pay_annual_url", None)
                st.session_state["screen"] = 1
                st.session_state["current_tab"] = 0
                st.query_params["screen"] = "1"
                st.query_params["tab"] = "0"
                st.session_state["entry_mode"] = "⚡ Instant Report Check"
                st.session_state["_payment_success"] = True
                try:
                    st.query_params.clear()
                except Exception:
                    pass
                st.rerun()
            else:
                # Paid, but no email on either side yet — defer until the
                # user signs in via the email gate, which completes this.
                st.session_state["pending_paystack_ref"] = _paystack_ref
                try:
                    st.query_params.clear()
                except Exception:
                    pass
        elif _pay_result.get("status") == "failed":
            st.warning("Payment didn't go through. Please try again.")
            try:
                st.query_params.clear()
            except Exception:
                pass
    # --- End Paystack handler ---

    _init_from_query_params()
    # Magic-link login confirm-click landing — shown regardless of screen
    _login_token = st.query_params.get("login_token", "")
    if _login_token:
        _render_login_link_landing(_login_token)
        return
    # Hidden admin view — usage metrics, passphrase-gated
    if st.query_params.get("admin") == "1":
        _render_admin_view()
        return
    # Pricing overlay — shown regardless of screen
    if st.session_state.get("_show_pricing"):
        render_pricing_page()
        return
    # Billing & account overlay — shown regardless of screen
    if st.session_state.get("_show_billing"):
        render_billing_page()
        return
    # My Audits overlay — shown regardless of screen
    if st.session_state.get("_show_my_audits"):
        render_my_audits_page()
        return
    try:
        screen = st.session_state["screen"]
        {0: render_screen_0, 1: render_screen_1, 2: render_screen_2, 3: render_screen_3}.get(
            screen, render_screen_0
        )()
    except Exception as _top_exc:
        import logging as _logging
        _logging.error("Unhandled top-level exception", exc_info=True)
        st.error(
            "Something went wrong rendering the app. Please refresh the page. "
            "If the problem persists, contact us: info@impact-receipts.com\n\n"
            f"Details: {type(_top_exc).__name__}: {_top_exc}"
        )
        # Error support WhatsApp CTA — server-side notification (council XXIV)
        _err_email  = st.session_state.get("user_email", "")
        _err_wa_key = "wa_error_support_clicked"
        if st.button("📱 Report this error on WhatsApp", key="wa_error_support_btn"):
            try:
                from utils.whatsapp import notify_founder
                notify_founder("error_support", user_email=_err_email,
                               result_data={"verdict": type(_top_exc).__name__})
            except Exception:
                pass
            st.session_state[_err_wa_key] = True
        if st.session_state.get(_err_wa_key):
            try:
                from utils.whatsapp import build_wa_url
                st.link_button("Open WhatsApp →",
                               build_wa_url("error_support", _err_email),
                               use_container_width=True)
                st.caption("✓ Error reported — we'll investigate and get back to you shortly.")
            except Exception:
                st.markdown("[WhatsApp +233 50 364 8195](https://wa.me/233503648195)")


if __name__ == "__main__":
    main()
